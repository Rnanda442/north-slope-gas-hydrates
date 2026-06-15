from __future__ import annotations

import csv
import hashlib
import math
from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[2]
BLUEPRINT_DIR = ROOT / "docs" / "project_blueprints"
ASSET_DIR = BLUEPRINT_DIR / "presentation_assets" / "stability_ml_remake_2026_06_15"
SOURCE_DECK = BLUEPRINT_DIR / "CURRENT_GMAIL_VISUAL_REVISION_9_SLIDE_North_Slope_Gas_Hydrate_Slides_2026-06-11.pptx"
OUT_DECK = BLUEPRINT_DIR / "STABILITY_ML_REMAKE_9_SLIDE_North_Slope_Gas_Hydrate_Slides_2026-06-15.pptx"
LOCKED_SLIDE_HASHES: dict[int, str] = {}

PUBLIC_PRODUCTS = ROOT / "data" / "public_stability_products"
REVISION_IMAGES = ROOT / "references" / "presentation-revision-2026-06-11" / "images"
GMAIL_IMAGES = ROOT / "references" / "presentation-revision-2026-06-11" / "gmail-2026-06-11"

W, H = 1600, 900

NAVY = (8, 33, 49)
DEEP = (4, 18, 28)
TEAL = (18, 124, 139)
ICE = (91, 203, 219)
ICE_LIGHT = (221, 244, 247)
GREEN = (34, 162, 127)
GREEN_LIGHT = (218, 245, 236)
AMBER = (219, 160, 60)
AMBER_LIGHT = (248, 235, 203)
RED = (203, 71, 70)
RED_LIGHT = (250, 229, 228)
BLUE = (55, 134, 205)
PURPLE = (125, 133, 213)
MUTED = (82, 103, 114)
LIGHT = (244, 249, 250)
PANEL = (231, 242, 244)
LINE = (184, 211, 218)
WHITE = (255, 255, 255)
BLACKISH = (21, 47, 60)


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        "arialbd.ttf" if bold else "arial.ttf",
        "segoeuib.ttf" if bold else "segoeui.ttf",
        "calibrib.ttf" if bold else "calibri.ttf",
    ]
    roots = [Path("C:/Windows/Fonts"), Path("/usr/share/fonts/truetype/dejavu")]
    for name in candidates:
        for root in roots:
            path = root / name
            if path.exists():
                return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def wrap_lines(draw: ImageDraw.ImageDraw, text: str, size: int, width: int, bold: bool = False) -> list[str]:
    f = font(size, bold)
    lines: list[str] = []
    for raw in text.splitlines():
        if not raw:
            lines.append("")
            continue
        current = ""
        for word in raw.split():
            test = f"{current} {word}".strip()
            if draw.textlength(test, font=f) <= width or not current:
                current = test
            else:
                lines.append(current)
                current = word
        if current:
            lines.append(current)
    return lines


def draw_text(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    size: int,
    fill: tuple[int, int, int] = NAVY,
    bold: bool = False,
    width: int | None = None,
    align: str = "left",
    line_gap: int = 6,
) -> int:
    f = font(size, bold)
    if width is None:
        draw.text(xy, text, font=f, fill=fill)
        return xy[1] + size

    lines = wrap_lines(draw, text, size, width, bold)
    y = xy[1]
    for line in lines:
        if align == "center":
            x = xy[0] + width // 2 - int(draw.textlength(line, font=f) // 2)
        elif align == "right":
            x = xy[0] + width - int(draw.textlength(line, font=f))
        else:
            x = xy[0]
        draw.text((x, y), line, font=f, fill=fill)
        y += size + line_gap
    return y


def title(draw: ImageDraw.ImageDraw, main: str, sub: str, dark: bool = False) -> None:
    draw_text(draw, (58, 42), main, 38, WHITE if dark else NAVY, True)
    if sub:
        draw_text(draw, (61, 92), sub, 17, (194, 229, 235) if dark else MUTED, width=1320)


def footer(draw: ImageDraw.ImageDraw, text: str, dark: bool = False) -> None:
    y = H - 52
    draw.line((54, y - 8, W - 54, y - 8), fill=(70, 94, 105) if dark else LINE, width=2)
    draw_text(draw, (58, y), text, 12, (166, 192, 201) if dark else MUTED, width=1420)


def new_canvas(dark: bool = False) -> Image.Image:
    img = Image.new("RGB", (W, H), DEEP if dark else WHITE)
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 18, H), fill=ICE if dark else TEAL)
    return img


def card(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    fill: tuple[int, int, int] = WHITE,
    outline: tuple[int, int, int] = LINE,
    radius: int = 16,
    width: int = 2,
) -> None:
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def chip(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    label: str,
    fill: tuple[int, int, int],
    text_fill: tuple[int, int, int] = NAVY,
    bold: bool = True,
) -> None:
    card(draw, box, fill=fill, outline=fill, radius=13, width=1)
    f = font(15, bold)
    x = (box[0] + box[2]) // 2 - int(draw.textlength(label, font=f) // 2)
    y = box[1] + (box[3] - box[1] - 15) // 2 - 2
    draw.text((x, y), label, font=f, fill=text_fill)


def arrow(
    draw: ImageDraw.ImageDraw,
    start: tuple[int, int],
    end: tuple[int, int],
    fill: tuple[int, int, int] = TEAL,
    width: int = 4,
) -> None:
    draw.line((start, end), fill=fill, width=width)
    angle = math.atan2(end[1] - start[1], end[0] - start[0])
    length = 15
    pts = [
        end,
        (int(end[0] - length * math.cos(angle - 0.45)), int(end[1] - length * math.sin(angle - 0.45))),
        (int(end[0] - length * math.cos(angle + 0.45)), int(end[1] - length * math.sin(angle + 0.45))),
    ]
    draw.polygon(pts, fill=fill)


def paste_cover(base: Image.Image, path: Path, box: tuple[int, int, int, int], tint: tuple[int, int, int] | None = None) -> None:
    draw = ImageDraw.Draw(base)
    if not path.exists():
        card(draw, box, fill=RED_LIGHT, outline=RED)
        draw_text(draw, (box[0] + 18, box[1] + 18), f"Missing asset: {path.name}", 16, RED, True, width=box[2] - box[0] - 36)
        return

    img = Image.open(path).convert("RGB")
    bw, bh = box[2] - box[0], box[3] - box[1]
    scale = max(bw / img.width, bh / img.height)
    resized = img.resize((max(1, int(img.width * scale)), max(1, int(img.height * scale))), Image.Resampling.LANCZOS)
    left = max(0, (resized.width - bw) // 2)
    top = max(0, (resized.height - bh) // 2)
    cropped = resized.crop((left, top, left + bw, top + bh))
    if tint:
        overlay = Image.new("RGB", cropped.size, tint)
        cropped = Image.blend(cropped, overlay, 0.18)
    base.paste(cropped, (box[0], box[1]))
    draw.rounded_rectangle(box, radius=18, outline=LINE, width=2)


def paste_contain(base: Image.Image, path: Path, box: tuple[int, int, int, int]) -> None:
    draw = ImageDraw.Draw(base)
    card(draw, box, fill=WHITE, outline=LINE)
    if not path.exists():
        draw_text(draw, (box[0] + 18, box[1] + 18), f"Missing asset: {path.name}", 16, RED, True, width=box[2] - box[0] - 36)
        return
    img = Image.open(path).convert("RGB")
    bw, bh = box[2] - box[0], box[3] - box[1]
    scale = min((bw - 14) / img.width, (bh - 14) / img.height)
    resized = img.resize((max(1, int(img.width * scale)), max(1, int(img.height * scale))), Image.Resampling.LANCZOS)
    x = box[0] + (bw - resized.width) // 2
    y = box[1] + (bh - resized.height) // 2
    base.paste(resized, (x, y))


def read_summary(path: Path) -> dict[str, str]:
    if not path.exists():
        return {}
    with path.open(newline="", encoding="utf-8") as fh:
        reader = csv.DictReader(fh)
        return {row["metric"]: row["value"] for row in reader}


def metric_value(summary: dict[str, str], key: str, fallback: str) -> str:
    value = summary.get(key, fallback)
    try:
        number = float(value)
    except ValueError:
        return value
    if number.is_integer():
        return f"{int(number):,}"
    return f"{number:,.1f}"


def draw_status_card(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    label: str,
    value: str,
    note: str,
    accent: tuple[int, int, int],
    fill: tuple[int, int, int] = WHITE,
) -> None:
    card(draw, box, fill=fill, outline=LINE)
    draw.rounded_rectangle((box[0], box[1], box[0] + 12, box[3]), radius=8, fill=accent)
    draw_text(draw, (box[0] + 28, box[1] + 20), label, 16, MUTED, True, width=box[2] - box[0] - 46)
    draw_text(draw, (box[0] + 28, box[1] + 50), value, 36, accent, True)
    draw_text(draw, (box[0] + 28, box[1] + 98), note, 14, NAVY, width=box[2] - box[0] - 46)


def draw_pipeline_node(
    draw: ImageDraw.ImageDraw,
    center: tuple[int, int],
    label: str,
    fill: tuple[int, int, int],
    w: int = 178,
    h: int = 74,
    text_fill: tuple[int, int, int] = NAVY,
) -> tuple[int, int, int, int]:
    x, y = center
    box = (x - w // 2, y - h // 2, x + w // 2, y + h // 2)
    card(draw, box, fill=fill, outline=fill, radius=20)
    draw_text(draw, (box[0] + 14, box[1] + 18), label, 16, text_fill, True, width=w - 28, align="center")
    return box


def slide_01() -> Path:
    img = new_canvas(False)
    draw = ImageDraw.Draw(img)
    title(draw, "Gas Hydrate Occurrence and Saturation Prediction", "A public-safe workflow from source-backed stability screening to approved-data ML.")

    paste_cover(img, BLUEPRINT_DIR / "presentation_assets" / "rohan_profile_photo.jpg", (1110, 136, 1518, 650))
    draw.rectangle((1110, 568, 1518, 650), fill=(255, 255, 255))
    draw_text(draw, (1136, 590), "Presenter / project lead", 18, NAVY, True)
    draw_text(draw, (1136, 620), "North Slope gas-hydrate research workflow", 14, MUTED, width=330)

    draw_text(draw, (70, 160), "What this deck now explains", 24, NAVY, True)
    steps = [
        ("Public sources", "GitHub + Streamlit explain the source-backed workflow.", TEAL, ICE_LIGHT),
        ("OSL workbench", "OpenScienceLab handles heavy source bundles and guarded calculations.", BLUE, (225, 239, 250)),
        ("Approved ML later", "Occurrence and saturation models wait for approved log/core validation.", GREEN, GREEN_LIGHT),
    ]
    x = 76
    prev = None
    for label, note, color, fill in steps:
        box = draw_pipeline_node(draw, (x + 120, 338), label, fill, w=240, h=82, text_fill=color)
        draw_text(draw, (box[0] + 18, box[3] + 16), note, 15, MUTED, width=210, align="center")
        if prev:
            arrow(draw, (prev[2] + 14, 338), (box[0] - 14, 338), fill=TEAL, width=4)
        prev = box
        x += 310

    chip(draw, (88, 565, 250, 606), "9-slide story", TEAL, WHITE)
    chip(draw, (278, 565, 475, 606), "stability screen", ICE_LIGHT, NAVY)
    chip(draw, (504, 565, 676, 606), "future ML", GREEN_LIGHT, NAVY)

    card(draw, (74, 665, 1006, 792), fill=LIGHT, outline=LINE)
    draw_text(draw, (104, 690), "Non-negotiable guardrail", 19, RED, True)
    draw_text(
        draw,
        (104, 724),
        "The current stability product is an admissibility screen. It is not hydrate proof, not saturation, and not a sweet-spot ranking.",
        21,
        NAVY,
        True,
        width=840,
    )
    footer(draw, "Deck authority: current Gmail 9-slide deck; remake generated as public-safe raster panels.")
    return save_slide(img, "slide_01_title_project_promise.png")


def slide_02() -> Path:
    img = new_canvas(False)
    draw = ImageDraw.Draw(img)
    title(draw, "Methane Hydrate Is A Stability-Gated Material", "Pressure-temperature conditions are necessary, but logs/core still decide occurrence.")

    left = (70, 155, 730, 705)
    paste_cover(img, REVISION_IMAGES / "usgs_gas_hydrate_crystals_sem_public_domain.jpg", left, tint=(28, 82, 103))
    draw.rounded_rectangle((left[0], left[1], left[2], left[1] + 86), radius=18, fill=(6, 31, 45))
    draw_text(draw, (left[0] + 30, left[1] + 28), "USGS hydrate crystals + clathrate concept", 24, WHITE, True)
    cx, cy = 470, 432
    for r in (138, 100):
        pts = []
        for i in range(8):
            a = 2 * math.pi * i / 8 + 0.18
            pts.append((cx + int(math.cos(a) * r), cy + int(math.sin(a) * r)))
        draw.line(pts + [pts[0]], fill=ICE, width=5)
        for px, py in pts:
            draw.ellipse((px - 10, py - 10, px + 10, py + 10), fill=ICE_LIGHT, outline=WHITE, width=2)
    draw.ellipse((cx - 58, cy - 58, cx + 58, cy + 58), fill=AMBER, outline=WHITE, width=5)
    draw_text(draw, (cx - 33, cy - 17), "CH4", 31, WHITE, True)
    draw_text(draw, (120, 605), "water cages + methane guest -> clathrate solid", 23, WHITE, True)

    right = (830, 150, 1505, 705)
    card(draw, right, fill=WHITE, outline=LINE)
    draw_text(draw, (862, 178), "P-T stability gate", 27, NAVY, True)
    # Axes and envelope.
    ax0, ay0, ax1, ay1 = 910, 620, 1460, 240
    draw.line((ax0, ay0, ax0, ay1), fill=MUTED, width=4)
    draw.line((ax0, ay0, ax1, ay0), fill=MUTED, width=4)
    arrow(draw, (ax0, ay1 + 24), (ax0, ay1), fill=MUTED, width=4)
    arrow(draw, (ax1 - 24, ay0), (ax1, ay0), fill=MUTED, width=4)
    draw_text(draw, (ax0 - 18, ay1 - 26), "P", 22, MUTED, True)
    draw_text(draw, (ax1 + 10, ay0 - 16), "T", 22, MUTED, True)
    envelope = [(ax0, 292), (1100, 300), (1255, 378), (1370, 575), (1000, 575)]
    draw.polygon(envelope, fill=ICE_LIGHT)
    curve = [(940, 580), (990, 520), (1048, 470), (1130, 420), (1225, 380), (1340, 342), (1440, 324)]
    draw.line(curve, fill=TEAL, width=5)
    geo = [(985, 620), (1090, 550), (1210, 470), (1335, 365), (1440, 290)]
    draw.line(geo, fill=RED, width=4)
    draw_text(draw, (1110, 365), "GHSZ", 34, TEAL, True)
    draw_text(draw, (1332, 386), "geotherm", 15, RED, True)
    draw_text(draw, (1110, 655), "too warm / unstable", 15, MUTED)
    chip(draw, (862, 655, 1248, 698), "baseline: 100% methane + 5 ppt salinity", ICE_LIGHT, NAVY)
    chip(draw, (1265, 655, 1485, 698), "necessary, not proof", RED_LIGHT, RED)

    footer(draw, "Sources: USGS SIR 2008-5175 methane 5 ppt phase boundary; USGS/NETL hydrate primers; USGS SEM image.")
    return save_slide(img, "slide_02_hydrate_stability_gate.png")


def slide_03() -> Path:
    img = new_canvas(False)
    draw = ImageDraw.Draw(img)
    title(draw, "The Model Sees Evidence Tiers, Not A Flat Parameter List", "The workflow checks can-exist, can-host, and behaves-like-hydrate in sequence.")

    tiers = [
        ("1", "Stability context", "depth | hydrostatic pressure | temperature model | methane 5 ppt phase curve", ICE, ICE_LIGHT),
        ("2", "Reservoir quality", "GR | porosity | RHOB | caliper QC | clean sand context", GREEN, GREEN_LIGHT),
        ("3", "Hydrate response", "Rt | Vp | Vs | Vp/Vs | AI | mu-rho | NMR/core support", TEAL, (222, 243, 239)),
    ]
    x0, y0, w, h = 105, 168, 950, 132
    previous = None
    for i, (num, label, items, accent, fill) in enumerate(tiers):
        y = y0 + i * 166
        card(draw, (x0, y, x0 + w, y + h), fill=fill, outline=accent, radius=22, width=3)
        draw.ellipse((x0 + 28, y + 31, x0 + 92, y + 95), fill=accent)
        draw_text(draw, (x0 + 52, y + 49), num, 25, WHITE, True)
        draw_text(draw, (x0 + 120, y + 27), label, 29, NAVY, True)
        draw_text(draw, (x0 + 120, y + 76), items, 20, MUTED, width=760)
        if previous:
            arrow(draw, (x0 + 475, previous[3] + 12), (x0 + 475, y - 12), fill=accent, width=4)
        previous = (x0, y, x0 + w, y + h)

    # Leakage rail.
    rail = (1125, 170, 1500, 636)
    card(draw, rail, fill=RED_LIGHT, outline=RED, radius=24, width=3)
    draw_text(draw, (1160, 202), "Target fields stay outside inputs", 23, RED, True, width=300)
    for j, field in enumerate(["S_h", "Sgh", "NMR_SAT", "phase labels", "final ranks"]):
        chip(draw, (1160, 285 + j * 54, 1460, 326 + j * 54), field, WHITE, RED)
    draw_text(draw, (1160, 575), "They supervise, calibrate, or review. They do not become predictors.", 17, NAVY, True, width=290)

    card(draw, (105, 704, 1500, 792), fill=LIGHT, outline=LINE)
    draw_text(draw, (136, 730), "Slide rule", 18, TEAL, True)
    draw_text(draw, (250, 730), "Ranges and screens guide review; no single parameter becomes a hydrate label.", 24, NAVY, True, width=1180)
    footer(draw, "Sources: WELL_LOG_REQUIREMENTS_MAP; SCIENCE_TO_ML_LOGIC_LADDER; ML pipeline source ledger.")
    return save_slide(img, "slide_03_evidence_tiers.png")


def slide_04() -> Path:
    img = new_canvas(True)
    draw = ImageDraw.Draw(img)
    title(draw, "Two Workspaces, One Guarded Pipeline", "Public delivery and heavy-data calculation stay connected, but the data boundary stays closed.", dark=True)

    public = (80, 160, 660, 500)
    osl = (940, 160, 1520, 500)
    card(draw, public, fill=(13, 54, 70), outline=ICE, radius=24, width=2)
    card(draw, osl, fill=(15, 57, 53), outline=GREEN, radius=24, width=2)
    draw_text(draw, (115, 192), "GitHub / Streamlit", 30, ICE, True)
    draw_text(draw, (115, 246), "public GIS\nsource docs\ndiagrams\npublic stability products", 22, WHITE, width=470, line_gap=11)
    draw_text(draw, (975, 192), "OpenScienceLab", 30, GREEN, True)
    draw_text(draw, (975, 246), "raw source bundles\napproved logs/core later\nruntime-only outputs\nreviewed derived products", 22, WHITE, width=470, line_gap=11)

    # Boundary.
    draw.rounded_rectangle((704, 152, 896, 510), radius=22, fill=(66, 25, 31), outline=RED, width=3)
    draw_text(draw, (728, 190), "Boundary", 25, WHITE, True)
    draw_text(draw, (728, 244), "No approved rows, restricted IDs, trained models, or sensitive outputs cross into public slides.", 18, (253, 224, 224), width=145, line_gap=6)
    arrow(draw, (660, 330), (704, 330), fill=ICE, width=4)
    arrow(draw, (896, 330), (940, 330), fill=GREEN, width=4)
    draw_text(draw, (695, 540), "reviewed summaries only", 18, ICE, True, width=230, align="center")

    # ML branch.
    draw_text(draw, (90, 575), "Future approved-data ML branch", 25, WHITE, True)
    nodes = [
        ("schema\n+ units", ICE_LIGHT),
        ("QC + depth\nalignment", ICE_LIGHT),
        ("features +\nstability", GREEN_LIGHT),
        ("target registry\n+ leakage lock", RED_LIGHT),
        ("occurrence\nclassifier", (226, 241, 255)),
        ("saturation\nregressor", (230, 246, 235)),
        ("uncertainty\n+ reason flags", AMBER_LIGHT),
    ]
    x = 125
    last = None
    for label, fill in nodes:
        box = draw_pipeline_node(draw, (x, 702), label, fill, w=165, h=78, text_fill=NAVY)
        if last:
            arrow(draw, (last[2] + 8, 702), (box[0] - 8, 702), fill=ICE, width=3)
        last = box
        x += 208
    footer(draw, "Guardrail: occurrence classification and saturation regression are linked outputs, not the same result.", dark=True)
    return save_slide(img, "slide_04_public_osl_guarded_pipeline.png")


def draw_log_track(draw: ImageDraw.ImageDraw, x: int, y: int, w: int, h: int, color: tuple[int, int, int], label: str, seed: float) -> None:
    card(draw, (x, y, x + w, y + h), fill=(249, 252, 253), outline=LINE, radius=10, width=1)
    draw_text(draw, (x + 8, y + 8), label, 13, color, True)
    draw.line((x + w // 2, y + 34, x + w // 2, y + h - 14), fill=(215, 230, 234), width=2)
    pts = []
    for i in range(36):
        py = y + 38 + int(i * (h - 58) / 35)
        val = 0.5 + 0.28 * math.sin(seed + i * 0.45) + 0.12 * math.sin(seed * 1.7 + i * 1.03)
        px = x + 12 + int(max(0.05, min(0.95, val)) * (w - 24))
        pts.append((px, py))
    draw.line(pts, fill=color, width=3)


def slide_05() -> Path:
    img = new_canvas(False)
    draw = ImageDraw.Draw(img)
    title(draw, "Hydrate Signals Have Mimics", "The model should learn multi-log agreement after stability, reservoir quality, QC, and false-positive checks.")

    # Synthetic log panel.
    panel = (80, 160, 920, 710)
    card(draw, panel, fill=WHITE, outline=LINE, radius=20)
    draw_text(draw, (112, 190), "One interval, competing explanations", 25, NAVY, True)
    track_labels = [("GR", AMBER, 0.4), ("Rt", RED, 1.5), ("RHOB", GREEN, 2.1), ("Vp", BLUE, 3.0), ("Vs", PURPLE, 4.1), ("NMR", TEAL, 5.0)]
    tx = 120
    for lab, col, seed in track_labels:
        draw_log_track(draw, tx, 245, 95, 395, col, lab, seed)
        tx += 125
    # Depth interval highlights.
    draw.rounded_rectangle((115, 357, 870, 430), radius=16, fill=(38, 174, 136, 52), outline=GREEN, width=3)
    draw_text(draw, (132, 372), "hydrate-supportive: clean sand + Rt + stiffness + context", 17, GREEN, True)
    draw.rounded_rectangle((115, 500, 870, 566), radius=16, fill=(230, 200, 90), outline=AMBER, width=2)
    draw_text(draw, (132, 514), "mimic risk: gas / ice / cement / shale / bad hole can copy one signal", 17, NAVY, True)
    draw_text(draw, (120, 660), "Synthetic pattern only: used to explain logic, not as project results.", 16, MUTED, width=730)

    # Mimic cards.
    mimic_cards = [
        ("High Rt", "hydrate, gas, ice, tight rock, salinity", RED),
        ("High Vp/Vs or stiffness", "hydrate, ice, cement, compaction", PURPLE),
        ("Low GR", "clean sand context, not hydrate proof", AMBER),
        ("Bad CAL", "washout can corrupt density, sonic, NMR, Rt", RED),
    ]
    y = 185
    for label, body, color in mimic_cards:
        draw_status_card(draw, (990, y, 1490, y + 106), label, "mimic check", body, color, fill=LIGHT)
        y += 128

    footer(draw, "Sources: Lee and Collett 2011; Haines et al. 2022; WELL_LOG_REQUIREMENTS_MAP; parameter logic ledger.")
    return save_slide(img, "slide_05_hydrate_signals_have_mimics.png")


def slide_06() -> Path:
    img = new_canvas(False)
    draw = ImageDraw.Draw(img)
    title(draw, "Physics Features Keep The ML Honest", "Feature equations create interpretable evidence; the stability screen stays a separate admissibility gate.")

    bands = [
        ("Raw logs", "RHOB | Vp | Vs | Rt | GR | CAL", TEAL, ICE_LIGHT),
        ("Derived features", "AI | Vp/Vs | G | K | mu-rho | lambda-rho", GREEN, GREEN_LIGHT),
        ("Stability screen", "hydrostatic pressure + temperature model + methane 5 ppt phase curve", BLUE, (226, 241, 255)),
    ]
    x0, y0 = 85, 170
    last_box = None
    for i, (label, items, accent, fill) in enumerate(bands):
        y = y0 + i * 148
        card(draw, (x0, y, 925, y + 104), fill=fill, outline=accent, radius=20, width=3)
        draw_text(draw, (x0 + 30, y + 20), label, 25, accent, True)
        draw_text(draw, (x0 + 300, y + 26), items, 24, NAVY, True, width=560)
        if last_box:
            arrow(draw, (505, last_box[3] + 10), (505, y - 10), fill=accent, width=4)
        last_box = (x0, y, 925, y + 104)

    card(draw, (1000, 158, 1508, 520), fill=WHITE, outline=LINE)
    draw_text(draw, (1032, 190), "Compact equation cluster", 25, NAVY, True)
    equations = [
        "AI = rho_b * Vp",
        "G = rho * Vs^2",
        "P_abs = 0.101325 + 0.00980665 * z_m",
        "stable when T_model <= T_eq(P_abs)",
    ]
    yy = 248
    for eq in equations:
        chip(draw, (1032, yy, 1468, yy + 42), eq, LIGHT, NAVY)
        yy += 62

    card(draw, (1000, 565, 1508, 724), fill=RED_LIGHT, outline=RED, radius=20, width=3)
    draw_text(draw, (1032, 595), "Equation caution", 23, RED, True)
    draw_text(draw, (1032, 633), "Equations create features and screens. They do not prove hydrate by themselves.", 22, NAVY, True, width=430)
    footer(draw, "Sources: stability calculation plan; runtime feature engineering docs; hydrate log-response sources.")
    return save_slide(img, "slide_06_physics_features_and_stability_screen.png")


def slide_07() -> Path:
    summaries = {
        "context": read_summary(PUBLIC_PRODUCTS / "north_slope_well_stability_context_summary_2026-06-14.csv"),
        "temp": read_summary(PUBLIC_PRODUCTS / "g10015_temperature_profile_summary_2026-06-14.csv"),
        "model": read_summary(PUBLIC_PRODUCTS / "stability_temperature_model_summary_2026-06-14.csv"),
        "screen": read_summary(PUBLIC_PRODUCTS / "stability_screen_summary_2026-06-14_methane_5ppt_v1.csv"),
    }
    img = new_canvas(False)
    draw = ImageDraw.Draw(img)
    title(draw, "The Public Site Shows Readiness And Context", "The app can display guarded stability status; approved logs/core still decide occurrence and saturation.")

    # Bridge.
    card(draw, (70, 160, 500, 455), fill=ICE_LIGHT, outline=TEAL, radius=22, width=3)
    draw_text(draw, (102, 192), "GitHub / Streamlit", 28, TEAL, True)
    draw_text(draw, (102, 246), "public summaries\nderived CSV products\nstatus charts\ndownloads", 22, NAVY, width=340, line_gap=10)
    card(draw, (1095, 160, 1520, 455), fill=GREEN_LIGHT, outline=GREEN, radius=22, width=3)
    draw_text(draw, (1127, 192), "OpenScienceLab", 28, GREEN, True)
    draw_text(draw, (1127, 246), "heavy source bundle\nG10015 profiles\nraw calculations\napproved data later", 22, NAVY, width=330, line_gap=10)
    draw.rounded_rectangle((585, 215, 1010, 398), radius=24, fill=RED_LIGHT, outline=RED, width=3)
    draw_text(draw, (625, 242), "public-safe bridge", 28, RED, True)
    draw_text(draw, (625, 292), "Only reviewed derived outputs move into the public deck/site.", 21, NAVY, True, width=340)
    arrow(draw, (500, 305), (585, 305), fill=TEAL, width=5)
    arrow(draw, (1010, 305), (1095, 305), fill=GREEN, width=5)

    # Metric cards.
    metrics = [
        ("scaffold wells", metric_value(summaries["context"], "Arctic Slope public wells", "8,084"), "current public stability scaffold", TEAL),
        (
            "G10015 profiles",
            metric_value(summaries["temp"], "G10015 profiles", "184"),
            f'{metric_value(summaries["temp"], "Unique well codes", "24")} well codes indexed',
            BLUE,
        ),
        ("key-depth rows", metric_value(summaries["model"], "Temperature model rows", "16,168"), "temperature-input product", ICE),
        ("calculated intervals", metric_value(summaries["screen"], "Calculated stability intervals", "22"), "baseline admissibility only", GREEN),
        ("blocked screen rows", metric_value(summaries["screen"], "Blocked rows", "8,054"), "source/calculation gates did not pass", RED),
    ]
    x, y = 75, 525
    for i, (label, value, note, color) in enumerate(metrics):
        draw_status_card(draw, (x + i * 300, y, x + i * 300 + 260, y + 150), label, value, note, color, fill=WHITE)

    card(draw, (215, 718, 1385, 792), fill=RED_LIGHT, outline=RED, radius=18, width=2)
    draw_text(draw, (255, 740), "All rows remain not hydrate proof.", 28, RED, True)
    draw_text(draw, (775, 741), "Context/readiness only: not occurrence, saturation, or sweet spots.", 18, NAVY, True, width=560)
    footer(draw, "Sources: data/public_stability_products summaries and guarded screen output.")
    return save_slide(img, "slide_07_public_stability_products.png")


def slide_08() -> Path:
    img = new_canvas(False)
    draw = ImageDraw.Draw(img)
    title(draw, "What Is Complete, Calculated, And Still Blocked", "The deliverable should show status honestly and leave ML results for approved validation.")

    columns = [
        ("Complete", GREEN, GREEN_LIGHT, ["public/OSL boundary", "methane 5 ppt baseline", "hydrostatic pressure model", "temperature-model logic", "guarded screen writer"]),
        ("Calculated", BLUE, (226, 241, 255), ["919 temperature key depths", "387 extrapolated key depths", "22 admissibility intervals", "source-control labels"]),
        ("Blocked / future", RED, RED_LIGHT, ["15,249 key depths blocked", "8,054 screen rows blocked", "approved log/core validation", "occurrence labels", "saturation targets", "trained ML metrics"]),
    ]
    x = 82
    for label, accent, fill, items in columns:
        card(draw, (x, 162, x + 450, 595), fill=fill, outline=accent, radius=24, width=3)
        draw_text(draw, (x + 30, 198), label, 30, accent, True)
        yy = 270
        for item in items:
            draw.ellipse((x + 34, yy + 6, x + 50, yy + 22), fill=accent)
            draw_text(draw, (x + 68, yy), item, 21, NAVY, True, width=330)
            yy += 58
        x += 510

    card(draw, (105, 655, 1495, 785), fill=LIGHT, outline=LINE, radius=22)
    draw_text(draw, (135, 685), "Future result slots after approved validation", 22, NAVY, True)
    slots = [("occurrence probability", TEAL), ("saturation estimate", GREEN), ("uncertainty", AMBER), ("review flags", RED)]
    sx = 140
    for label, color in slots:
        chip(draw, (sx, 730, sx + 300, 775), label, color, WHITE)
        sx += 330
    footer(draw, "Do not show fake confusion matrices, saturation tracks, or project performance metrics before approved validation.")
    return save_slide(img, "slide_08_readiness_and_validation_slots.png")


def slide_09() -> Path:
    img = new_canvas(False)
    draw = ImageDraw.Draw(img)
    title(draw, "A Defensible Workflow Before A Prediction Claim", "Stability, occurrence, saturation, uncertainty, and producibility stay separate.")

    center = (800, 420)
    draw.ellipse((center[0] - 205, center[1] - 155, center[0] + 205, center[1] + 155), fill=ICE_LIGHT, outline=TEAL, width=4)
    draw_text(draw, (center[0] - 145, center[1] - 52), "guarded workflow", 32, NAVY, True)
    draw_text(draw, (center[0] - 145, center[1] + 2), "ready for approved\nvalidation", 24, TEAL, True, width=290, align="center")

    decisions = [
        ("baseline curve?", "100% methane + 5 ppt", (270, 190), ICE),
        ("mixed-gas sensitivity?", "Collett/Holder curve only as sensitivity", (1160, 190), AMBER),
        ("confidence thresholds?", "high / medium / low source control", (220, 505), BLUE),
        ("OM-222 digitization?", "permafrost base surface next?", (1180, 505), GREEN),
        ("approved validation fields?", "which logs/core targets later", (400, 685), RED),
        ("two ML outputs?", "occurrence + saturation", (980, 685), TEAL),
    ]
    for label, note, pos, color in decisions:
        x, y = pos
        card(draw, (x - 190, y - 62, x + 190, y + 62), fill=WHITE, outline=color, radius=20, width=3)
        draw_text(draw, (x - 155, y - 42), label, 21, color, True, width=310, align="center")
        draw_text(draw, (x - 155, y - 4), note, 15, MUTED, width=310, align="center")
        arrow(draw, (x, y + (62 if y < center[1] else -62)), (center[0] + int((x - center[0]) * 0.42), center[1] + int((y - center[1]) * 0.42)), fill=color, width=3)

    draw_text(
        draw,
        (220, 792),
        "Final message: validation-ready workflow, not a hydrate claim.",
        26,
        NAVY,
        True,
        width=1160,
        align="center",
    )
    footer(draw, "Next decisions: phase curve baseline, mixed-gas sensitivity, control-distance thresholds, OM-222, validation fields, ML output framing.")
    return save_slide(img, "slide_09_conclusion_mentor_decisions.png")


def save_slide(img: Image.Image, name: str) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / name
    img.save(path, quality=94)
    return path


def preserve_source_slide_panel(slide_number: int, name: str) -> Path:
    """Copy a locked raster panel out of the current Gmail authority deck."""
    if not SOURCE_DECK.exists():
        raise FileNotFoundError(SOURCE_DECK)

    prs = Presentation(SOURCE_DECK)
    slide = prs.slides[slide_number - 1]
    pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
    if len(pictures) != 1:
        raise ValueError(f"Source slide {slide_number} should have exactly one raster panel, found {len(pictures)}")

    blob = pictures[0].image.blob
    with Image.open(BytesIO(blob)) as image:
        if image.size != (W, H):
            raise ValueError(f"Source slide {slide_number} image size is {image.size}, expected {(W, H)}")

    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / name
    path.write_bytes(blob)
    LOCKED_SLIDE_HASHES[slide_number] = hashlib.sha256(blob).hexdigest()
    return path


def build_slides() -> list[Path]:
    return [
        preserve_source_slide_panel(1, "slide_01_locked_from_current_gmail_deck.png"),
        preserve_source_slide_panel(2, "slide_02_locked_from_current_gmail_deck.png"),
        slide_03(),
        slide_04(),
        slide_05(),
        slide_06(),
        slide_07(),
        slide_08(),
        slide_09(),
    ]


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        element = shape._element  # noqa: SLF001 - python-pptx has no public remove API.
        element.getparent().remove(element)


def rebuild_deck(panel_paths: list[Path]) -> Path:
    if not SOURCE_DECK.exists():
        raise FileNotFoundError(SOURCE_DECK)
    prs = Presentation(SOURCE_DECK)
    if len(prs.slides) != 9:
        raise ValueError(f"Expected 9 slides in {SOURCE_DECK}, found {len(prs.slides)}")
    if len(panel_paths) != 9:
        raise ValueError(f"Expected 9 generated panels, found {len(panel_paths)}")

    for slide, panel in zip(prs.slides, panel_paths, strict=True):
        clear_slide(slide)
        slide.shapes.add_picture(str(panel), 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(OUT_DECK)
    return OUT_DECK


def verify_deck(path: Path) -> None:
    prs = Presentation(path)
    if len(prs.slides) != 9:
        raise ValueError(f"Expected 9 slides, found {len(prs.slides)}")
    for idx, slide in enumerate(prs.slides, start=1):
        pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
        if len(pictures) != 1:
            raise ValueError(f"Slide {idx} should have exactly one raster panel, found {len(pictures)}")
        blob = pictures[0].image.blob
        with Image.open(BytesIO(blob)) as image:
            if image.size != (W, H):
                raise ValueError(f"Slide {idx} image size is {image.size}, expected {(W, H)}")
        if idx in LOCKED_SLIDE_HASHES:
            actual = hashlib.sha256(blob).hexdigest()
            if actual != LOCKED_SLIDE_HASHES[idx]:
                raise ValueError(f"Slide {idx} no longer matches the locked source panel")


def main() -> None:
    panels = build_slides()
    deck = rebuild_deck(panels)
    verify_deck(deck)
    print(f"Wrote {deck}")
    print(f"Wrote panels to {ASSET_DIR}")


if __name__ == "__main__":
    main()
