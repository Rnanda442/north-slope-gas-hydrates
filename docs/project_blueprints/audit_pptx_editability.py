"""Audit whether a PPTX uses editable slide objects or flat raster slides.

This script reads a PPTX and reports per-slide counts for shapes, pictures,
text-bearing shapes, connector-like shapes, and likely full-slide raster panels.
It does not modify the source deck.
"""

from __future__ import annotations

import argparse
import csv
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_TOLERANCE_RATIO = 0.02


@dataclass(frozen=True)
class SlideAudit:
    slide_number: int
    shape_count: int
    picture_count: int
    text_shape_count: int
    connector_count: int
    full_slide_picture_count: int
    one_picture_full_slide_raster: bool
    passes_editability: bool


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()


def _is_picture(shape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE


def _is_connector_like(shape) -> bool:
    return shape.shape_type in {
        MSO_SHAPE_TYPE.LINE,
        MSO_SHAPE_TYPE.FREEFORM,
    }


def _is_full_slide_picture(shape, slide_width: int, slide_height: int) -> bool:
    if not _is_picture(shape):
        return False
    left_tol = int(slide_width * EMU_TOLERANCE_RATIO)
    top_tol = int(slide_height * EMU_TOLERANCE_RATIO)
    width_min = int(slide_width * (1 - EMU_TOLERANCE_RATIO))
    height_min = int(slide_height * (1 - EMU_TOLERANCE_RATIO))
    return (
        shape.left <= left_tol
        and shape.top <= top_tol
        and shape.width >= width_min
        and shape.height >= height_min
    )


def audit_presentation(
    pptx_path: Path,
    allowed_full_raster_slides: Iterable[int] = (),
) -> list[SlideAudit]:
    presentation = Presentation(str(pptx_path))
    allowed = set(allowed_full_raster_slides)
    results: list[SlideAudit] = []

    for index, slide in enumerate(presentation.slides, start=1):
        shapes = list(slide.shapes)
        picture_count = sum(1 for shape in shapes if _is_picture(shape))
        text_shape_count = sum(1 for shape in shapes if _shape_text(shape))
        connector_count = sum(1 for shape in shapes if _is_connector_like(shape))
        full_slide_picture_count = sum(
            1
            for shape in shapes
            if _is_full_slide_picture(shape, presentation.slide_width, presentation.slide_height)
        )
        one_picture_full_slide_raster = (
            len(shapes) == 1 and picture_count == 1 and full_slide_picture_count == 1
        )
        passes_editability = (
            (not one_picture_full_slide_raster or index in allowed)
            and text_shape_count > 0
        )
        results.append(
            SlideAudit(
                slide_number=index,
                shape_count=len(shapes),
                picture_count=picture_count,
                text_shape_count=text_shape_count,
                connector_count=connector_count,
                full_slide_picture_count=full_slide_picture_count,
                one_picture_full_slide_raster=one_picture_full_slide_raster,
                passes_editability=passes_editability,
            )
        )
    return results


def write_csv(results: list[SlideAudit], output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with output_path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(SlideAudit.__dataclass_fields__))
        writer.writeheader()
        for row in results:
            writer.writerow(row.__dict__)


def write_markdown(results: list[SlideAudit], output_path: Path, pptx_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    lines = [
        "# PPTX Editability Audit",
        "",
        f"Source deck: `{pptx_path}`",
        "",
        "| slide | shapes | pictures | text shapes | connectors | full-slide pictures | one-picture full-slide raster | pass |",
        "|---:|---:|---:|---:|---:|---:|---|---|",
    ]
    for row in results:
        lines.append(
            "| {slide} | {shapes} | {pictures} | {text} | {connectors} | {full} | {raster} | {passed} |".format(
                slide=row.slide_number,
                shapes=row.shape_count,
                pictures=row.picture_count,
                text=row.text_shape_count,
                connectors=row.connector_count,
                full=row.full_slide_picture_count,
                raster="yes" if row.one_picture_full_slide_raster else "no",
                passed="PASS" if row.passes_editability else "FAIL",
            )
        )
    output_path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def _parse_slide_list(raw: str) -> list[int]:
    if not raw:
        return []
    return [int(part.strip()) for part in raw.split(",") if part.strip()]


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--out-csv", type=Path)
    parser.add_argument("--out-md", type=Path)
    parser.add_argument("--allow-full-raster-slides", default="")
    args = parser.parse_args()

    results = audit_presentation(
        args.pptx,
        allowed_full_raster_slides=_parse_slide_list(args.allow_full_raster_slides),
    )
    if args.out_csv:
        write_csv(results, args.out_csv)
    if args.out_md:
        write_markdown(results, args.out_md, args.pptx)
    for row in results:
        status = "PASS" if row.passes_editability else "FAIL"
        print(
            f"slide {row.slide_number}: {status}; "
            f"shapes={row.shape_count}; pictures={row.picture_count}; "
            f"text_shapes={row.text_shape_count}; full_slide_raster={row.one_picture_full_slide_raster}"
        )
    return 0 if all(row.passes_editability for row in results) else 1


if __name__ == "__main__":
    raise SystemExit(main())
