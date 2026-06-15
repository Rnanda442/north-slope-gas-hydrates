from __future__ import annotations

import sys
from pathlib import Path


REPO_ROOT = Path(__file__).resolve().parents[2]
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

import pandas as pd
import plotly.graph_objects as go
from PIL import Image, ImageChops, ImageDraw, ImageFont
from plotly.subplots import make_subplots
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptInches, Pt as PptPt

from dashboard.well_log_engine import (
    generate_synthetic_logs,
    model_placeholder_figures,
    screen_intervals,
    sweet_spot_review_table,
    well_log_panel,
)
from dashboard.app import build_geographic_structural_figure
from dashboard.visual_story_data import HEADER_DERIVED_SYNTHETIC_NOTE, SOURCE_ANCHORS
OUT_DIR = REPO_ROOT / "docs" / "project_blueprints"
DOCX_OUT = OUT_DIR / "North_Slope_Gas_Hydrate_Reservoir_Characterization_Research_Overview.docx"
PPTX_OUT = OUT_DIR / "North_Slope_Gas_Hydrate_Reservoir_Characterization_Research_Overview.pptx"
ASSET_DIR = OUT_DIR / "presentation_assets"
PARAMETER_CSV = OUT_DIR / "ml_parameter_effect_tree.csv"
MASTER_3D = REPO_ROOT / "03_data_final" / "master_layers" / "north_slope_master_3d_surfaces.parquet"
MASTER_2D = REPO_ROOT / "03_data_final" / "master_layers" / "north_slope_master_2d_layers.parquet"
STREAMLIT_URL = "https://north-slope-gas-hydrates-vj67xkke9ksfzveon8ldt2.streamlit.app/"

NAVY = "123447"
TEAL = "167D8D"
ICE = "EAF5F6"
SAND = "F4EFE6"
GRAY = "EEF2F4"
INK = "1F2933"
MUTED = "5E6A71"
WHITE = "FFFFFF"


REFERENCES = [
    "Collett, T.S., Lewis, K.A., Zyrianova, M.V., and others. 2019. Assessment of undiscovered gas hydrate resources in the North Slope of Alaska, 2018. USGS Fact Sheet 2019-3037. DOI: 10.3133/fs20193037. https://pubs.usgs.gov/publication/fs20193037",
    "National Energy Technology Laboratory. Alaska North Slope Gas Hydrate Reservoir Characterization. Project summary. https://netl.doe.gov/node/6846",
    "U.S. Department of Energy. 2024. DOE and international partners complete gas hydrates production testing on Alaska North Slope. https://www.energy.gov/hgeo/articles/doe-and-international-partners-complete-gas-hydrates-production-testing-alaska-north",
    "Chong, L.B., Singh, H., Creason, C.G., and others. 2022. Application of machine learning to characterize gas hydrate reservoirs in Mackenzie Delta (Canada) and on the Alaska North Slope (USA). Computational Geosciences, 26, 991-1006. DOI: 10.1007/s10596-022-10151-9. https://link.springer.com/article/10.1007/s10596-022-10151-9",
    "Singh, H., Seol, Y., and Myshakin, E.M. 2021. Prediction of gas hydrate saturation using machine learning and optimal set of well-logs. Computational Geosciences, 25, 267-283. DOI: 10.1007/s10596-020-10004-3.",
    "Chong, L., Collett, T., Creason, C.G., Seol, Y., and Myshakin, E. 2024. Machine learning application to assess occurrence and saturations of methane hydrate in marine deposits offshore India. Interpretation, 12, T63-T75. DOI: 10.1190/int-2023-0056.1.",
    "Lee, M.W., and Collett, T.S. 2011. In-situ gas hydrate saturation estimated from various well logs at the Mount Elbert Gas Hydrate Stratigraphic Test Well, Alaska North Slope. Marine and Petroleum Geology, 28, 439-449. https://pubs.usgs.gov/publication/70036903",
    "Haines, S.S., Collett, T.S., Yoneda, J., Shimoda, N., Boswell, R., and Okinaka, N. 2022. Gas hydrate saturation estimates, gas hydrate occurrence, and reservoir characteristics based on well log data from the Hydrate-01 stratigraphic test well, Alaska North Slope. Energy & Fuels, 36, 3040-3050. DOI: 10.1021/acs.energyfuels.1c04100. https://pubs.acs.org/doi/10.1021/acs.energyfuels.1c04100",
    "Zyrianova, M.V., Collett, T.S., and Boswell, R. 2024. Characterization of structural, stratigraphic, and reservoir controls on gas hydrate occurrence in the Eileen Gas Hydrate Trend, Alaska North Slope. https://www.mdpi.com/2077-1312/12/3/472",
    "Aung, T.T., Naito, K., Tano, K., Tamaki, M., and Boswell, R. 2026. Alaska North Slope Extended-Duration Gas Hydrate Production Test Site Logging-While-Drilling Data Acquisition. Energy & Fuels. DOI: 10.1021/acs.energyfuels.5c06115.",
    "Yoneda, J., Hiruta, A., Oshima, M., Jin, Y., Ohtsuki, S., Arima, Y., Nakatsuka, Y., and Okinaka, N. 2026. Permeability Evaluation of Hydrate Reservoirs Based on NMR T2 Relaxation Time from Both Log and Laboratory Data, Alaska North Slope HYDRATE 02 Geo Data Well. Energy & Fuels. DOI: 10.1021/acs.energyfuels.5c05321.",
    "Tian, D., Yang, S., Gong, Y., Geng, M., Li, Y., and Hu, G. 2023. A Comparative Study of Machine Learning Methods for Gas Hydrate Identification. Geoenergy Science and Engineering, 223, 211564. DOI: 10.1016/j.geoen.2023.211564.",
    "Li, C., and Liu, X. 2020. Research on the Estimate of Gas Hydrate Saturation Based on LSTM Recurrent Neural Network. Energies, 13, 6536. DOI: 10.3390/en13246536.",
    "Naim, F., Cook, A.E., and Moortgat, J. 2023. Estimating Compressional Velocity and Bulk Density Logs in Marine Gas Hydrates Using Machine Learning. Energies, 16, 7709. DOI: 10.3390/en16237709.",
    "Cook, A.E., and Waite, W.F. 2018. Archie's saturation exponent for natural gas hydrate in coarse-grained reservoirs. Journal of Geophysical Research: Solid Earth, 123.",
    "Rohan Nanda. 2026. ML Project Reference and CreditScoreV4 Case Notes. User-supplied general ML methodology notes recovered from Gmail on 2026-06-11 and stored in references/ml-sources/2026-06-11/.",
]


def export_plotly(figure: go.Figure, path: Path, width=1200, height=720) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    figure.write_image(str(path), width=width, height=height, scale=1.7)
    trim_white_margin(path)
    return path


def trim_white_margin(path: Path, padding: int = 28) -> None:
    image = Image.open(path).convert("RGB")
    background = Image.new("RGB", image.size, (255, 255, 255))
    diff = ImageChops.difference(image, background)
    bbox = diff.getbbox()
    if not bbox:
        image.save(path)
        return
    left = max(0, bbox[0] - padding)
    top = max(0, bbox[1] - padding)
    right = min(image.size[0], bbox[2] + padding)
    bottom = min(image.size[1], bbox[3] + padding)
    image.crop((left, top, right, bottom)).save(path)


def project_asset_path(name: str) -> Path:
    return ASSET_DIR / name


def build_regional_3d_asset() -> Path:
    figure = build_geographic_structural_figure(
        ["NStopo", "NSLCU", "NSshublik", "NSbasement"],
        1800,
        [
            "North Slope study-area boundary",
            "Assessment-unit outlines",
            "North Slope public wells",
        ],
    )
    figure.update_traces(selector={"type": "surface"}, opacity=0.66)
    figure.update_layout(
        title="Current Streamlit structural explorer view: public wells, boundaries, and horizons",
        paper_bgcolor="white",
        plot_bgcolor="white",
        margin={"l": 0, "r": 0, "t": 58, "b": 0},
        legend={"orientation": "h", "y": 1.02, "x": 0.0},
        scene={
            "xaxis_title": "Longitude",
            "yaxis_title": "Latitude",
            "zaxis_title": "Depth (m)",
            "zaxis": {"autorange": "reversed"},
            "aspectmode": "manual",
            "aspectratio": {"x": 1.8, "y": 1, "z": 0.55},
            "camera": {"eye": {"x": 1.55, "y": -1.7, "z": 0.85}},
        },
    )
    return export_plotly(figure, project_asset_path("streamlit_3d_context_v2.png"), 1400, 820)


def build_well_log_asset(logs: pd.DataFrame) -> Path:
    figure = well_log_panel(logs, "SYNTH-WELL-01")
    figure.update_layout(title="Synthetic well-log scaffold: multi-log hydrate interpretation")
    return export_plotly(figure, project_asset_path("synthetic_well_log_panel.png"), 1280, 760)


def build_ml_metrics_asset() -> Path:
    confusion, calibration = model_placeholder_figures()
    figure = make_subplots(
        rows=1,
        cols=2,
        subplot_titles=("Classification error view", "Saturation calibration view"),
        horizontal_spacing=0.14,
    )
    for trace in confusion.data:
        figure.add_trace(trace, row=1, col=1)
    for trace in calibration.data:
        figure.add_trace(trace, row=1, col=2)
    figure.update_xaxes(title_text="Predicted class", row=1, col=1)
    figure.update_yaxes(title_text="Reference class", row=1, col=1)
    figure.update_xaxes(title_text="Predicted probability", row=1, col=2)
    figure.update_yaxes(title_text="Observed frequency", row=1, col=2)
    figure.update_layout(
        title="Model validation outputs to replace placeholders after approved-data execution",
        showlegend=False,
        paper_bgcolor="white",
        plot_bgcolor="white",
        margin={"l": 58, "r": 30, "t": 70, "b": 52},
    )
    return export_plotly(figure, project_asset_path("model_validation_placeholders.png"), 1280, 640)


def build_sweet_spot_asset(logs: pd.DataFrame) -> Path:
    intervals = screen_intervals(logs)
    ranked = sweet_spot_review_table(intervals).head(8).sort_values("Synthetic review priority")
    figure = go.Figure(
        go.Bar(
            x=ranked["Synthetic review priority"],
            y=ranked["Well alias"] + " " + ranked["Top depth (m)"].astype(str) + "-" + ranked["Base depth (m)"].astype(str) + " m",
            orientation="h",
            marker={"color": ranked["Hydrate-saturation proxy"], "colorscale": "Tealgrn", "showscale": True, "colorbar": {"title": "Sh proxy"}},
            text=ranked["Evidence domains passed"],
            textposition="outside",
        )
    )
    figure.update_layout(
        title="Synthetic sweet-spot ranking scaffold",
        xaxis_title="Review priority score",
        yaxis_title="Synthetic interval",
        margin={"l": 190, "r": 80, "t": 70, "b": 50},
        paper_bgcolor="white",
        plot_bgcolor="white",
    )
    return export_plotly(figure, project_asset_path("sweet_spot_ranking.png"), 1280, 640)


def build_website_flow_asset() -> Path:
    labels = ["Overview", "Explore North Slope", "Analyze Hydrates", "Project Plan"]
    subtitles = [
        "project goal + evidence flow",
        "public maps and structure",
        "header-derived synthetic log workflow",
        "built now / activate with data",
    ]
    colors = ["#123447", "#167d8d", "#d8a24a", "#8ea7ff"]
    figure = go.Figure()
    for index, (label, subtitle, color) in enumerate(zip(labels, subtitles, colors)):
        x0 = 0.06 + index * 0.235
        x1 = x0 + 0.18
        figure.add_shape(
            type="rect",
            xref="paper",
            yref="paper",
            x0=x0,
            x1=x1,
            y0=0.34,
            y1=0.74,
            fillcolor=color,
            line={"color": color, "width": 2},
        )
        figure.add_annotation(
            x=(x0 + x1) / 2,
            y=0.58,
            xref="paper",
            yref="paper",
            text=f"<b>{label}</b><br><span style='font-size:12px'>{subtitle}</span>",
            showarrow=False,
            font={"color": "white", "size": 16},
            align="center",
        )
        if index < len(labels) - 1:
            figure.add_annotation(
                x=x1 + 0.03,
                y=0.54,
                xref="paper",
                yref="paper",
                text="->",
                showarrow=False,
                font={"color": "#167d8d", "size": 24},
            )
    figure.add_annotation(
        x=0.5,
        y=0.16,
        xref="paper",
        yref="paper",
        text="Public website = visual prototype surface. Approved well/core rows remain runtime-only.",
        showarrow=False,
        font={"color": "#5E6A71", "size": 15},
    )
    figure.update_layout(
        title="Implemented Streamlit website structure",
        width=1280,
        height=520,
        margin={"l": 30, "r": 30, "t": 75, "b": 35},
        paper_bgcolor="white",
        plot_bgcolor="white",
        xaxis={"visible": False},
        yaxis={"visible": False},
    )
    return export_plotly(figure, project_asset_path("website_four_page_flow.png"), 1280, 520)


def build_evidence_stack_asset() -> Path:
    path = project_asset_path("subsurface_evidence_stack.png")
    path.parent.mkdir(parents=True, exist_ok=True)
    rows = [
        ("Regional context", "USGS/Collett assessments; public GIS layers", "#67d0df"),
        ("Stability + reservoir", "P-T admissibility plus sand quality and continuity", "#25b99a"),
        ("Logs + core", "GR, Rt, density, Vp, Vs, NMR, core labels", "#d8a24a"),
        ("Decision", "phase, saturation, uncertainty, sweet-spot rank", "#8ea7ff"),
    ]
    image = Image.new("RGB", (1280, 620), "white")
    draw = ImageDraw.Draw(image)
    try:
        title_font = ImageFont.truetype("arial.ttf", 34)
        label_font = ImageFont.truetype("arialbd.ttf", 24)
        detail_font = ImageFont.truetype("arial.ttf", 22)
        note_font = ImageFont.truetype("arial.ttf", 21)
    except OSError:
        title_font = label_font = detail_font = note_font = ImageFont.load_default()

    draw.text((55, 36), "Source-backed subsurface evidence stack", fill=f"#{NAVY}", font=title_font)
    y = 115
    for label, detail, color in rows:
        rgb = tuple(int(color.lstrip("#")[i : i + 2], 16) for i in (0, 2, 4))
        fill = tuple(int(255 * 0.78 + channel * 0.22) for channel in rgb)
        draw.rounded_rectangle((95, y, 1185, y + 82), radius=12, fill=fill, outline=rgb, width=3)
        draw.text((130, y + 27), label, fill=f"#{NAVY}", font=label_font)
        draw.text((450, y + 29), detail, fill=f"#{INK}", font=detail_font)
        if y < 430:
            draw.line((640, y + 88, 640, y + 112), fill=f"#{TEAL}", width=4)
            draw.polygon([(640, y + 120), (628, y + 104), (652, y + 104)], fill=f"#{TEAL}")
        y += 112
    draw.text(
        (250, 565),
        "Maps constrain confidence; logs/core/labels determine interval decisions.",
        fill=f"#{MUTED}",
        font=note_font,
    )
    image.save(path)
    return path


def build_visual_assets() -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    logs = generate_synthetic_logs()
    return {
        "regional_3d": build_regional_3d_asset(),
        "well_log": build_well_log_asset(logs),
        "metrics": build_ml_metrics_asset(),
        "sweet_spot": build_sweet_spot_asset(logs),
        "website_flow": build_website_flow_asset(),
        "evidence_stack": build_evidence_stack_asset(),
    }


def shade_cell(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), fill)
    tc_pr.append(shading)


def set_borders(table, color="B7C9D3") -> None:
    tbl_pr = table._tbl.tblPr
    borders = OxmlElement("w:tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        element = OxmlElement(f"w:{edge}")
        element.set(qn("w:val"), "single")
        element.set(qn("w:sz"), "4")
        element.set(qn("w:space"), "0")
        element.set(qn("w:color"), color)
        borders.append(element)
    tbl_pr.append(borders)


def write_cell(cell, text: str, *, bold=False, color=INK, size=8.6, align=WD_ALIGN_PARAGRAPH.LEFT) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = align
    p.paragraph_format.space_after = Pt(0)
    p.paragraph_format.line_spacing = 1.05
    run = p.add_run(text)
    run.bold = bold
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string(color)


def configure_doc(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)
    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Aptos"
    normal.font.size = Pt(10.2)
    normal.paragraph_format.space_after = Pt(5)
    for style_name, size, color in [
        ("Heading 1", 15, NAVY),
        ("Heading 2", 12.5, TEAL),
        ("Heading 3", 11, NAVY),
    ]:
        style = styles[style_name]
        style.font.name = "Aptos Display"
        style.font.size = Pt(size)
        style.font.color.rgb = RGBColor.from_string(color)
        style.paragraph_format.space_before = Pt(10)
        style.paragraph_format.space_after = Pt(4)


def add_body(doc: Document, text: str, size=10.2) -> None:
    p = doc.add_paragraph(text)
    p.paragraph_format.space_after = Pt(5)
    p.paragraph_format.line_spacing = 1.1
    for run in p.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor.from_string(INK)


def add_note(doc: Document, text: str) -> None:
    table = doc.add_table(rows=1, cols=1)
    table.style = "Table Grid"
    set_borders(table, "D2C8B8")
    cell = table.cell(0, 0)
    shade_cell(cell, SAND)
    write_cell(cell, text, bold=True, size=9.2, color=NAVY)
    doc.add_paragraph()


def add_process_sketch(doc: Document, title: str, labels: list[str], fill=SAND) -> None:
    doc.add_heading(title, level=3)
    cols = len(labels) * 2 - 1
    table = doc.add_table(rows=1, cols=cols)
    table.style = "Table Grid"
    set_borders(table, "C5D3D8")
    widths = [1.0 if i % 2 == 0 else 0.25 for i in range(cols)]
    for idx, width in enumerate(widths):
        table.cell(0, idx).width = Inches(width)
    for idx in range(cols):
        cell = table.cell(0, idx)
        if idx % 2 == 0:
            shade_cell(cell, fill if idx in {0, cols - 1} else ICE)
            write_cell(cell, labels[idx // 2], bold=True, size=7.8, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER)
        else:
            shade_cell(cell, WHITE)
            write_cell(cell, "->", bold=True, size=9.5, color=TEAL, align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()


def read_parameter_matrix() -> list[dict[str, str]]:
    if not PARAMETER_CSV.exists():
        return []
    matrix = pd.read_csv(PARAMETER_CSV).fillna("")
    return [{key: str(value) for key, value in row.items()} for row in matrix.to_dict("records")]


def add_parameter_matrix_table(doc: Document, params: list[dict[str, str]]) -> None:
    if not params:
        add_note(doc, "Parameter matrix not found; regenerate docs/project_blueprints/ml_parameter_effect_tree.csv before finalizing this section.")
        return

    doc.add_heading("Current parameter and masking matrix", level=2)
    add_body(
        doc,
        "The matrix below matches the current PowerPoint revamp plan. The percentages are visual planning priors, not trained feature importance, and will be replaced or recalibrated only after approved data, labels, and held-out-well results are available.",
    )
    table = doc.add_table(rows=1, cols=5)
    table.style = "Table Grid"
    set_borders(table, "C5D3D8")
    headers = ["Parameter family", "Weight", "Hydrate-supportive evidence", "Major masks", "Model role"]
    widths = [1.55, 0.55, 2.15, 2.25, 1.65]
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.width = Inches(widths[idx])
        shade_cell(cell, NAVY)
        write_cell(cell, header, bold=True, size=7.3, color=WHITE, align=WD_ALIGN_PARAGRAPH.CENTER)

    for row in params:
        cells = table.add_row().cells
        values = [
            row.get("parameter", ""),
            f"{row.get('planned_importance_percent', '').strip()}%",
            row.get("primary_hydrate_effect", ""),
            row.get("major_masks_or_false_positives", "").replace("; ", "; "),
            row.get("model_role", ""),
        ]
        for idx, value in enumerate(values):
            cells[idx].width = Inches(widths[idx])
            shade_cell(cells[idx], ICE if idx in {0, 1} else WHITE)
            write_cell(
                cells[idx],
                value,
                bold=idx == 0,
                size=6.6 if idx in {2, 3, 4} else 6.9,
                color=NAVY if idx in {0, 1} else INK,
                align=WD_ALIGN_PARAGRAPH.CENTER if idx == 1 else WD_ALIGN_PARAGRAPH.LEFT,
            )
    doc.add_paragraph()


def add_placeholder_section(doc: Document, title: str, sentence: str) -> None:
    doc.add_heading(title, level=1)
    add_body(doc, sentence)


def add_compact_table(
    doc: Document,
    headers: list[str],
    rows: list[list[str]],
    *,
    widths: list[float] | None = None,
    header_fill: str = NAVY,
    body_size: float = 7.8,
) -> None:
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    set_borders(table, "C5D3D8")
    if widths is None:
        widths = [1.0 for _ in headers]
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.width = Inches(widths[idx])
        shade_cell(cell, header_fill)
        write_cell(cell, header, bold=True, size=7.5, color=WHITE, align=WD_ALIGN_PARAGRAPH.CENTER)
    for row_index, values in enumerate(rows):
        cells = table.add_row().cells
        for idx, value in enumerate(values):
            cells[idx].width = Inches(widths[idx])
            shade_cell(cells[idx], ICE if row_index % 2 == 0 and idx == 0 else WHITE)
            write_cell(
                cells[idx],
                value,
                bold=idx == 0,
                size=body_size,
                color=NAVY if idx == 0 else INK,
                align=WD_ALIGN_PARAGRAPH.LEFT,
            )
    doc.add_paragraph()


def build_legacy_docx() -> None:
    doc = Document()
    configure_doc(doc)
    parameter_matrix = read_parameter_matrix()

    title = doc.add_paragraph()
    title.paragraph_format.space_after = Pt(2)
    run = title.add_run("Gas Hydrate Occurrence and Saturation Prediction in Permafrost Sediments on the Alaska North Slope Using AI/ML")
    run.bold = True
    run.font.name = "Aptos Display"
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor.from_string(NAVY)

    subtitle = doc.add_paragraph("Research-paper outline updated with ML source specifics, parameter logic, and pipeline controls")
    subtitle.paragraph_format.space_after = Pt(8)
    for run in subtitle.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(10.5)
        run.font.color.rgb = RGBColor.from_string(MUTED)

    add_note(
        doc,
        "Boundary: this public-source planning document contains no restricted well logs, approved-runtime data, named restricted well identifiers, populated model results, credentials, or sensitive derived outputs.",
    )
    add_note(
        doc,
        "Source update: this revision incorporates the June 11 Gmail ML sources now stored under references/ml-sources/2026-06-11/. Chong et al. (2022) supplies the gas-hydrate ANN/well-log anchor; the ML Project Reference Notes supply general leakage-safe pipeline, data-quality, validation, and monitoring controls.",
    )
    add_note(doc, HEADER_DERIVED_SYNTHETIC_NOTE)

    doc.add_heading("Abstract", level=1)
    add_body(
        doc,
        "Gas hydrates are solid crystalline compounds in which gas molecules, dominantly methane, are trapped inside cages of water molecules under cold and high-pressure conditions. On the Alaska North Slope, permafrost and shallow subsurface pressure-temperature conditions create a setting where hydrate can accumulate in sand-rich reservoirs and represent a significant natural-gas resource. The United States Geological Survey assessed the North Slope gas hydrate system as containing a mean of approximately 53.8 trillion cubic feet of undiscovered, technically recoverable gas in hydrate accumulations, making the region an important target for energy-resource characterization and long-term energy security planning. This project develops a physics-constrained machine-learning workflow for predicting gas hydrate occurrence and saturation from approved well-log and core-analysis data. The workflow adapts the Chong et al. ANN saturation approach, builds a depth-indexed feature table from density, porosity, natural gamma ray, resistivity, Vp, Vs, NMR, caliper, and core measurements, keeps NMR/core saturation fields as labels rather than predictors, compares baseline and nonlinear models, and validates by complete wells. The expected outcome is an auditable reservoir characterization framework that separates hydrate occurrence, saturation, reservoir quality, uncertainty, data-quality state, and later producibility screening.",
    )

    doc.add_heading("Introduction", level=1)
    add_body(
        doc,
        "Gas hydrates are ice-like solids that store methane in the pore systems of sediment. They are not conventional gas reservoirs, because the gas is held in a solid water-gas structure rather than existing as a free gas phase. Their formation requires pressure-temperature stability, but stability alone is not enough to produce an accumulation. A reservoir interval must also have pore space, a supply of gas, migration pathways, water, and geologic architecture that allows hydrate to occupy and remain within the sediment. This distinction is central to the Alaska North Slope project. The purpose is not only to ask where hydrate could be stable, but to identify which intervals are most likely to contain hydrate, how much pore volume is occupied, and how reliable that interpretation is."
    )
    add_body(
        doc,
        "The Alaska North Slope is a strong study area because it combines permafrost-associated hydrate stability, documented hydrate-bearing test wells, public regional geologic context, and a history of DOE, NETL, USGS, and industry-supported hydrate research. Mount Elbert, Hydrate-01, and Eileen-trend studies show that hydrate-bearing sands can be investigated using a multi-log interpretation that includes resistivity, density, neutron and NMR porosity, sonic velocity, and core calibration. This project uses that research base to build a practical workflow for approved data. The final scientific question is whether well-log and core measurements can be processed into reliable machine-learning features that predict occurrence and saturation across known and prediction wells."
    )
    add_body(
        doc,
        "The closest published machine-learning anchor is Chong et al. (2022), which used permafrost-associated gas hydrate wells from the Alaska North Slope and Mallik, more than 10,000 depth points, and Keras/TensorFlow artificial neural networks to predict NMR-density-derived gas hydrate saturation from combinations of density, density porosity, gamma ray, resistivity, Vp, and Vs. Their workflow included preprocessing, hyperparameter tuning, well-log combination optimization, model validation, caliper-based washout removal, multivariate outlier review, min-max normalization, and R2-style saturation accuracy scoring. This project uses that paper as a methodological analogue, but it upgrades the validation design for the current DOE-style workflow by holding out complete wells or compartments rather than relying on neighboring random depth rows.",
    )
    add_body(
        doc,
        "Characterizing a gas hydrate reservoir matters because a regional resource estimate does not identify the best intervals for detailed evaluation. Characterization translates resource potential into interval-scale evidence. It asks whether the rock is a reservoir, whether the log response is consistent with hydrate rather than gas, ice, shale, carbonate, coal, or borehole effects, whether NMR or core evidence supports a saturation estimate, and whether the uncertainty is low enough for the result to be used in a technical decision. High resistivity alone is not a hydrate label, and high saturation alone is not a producibility ranking. Occurrence, saturation, reservoir quality, uncertainty, and producibility must remain separate outputs."
    )
    add_body(
        doc,
        "The current public scaffold is based on three Excel header/schema references, not real sample rows. Those headers indicate the expected input families: depth, location, gamma ray, density, density porosity, neutron porosity, NMR porosity, resistivity, caliper, Vp, Vs, Vp/Vs, acoustic impedance, interpreted hydrate saturation fields such as Sgh or S_h, NMR-derived saturation where supplied, core porosity, permeability, lithology, and quality flags. Synthetic rows may be generated from those headers to test layout, validation, and visual behavior, but the generated values are not scientific measurements. These variables are scientifically useful only when their roles are controlled. Measured curves, derived features, quality-control fields, alignment fields, and supervised targets must be separated so that the model does not accidentally learn from the answer column it is supposed to predict."
    )
    add_process_sketch(
        doc,
        "Conceptual hydrate-reservoir sketch",
        ["stability", "reservoir sand", "gas charge", "multi-log evidence", "saturation", "sweet spot"],
    )
    add_body(
        doc,
        "The machine-learning component will follow a staged workflow rather than immediately relying on a black-box model. Raw well-log data must first be compared, cleaned, depth-aligned, and screened for outliers and borehole-quality problems. Relevant depth intervals will then be selected, processed into a consistent dataset, and transformed into petrophysical and geomechanical features. The project will compare historical AI/ML approaches for gas hydrate occurrence and saturation prediction, then use baseline classification and regression models before testing nonlinear models such as artificial neural networks. Keras-based neural-network workflows are appropriate only after the training labels, curve coverage, and validation design are defensible."
    )
    add_process_sketch(
        doc,
        "Approved-data processing sketch",
        ["raw logs", "QC", "features", "labels", "ML models", "validated outputs"],
        fill=ICE,
    )
    add_process_sketch(
        doc,
        "Website visual workflow now implemented",
        ["Overview", "Explore", "Analyze", "Project Plan"],
        fill=ICE,
    )
    add_body(
        doc,
        "The Streamlit website has been reduced to four visual-first pages. It is a prototype and figure-generation surface for the Word document and PowerPoint. It should continue to show only public GIS layers and header-derived synthetic examples until approved data are loaded inside the authorized runtime environment."
    )
    add_body(
        doc,
        "A key methodological control is validation by complete well. Randomly splitting neighboring depth samples can make performance appear stronger than it will be on an unseen well because adjacent samples are highly correlated. The planned workflow will reserve complete wells for validation and testing, report classification and regression error, and compare predicted saturation against core and interpreted reference measurements. The final deliverables will present the work in graphical and tabular form so that conclusions, recommendations, and uncertainty can be reviewed by geoscientists and energy-resource decision makers."
    )

    doc.add_heading("Parameters", level=1)
    add_body(
        doc,
        "The approved-data workflow will treat each well-log family as physical evidence with limits, not as an isolated hydrate label. Density, porosity, gamma ray, resistivity, Vp, Vs, NMR, caliper, core porosity, permeability, lithology, and pressure-temperature context all support different parts of the interpretation. The key rule is parameter -> hydrate effect -> masking condition -> QC or context decision.",
    )
    add_body(
        doc,
        "Depth remains an alignment, stratigraphic, and pressure-temperature context variable. It should not be normalized in the same way as the privacy-protected curve values. Non-depth log and core variables may be standardized or normalized for modeling, but those transformations must be fitted on training wells only and then applied unchanged to validation, locked-test, and prediction wells.",
    )
    add_body(
        doc,
        "Target columns such as Sgh, S_h, NMR_SAT, phase labels, final hydrate calls, and sweet-spot rankings are supervised labels or outputs. They must be kept out of the predictor matrix except when they are explicitly used as calibration, training target, validation reference, or post-prediction review evidence.",
    )
    add_body(
        doc,
        "For each approved-data curve, the feature table should retain the original mnemonic, canonical field, source unit, standardized unit, schema role, missingness flag, QC state, and provenance. That level of metadata is necessary because the same physical parameter can serve different roles depending on source: measured NMR porosity can be an input, while NMR-derived saturation is a label; density can be a measured input, while density porosity or acoustic impedance may be derived features.",
    )
    add_parameter_matrix_table(doc, parameter_matrix)
    add_process_sketch(
        doc,
        "Parameter-role sketch",
        ["measured logs", "derived physics", "QC flags", "targets", "outputs"],
        fill=SAND,
    )

    doc.add_heading("Methodology", level=1)
    add_body(
        doc,
        "The methodology begins with source classification and data intake. Public GIS and source documents remain in the repository, while approved LAS, CSV, Excel, and core-analysis files remain in the authorized runtime. The runtime workflow standardizes curve aliases, verifies units, checks monotonic depth, flags missing or out-of-range fields, and records whether each output is ready, partial, or blocked.",
    )
    add_body(
        doc,
        "Quality control happens before feature generation. Caliper and differential-caliper fields are used to identify washout, rugosity, and tool-standoff risk; missing-curve routing prevents an absent measurement from becoming a silent zero; and depth matching controls how log curves, NMR-derived values, and core intervals are aligned. Each derived feature must retain provenance from the curves used to calculate it.",
    )
    add_body(
        doc,
        "The approved-data QC sequence should be explicit and reproducible: verify schema and units; standardize depth; drop or quarantine rows where a required feature or label is missing for the intended task; flag the upper-tail caliper washout intervals before model fitting; run multivariate outlier screening on the feature space; and record row loss after every step. This mirrors the source-paper logic while leaving real logs and row-level outputs inside the authorized runtime.",
    )
    add_body(
        doc,
        "The geomechanical feature set will include dynamic Young's modulus, Poisson's ratio, brittleness, lambda-rho, and mu-rho where density, Vp, and Vs are available. These features help place hydrate interpretation in a rock-physics context because stiffness, rigidity, and incompressibility can help distinguish hydrate-supported responses from gas, lithology, or stress effects."
    )
    add_body(
        doc,
        "The future overburden map is a context layer, not a hydrate classifier. Its role is to explain how burial pressure, stratigraphic load, structural position, and compaction can shift density, porosity, Vp, Vs, acoustic impedance, and geomechanical baselines across the Alaska North Slope. That context helps prevent one universal threshold from being applied across wells with different burial histories.",
    )

    doc.add_heading("Machine-Learning Framework", level=1)
    add_body(
        doc,
        "The model architecture uses shared physics-backed inputs and then separates into two parallel heads: a classification branch for hydrate phase or occurrence and a saturation-regression branch for continuous hydrate saturation. The branches share QC-reviewed logs, derived petrophysical features, rock-physics attributes, and approved core or interpretation context, but they optimize different outputs and should report different uncertainty measures.",
    )
    add_body(
        doc,
        "The planned model ladder is baseline-first. Rule-based screens and simple linear or logistic baselines establish defensible reference behavior; tree-based models then test nonlinear interactions; artificial neural networks are appropriate after label quality, curve coverage, missingness, and complete-well validation are understood. Chong et al. (2022) is the closest source anchor for ANN-style hydrate-saturation modeling, but this project will use held-out wells rather than random neighboring depth samples for field-generalization checks.",
    )
    add_body(
        doc,
        "The general ML notes add three implementation controls to that ladder. First, the validation scheme must mimic the production use case: if the model will predict unseen wells, the evaluation must withhold unseen wells. Second, train-only preprocessing must be leakage-safe: imputation, scaling, min-max normalization, feature selection, and dimensionality reduction are fitted on training wells only and then applied unchanged to validation, locked-test, and prediction wells. Third, flexible models such as gradient boosting or ANN/Keras should be adopted only if they materially improve held-out performance over simpler baselines without increasing false positives, drift sensitivity, or interpretability risk.",
    )
    add_body(
        doc,
        "A visible target-leakage barrier is required in every modeling implementation. Hydrate saturation labels, NMR-derived target fields, phase labels, and final sweet-spot rankings may supervise or calibrate the model, but they cannot be included as predictors. The exported feature matrix should show which columns are measured inputs, derived inputs, QC/context variables, calibration fields, labels, and blocked outputs.",
    )
    add_process_sketch(
        doc,
        "ML architecture sketch",
        ["features", "classification", "saturation regression", "ANN test", "held-out wells"],
        fill=ICE,
    )

    doc.add_heading("Error and Validation", level=1)
    add_body(
        doc,
        "Validation must be performed by complete well. Random row splitting is not sufficient because adjacent depth samples share geology, tool response, and labeling assumptions. The known-well cohort should be divided into training, validation, and locked-test wells, and any prediction-well outcomes should remain hidden until predictions are frozen.",
    )
    add_body(
        doc,
        "Classification reporting should include per-class precision, recall, F1, confusion matrices, abstention or expert-review counts, and false-positive review by masking condition. Saturation reporting should include MAE, RMSE, R2, calibration by saturation band, residuals by well, and sensitivity to missing curves. Both branches should report where the model is blocked, partial, or uncertain rather than forcing an unsupported hydrate call.",
    )
    add_body(
        doc,
        "Pipeline validation should also report data quality, not just model metrics. Required checks include completeness, uniqueness, validity, cross-field consistency, row-count reconciliation, source-to-feature row loss, out-of-range values, duplicate depth keys, unit mismatches, core-log depth offsets, and feature drift by well, unit, or compartment. Distribution-shift triage should separate sampling noise, pipeline changes, and real geologic or acquisition differences before the model is blamed or trusted.",
    )
    doc.add_heading("Discussion", level=1)
    add_body(
        doc,
        "The discussion will interpret predicted hydrate intervals only after approved-data execution. Until then, the public deliverables should emphasize the review logic: high resistivity must be checked against lithology, gas, ice, carbonate, porosity, and water-salinity assumptions; density and velocity features must be read with overburden and mineralogy context; and NMR or core evidence must be checked for depth alignment and target provenance.",
    )
    add_body(
        doc,
        "Sweet-spot ranking should combine occurrence probability, predicted saturation, reservoir quality, uncertainty, and masking explanations. It should not collapse hydrate presence, saturation, producibility, and confidence into a single hidden score. The final ranking should remain traceable back to measured curves, derived physics, core or interpretation support, and the validation results for unseen wells.",
    )
    add_body(
        doc,
        "The final discussion should therefore read the model as an evidence system: which parameters supported the prediction, which masking conditions were ruled out or left open, which QC checks passed, which wells or compartments drove the error, and whether the predicted interval should be accepted, reviewed, or held back for more data.",
    )
    doc.add_heading("Current Evidence Anchors", level=1)
    for anchor in SOURCE_ANCHORS:
        add_body(doc, f"{anchor['claim']}: {anchor['source']} - {anchor['use']}")
    add_placeholder_section(
        doc,
        "Conclusion",
        "This section will summarize what the workflow demonstrates, what the final approved-data results show, and how the project supports North Slope gas hydrate reservoir characterization for future energy-resource evaluation.",
    )
    doc.add_heading("Current Website and Source Integration", level=1)
    add_body(
        doc,
        "The Streamlit website now uses a four-page visual workflow: Overview, Explore North Slope, Analyze Hydrates, and Project Plan. The visuals are public/synthetic communication surfaces, while the real-data execution path remains the approved runtime.",
    )
    for item in [
        "ML inputs and saturation target design are anchored to Chong et al. (2022): density, porosity, resistivity, gamma ray, Vp, Vs, and NMR/core-calibrated saturation references.",
        "NMR, sonic, density, and resistivity comparison logic is cross-checked against Lee and Collett (2011) and Haines et al. (2022).",
        "Reservoir quality and hydrate presence remain separate interpretations; clean sand, stability, and gas charge can still produce no-hydrate cases.",
        "Public maps constrain structural and regional context, but they do not classify hydrate occurrence without direct log/core evidence.",
        "Generated rows in the public website and tests are header-derived synthetic records, not user-supplied sample data.",
    ]:
        add_body(doc, f"- {item}", size=9.8)

    doc.add_heading("References", level=1)
    for ref in REFERENCES:
        add_body(doc, ref, size=9)

    doc.save(DOCX_OUT)


def ppt_color(hex_color: str) -> PptRGB:
    return PptRGB(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = PptInches(0.05)
    frame.margin_right = PptInches(0.05)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = ppt_color(color)
    return box


def add_link_textbox(slide, x, y, w, h, text, url, size=11.5):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = PptInches(0.06)
    frame.margin_right = PptInches(0.06)
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = PptPt(size)
    run.font.bold = True
    run.font.color.rgb = ppt_color(TEAL)
    run.hyperlink.address = url
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.24, 12.1, 0.45, title, size=24, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, 0.58, 0.76, 11.8, 0.35, subtitle, size=11.5, color=MUTED)


def add_panel(slide, x, y, w, h, fill=ICE, line=TEAL, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_color(fill)
    shape.line.color.rgb = ppt_color(line)
    shape.line.width = PptPt(1.1)
    return shape


def add_label(slide, x, y, w, h, label, body="", fill=ICE):
    add_panel(slide, x, y, w, h, fill=fill)
    add_textbox(slide, x + 0.12, y + 0.1, w - 0.24, 0.27, label, size=12.2, bold=True, color=NAVY)
    if body:
        add_textbox(slide, x + 0.12, y + 0.42, w - 0.24, h - 0.48, body, size=8.9, color=INK)


def add_footer(slide, text="Sources: USGS 2019; DOE/NETL; Chong et al. 2022; Lee and Collett 2011; Haines et al. 2022"):
    add_textbox(slide, 0.55, 7.18, 12.1, 0.22, text, size=7.3, color=MUTED)


def add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, PptInches(x1), PptInches(y1), PptInches(x2), PptInches(y2))
    line.line.color.rgb = ppt_color(TEAL)
    line.line.width = PptPt(2)
    return line


def add_bullet_list(slide, x, y, w, h, items, size=11.5):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = PptInches(0.08)
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = PptPt(size)
        p.font.color.rgb = ppt_color(INK)
    return box


def add_flow(slide, labels: list[str], x=0.8, y=2.0, box_w=1.45, box_h=0.68, gap=0.38):
    for idx, label in enumerate(labels):
        bx = x + idx * (box_w + gap)
        add_label(slide, bx, y, box_w, box_h, label, "", SAND if idx in {0, len(labels) - 1} else ICE)
        if idx < len(labels) - 1:
            add_arrow(slide, bx + box_w + 0.02, y + box_h / 2, bx + box_w + gap - 0.08, y + box_h / 2)


def add_image(slide, path: Path, x, y, w, h):
    return slide.shapes.add_picture(str(path), PptInches(x), PptInches(y), PptInches(w), PptInches(h))


def build_pptx() -> None:
    assets = build_visual_assets()
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Gas Hydrate Occurrence and Saturation Prediction", "Permafrost sediments on the Alaska North Slope using AI/ML")
    add_image(slide, assets["regional_3d"], 7.05, 1.08, 5.65, 3.35)
    photo = slide.shapes.add_shape(MSO_SHAPE.OVAL, PptInches(0.78), PptInches(1.42), PptInches(1.85), PptInches(1.85))
    photo.fill.solid()
    photo.fill.fore_color.rgb = ppt_color(GRAY)
    photo.line.color.rgb = ppt_color(TEAL)
    add_textbox(slide, 1.02, 2.1, 1.35, 0.3, "photo", size=16, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_label(slide, 2.95, 1.4, 3.55, 1.2, "Rohan Nanda", "Lead Geologist\nM.S. Geosciences, University of Texas at Dallas", SAND)
    add_label(slide, 0.8, 3.75, 5.7, 1.35, "Project focus", "Use approved well-log, NMR, and core-analysis data to predict gas hydrate occurrence and saturation, then communicate results as reservoir-characterization evidence.")
    add_label(slide, 7.25, 4.72, 5.2, 0.9, "Website asset embedded", "Public North Slope structural scene and well inventory generated from the project atlas.", ICE)
    add_footer(slide, "Deck structure follows project email instructions received 2026-06-09.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Introduction: What and Why Gas Hydrates", "A potential natural-gas resource that requires reservoir-scale characterization")
    add_label(slide, 0.8, 1.15, 3.45, 1.05, "Gas hydrate", "Methane is held inside crystalline water cages in sediment pore space.", SAND)
    add_label(slide, 0.8, 2.45, 3.45, 1.05, "North Slope value", "USGS assessed a mean of about 53.8 TCF of technically recoverable gas in North Slope hydrate accumulations.")
    add_label(slide, 0.8, 3.75, 3.45, 1.05, "Characterization goal", "Translate resource potential into interval-scale occurrence, saturation, uncertainty, and sweet-spot evidence.")
    add_image(slide, assets["evidence_stack"], 4.75, 1.15, 7.25, 3.35)
    add_label(slide, 4.95, 4.72, 6.85, 0.9, "Core message", "Stability and maps constrain confidence. Logs, NMR/core labels, and uncertainty determine interval decisions.", SAND)
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Parameters: Header-Derived Well-Log Scaffold", "Only Excel header/schema references are available; generated rows are synthetic")
    add_image(slide, assets["well_log"], 5.25, 1.08, 7.45, 4.65)
    for i, (label, body) in enumerate([
        ("Lithology", "GR, core lithology, clean-sand screening"),
        ("Porosity", "Density, neutron, NMR, core porosity"),
        ("Electrical", "Deep resistivity and source mnemonics"),
        ("Elastic", "Vp, Vs, Vp/Vs, impedance, moduli"),
        ("Targets", "Sgh, S_h, NMR_SAT, phase classes"),
        ("QC", "Caliper, outliers, depth match, missing curves"),
    ]):
        add_label(slide, 0.65, 1.12 + i * 0.78, 4.25, 0.62, label, body, SAND if label == "Targets" else ICE)
    add_label(slide, 5.45, 5.88, 6.9, 0.65, "Data provenance", "Header-derived synthetic records only. Values test layout and QA; they are not user-supplied sample data.")
    add_footer(slide, "Sources: three Excel header references; project Q&A update; Lee and Collett 2011; Haines et al. 2022")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "ML Methodology: Architecture", "ANN paper logic adapted into a leakage-controlled well-log workflow")
    add_label(slide, 0.75, 1.18, 2.35, 0.9, "Chong et al. anchor", "5 permafrost wells\n>10,000 depth points\nNMR-derived Sh target", SAND)
    add_label(slide, 3.45, 1.18, 2.35, 0.9, "Input log families", "density, porosity, Rt, GR,\nVp and Vs candidates", ICE)
    add_label(slide, 6.15, 1.18, 2.35, 0.9, "Project target control", "NMR/core saturation is a label,\nnot an input when target-derived", ICE)
    add_label(slide, 8.85, 1.18, 2.35, 0.9, "Validation upgrade", "complete-well holdout\ninstead of random depth rows", SAND)
    add_flow(slide, ["approved logs", "QC", "feature sets", "model ladder", "held-out wells", "outputs"], x=0.75, y=2.48, box_w=1.55)
    add_label(slide, 0.75, 3.75, 2.95, 1.15, "1. QC and alignment", "Caliper washout, depth matching, missing-curve routing, unit standardization, multivariate outlier review.", SAND)
    add_label(slide, 3.95, 3.75, 2.95, 1.15, "2. Physics features", "Separate measured logs, derived elastic/geomechanical features, and target-derived saturation fields.")
    add_label(slide, 7.15, 3.75, 2.95, 1.15, "3. Model ladder", "Rules and linear baselines first; tree/GBM for feature behavior; ANN/Keras for nonlinear saturation response.")
    add_label(slide, 10.35, 3.75, 2.2, 1.15, "4. Outputs", "Occurrence class\nSh regression\nuncertainty flags", ICE)
    add_image(slide, assets["metrics"], 2.25, 5.33, 8.8, 1.18)
    add_footer(slide, "Source anchor: Chong et al. 2022, DOI: 10.1007/s10596-022-10151-9; project runtime requirements map.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Website Integration: Four-Page Visual Workflow", "Streamlit is a public-source prototype and figure-generation surface")
    add_image(slide, assets["website_flow"], 0.95, 1.12, 11.25, 3.25)
    add_label(slide, 0.95, 4.82, 3.45, 0.95, "Overview", "Project goal, data-to-decision pipeline, and evidence stack.", SAND)
    add_label(slide, 4.72, 4.82, 3.45, 0.95, "Analyze Hydrates", "Header-derived synthetic log board, target guardrail, runtime readiness.", ICE)
    add_label(slide, 8.48, 4.82, 3.45, 0.95, "Project Plan", "Built now versus approved-data activation, blockers, deliverables.", ICE)
    add_footer(slide, "Public website remains synthetic/public-source only; legacy eight-page links route into the four current pages.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Geomechanical Feature Sketch", "Rock-physics parameters help evaluate hydrate-consistent stiffness")
    add_panel(slide, 0.8, 1.1, 11.75, 4.55, fill=WHITE, line=TEAL, radius=False)
    add_textbox(slide, 1.1, 1.35, 4.15, 0.35, "Inputs from the well-log scaffold", 16, True, NAVY)
    add_textbox(slide, 1.1, 1.88, 3.9, 1.0, "rho, Vp, Vs, GR, porosity, NMR", 20, True, TEAL, PP_ALIGN.CENTER)
    add_textbox(slide, 5.1, 1.35, 3.8, 0.35, "Derived rock-physics terms", 16, True, NAVY)
    add_textbox(slide, 5.1, 1.9, 3.8, 1.95, "YM = f(rho, Vp, Vs)\nPR = f(Vp, Vs)\nMR = (Vs * rho)^2\nLR = (Vp * rho)^2 - 2(Vs * rho)^2", 17, True, INK, PP_ALIGN.CENTER)
    add_textbox(slide, 9.15, 1.35, 2.95, 0.35, "Interpretation use", 16, True, NAVY)
    add_textbox(slide, 9.15, 1.92, 2.9, 1.35, "Compare hydrate, gas, water, shale, ice, and stress effects.", 15, True, TEAL, PP_ALIGN.CENTER)
    add_arrow(slide, 4.65, 2.55, 5.0, 2.55)
    add_arrow(slide, 8.95, 2.55, 9.1, 2.55)
    add_label(slide, 1.1, 4.35, 10.85, 0.82, "Purpose", "Integrate mechanical behavior with log response so high resistivity or high velocity is not treated as a hydrate label by itself.", SAND)
    add_footer(slide, "Source anchor: project geomechanics screenshots and local equation-map documents.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "3D Map and Well Context", "Current Streamlit structural explorer view with a live-map link")
    add_image(slide, assets["regional_3d"], 0.75, 1.05, 7.35, 4.7)
    add_label(slide, 8.45, 1.35, 3.75, 1.2, "Map use", "Show well positions, structural context, data coverage, and candidate intervals.", SAND)
    add_label(slide, 8.45, 3.0, 3.75, 1.05, "Live interaction", "Google Slides stores a static image; use the linked Streamlit app for rotation, zoom, and layer toggles.")
    add_link_textbox(slide, 8.55, 4.2, 3.55, 0.3, "Open live Streamlit 3D explorer", STREAMLIT_URL)
    add_label(slide, 8.45, 4.82, 3.75, 0.82, "Data boundary", "Public map uses regional/public layers; approved well results stay runtime-only.", ICE)
    add_footer(slide, "Source anchor: project GIS atlas and future approved-runtime outputs.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Results and Discussion Plan", "Website placeholders now show what approved-data outputs will replace")
    add_image(slide, assets["metrics"], 0.78, 1.05, 5.85, 2.8)
    add_image(slide, assets["sweet_spot"], 6.9, 1.05, 5.45, 2.8)
    add_label(slide, 0.9, 4.35, 3.5, 1.0, "Expected figures", "Well-log panel, saturation plot, confusion matrix, calibration plot.", SAND)
    add_label(slide, 4.85, 4.35, 3.5, 1.0, "Discussion", "Explain model behavior, false positives, uncertain intervals, and geologic controls.")
    add_label(slide, 8.8, 4.35, 3.5, 1.0, "Outputs", "Occurrence class, saturation estimate, uncertainty, sweet-spot rank, review flags.")
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Conclusion", "Final project message")
    add_label(slide, 0.95, 1.35, 3.2, 1.45, "Scientific value", "Hydrate occurrence, saturation, reservoir quality, and producibility are separated.")
    add_label(slide, 4.6, 1.35, 3.2, 1.45, "ML value", "Models are trained on physics-backed features and tested by complete well.")
    add_label(slide, 8.25, 1.35, 3.2, 1.45, "Energy value", "Characterization supports future natural-gas resource evaluation and energy security.")
    add_textbox(slide, 1.2, 4.25, 10.7, 0.7, "The goal is a defensible, explainable workflow for predicting gas hydrate occurrence and saturation from approved North Slope well-log and core data.", 20, True, NAVY, PP_ALIGN.CENTER)
    add_footer(slide)

    prs.save(PPTX_OUT)


def build_docx() -> None:
    doc = Document()
    configure_doc(doc)

    title = doc.add_paragraph()
    title.paragraph_format.space_after = Pt(2)
    run = title.add_run("North Slope Gas Hydrate Occurrence and Saturation ML Workflow")
    run.bold = True
    run.font.name = "Aptos Display"
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor.from_string(NAVY)

    subtitle = doc.add_paragraph("Science-to-ML methods draft rebuilt from the June 13 source ledger")
    subtitle.paragraph_format.space_after = Pt(8)
    for run in subtitle.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(10.5)
        run.font.color.rgb = RGBColor.from_string(MUTED)

    add_note(
        doc,
        "Public-safe boundary: this document describes a source-backed workflow. It does not include approved runtime well-log rows, core rows, restricted well identifiers, trained models, populated runtime configs, or project model results.",
    )
    add_note(
        doc,
        "Current status: the numeric ranges are first-pass screening envelopes for QC, feature engineering, and crossplots. They are not final DOE thresholds and must be calibrated against approved well, NMR, and core targets before results claims.",
    )

    doc.add_heading("Abstract", level=1)
    add_body(
        doc,
        "This project designs a physics-constrained machine-learning workflow for predicting gas hydrate occurrence and saturation in sand-rich, permafrost-associated reservoirs on the Alaska North Slope. The workflow is built for pore-filling methane hydrate rather than every possible hydrate habit, because pore-filling hydrate in reservoir-quality sand has the clearest connection to well-log response, NMR/core calibration, and saturation modeling. The method follows a science-to-ML ladder: first test whether hydrate can exist, then whether the rock can host hydrate, and finally whether the measured logs show an electrical, elastic, and porosity response consistent with hydrate-bearing sediment.",
    )
    add_body(
        doc,
        "The planned runtime preserves raw headers and units, converts to canonical fields, aligns depth, runs QC and missingness checks, computes measured and derived physics features, blocks target leakage, and trains separate occurrence-classification and saturation-regression outputs. Published gas-hydrate ML studies support this architecture, especially Chong et al. (2022) for permafrost-associated ANS and Mallik saturation prediction, Singh et al. (2021) for selected well-log saturation features, and Chong et al. (2024) for separate occurrence and saturation heads. Direct ANS sources such as Aung et al. (2026), Yoneda et al. (2026), Lee and Collett (2011), and Haines et al. (2022) anchor the logging, NMR, core, QC, and reservoir interpretation logic.",
    )

    doc.add_heading("1. Project Claim and Source Boundary", level=1)
    add_body(
        doc,
        "The current claim is a workflow claim, not a performance claim. The project is not reporting trained North Slope model accuracy yet. It is claiming that a defensible DOE-style pipeline should teach the model the same sequence a geoscientist would use: pressure-temperature admissibility, sand-reservoir quality, multi-log hydrate response, target-safe model training, and complete-well validation.",
    )
    add_body(
        doc,
        "The source hierarchy prevents overclaiming. Direct ANS logging and core papers support North Slope field logic. Chong et al. (2022) supports the direct permafrost hydrate ML analogy. Comparative ML papers support model choices, but they do not become North Slope field truth. Project synthesis documents such as the science-to-ML ladder and baseline source ledger organize the workflow; they are not independent scientific proof.",
    )
    add_compact_table(
        doc,
        ["Evidence tier", "Main sources", "Use in this project", "Guardrail"],
        [
            ["Direct ANS logging/core", "Aung 2026; Yoneda 2026; Lee and Collett 2011; Haines 2022", "Log suite, QC, NMR/core calibration, reservoir vs seal, producibility context", "Do not cite source quick-look values as this project's results"],
            ["Direct permafrost ML", "Chong et al. 2022", "Feature families, NMR-derived saturation target concept, ANN saturation workflow", "Use as analogue until approved project validation exists"],
            ["Comparative ML", "Singh 2021; Chong 2024; Tian 2023; Li and Liu 2020; Naim 2023", "Model ladder, occurrence plus saturation, optional sequence and missing-log paths", "Method support only, not ANS calibration"],
            ["Project synthesis", "Logic ladder; baseline source ledger; requirements map", "Parameter grammar, leakage barrier, feature roles, public/runtime boundary", "Planning support, not source-paper evidence"],
        ],
        widths=[1.35, 1.9, 2.7, 2.0],
        body_size=6.9,
    )

    doc.add_heading("2. Hydrate System First", level=1)
    add_body(
        doc,
        "The DOE workflow is mainly built for pore-filling methane hydrate in clean to sand-rich permafrost reservoirs. That focus matters because pore-filling hydrate changes the pore-fluid system and the sediment frame in ways that logs can detect: resistivity tends to rise as conductive pore water is displaced, sonic stiffness can increase as the frame is stiffened, and NMR/core evidence can independently support saturation or reservoir quality. The model should therefore learn hydrate behavior in a defined physical system instead of scanning a flat list of variables.",
    )
    add_compact_table(
        doc,
        ["Hydrate habit", "Meaning for logs and ML"],
        [
            ["Pore-filling sand hydrate", "Best fit for this workflow. Logs can combine resistivity, porosity, Vp, Vs, impedance, NMR, and core context when the interval is clean enough."],
            ["Fracture or vein hydrate", "Harder for ML because thin or anisotropic hydrate may be averaged by log-tool response and may not follow clean-sand reservoir logic."],
            ["Massive or nodular hydrate", "Can create strong local responses, but it is less cleanly tied to the sand-reservoir saturation workflow and should be treated as a separate habit if labels support it."],
        ],
        widths=[2.0, 5.6],
        body_size=7.4,
    )
    add_process_sketch(
        doc,
        "Science-to-ML logic ladder",
        ["can hydrate exist?", "can rock host it?", "hydrate-like logs?", "ML prediction", "review output"],
        fill=ICE,
    )

    doc.add_heading("3. Parameters as Physical Evidence", level=1)
    add_body(
        doc,
        "Parameters are organized into three tiers: stability context, reservoir quality, and hydrate response. This is the core change from the earlier flat parameter list. Each parameter must be explained through the same grammar: physical reason, hydrate signal, false positives, and ML role.",
    )
    add_compact_table(
        doc,
        ["Parameter family", "Physical reason", "Hydrate signal", "False positives or masks", "ML role"],
        [
            ["Depth, pressure, temperature", "Hydrate has a stability window", "Interval is admissible", "Thermal gradient, pressure, depth-unit errors", "Context feature or mask"],
            ["GR and lithology", "Separates cleaner sand from shale-prone intervals", "Low GR supports reservoir-quality sand", "Clean sand without hydrate; radioactive minerals; laminations", "Reservoir gate, not hydrate label"],
            ["Porosity and density", "Defines pore volume and supports elastic calculations", "Porous sand can host pore-filling hydrate", "Shale, carbonate, coal, compaction, washout", "Input feature and equation input"],
            ["NMR and core", "Independent pore-fluid and reservoir evidence", "NMR-density separation or calibrated saturation support", "Clay-bound water, tool settings, sparse core, depth mismatch", "Measured NMRPHI can be input; Sgh/NMR_SAT are targets"],
            ["Resistivity Rt", "Hydrate displaces conductive pore water", "Rt rises in clean porous sand", "Free gas, ice, tight rock, cement, salinity, bad hole", "Log-transform and combine with context"],
            ["Vp and Vs", "Hydrate can stiffen the sediment frame", "Vp and especially Vs may rise", "Ice, cement, carbonate, stress, lithology", "Input for stiffness and mimic separation"],
            ["Vp/Vs, impedance, lambda-rho, mu-rho", "Elastic attributes combine density and velocities", "Hydrate-consistent stiffness pattern", "Input errors, compaction, gas, shale, ice", "Derived features and crossplot checks"],
            ["Caliper and missingness", "Borehole condition controls log trust", "In-gauge intervals support other logs", "Washout, tool standoff, unknown bit size", "QC exclusion, downweighting, or confidence flag"],
        ],
        widths=[1.35, 1.55, 1.55, 2.0, 1.45],
        body_size=5.9,
    )
    add_body(
        doc,
        "The most important false positive is high resistivity without frame-stiffening and reservoir support. Free gas, ice, tight rock, cemented intervals, carbonate, salinity assumptions, and bad-hole conditions can all look hydrate-supportive in one curve. The model should see these as mimic risks, not as hidden rules to ignore.",
    )

    doc.add_heading("4. Parameter Movement Patterns", level=1)
    add_body(
        doc,
        "The model should learn how hydrate changes the local pocket of sediment. The setting matters because it changes the background, but the ML-relevant signal is the way the hydrate-bearing interval moves relative to water sand, gas sand, frozen sediment, shale, and tight rock.",
    )
    add_compact_table(
        doc,
        ["Scenario", "Expected parameter movement", "Why it matters for ML"],
        [
            ["Pore-filling hydrate in clean sand", "Low GR; porosity present; Rt up; Vp up; Vs up; mu-rho and impedance up; NMR mobile-fluid response reduced or NMR-density separation possible", "Main target pattern; needs multi-log agreement"],
            ["Water-bearing clean sand", "Low GR; porosity present; Rt low to moderate; velocities lower than hydrate sand; NMR mobile fluid present", "No-hydrate reservoir analogue"],
            ["Free gas in sand", "Rt can be high; Vp often lower; Vs does not rise like hydrate; Vp/Vs can rise", "Major resistivity mimic; elastic features matter"],
            ["Ice or frozen sediment", "Rt very high; Vp and Vs high; may overlap with hydrate stiffness", "Needs permafrost, depth, and reservoir-quality context"],
            ["Tight/cemented/carbonate rock", "Rt, Vp, Vs, and impedance can be high; porosity and NMR mobile fluid low", "Electrical and mechanical mimic; porosity gate matters"],
            ["Shale or clay-rich interval", "GR high; neutron porosity can read high; NMR may include clay-bound effects", "Usually not target reservoir logic"],
        ],
        widths=[1.8, 3.9, 2.0],
        body_size=6.3,
    )

    doc.add_heading("5. Feature Engineering and Screening Envelopes", level=1)
    add_body(
        doc,
        "Measured features should be the first block: GR, Rt, RHOB, porosity fields, NPHI, measured NMRPHI, Vp, Vs, depth/context fields, and available core or lithology fields. Derived features should be computed only after unit checks: Vp/Vs, acoustic impedance, shear modulus, bulk modulus, Young's modulus, Poisson's ratio, lambda-rho, and mu-rho. Archie-style saturation belongs as a physics baseline or review reference after Rw, Archie exponents, shale correction, and calibration are documented; it should not leak into predictors if it defines the target.",
    )
    add_body(
        doc,
        "The working ranges below are screening envelopes, not cutoffs. They are useful for QC, crossplots, and feature sanity checks, but the final model must learn from calibrated Sgh, NMR, core, or interpreted labels inside the authorized workflow.",
    )
    add_compact_table(
        doc,
        ["Feature", "Hydrate sand envelope", "Main use", "Main issue"],
        [
            ["Rt", "10-100+ ohm-m", "Electrical hydrate response", "Gas, ice, tight rock, salinity, cement, bad hole"],
            ["Vp", "2.5-4.0 km/s", "Compressional stiffness", "Ice, cement, carbonate, compaction"],
            ["Vs", "1.0-2.5 km/s", "Rigidity and gas separation", "Missing shear, ice, stress, lithology"],
            ["Vp/Vs", "1.6-2.4 broad; 1.4-1.6 crossplot hypothesis", "Elastic crossplot feature", "Strong overlap; never standalone"],
            ["mu-rho", "10-55 GPa*g/cc", "Rigidity discriminator", "Inherited density/Vs errors"],
            ["lambda-rho", "12-55 GPa*g/cc", "Incompressibility-sensitive feature", "Input units and gas/lithology overlap"],
            ["Porosity phi", "0.20-0.35", "Host-rock capacity", "Shale, compaction, density/NMR provenance"],
        ],
        widths=[1.35, 1.7, 1.65, 2.6],
        body_size=6.8,
    )

    doc.add_heading("6. Target Registry and Leakage Barrier", level=1)
    add_body(
        doc,
        "Before modeling, every candidate column must be assigned a role: measured input, derived feature, QC/context field, alignment field, target, calibration reference, or output. The target registry should fail closed. If a field might be target-derived, it should be excluded from predictors until provenance is confirmed.",
    )
    add_body(
        doc,
        "S_h, Sgh, NMR_SAT, phase labels, final rankings, and known/prediction outcomes are labels or outputs. They can supervise, calibrate, or score the model, but they cannot appear in the feature matrix. Measured NMRPHI can be an input if it is a measured porosity curve; NMR-derived saturation is a target. This distinction is essential for both the Word document and the runtime pipeline.",
    )
    add_process_sketch(
        doc,
        "Leakage-safe runtime flow",
        ["raw headers", "unit map", "QC", "features", "target registry", "split by well", "models", "outputs"],
        fill=SAND,
    )

    doc.add_heading("7. ML Pipeline and Model Ladder", level=1)
    add_body(
        doc,
        "The recommended implementation is baseline-first. Start with a physics/rule baseline and a simple supervised baseline, then test tree-based tabular models and ANN/MLP models only after the feature table, target registry, and complete-well split are working. LSTM sequence models and missing-log adapters should be optional future paths, not the first deliverable, because they increase leakage and provenance risk.",
    )
    add_compact_table(
        doc,
        ["Model option", "Why include it", "Current role"],
        [
            ["Physics/rule baseline", "Transparent comparison to petrophysical logic", "Required baseline and sanity check"],
            ["Logistic/linear/SGD", "Simple supervised benchmark", "Leakage check and first classifier/regressor"],
            ["Random forest or gradient boosting", "Strong nonlinear tabular model", "Main tabular candidate after grouped splits"],
            ["ANN/MLP", "Matches direct hydrate ML analogue", "Advanced saturation candidate after validation controls"],
            ["LSTM", "Uses depth sequence", "Optional only if continuous grouped sequences support it"],
            ["Missing-log model", "Can fill blocked curves", "Optional adapter with measured/estimated provenance flags"],
        ],
        widths=[1.8, 3.0, 2.7],
        body_size=7.0,
    )

    doc.add_heading("8. Validation and Outputs", level=1)
    add_body(
        doc,
        "Final validation should use complete-well or compartment holdouts. Random depth-row splits can be used only for debugging because adjacent samples from the same well share geology, tool response, and labeling assumptions. Preprocessing, imputation, scaling, feature selection, and calibration must be fitted on training wells only, then applied unchanged to validation, locked-test, and prediction wells.",
    )
    add_body(
        doc,
        "Outputs should remain separate: occurrence probability, saturation estimate, uncertainty or confidence, QC status, reason flags, mimic flags, and producibility or permeability context. Permeability is important for reservoir development, but it is not an occurrence label. The first results section should show planned output slots rather than fake metrics: well-log panel, occurrence probability track, saturation track, target overlay where allowed, residuals, calibration plot, QC/reason flags, and public-safe maps or interval summaries.",
    )

    doc.add_heading("9. Open Decisions Before Results Claims", level=1)
    for item in [
        "Which field is the authoritative hydrate-saturation target?",
        "Are saturation values fractions or percentages?",
        "Which fields define hydrate occurrence, no hydrate, uncertain hydrate, and possible hydrate habit?",
        "Which wells are training, validation, locked-test, and prediction wells?",
        "Which intervals are excluded for bad hole, missing curves, depth mismatch, or outliers?",
        "Which resistivity mnemonic is the preferred deep-resistivity predictor?",
        "Are Vp and Vs supplied directly or calculated from slowness curves?",
        "Which NMR fields are measured inputs and which are derived targets?",
        "Which caliper thresholds are valid after bit size and unit confirmation?",
    ]:
        add_body(doc, f"- {item}", size=9.6)

    doc.add_heading("10. Deliverable Implications", level=1)
    add_body(
        doc,
        "The Word document should explain the logic in prose and diagrams, not as a wall of tables. Tables should be used only where they clarify evidence tiers, parameter roles, model choices, or screening-envelope caveats. The PowerPoint should keep exactly nine slides and turn this document into left-to-right visual logic: hydrate definition, three-tier parameter ladder, raw headers to ML architecture, parameter movement, geomechanics, map context, output plan, and conclusion.",
    )

    doc.add_heading("References", level=1)
    for ref in REFERENCES:
        add_body(doc, ref, size=8.7)

    doc.save(DOCX_OUT)


def main() -> None:
    build_docx()
    print(DOCX_OUT)
    print("Current PPTX is built by docs/project_blueprints/build_ml_revamp_powerpoint.py")


if __name__ == "__main__":
    main()
