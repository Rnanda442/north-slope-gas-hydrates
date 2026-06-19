from __future__ import annotations

import csv
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
BLUEPRINT_DIR = ROOT / "docs" / "project_blueprints"
ASSET_DIR = BLUEPRINT_DIR / "presentation_assets" / "slide2_methane_context_2026_06_19"
SOURCE_DECK = BLUEPRINT_DIR / "V5_5_SLIDE2_SOURCE_UPDATE_North_Slope_Gas_Hydrate_ML_Workflow_Slides_2026-06-17.pptx"
OUT_DECK = BLUEPRINT_DIR / "V5_5_SLIDE2_METHANE_CONTEXT_REBUILD_North_Slope_Gas_Hydrate_ML_Workflow_Slides_2026-06-19.pptx"
OUT_ONE_SLIDE = ASSET_DIR / "slide_02_methane_hydrate_context_editable_2026_06_19.pptx"
OUT_PREVIEW = ASSET_DIR / "slide_02_methane_hydrate_context_rebuild_2026_06_19.png"
MAP_SOURCE = (
    BLUEPRINT_DIR
    / "presentation_assets"
    / "website_well_maps_2026_06_18"
    / "unified_north_slope_well_stability_context_map_2026_06_18.png"
)
MAP_CROP = ASSET_DIR / "unified_north_slope_map_crop_slide2_2026_06_19.png"
STRUCTURE_SOURCE = (
    ROOT
    / "docs"
    / "evidence"
    / "slide02_source_bundle_2026_06_17"
    / "slide02_selected_14_world_atlas_fig1_1_structure_types_clean.png"
)
STRUCTURE_CROP = ASSET_DIR / "hydrate_structure_i_ii_h_crop_slide2_2026_06_19.png"
CROSS_SECTION_SOURCE = (
    ROOT
    / "docs"
    / "evidence"
    / "slide02_source_bundle_2026_06_17"
    / "slide02_selected_06_usgs_arctic_alaska_cross_section_fig2_crop.png"
)
PHASE_CSV = ROOT / "data" / "public_stability_products" / "phase_curve_methane_5ppt_sir2008_csmhyd_digitized_v1.csv"
PT_DIAGRAM = ASSET_DIR / "methane_5ppt_pt_diagram_from_csv_2026_06_19.png"

W, H = 1600, 900
NAVY = (12, 34, 49)
INK = (24, 45, 58)
MUTED = (90, 110, 124)
LINE = (178, 197, 207)
PALE = (246, 250, 251)
WHITE = (255, 255, 255)
TEAL = (10, 118, 140)
BLUE = (37, 99, 235)
GREEN = (37, 154, 113)
AMBER = (217, 119, 6)
PURPLE = (124, 58, 237)
RED = (185, 44, 55)


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    names = ["arialbd.ttf" if bold else "arial.ttf", "segoeuib.ttf" if bold else "segoeui.ttf"]
    for root in [Path("C:/Windows/Fonts"), Path("/usr/share/fonts/truetype/dejavu")]:
        for name in names:
            path = root / name
            if path.exists():
                return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def wrap(draw: ImageDraw.ImageDraw, value: str, fnt: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    lines: list[str] = []
    for raw in value.splitlines():
        words = raw.split()
        line = ""
        for word in words:
            candidate = word if not line else f"{line} {word}"
            if draw.textlength(candidate, font=fnt) <= max_width:
                line = candidate
            else:
                if line:
                    lines.append(line)
                line = word
        if line:
            lines.append(line)
    return lines


def draw_text(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    value: str,
    size: int,
    fill: tuple[int, int, int] = INK,
    bold: bool = False,
    width: int | None = None,
    gap: int = 5,
    align: str = "left",
) -> int:
    fnt = font(size, bold)
    if width is None:
        draw.text(xy, value, font=fnt, fill=fill)
        return xy[1] + size + gap
    y = xy[1]
    for line in wrap(draw, value, fnt, width):
        line_width = draw.textlength(line, font=fnt)
        x = xy[0]
        if align == "center":
            x += max(0, int((width - line_width) / 2))
        elif align == "right":
            x += max(0, int(width - line_width))
        draw.text((x, y), line, font=fnt, fill=fill)
        y += size + gap
    return y


def card(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], fill=WHITE, outline=LINE, width: int = 2) -> None:
    draw.rounded_rectangle(box, radius=10, fill=fill, outline=outline, width=width)


def paste_contain(img: Image.Image, path: Path, box: tuple[int, int, int, int]) -> None:
    draw = ImageDraw.Draw(img)
    if not path.exists():
        card(draw, box, fill=(240, 246, 248), outline=RED, width=2)
        draw_text(draw, (box[0] + 18, box[1] + 20), f"Missing source: {path.name}", 18, RED, True, box[2] - box[0] - 36)
        return
    src = Image.open(path).convert("RGB")
    bw, bh = box[2] - box[0], box[3] - box[1]
    scale = min(bw / src.width, bh / src.height)
    new_size = (max(1, int(src.width * scale)), max(1, int(src.height * scale)))
    src = src.resize(new_size, Image.Resampling.LANCZOS)
    x = box[0] + (bw - new_size[0]) // 2
    y = box[1] + (bh - new_size[1]) // 2
    img.paste(src, (x, y))


def prepare_source_crops() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)

    if MAP_SOURCE.exists():
        src = Image.open(MAP_SOURCE).convert("RGB")
        # Crop to the map body and legend row; remove title/sidebar/footer text.
        crop = src.crop((125, 290, 2860, 1850))
        crop.save(MAP_CROP)

    if STRUCTURE_SOURCE.exists():
        src = Image.open(STRUCTURE_SOURCE).convert("RGB")
        # Keep source structure labels sI/sII/sH visible and remove blank margins.
        crop = src.crop((650, 0, 1278, 900))
        crop.save(STRUCTURE_CROP)


def build_pt_diagram() -> None:
    rows: list[tuple[float, float]] = []
    with PHASE_CSV.open(newline="", encoding="utf-8") as fh:
        for row in csv.DictReader(fh):
            try:
                rows.append((float(row["equilibrium_temperature_c"]), float(row["pressure_mpa_absolute"])))
            except (KeyError, ValueError):
                continue
    if not rows:
        raise ValueError(f"No plottable rows in {PHASE_CSV}")

    img = Image.new("RGB", (760, 470), WHITE)
    draw = ImageDraw.Draw(img)
    card(draw, (0, 0, 759, 469), fill=WHITE, outline=LINE, width=2)
    draw_text(draw, (28, 22), "Methane 5 ppt P-T diagram", 26, NAVY, True)
    draw_text(draw, (30, 56), "CSV-derived phase curve; context only", 16, MUTED, True)

    ax = (82, 104, 708, 388)
    t_values = [t for t, _ in rows]
    p_values = [p for _, p in rows]
    t_min, t_max = min(t_values) - 1.5, max(t_values) + 1.5
    p_min, p_max = 0.0, max(p_values) + 1.0

    def project(t: float, p: float) -> tuple[int, int]:
        x = ax[0] + int((t - t_min) / (t_max - t_min) * (ax[2] - ax[0]))
        y = ax[3] - int((p - p_min) / (p_max - p_min) * (ax[3] - ax[1]))
        return x, y

    draw.rectangle(ax, fill=(244, 249, 251), outline=LINE, width=1)
    for t_tick in [-10, -5, 0, 5, 10, 15]:
        if t_min <= t_tick <= t_max:
            x, _ = project(t_tick, p_min)
            draw.line((x, ax[1], x, ax[3]), fill=(220, 232, 238), width=1)
            draw_text(draw, (x - 18, ax[3] + 8), str(t_tick), 13, MUTED, True, width=36, align="center")
    for p_tick in [0, 4, 8, 12]:
        if p_min <= p_tick <= p_max:
            _, y = project(t_min, p_tick)
            draw.line((ax[0], y, ax[2], y), fill=(220, 232, 238), width=1)
            draw_text(draw, (24, y - 8), str(p_tick), 13, MUTED, True, width=46, align="right")

    points = [project(t, p) for t, p in rows]
    # Stable side is conceptual; the curve itself is the auditable CSV product.
    stable_poly = [(ax[0], ax[1]), *points, (ax[2], ax[1])]
    draw.polygon(stable_poly, fill=(224, 246, 236))
    draw.line(points, fill=TEAL, width=5)
    draw.line((ax[0], ax[3], ax[2], ax[3]), fill=MUTED, width=2)
    draw.line((ax[0], ax[3], ax[0], ax[1]), fill=MUTED, width=2)
    draw_text(draw, (ax[0] + 210, ax[3] + 30), "Temperature (deg C)", 15, NAVY, True, width=250, align="center")
    draw_text(draw, (8, ax[1] + 98), "Pressure\n(MPa abs)", 15, NAVY, True, width=68, align="center")
    draw_text(draw, (ax[0] + 78, ax[1] + 42), "hydrate-stable side\nunder assumptions", 15, GREEN, True, width=220)
    draw_text(draw, (ax[2] - 230, ax[3] - 70), "warmer / not stable\nfor this curve", 14, RED, True, width=210)
    img.save(PT_DIAGRAM)


def build_preview() -> None:
    prepare_source_crops()
    build_pt_diagram()
    img = Image.new("RGB", (W, H), (248, 251, 252))
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, W, 112), fill=(236, 248, 250))
    draw_text(draw, (54, 30), "Methane Hydrate and the Alaska North Slope", 37, NAVY, True)
    draw_text(draw, (58, 78), "Source-backed context for what hydrate is, why the region matters, and why P-T conditions vary.", 17, MUTED, True)

    map_box = (54, 132, 704, 476)
    card(draw, (46, 122, 714, 522), fill=WHITE, outline=LINE, width=2)
    paste_contain(img, MAP_CROP, map_box)
    draw_text(draw, (70, 486), "Combined 2D North Slope map: geology, public wells, permafrost controls, hydrate AUs, roads/TAPS/fields.", 14, NAVY, True, 616)

    pt_box = (748, 132, 1164, 390)
    card(draw, (738, 122, 1174, 430), fill=WHITE, outline=LINE, width=2)
    paste_contain(img, PT_DIAGRAM, pt_box)
    draw_text(draw, (762, 398), "P-T diagram varies with pressure, temperature, salinity, and gas mix.", 14, NAVY, True, 382)

    struct_box = (1206, 132, 1538, 390)
    card(draw, (1196, 122, 1548, 430), fill=WHITE, outline=LINE, width=2)
    paste_contain(img, STRUCTURE_CROP, struct_box)
    draw.rounded_rectangle((1280, 146, 1482, 228), radius=41, outline=TEAL, width=5)
    draw_text(draw, (1214, 398), "Structure I circled: methane-dominant baseline.", 14, TEAL, True, 316)

    cross_box = (54, 552, 704, 760)
    card(draw, (46, 540, 714, 814), fill=WHITE, outline=LINE, width=2)
    paste_contain(img, CROSS_SECTION_SOURCE, cross_box)
    draw_text(draw, (70, 768), "Regional cross-section context: geology and permafrost/stability setting explain why map context matters.", 13, NAVY, True, 616)

    cards = [
        ((748, 462, 1002, 584), "sI", "methane-rich baseline", TEAL),
        ((1022, 462, 1276, 584), "sII", "larger gases / mixed gas", AMBER),
        ((1296, 462, 1550, 584), "sH", "large hydrocarbon + methane", PURPLE),
        ((748, 610, 1002, 748), "Gas origin", "biogenic CH4 vs thermogenic C2+ shifts stability", BLUE),
        ((1022, 610, 1276, 748), "Why it matters", "USGS 2018 mean estimate: about 54 Tcf technically recoverable hydrate gas", GREEN),
        ((1296, 610, 1550, 748), "Guardrail", "stability context only: not proof, saturation, producibility, or ranking", RED),
    ]
    for box, head, body, color in cards:
        card(draw, box, fill=WHITE, outline=color, width=2)
        draw_text(draw, (box[0] + 16, box[1] + 14), head, 20, color, True, box[2] - box[0] - 32)
        draw_text(draw, (box[0] + 16, box[1] + 48), body, 15, NAVY, True, box[2] - box[0] - 32, gap=4)

    draw_text(draw, (60, 842), "Slide labels are concise; detailed citations stay in the companion notes.", 13, MUTED, True, 900)
    draw_text(draw, (1190, 842), "No hydrate proof. No occurrence/saturation claim.", 14, RED, True, 340, align="right")
    img.save(OUT_PREVIEW)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def add_textbox(slide, left, top, width, height, text, size=14, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_card(slide, left, top, width, height, outline, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*outline)
    shape.line.width = Pt(1.5)
    return shape


def add_panel(slide, title: str = "Methane Hydrate and the Alaska North Slope") -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 251, 252)
    add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.02), outline=(236, 248, 250), fill=(236, 248, 250))
    add_textbox(slide, Inches(0.42), Inches(0.20), Inches(10.7), Inches(0.42), title, 26, NAVY, True)
    add_textbox(slide, Inches(0.45), Inches(0.66), Inches(11.6), Inches(0.28), "Source-backed context: what hydrate is, why the region matters, and why P-T conditions vary.", 11, MUTED)

    # Source figures remain images; labels and callouts below are editable.
    add_card(slide, Inches(0.38), Inches(1.18), Inches(5.56), Inches(3.47), LINE)
    slide.shapes.add_picture(str(MAP_CROP), Inches(0.48), Inches(1.30), width=Inches(5.34), height=Inches(2.55))
    add_textbox(slide, Inches(0.56), Inches(3.93), Inches(5.04), Inches(0.48), "Combined 2D North Slope map: geology, public wells, permafrost controls, hydrate AUs, roads/TAPS/fields.", 8.5, NAVY, True)

    add_card(slide, Inches(6.18), Inches(1.18), Inches(3.52), Inches(2.74), LINE)
    slide.shapes.add_picture(str(PT_DIAGRAM), Inches(6.30), Inches(1.30), width=Inches(3.26), height=Inches(2.02))
    add_textbox(slide, Inches(6.42), Inches(3.34), Inches(3.02), Inches(0.36), "P-T diagram: pressure, temperature, salinity, and gas mix control admissibility.", 8.5, NAVY, True)

    add_card(slide, Inches(10.02), Inches(1.18), Inches(2.88), Inches(2.74), LINE)
    slide.shapes.add_picture(str(STRUCTURE_CROP), Inches(10.70), Inches(1.24), width=Inches(1.36), height=Inches(2.10))
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.67), Inches(1.29), Inches(1.82), Inches(0.72))
    oval.fill.background()
    oval.line.color.rgb = RGBColor(*TEAL)
    oval.line.width = Pt(2.5)
    add_textbox(slide, Inches(10.20), Inches(3.34), Inches(2.46), Inches(0.36), "Structure I circled: methane-dominant baseline.", 8.5, TEAL, True)

    add_card(slide, Inches(0.38), Inches(4.82), Inches(5.56), Inches(2.06), LINE)
    slide.shapes.add_picture(str(CROSS_SECTION_SOURCE), Inches(0.52), Inches(4.96), width=Inches(5.26), height=Inches(1.26))
    add_textbox(slide, Inches(0.58), Inches(6.18), Inches(5.00), Inches(0.45), "Regional cross-section context explains why permafrost, structure, and stability vary across the map.", 8.3, NAVY, True)

    native_cards = [
        (6.18, 4.18, 1.68, 0.92, "sI", "methane-rich baseline", TEAL),
        (8.02, 4.18, 1.68, 0.92, "sII", "larger gases / mixed gas", AMBER),
        (9.86, 4.18, 1.68, 0.92, "sH", "large hydrocarbon + methane", PURPLE),
        (11.70, 4.18, 1.20, 0.92, "noted", "scenario only", RED),
        (6.18, 5.40, 2.13, 1.16, "Gas origin", "biogenic CH4 vs thermogenic C2+ shifts stability", BLUE),
        (8.52, 5.40, 2.13, 1.16, "Why it matters", "USGS 2018 mean: about 54 Tcf technically recoverable hydrate gas", GREEN),
        (10.86, 5.40, 2.04, 1.16, "Guardrail", "context only; not proof, saturation, producibility, or ranking", RED),
    ]
    for x, y, w, h, head, body, color in native_cards:
        add_card(slide, Inches(x), Inches(y), Inches(w), Inches(h), color)
        add_textbox(slide, Inches(x + 0.08), Inches(y + 0.08), Inches(w - 0.16), Inches(0.22), head, 10, color, True)
        add_textbox(slide, Inches(x + 0.08), Inches(y + 0.36), Inches(w - 0.16), Inches(h - 0.42), body, 8.5, NAVY, True)

    add_textbox(slide, Inches(0.46), Inches(7.16), Inches(8.3), Inches(0.30), "Slide labels are concise; detailed citations stay in the companion notes.", 8, MUTED)
    add_textbox(slide, Inches(10.1), Inches(7.14), Inches(2.72), Inches(0.34), "No hydrate proof. No occurrence/saturation claim.", 8.5, RED, True, PP_ALIGN.RIGHT)

    # Small editable leader line from structure card to the Structure I oval.
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10.24), Inches(3.30), Inches(10.88), Inches(2.08))
    line.line.color.rgb = RGBColor(*TEAL)
    line.line.width = Pt(1.5)


def build_pptx() -> None:
    build_preview()
    prs = Presentation(SOURCE_DECK)
    slide = prs.slides[1]
    clear_slide(slide)
    add_panel(slide)
    prs.save(OUT_DECK)

    one = Presentation()
    one.slide_width = prs.slide_width
    one.slide_height = prs.slide_height
    blank = one.slide_layouts[6]
    slide = one.slides.add_slide(blank)
    add_panel(slide)
    one.save(OUT_ONE_SLIDE)


def main() -> None:
    build_pptx()
    print(f"Wrote {OUT_PREVIEW}")
    print(f"Wrote {OUT_ONE_SLIDE}")
    print(f"Wrote {OUT_DECK}")


if __name__ == "__main__":
    main()
