from __future__ import annotations

import math
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont


ROOT = Path(__file__).resolve().parents[2]
ASSET_DIR = (
    ROOT
    / "docs"
    / "project_blueprints"
    / "presentation_assets"
    / "slide3_signal_response_geopackage_update_2026_06_18"
)
MAP_PATH = (
    ROOT
    / "docs"
    / "project_blueprints"
    / "presentation_assets"
    / "website_well_maps_2026_06_18"
    / "slide3_correct_2d_well_stability_map_2026_06_18.png"
)
PNG_PATH = ASSET_DIR / "slide_03_signal_response_geopackage_update_2026_06_18.png"
PPTX_PATH = (
    ROOT
    / "docs"
    / "project_blueprints"
    / "North_Slope_Gas_Hydrate_Slide3_Signal_Response_GeoPackage_Update_2026-06-18.pptx"
)

W, H = 1920, 1080

NAVY = (11, 42, 57)
INK = (24, 54, 65)
MUTED = (89, 111, 121)
PALE = (244, 248, 249)
LINE = (194, 213, 219)
WHITE = (255, 255, 255)
BLUE = (45, 125, 190)
BLUE_LIGHT = (226, 241, 250)
PURPLE = (116, 104, 190)
PURPLE_LIGHT = (237, 233, 249)
GREEN = (35, 155, 126)
GREEN_LIGHT = (220, 245, 236)
AMBER = (218, 151, 61)
AMBER_LIGHT = (250, 238, 215)
RED = (188, 62, 68)
RED_LIGHT = (250, 229, 229)
GRAY = (122, 135, 141)


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        "arialbd.ttf" if bold else "arial.ttf",
        "segoeuib.ttf" if bold else "segoeui.ttf",
        "calibrib.ttf" if bold else "calibri.ttf",
    ]
    roots = [Path("C:/Windows/Fonts"), Path("/usr/share/fonts/truetype/dejavu")]
    for root in roots:
        for name in candidates:
            path = root / name
            if path.exists():
                return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def text(
    d: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    value: str,
    size: int,
    fill: tuple[int, int, int] = INK,
    bold: bool = False,
    width: int | None = None,
    align: str = "left",
    anchor: str | None = None,
) -> None:
    f = font(size, bold)
    if width is None:
        d.text(xy, value, font=f, fill=fill, anchor=anchor)
        return

    lines: list[str] = []
    for raw in value.splitlines():
        current = ""
        for word in raw.split():
            proposed = f"{current} {word}".strip()
            if d.textlength(proposed, font=f) <= width or not current:
                current = proposed
            else:
                lines.append(current)
                current = word
        if current:
            lines.append(current)
    y = xy[1]
    line_h = int(size * 1.18)
    for line in lines:
        if align == "center":
            d.text((xy[0] + width // 2, y), line, font=f, fill=fill, anchor="ma")
        else:
            d.text((xy[0], y), line, font=f, fill=fill)
        y += line_h


def rounded(
    d: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    fill: tuple[int, int, int],
    outline: tuple[int, int, int] = LINE,
    radius: int = 8,
    width: int = 2,
) -> None:
    d.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def alpha_rect(base: Image.Image, box: tuple[int, int, int, int], fill: tuple[int, int, int], alpha: int) -> None:
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    od.rectangle(box, fill=fill + (alpha,))
    base.alpha_composite(overlay)


def arrow(d: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], fill: tuple[int, int, int], width: int = 4) -> None:
    d.line([start, end], fill=fill, width=width)
    angle = math.atan2(end[1] - start[1], end[0] - start[0])
    length = 16
    p1 = (
        int(end[0] - length * math.cos(angle - 0.48)),
        int(end[1] - length * math.sin(angle - 0.48)),
    )
    p2 = (
        int(end[0] - length * math.cos(angle + 0.48)),
        int(end[1] - length * math.sin(angle + 0.48)),
    )
    d.polygon([end, p1, p2], fill=fill)


def marker(d: ImageDraw.ImageDraw, xy: tuple[int, int], label: str, fill: tuple[int, int, int]) -> None:
    x, y = xy
    d.ellipse((x - 18, y - 18, x + 18, y + 18), fill=fill, outline=WHITE, width=4)
    text(d, (x, y - 12), label, 21, WHITE, True, anchor="ma")


def log_curve_points(
    x: int,
    y: int,
    w: int,
    h: int,
    phase: float,
    high_shift: float = 0.0,
    low_shift: float = 0.0,
) -> list[tuple[int, int]]:
    points: list[tuple[int, int]] = []
    for i in range(90):
        t = i / 89
        wave = 0.5 + 0.28 * math.sin(phase + 10.5 * t) + 0.13 * math.sin(phase * 1.7 + 22 * t)
        if 0.26 < t < 0.39:
            wave += high_shift
        if 0.58 < t < 0.72:
            wave += low_shift
        wave = max(0.05, min(0.95, wave))
        points.append((x + int(18 + wave * (w - 36)), y + int(t * h)))
    return points


def draw_log_stack(img: Image.Image) -> None:
    d = ImageDraw.Draw(img)
    x0, y0 = 74, 170
    w, h = 860, 610
    rounded(d, (x0, y0, x0 + w, y0 + h), WHITE, LINE, radius=8)
    text(d, (x0 + 20, y0 + 16), "QC-cleaned normalized well-log scaffold", 28, NAVY, True)
    text(d, (x0 + 20, y0 + 53), "Depth aligned | schematic only | no approved row values", 18, MUTED)

    plot_y = y0 + 96
    plot_h = 430
    label_y = plot_y - 32
    track_w = 76
    gap = 22
    depth_x = x0 + 40
    track_x0 = x0 + 96

    alpha_rect(img, (track_x0 - 16, plot_y + 112, x0 + w - 34, plot_y + 188), GREEN, 58)
    alpha_rect(img, (track_x0 - 16, plot_y + 276, x0 + w - 34, plot_y + 356), AMBER, 54)

    d.line((depth_x, plot_y, depth_x, plot_y + plot_h), fill=INK, width=3)
    for offset, label in [(0, "shallow"), (plot_h // 2, "depth"), (plot_h, "deep")]:
        y = plot_y + offset
        d.line((depth_x - 8, y, depth_x + 8, y), fill=INK, width=2)
        text(d, (x0 + 18, y - 10), label, 14, MUTED)

    tracks = [
        ("GR", BLUE, -0.28, 0.10, "clean"),
        ("Phi/RHOB", BLUE, 0.22, -0.05, "pore"),
        ("Rt", BLUE, 0.34, 0.25, "resistive"),
        ("Vp", BLUE, 0.26, -0.12, "fast"),
        ("Vs/mu-rho", PURPLE, 0.31, -0.18, "stiff"),
        ("NMR/core", BLUE, -0.10, 0.18, "calib"),
        ("Stability", AMBER, 0.36, -0.30, "context"),
    ]
    for idx, (label, color, high, low, note) in enumerate(tracks):
        x = track_x0 + idx * (track_w + gap)
        d.rectangle((x, plot_y, x + track_w, plot_y + plot_h), fill=(252, 254, 254), outline=LINE)
        d.line((x + track_w // 2, plot_y, x + track_w // 2, plot_y + plot_h), fill=(226, 236, 239), width=2)
        text(d, (x + track_w // 2, label_y), label, 16, color, True, anchor="ma")
        pts = log_curve_points(x, plot_y, track_w, plot_h, phase=idx * 0.83 + 0.4, high_shift=high, low_shift=low)
        d.line(pts, fill=color, width=4, joint="curve")
        text(d, (x + track_w // 2, plot_y + plot_h + 10), note, 14, MUTED, True, anchor="ma")

    # Interval labels live outside the curves so highlighted bands stay visible.
    marker(d, (x0 + w - 27, plot_y + 150), "1", GREEN)
    marker(d, (x0 + w - 27, plot_y + 315), "2", AMBER)
    # Lanes under the log stack.
    lane_y = y0 + h - 72
    lane_h = 52
    lanes = [
        (x0 + 22, "Measured logs\nDOE inputs", BLUE, BLUE_LIGHT),
        (x0 + 234, "Derived\nfeatures", PURPLE, PURPLE_LIGHT),
        (x0 + 446, "QC / context", AMBER, AMBER_LIGHT),
        (x0 + 658, "Y-only targets", RED, RED_LIGHT),
    ]
    for x, label, color, fill in lanes:
        rounded(d, (x, lane_y, x + 184, lane_y + lane_h), fill, color, radius=8, width=2)
        text(d, (x + 92, lane_y + 8), label, 14, color, True, width=164, align="center")


def draw_map_inset(img: Image.Image) -> None:
    d = ImageDraw.Draw(img)
    x1, y1, x2, y2 = 995, 170, 1844, 555
    rounded(d, (x1, y1, x2, y2), WHITE, LINE, radius=8)
    text(d, (x1 + 18, y1 + 16), "Latest public website stability map", 25, NAVY, True)
    text(d, (x1 + 18, y1 + 47), "GeoPackage/GIS orientation layers visible", 17, MUTED)

    if MAP_PATH.exists():
        src = Image.open(MAP_PATH).convert("RGB")
        sw, sh = src.size
        # Focus on the North Slope trend so field labels, unit outlines, roads,
        # and TAPS remain readable in the slide inset.
        crop = src.crop((int(sw * 0.135), int(sh * 0.25), int(sw * 0.87), int(sh * 0.63)))
        box = (x1 + 20, y1 + 78, x2 - 20, y2 - 58)
        bw, bh = box[2] - box[0], box[3] - box[1]
        scale = min(bw / crop.width, bh / crop.height)
        resized = crop.resize((int(crop.width * scale), int(crop.height * scale)), Image.Resampling.LANCZOS)
        px = box[0] + (bw - resized.width) // 2
        py = box[1] + (bh - resized.height) // 2
        img.paste(resized, (px, py))
        d.rectangle((box[0], box[1], box[2], box[3]), outline=LINE, width=2)
    else:
        d.rectangle((x1 + 20, y1 + 78, x2 - 20, y2 - 58), fill=(252, 239, 239), outline=RED, width=2)
        text(d, (x1 + 38, y1 + 110), f"Missing map:\n{MAP_PATH.name}", 24, RED, True, width=740)

    d.rectangle((x1 + 20, y2 - 49, x2 - 20, y2 - 18), fill=AMBER_LIGHT, outline=None)
    text(d, (x1 + 34, y2 - 43), "Public stability-admissibility context, not hydrate proof.", 19, NAVY, True)


def draw_callouts(img: Image.Image) -> None:
    d = ImageDraw.Draw(img)
    x1, y1, x2, y2 = 995, 585, 1844, 812
    rounded(d, (x1, y1, x2, y2), WHITE, LINE, radius=8)
    text(d, (x1 + 18, y1 + 15), "Visual flow and leakage guardrail", 25, NAVY, True)

    rows = [
        ("1", GREEN, "Host + response", "Low GR / pore space plus Rt and velocity stiffness make a stronger candidate stack."),
        ("2", AMBER, "Mimic review", "A resistive interval alone can be shale, gas, cement, or poor-hole response."),
        ("3", AMBER, "QC/context", "Caliper and washout are trust gates; stability is admissible-under-assumptions only."),
        ("4", RED, "Targets locked", "S_h, Sgh, Sh, NMR_SAT, Hydrate Saturation, Swr/S_wr, phase labels never enter X."),
    ]
    y = y1 + 56
    for label, color, head, body in rows:
        marker(d, (x1 + 29, y + 14), label, color)
        text(d, (x1 + 60, y), head, 18, color, True)
        text(d, (x1 + 60, y + 22), body, 15, INK, width=760)
        y += 42

    # Leaders to the scaffold interval markers.
    arrow(d, (936, 416), (994, 668), GREEN, 3)
    arrow(d, (936, 581), (994, 700), AMBER, 3)


def draw_equations(img: Image.Image) -> None:
    d = ImageDraw.Draw(img)
    x1, y1, x2, y2 = 74, 820, 1844, 1010
    rounded(d, (x1, y1, x2, y2), WHITE, LINE, radius=8)
    text(d, (x1 + 18, y1 + 14), "Equation lane: derived features are built after units/QC, before labels are joined", 24, NAVY, True)

    cards = [
        ("Density porosity", "phi_D = (rho_ma - RHOB) / (rho_ma - rho_f)", PURPLE),
        ("Archie water/hydrate", "S_w^n = a R_w / (phi^m R_t)  |  S_h = 1 - S_w", PURPLE),
        ("Sonic + elastic", "Vp, Vs -> Vp/Vs, AI = rho Vp, mu-rho = rho Vs^2, lambda-rho = rho(Vp^2 - 2Vs^2)", PURPLE),
        ("Stability context", "P(z) = P0 + rho_w g z  |  T_model <= T_eq(P, gas, salinity)", AMBER),
    ]
    cx = x1 + 18
    cy = y1 + 56
    widths = [325, 380, 510, 350]
    for idx, (head, body, color) in enumerate(cards):
        fill = PURPLE_LIGHT if color == PURPLE else AMBER_LIGHT
        rounded(d, (cx, cy, cx + widths[idx], y2 - 20), fill, color, radius=8, width=2)
        text(d, (cx + 16, cy + 15), head, 19, color, True)
        text(d, (cx + 16, cy + 48), body, 15, INK, width=widths[idx] - 32)
        if idx < len(cards) - 1:
            arrow(d, (cx + widths[idx] + 6, cy + 59), (cx + widths[idx] + 28, cy + 59), MUTED, 3)
        cx += widths[idx] + 36


def draw_legend(img: Image.Image) -> None:
    d = ImageDraw.Draw(img)
    x, y = 1015, 84
    items = [
        ("Measured log/input", BLUE),
        ("Derived feature", PURPLE),
        ("Assumption/context", AMBER),
        ("Target-only field", RED),
    ]
    for label, color in items:
        d.rectangle((x, y + 7, x + 22, y + 29), fill=color)
        text(d, (x + 32, y + 5), label, 16, INK)
        x += 220


def build_slide() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGBA", (W, H), WHITE + (255,))
    d = ImageDraw.Draw(img)
    d.rectangle((0, 0, W, 108), fill=PALE)
    d.rectangle((0, 0, 16, H), fill=GREEN)
    text(d, (74, 36), "Slide 3 - Equation / Well-Log Scaffold", 42, NAVY, True)
    text(
        d,
        (76, 86),
        "Cleaned signals -> equation features -> QC/context -> leakage-safe occurrence and saturation labels",
        20,
        MUTED,
    )
    draw_legend(img)
    draw_log_stack(img)
    draw_map_inset(img)
    draw_callouts(img)
    draw_equations(img)
    text(
        d,
        (74, 1034),
        "Public-safe schematic. No approved rows, well identifiers, row-level predictions, fitted models, or final hydrate claims are shown.",
        17,
        MUTED,
    )
    img.convert("RGB").save(PNG_PATH, quality=95)


def build_pptx() -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:  # pragma: no cover - environment dependent
        raise SystemExit("python-pptx is required to build the PPTX export.") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(PNG_PATH), 0, 0, width=prs.slide_width, height=prs.slide_height)
    PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)


def main() -> None:
    build_slide()
    build_pptx()
    print(PNG_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
