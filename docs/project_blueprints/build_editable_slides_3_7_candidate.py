from __future__ import annotations

import csv
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
BLUEPRINT_DIR = ROOT / "docs" / "project_blueprints"
ASSET_DIR = BLUEPRINT_DIR / "presentation_assets" / "editable_slides_3_7_2026_06_19"
SLIDE2_CANDIDATE_DECK = (
    BLUEPRINT_DIR
    / "V5_5_SLIDE2_METHANE_CONTEXT_REBUILD_North_Slope_Gas_Hydrate_ML_Workflow_Slides_2026-06-19.pptx"
)
SOURCE_UPDATE_DECK = (
    BLUEPRINT_DIR
    / "V5_5_SLIDE2_SOURCE_UPDATE_North_Slope_Gas_Hydrate_ML_Workflow_Slides_2026-06-17.pptx"
)
OUT_DECK = (
    BLUEPRINT_DIR
    / "V5_5_EDITABLE_SLIDES_3_7_North_Slope_Gas_Hydrate_ML_Workflow_Slides_2026-06-19.pptx"
)
MAP_SOURCE = (
    BLUEPRINT_DIR
    / "presentation_assets"
    / "website_well_maps_2026_06_18"
    / "unified_north_slope_slide_export_callout_space_2026_06_18.png"
)
AUDIT_CSV = ASSET_DIR / "editable_slides_3_7_editability_audit_2026_06_19.csv"
SOURCE_NOTES = ASSET_DIR / "editable_slides_3_7_source_claim_notes_2026_06_19.md"
CONTACT_SHEET = ASSET_DIR / "editable_slides_3_7_contact_sheet_2026_06_19.png"

PREVIEW_FILES = {
    3: ASSET_DIR / "slide_03_log_lithology_core_editable_2026_06_19.png",
    4: ASSET_DIR / "slide_04_simplified_ml_architecture_editable_2026_06_19.png",
    5: ASSET_DIR / "slide_05_equation_cards_editable_2026_06_19.png",
    6: ASSET_DIR / "slide_06_evidence_review_board_editable_2026_06_19.png",
    7: ASSET_DIR / "slide_07_stability_context_map_editable_2026_06_19.png",
}

W, H = 1600, 900
NAVY = (12, 34, 49)
INK = (25, 45, 58)
MUTED = (86, 106, 121)
LINE = (177, 196, 206)
PALE = (246, 250, 251)
WHITE = (255, 255, 255)
TEAL = (10, 118, 140)
BLUE = (37, 99, 235)
GREEN = (35, 148, 103)
AMBER = (217, 119, 6)
PURPLE = (124, 58, 237)
RED = (185, 44, 55)
GRAY = (98, 115, 129)
YELLOW = (250, 204, 21)


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    names = ["arialbd.ttf" if bold else "arial.ttf", "segoeuib.ttf" if bold else "segoeui.ttf"]
    for root in [Path("C:/Windows/Fonts"), Path("/usr/share/fonts/truetype/dejavu")]:
        for name in names:
            path = root / name
            if path.exists():
                return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def wrap(draw: ImageDraw.ImageDraw, value: str, fnt: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    lines: list[str] = []
    for raw in value.splitlines():
        words = raw.split()
        if not words:
            lines.append("")
            continue
        line = ""
        for word in words:
            candidate = word if not line else f"{line} {word}"
            if draw.textlength(candidate, font=fnt) <= max_width:
                line = candidate
            else:
                if line:
                    lines.append(line)
                line = word
        if line:
            lines.append(line)
    return lines


def draw_text(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    value: str,
    size: int,
    fill: tuple[int, int, int] = INK,
    bold: bool = False,
    width: int | None = None,
    gap: int = 5,
    align: str = "left",
) -> int:
    fnt = font(size, bold)
    if width is None:
        draw.text(xy, value, font=fnt, fill=fill)
        return xy[1] + size + gap
    y = xy[1]
    for line in wrap(draw, value, fnt, width):
        line_width = draw.textlength(line, font=fnt)
        x = xy[0]
        if align == "center":
            x += max(0, int((width - line_width) / 2))
        elif align == "right":
            x += max(0, int(width - line_width))
        draw.text((x, y), line, font=fnt, fill=fill)
        y += size + gap
    return y


def draw_card(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    fill: tuple[int, int, int] = WHITE,
    outline: tuple[int, int, int] = LINE,
    width: int = 2,
    radius: int = 10,
) -> None:
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def paste_contain(img: Image.Image, path: Path, box: tuple[int, int, int, int]) -> None:
    draw = ImageDraw.Draw(img)
    if not path.exists():
        draw_card(draw, box, fill=(249, 242, 242), outline=RED, width=2)
        draw_text(draw, (box[0] + 16, box[1] + 18), f"Missing map source: {path.name}", 22, RED, True, box[2] - box[0] - 32)
        return
    src = Image.open(path).convert("RGB")
    bw, bh = box[2] - box[0], box[3] - box[1]
    scale = min(bw / src.width, bh / src.height)
    new_size = (max(1, int(src.width * scale)), max(1, int(src.height * scale)))
    src = src.resize(new_size, Image.Resampling.LANCZOS)
    x = box[0] + (bw - new_size[0]) // 2
    y = box[1] + (bh - new_size[1]) // 2
    img.paste(src, (x, y))


def preview_header(draw: ImageDraw.ImageDraw, title: str, subtitle: str) -> None:
    draw.rectangle((0, 0, W, 112), fill=(236, 248, 250))
    draw_text(draw, (54, 27), title, 35, NAVY, True)
    draw_text(draw, (58, 75), subtitle, 18, MUTED, True)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    size: float = 14,
    color: tuple[int, int, int] = NAVY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_card(
    slide,
    left,
    top,
    width,
    height,
    outline: tuple[int, int, int] = LINE,
    fill: tuple[int, int, int] = WHITE,
    line_width: float = 1.2,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*outline)
    shape.line.width = Pt(line_width)
    return shape


def add_rect(slide, left, top, width, height, fill, outline=LINE, line_width: float = 0.8):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*outline)
    shape.line.width = Pt(line_width)
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, width: float = 1.2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = RGBColor(*color)
    line.line.width = Pt(width)
    return line


def add_ppt_header(slide, title: str, subtitle: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 251, 252)
    add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.02), outline=(236, 248, 250), fill=(236, 248, 250))
    add_textbox(slide, Inches(0.42), Inches(0.20), Inches(10.9), Inches(0.42), title, 24, NAVY, True)
    add_textbox(slide, Inches(0.45), Inches(0.66), Inches(11.8), Inches(0.28), subtitle, 10.5, MUTED)


def add_chip(slide, x, y, w, label: str, color: tuple[int, int, int]) -> None:
    add_card(slide, Inches(x), Inches(y), Inches(w), Inches(0.25), outline=color, fill=(248, 251, 252), line_width=1)
    add_textbox(slide, Inches(x + 0.04), Inches(y + 0.04), Inches(w - 0.08), Inches(0.16), label, 6.8, color, True, PP_ALIGN.CENTER)


def curve_points(track_left: float, track_top: float, track_w: float, track_h: float, values: list[float]) -> list[tuple[float, float]]:
    points = []
    n = len(values) - 1
    for i, value in enumerate(values):
        x = track_left + 0.08 + value * (track_w - 0.16)
        y = track_top + 0.18 + (i / n) * (track_h - 0.36)
        points.append((x, y))
    return points


def add_curve(slide, points: list[tuple[float, float]], color: tuple[int, int, int], width: float = 1.7) -> None:
    for (x1, y1), (x2, y2) in zip(points, points[1:]):
        add_line(slide, Inches(x1), Inches(y1), Inches(x2), Inches(y2), color=color, width=width)


def draw_curve(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], values: list[float], color: tuple[int, int, int]) -> None:
    x0, y0, x1, y1 = box
    pts = []
    n = len(values) - 1
    for i, value in enumerate(values):
        x = int(x0 + 14 + value * (x1 - x0 - 28))
        y = int(y0 + 28 + (i / n) * (y1 - y0 - 44))
        pts.append((x, y))
    draw.line(pts, fill=color, width=4)


def draw_preview_slide3() -> None:
    img = Image.new("RGB", (W, H), (248, 251, 252))
    draw = ImageDraw.Draw(img)
    preview_header(draw, "Slide 3: Log Signals, Lithology, And Calibration", "Editable scaffold; DOE export placeholder, no private rows.")
    panel = (60, 138, 1160, 738)
    draw_card(draw, panel)
    draw_text(draw, (84, 156), "Well-log style evidence panel", 20, NAVY, True)
    draw_text(draw, (84, 184), "schematic only - replace with reviewed row-free export", 14, RED, True)

    track_names = ["Lith", "GR", "Rt", "Phi/NMR", "Vp", "Vs", "Vp/Vs", "AI", "Caliper"]
    colors = [GRAY, GREEN, BLUE, TEAL, PURPLE, PURPLE, AMBER, AMBER, RED]
    values = [
        [0.3, 0.35, 0.4, 0.7, 0.8, 0.45, 0.3, 0.35],
        [0.55, 0.42, 0.28, 0.18, 0.22, 0.46, 0.64, 0.58],
        [0.25, 0.38, 0.65, 0.88, 0.78, 0.48, 0.30, 0.27],
        [0.65, 0.58, 0.48, 0.26, 0.34, 0.54, 0.67, 0.72],
        [0.36, 0.44, 0.60, 0.73, 0.68, 0.52, 0.40, 0.35],
        [0.32, 0.40, 0.56, 0.78, 0.70, 0.48, 0.38, 0.34],
        [0.46, 0.50, 0.62, 0.72, 0.66, 0.56, 0.44, 0.42],
        [0.34, 0.44, 0.61, 0.80, 0.74, 0.50, 0.36, 0.32],
        [0.22, 0.24, 0.30, 0.52, 0.30, 0.32, 0.24, 0.22],
    ]
    x = 92
    y = 226
    track_w = 110
    track_h = 430
    for idx, name in enumerate(track_names):
        box = (x + idx * 114, y, x + idx * 114 + track_w, y + track_h)
        draw.rectangle(box, fill=(250, 252, 253), outline=LINE, width=1)
        draw_text(draw, (box[0] + 5, box[1] - 24), name, 13, NAVY, True, track_w - 10, align="center")
        if idx == 0:
            draw.rectangle((box[0] + 16, box[1] + 82, box[2] - 16, box[1] + 174), fill=(226, 247, 236), outline=GREEN)
            draw.rectangle((box[0] + 16, box[1] + 176, box[2] - 16, box[1] + 256), fill=(239, 243, 245), outline=GRAY)
            draw.rectangle((box[0] + 16, box[1] + 258, box[2] - 16, box[1] + 328), fill=(255, 247, 225), outline=AMBER)
            draw_text(draw, (box[0] + 24, box[1] + 102), "clean", 12, GREEN, True)
            draw_text(draw, (box[0] + 24, box[1] + 204), "mixed", 12, GRAY, True)
            draw_text(draw, (box[0] + 24, box[1] + 282), "uncertain", 10, AMBER, True)
        else:
            draw.rectangle((box[0] + 2, box[1] + 112, box[2] - 2, box[1] + 210), fill=(229, 247, 238), outline=None)
            draw_curve(draw, box, values[idx], colors[idx])
    draw_text(draw, (96, 682), "Depth axis is schematic. Depth is not normalized; log values are future 0-1 scaled views after train split.", 15, MUTED, True, 930)

    target_box = (1212, 138, 1538, 584)
    draw_card(draw, target_box, outline=RED)
    draw_text(draw, (1234, 160), "Y-only calibration rail", 21, RED, True, 270)
    for i, text in enumerate(["S_h / Sgh / Sh", "NMR_SAT", "Hydrate saturation", "Swr / S_wr", "phase labels"]):
        y0 = 218 + i * 52
        draw_card(draw, (1234, y0, 1510, y0 + 34), fill=(255, 247, 247), outline=RED, width=1)
        draw_text(draw, (1248, y0 + 8), text, 14, RED, True)
    draw_text(draw, (1234, 502), "Targets supervise/review only. They never enter X predictors.", 16, NAVY, True, 270)

    callouts = [
        ((1212, 620, 1538, 688), "stack evidence", "clean + resistive + stiff + stable"),
        ((1212, 704, 1538, 772), "caliper caution", "washout is QC, not hydrate evidence"),
        ((1212, 788, 1538, 842), "public-safe", "no approved rows in GitHub"),
    ]
    for box, head, body in callouts:
        draw_card(draw, box, outline=TEAL if head != "caliper caution" else RED)
        draw_text(draw, (box[0] + 14, box[1] + 8), head, 15, TEAL if head != "caliper caution" else RED, True, box[2] - box[0] - 28)
        draw_text(draw, (box[0] + 14, box[1] + 32), body, 13, NAVY, True, box[2] - box[0] - 28)
    img.save(PREVIEW_FILES[3])


def build_slide3(slide) -> None:
    add_ppt_header(
        slide,
        "Log Signals, Lithology, And Calibration",
        "Editable well-log scaffold; real four-well/core data remain DOE-runtime exports until reviewed.",
    )
    add_card(slide, Inches(0.42), Inches(1.18), Inches(9.22), Inches(5.72), LINE)
    add_textbox(slide, Inches(0.62), Inches(1.34), Inches(3.1), Inches(0.25), "Well-log style evidence panel", 12, NAVY, True)
    add_textbox(slide, Inches(3.50), Inches(1.34), Inches(5.7), Inches(0.25), "DOE export placeholder: schematic only, no approved rows or private depths.", 8.5, RED, True, PP_ALIGN.RIGHT)

    track_top = 1.85
    track_h = 3.78
    track_w = 0.83
    start_x = 0.70
    gap = 0.12
    tracks = [
        ("Lith", GRAY, []),
        ("GR", GREEN, [0.56, 0.42, 0.25, 0.18, 0.22, 0.46, 0.64, 0.58]),
        ("Rt", BLUE, [0.25, 0.38, 0.65, 0.88, 0.78, 0.48, 0.30, 0.27]),
        ("Phi/NMR", TEAL, [0.65, 0.58, 0.48, 0.26, 0.34, 0.54, 0.67, 0.72]),
        ("Vp", PURPLE, [0.36, 0.44, 0.60, 0.73, 0.68, 0.52, 0.40, 0.35]),
        ("Vs", PURPLE, [0.32, 0.40, 0.56, 0.78, 0.70, 0.48, 0.38, 0.34]),
        ("Vp/Vs", AMBER, [0.46, 0.50, 0.62, 0.72, 0.66, 0.56, 0.44, 0.42]),
        ("AI", AMBER, [0.34, 0.44, 0.61, 0.80, 0.74, 0.50, 0.36, 0.32]),
        ("Caliper", RED, [0.22, 0.24, 0.30, 0.52, 0.30, 0.32, 0.24, 0.22]),
    ]
    for i, (name, color, vals) in enumerate(tracks):
        x = start_x + i * (track_w + gap)
        add_textbox(slide, Inches(x), Inches(track_top - 0.30), Inches(track_w), Inches(0.22), name, 7.2, NAVY, True, PP_ALIGN.CENTER)
        add_rect(slide, Inches(x), Inches(track_top), Inches(track_w), Inches(track_h), fill=(250, 252, 253), outline=LINE)
        if name == "Lith":
            lith = [
                (track_top + 0.58, 0.82, "clean", GREEN, (226, 247, 236)),
                (track_top + 1.42, 0.78, "mixed", GRAY, (239, 243, 245)),
                (track_top + 2.24, 0.62, "uncertain", AMBER, (255, 247, 225)),
            ]
            for y, h, label, line_color, fill in lith:
                add_rect(slide, Inches(x + 0.10), Inches(y), Inches(track_w - 0.20), Inches(h), fill=fill, outline=line_color)
                add_textbox(slide, Inches(x + 0.12), Inches(y + 0.22), Inches(track_w - 0.24), Inches(0.24), label, 6.4, line_color, True, PP_ALIGN.CENTER)
        else:
            add_rect(slide, Inches(x + 0.02), Inches(track_top + 1.08), Inches(track_w - 0.04), Inches(0.95), fill=(229, 247, 238), outline=(229, 247, 238))
            add_curve(slide, curve_points(x, track_top, track_w, track_h, vals), color, width=1.5)

    add_textbox(
        slide,
        Inches(0.72),
        Inches(5.88),
        Inches(8.40),
        Inches(0.34),
        "Depth axis is schematic. Depth stays as alignment/context; log values are future train-only 0-1 scaled views.",
        8.4,
        MUTED,
        True,
    )
    add_card(slide, Inches(0.70), Inches(6.36), Inches(8.72), Inches(0.34), outline=LINE, fill=(248, 251, 252), line_width=1)
    legend = [
        ("measured log", GREEN),
        ("derived signal", AMBER),
        ("QC only", RED),
        ("target/reference", RED),
        ("placeholder", GRAY),
    ]
    lx = 0.82
    for label, color in legend:
        add_rect(slide, Inches(lx), Inches(6.45), Inches(0.10), Inches(0.10), fill=color, outline=color)
        add_textbox(slide, Inches(lx + 0.14), Inches(6.38), Inches(1.18), Inches(0.22), label, 6.6, NAVY, True)
        lx += 1.45

    add_card(slide, Inches(10.02), Inches(1.18), Inches(2.86), Inches(3.72), outline=RED)
    add_textbox(slide, Inches(10.22), Inches(1.38), Inches(2.40), Inches(0.28), "Y-only calibration rail", 13, RED, True)
    target_labels = ["S_h / Sgh / Sh", "NMR_SAT", "Hydrate saturation", "Swr / S_wr", "phase labels"]
    for idx, label in enumerate(target_labels):
        y = 1.86 + idx * 0.43
        add_card(slide, Inches(10.24), Inches(y), Inches(2.36), Inches(0.29), outline=RED, fill=(255, 247, 247), line_width=0.8)
        add_textbox(slide, Inches(10.32), Inches(y + 0.06), Inches(2.16), Inches(0.14), label, 7.2, RED, True)
    add_textbox(slide, Inches(10.24), Inches(4.18), Inches(2.34), Inches(0.40), "Targets supervise/review only. They never enter X predictors.", 7.8, NAVY, True)

    callouts = [
        (10.02, 5.18, 2.86, 0.56, "stack evidence", "clean + resistive + stiff + stable", TEAL),
        (10.02, 5.90, 2.86, 0.56, "caliper caution", "washout is QC, not hydrate evidence", RED),
        (10.02, 6.62, 2.86, 0.42, "public-safe", "no approved rows in GitHub", GRAY),
    ]
    for x, y, w, h, head, body, color in callouts:
        add_card(slide, Inches(x), Inches(y), Inches(w), Inches(h), outline=color, fill=WHITE)
        add_textbox(slide, Inches(x + 0.10), Inches(y + 0.07), Inches(w - 0.20), Inches(0.18), head, 7.8, color, True)
        add_textbox(slide, Inches(x + 0.10), Inches(y + 0.28), Inches(w - 0.20), Inches(h - 0.30), body, 6.8, NAVY, True)
    add_line(slide, Inches(9.25), Inches(3.28), Inches(10.02), Inches(5.46), color=TEAL, width=1.2)
    add_line(slide, Inches(9.10), Inches(4.12), Inches(10.02), Inches(6.18), color=RED, width=1.2)


def draw_preview_slide4() -> None:
    img = Image.new("RGB", (W, H), (248, 251, 252))
    draw = ImageDraw.Draw(img)
    preview_header(draw, "Slide 4: Simplified ML Architecture", "Audience-facing workflow: inputs, leakage barrier, model outputs, validation.")
    boxes = [
        ((64, 160, 330, 330), "Inputs", "logs\ncore/NMR\nlithology/QC\nstability context", TEAL),
        ((398, 160, 664, 330), "Prepare", "original headers\nunits\nalign by depth\nvariable fingerprints", BLUE),
        ((732, 160, 998, 330), "Leakage barrier", "X inputs only\nY targets locked away", RED),
        ((1066, 160, 1532, 330), "Feature matrix", "measured logs + valid derived features\ncontext flags only", GREEN),
        ((398, 470, 664, 638), "Baselines", "physics/rules\ntree/boosting", GRAY),
        ((732, 470, 998, 638), "ANN / Keras later", "learned patterns after baseline checks", PURPLE),
        ((1066, 470, 1532, 638), "Two outputs", "occurrence classification\nsaturation regression", AMBER),
    ]
    for box, head, body, color in boxes:
        draw_card(draw, box, outline=color)
        draw_text(draw, (box[0] + 22, box[1] + 22), head, 24, color, True, box[2] - box[0] - 44)
        draw_text(draw, (box[0] + 22, box[1] + 68), body, 18, NAVY, True, box[2] - box[0] - 44)
    for x1, y1, x2, y2 in [(330, 244, 398, 244), (664, 244, 732, 244), (998, 244, 1066, 244), (664, 552, 732, 552), (998, 552, 1066, 552)]:
        draw.line((x1, y1, x2, y2), fill=NAVY, width=4)
        draw.polygon([(x2, y2), (x2 - 14, y2 - 8), (x2 - 14, y2 + 8)], fill=NAVY)
    draw.line((864, 330, 864, 470), fill=NAVY, width=4)
    draw.polygon([(864, 470), (856, 456), (872, 456)], fill=NAVY)
    rail = (64, 686, 1532, 784)
    draw_card(draw, rail, fill=(255, 247, 247), outline=RED)
    draw_text(draw, (90, 710), "Target-only rail: occurrence labels + S_h/Sgh/Sh/NMR_SAT/hydrate saturation supervise training and validation only.", 22, RED, True, 1280)
    draw_text(draw, (90, 810), "Validation before metrics. Whole-well/geographic split before scaling. Reviewed outputs only.", 19, NAVY, True, 1340)
    img.save(PREVIEW_FILES[4])


def build_slide4(slide) -> None:
    add_ppt_header(slide, "Simplified ML Architecture", "Inputs become features only after schema, unit, QC, and leakage checks.")
    flow = [
        (0.54, 1.32, 2.22, 1.28, "Inputs", "logs\ncore/NMR\nlithology/QC\nstability context", TEAL),
        (3.20, 1.32, 2.22, 1.28, "Prepare", "headers\nunits\nalign by depth\nfingerprints", BLUE),
        (5.86, 1.32, 2.22, 1.28, "Leakage barrier", "X inputs only\nY targets locked away", RED),
        (8.52, 1.32, 3.82, 1.28, "Feature matrix", "measured logs + valid derived features\ncontext flags only", GREEN),
        (3.20, 4.15, 2.22, 1.20, "Baselines", "physics/rules\ntree/boosting", GRAY),
        (5.86, 4.15, 2.22, 1.20, "ANN / Keras later", "learned patterns after baseline checks", PURPLE),
        (8.52, 4.15, 3.82, 1.20, "Two outputs", "occurrence classification\nsaturation regression", AMBER),
    ]
    for x, y, w, h, head, body, color in flow:
        add_card(slide, Inches(x), Inches(y), Inches(w), Inches(h), outline=color)
        add_textbox(slide, Inches(x + 0.14), Inches(y + 0.15), Inches(w - 0.28), Inches(0.25), head, 12.5, color, True)
        add_textbox(slide, Inches(x + 0.14), Inches(y + 0.48), Inches(w - 0.28), Inches(h - 0.54), body, 9.3, NAVY, True)
    for x1, y1, x2, y2 in [(2.76, 1.96, 3.20, 1.96), (5.42, 1.96, 5.86, 1.96), (8.08, 1.96, 8.52, 1.96), (5.42, 4.75, 5.86, 4.75), (8.08, 4.75, 8.52, 4.75)]:
        add_line(slide, Inches(x1), Inches(y1), Inches(x2), Inches(y2), NAVY, 1.7)
    add_line(slide, Inches(6.97), Inches(2.60), Inches(6.97), Inches(4.15), NAVY, 1.7)
    add_card(slide, Inches(0.54), Inches(6.04), Inches(11.80), Inches(0.76), outline=RED, fill=(255, 247, 247))
    add_textbox(
        slide,
        Inches(0.74),
        Inches(6.25),
        Inches(11.40),
        Inches(0.28),
        "Target-only rail: occurrence labels + S_h/Sgh/Sh/NMR_SAT/hydrate saturation supervise training and validation only.",
        10.7,
        RED,
        True,
    )
    add_textbox(
        slide,
        Inches(0.74),
        Inches(7.03),
        Inches(11.40),
        Inches(0.30),
        "Validation before metrics. Whole-well/geographic split before scaling. Reviewed outputs only.",
        10.3,
        NAVY,
        True,
    )


def draw_preview_slide5() -> None:
    img = Image.new("RGB", (W, H), (248, 251, 252))
    draw = ImageDraw.Draw(img)
    preview_header(draw, "Slide 5: Equation Transformations", "Equations explain feature creation and stability context; they are not final proof.")
    cards = [
        ("Pressure-depth", "P_abs(z) = P_surface + rho_w g z / 1e6", "context", BLUE),
        ("P-T stability", "stable if T_model(z) <= T_eq(P_abs, CH4, salinity)", "context", TEAL),
        ("Density porosity", "phi_D = (rho_ma - RHOB) / (rho_ma - rho_f)", "log-derived", GREEN),
        ("Vp/Vs", "Vp/Vs = Vp / Vs", "elastic", PURPLE),
        ("Acoustic impedance", "AI = RHOB * Vp", "elastic", PURPLE),
        ("mu-rho", "mu-rho = RHOB * Vs^2", "elastic", AMBER),
        ("lambda-rho", "lambda-rho = RHOB * (Vp^2 - 2Vs^2)", "elastic", AMBER),
        ("Archie-style check", "Sw^n ~ a Rw / (phi^m Rt)", "reference only", RED),
    ]
    for idx, (head, eq, role, color) in enumerate(cards):
        col = idx % 4
        row = idx // 4
        x0 = 72 + col * 382
        y0 = 160 + row * 230
        draw_card(draw, (x0, y0, x0 + 336, y0 + 178), outline=color)
        draw_text(draw, (x0 + 18, y0 + 18), head, 21, color, True, 300)
        draw_text(draw, (x0 + 18, y0 + 62), eq, 18, NAVY, True, 300)
        draw_card(draw, (x0 + 18, y0 + 128, x0 + 150, y0 + 154), fill=(248, 251, 252), outline=color, width=1)
        draw_text(draw, (x0 + 30, y0 + 133), role, 13, color, True)
    draw_text(draw, (74, 678), "Readable equation cards; detailed citations and derivations stay in Word/end material.", 20, NAVY, True, 880)
    draw_text(draw, (74, 728), "Archie-style saturation relation is kept as a source-check/reference card until target role and units are verified.", 17, RED, True, 1120)
    img.save(PREVIEW_FILES[5])


def build_slide5(slide) -> None:
    add_ppt_header(slide, "Equation Transformations", "Feature engineering and stability context equations; no map, no ML diagram, no log panel.")
    cards = [
        ("Pressure-depth", "P_abs(z) = P_surface + rho_w g z / 1e6", "context", BLUE),
        ("P-T stability", "stable if T_model(z) <= T_eq(P_abs, CH4, salinity)", "context", TEAL),
        ("Density porosity", "phi_D = (rho_ma - RHOB) / (rho_ma - rho_f)", "log-derived", GREEN),
        ("Vp/Vs", "Vp/Vs = Vp / Vs", "elastic", PURPLE),
        ("Acoustic impedance", "AI = RHOB * Vp", "elastic", PURPLE),
        ("mu-rho", "mu-rho = RHOB * Vs^2", "elastic", AMBER),
        ("lambda-rho", "lambda-rho = RHOB * (Vp^2 - 2Vs^2)", "elastic", AMBER),
        ("Archie-style check", "Sw^n ~ a Rw / (phi^m Rt)", "reference only", RED),
    ]
    for idx, (head, eq, role, color) in enumerate(cards):
        col = idx % 4
        row = idx // 4
        x = 0.58 + col * 3.10
        y = 1.34 + row * 1.90
        add_card(slide, Inches(x), Inches(y), Inches(2.72), Inches(1.48), outline=color)
        add_textbox(slide, Inches(x + 0.12), Inches(y + 0.12), Inches(2.42), Inches(0.24), head, 11.5, color, True)
        add_textbox(slide, Inches(x + 0.12), Inches(y + 0.52), Inches(2.46), Inches(0.48), eq, 9.0, NAVY, True)
        add_chip(slide, x + 0.12, y + 1.10, 1.08, role, color)
    add_textbox(
        slide,
        Inches(0.62),
        Inches(5.46),
        Inches(7.50),
        Inches(0.28),
        "Readable equation cards; detailed citations and derivations stay in Word/end material.",
        10.2,
        NAVY,
        True,
    )
    add_textbox(
        slide,
        Inches(0.62),
        Inches(5.98),
        Inches(10.10),
        Inches(0.34),
        "Archie-style saturation relation is a source-check/reference card until target role and units are verified; saturation labels remain Y-only.",
        9.0,
        RED,
        True,
    )
    add_card(slide, Inches(9.18), Inches(5.50), Inches(3.10), Inches(0.92), outline=LINE, fill=(248, 251, 252))
    add_textbox(slide, Inches(9.34), Inches(5.66), Inches(2.78), Inches(0.24), "Color key", 9.5, NAVY, True)
    for i, (label, color) in enumerate([("logs/core", GREEN), ("context", TEAL), ("elastic", PURPLE), ("reference", RED)]):
        add_rect(slide, Inches(9.36 + (i % 2) * 1.38), Inches(5.98 + (i // 2) * 0.22), Inches(0.10), Inches(0.10), fill=color, outline=color)
        add_textbox(slide, Inches(9.50 + (i % 2) * 1.38), Inches(5.90 + (i // 2) * 0.22), Inches(1.12), Inches(0.20), label, 6.4, NAVY, True)


def draw_preview_slide6() -> None:
    img = Image.new("RGB", (W, H), (248, 251, 252))
    draw = ImageDraw.Draw(img)
    preview_header(draw, "Slide 6: Four-Well Evidence Review Board", "Review lanes show what must agree before any occurrence or saturation claim.")
    lanes = ["Lithology/core", "Logs", "QC", "Stability context", "Target-only review"]
    wells = ["Well A", "Well B", "Well C", "Well D"]
    x0, y0 = 190, 170
    cell_w, cell_h = 270, 94
    for i, lane in enumerate(lanes):
        y = y0 + i * cell_h
        draw_text(draw, (52, y + 26), lane, 17, NAVY, True, 116, align="right")
        for j, well in enumerate(wells):
            x = x0 + j * cell_w
            if i == 0:
                draw_text(draw, (x + 26, 132), well, 18, NAVY, True, 180, align="center")
            color = [GREEN, BLUE, RED, TEAL, AMBER][i]
            draw_card(draw, (x, y, x + 232, y + 68), fill=WHITE, outline=color, width=2)
            status = ["pending", "needed", "QC", "context", "Y-only"][i]
            draw_text(draw, (x + 18, y + 20), status, 17, color, True)
    side = (1290, 170, 1530, 640)
    draw_card(draw, side, outline=RED)
    draw_text(draw, (1312, 194), "Guardrail", 23, RED, True, 190)
    draw_text(draw, (1312, 244), "This board is a review checklist, not a result table.", 18, NAVY, True, 184)
    draw_text(draw, (1312, 390), "No final occurrence, saturation, ranking, or producibility claim.", 18, RED, True, 184)
    draw_card(draw, (60, 706, 1532, 790), fill=(236, 248, 250), outline=TEAL)
    draw_text(draw, (92, 732), "Takeaway: hydrate interpretation becomes stronger only when lithology, logs, QC, stability context, and target-only labels are reviewed together.", 22, NAVY, True, 1350)
    img.save(PREVIEW_FILES[6])


def build_slide6(slide) -> None:
    add_ppt_header(slide, "Four-Well Evidence Review Board", "High-level review lanes; details and citations stay in Word/speaker notes.")
    lanes = [
        ("Lithology/core", GREEN, "clean sand / core support"),
        ("Logs", BLUE, "GR + Rt + sonic + NMR/density"),
        ("QC", RED, "caliper/washout caution"),
        ("Stability context", TEAL, "admissible under assumptions"),
        ("Target-only review", AMBER, "S_h / NMR / phase labels"),
    ]
    wells = ["Well A", "Well B", "Well C", "Well D"]
    x0, y0 = 1.58, 1.42
    cell_w, cell_h = 2.12, 0.68
    for j, well in enumerate(wells):
        add_textbox(slide, Inches(x0 + j * cell_w + 0.08), Inches(1.14), Inches(1.70), Inches(0.22), well, 9.2, NAVY, True, PP_ALIGN.CENTER)
    for i, (lane, color, status) in enumerate(lanes):
        y = y0 + i * 0.76
        add_textbox(slide, Inches(0.34), Inches(y + 0.20), Inches(1.00), Inches(0.20), lane, 7.8, NAVY, True, PP_ALIGN.RIGHT)
        for j in range(4):
            x = x0 + j * cell_w
            add_card(slide, Inches(x), Inches(y), Inches(1.82), Inches(0.54), outline=color, fill=WHITE)
            add_textbox(slide, Inches(x + 0.10), Inches(y + 0.15), Inches(1.52), Inches(0.18), status, 6.9, color, True, PP_ALIGN.CENTER)
    add_card(slide, Inches(10.68), Inches(1.42), Inches(2.20), Inches(3.52), outline=RED)
    add_textbox(slide, Inches(10.86), Inches(1.62), Inches(1.84), Inches(0.26), "Guardrail", 12.5, RED, True)
    add_textbox(slide, Inches(10.86), Inches(2.08), Inches(1.82), Inches(0.82), "This board is a review checklist, not a result table.", 9.0, NAVY, True)
    add_textbox(slide, Inches(10.86), Inches(3.44), Inches(1.82), Inches(0.90), "No final occurrence, saturation, ranking, or producibility claim.", 9.0, RED, True)
    add_card(slide, Inches(0.54), Inches(6.02), Inches(11.90), Inches(0.70), outline=TEAL, fill=(236, 248, 250))
    add_textbox(
        slide,
        Inches(0.74),
        Inches(6.22),
        Inches(11.46),
        Inches(0.24),
        "Takeaway: hydrate interpretation becomes stronger only when lithology, logs, QC, stability context, and target-only labels are reviewed together.",
        10.3,
        NAVY,
        True,
    )


def draw_preview_slide7() -> None:
    img = Image.new("RGB", (W, H), (248, 251, 252))
    draw = ImageDraw.Draw(img)
    preview_header(draw, "Slide 7: North Slope Stability Context Map", "Unified 2D public context map for slide callouts; not an ML overlay.")
    map_box = (60, 142, 1138, 746)
    draw_card(draw, (48, 130, 1150, 780), outline=LINE)
    paste_contain(img, MAP_SOURCE, map_box)
    side = [
        ((1190, 148, 1534, 236), "Public layers", "geology, borough edge, roads/TAPS/fields", TEAL),
        ((1190, 260, 1534, 348), "Controls", "GGD223 permafrost + public stability-screen status", BLUE),
        ((1190, 372, 1534, 460), "Hydrate AUs", "regional assessment context only", GREEN),
        ((1190, 484, 1534, 572), "Use", "orientation and caveat support", AMBER),
        ((1190, 596, 1534, 704), "Not a result", "not occurrence, saturation, producibility, or ranking", RED),
    ]
    for box, head, body, color in side:
        draw_card(draw, box, outline=color)
        draw_text(draw, (box[0] + 16, box[1] + 10), head, 18, color, True, box[2] - box[0] - 32)
        draw_text(draw, (box[0] + 16, box[1] + 38), body, 14, NAVY, True, box[2] - box[0] - 32)
    draw_text(draw, (68, 806), "Caption: public stability-admissibility context under explicit assumptions; source layers are public/regional and do not prove hydrate.", 18, NAVY, True, 1210)
    img.save(PREVIEW_FILES[7])


def build_slide7(slide) -> None:
    add_ppt_header(slide, "North Slope Stability Context Map", "Unified 2D public context map for slide callouts; not an ML overlay.")
    add_card(slide, Inches(0.42), Inches(1.16), Inches(9.30), Inches(5.50), outline=LINE)
    if MAP_SOURCE.exists():
        slide.shapes.add_picture(str(MAP_SOURCE), Inches(0.54), Inches(1.28), width=Inches(9.06), height=Inches(4.98))
    else:
        add_textbox(slide, Inches(0.84), Inches(3.22), Inches(8.40), Inches(0.32), f"Missing map source: {MAP_SOURCE.name}", 13, RED, True, PP_ALIGN.CENTER)
    side = [
        (10.06, 1.28, 2.74, 0.68, "Public layers", "geology, borough edge, roads/TAPS/fields", TEAL),
        (10.06, 2.16, 2.74, 0.68, "Controls", "GGD223 permafrost + public screen status", BLUE),
        (10.06, 3.04, 2.74, 0.68, "Hydrate AUs", "regional assessment context only", GREEN),
        (10.06, 3.92, 2.74, 0.68, "Use", "orientation and caveat support", AMBER),
        (10.06, 4.80, 2.74, 0.88, "Not a result", "not occurrence, saturation, producibility, or ranking", RED),
    ]
    for x, y, w, h, head, body, color in side:
        add_card(slide, Inches(x), Inches(y), Inches(w), Inches(h), outline=color)
        add_textbox(slide, Inches(x + 0.10), Inches(y + 0.08), Inches(w - 0.20), Inches(0.18), head, 8.6, color, True)
        add_textbox(slide, Inches(x + 0.10), Inches(y + 0.30), Inches(w - 0.20), Inches(h - 0.32), body, 7.1, NAVY, True)
    add_textbox(
        slide,
        Inches(0.54),
        Inches(6.86),
        Inches(10.50),
        Inches(0.34),
        "Caption: public stability-admissibility context under explicit assumptions; source layers are public/regional and do not prove hydrate.",
        8.8,
        NAVY,
        True,
    )
    add_textbox(slide, Inches(10.92), Inches(6.86), Inches(1.88), Inches(0.30), "callouts are editable", 8.0, MUTED, True, PP_ALIGN.RIGHT)


def build_previews() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    draw_preview_slide3()
    draw_preview_slide4()
    draw_preview_slide5()
    draw_preview_slide6()
    draw_preview_slide7()
    thumbs = []
    for slide_num in range(3, 8):
        img = Image.open(PREVIEW_FILES[slide_num]).convert("RGB")
        thumbs.append(img.resize((480, 270), Image.Resampling.LANCZOS))
    sheet = Image.new("RGB", (1540, 690), (248, 251, 252))
    draw = ImageDraw.Draw(sheet)
    draw_text(draw, (34, 24), "Editable Slides 3-7 Candidate Contact Sheet", 32, NAVY, True)
    for idx, thumb in enumerate(thumbs):
        col = idx % 3
        row = idx // 3
        x = 34 + col * 500
        y = 86 + row * 310
        sheet.paste(thumb, (x, y))
        draw.rectangle((x, y, x + 480, y + 270), outline=LINE, width=2)
        draw_text(draw, (x, y + 276), f"Slide {idx + 3}", 18, NAVY, True)
    sheet.save(CONTACT_SHEET)


def build_deck() -> None:
    source = SLIDE2_CANDIDATE_DECK if SLIDE2_CANDIDATE_DECK.exists() else SOURCE_UPDATE_DECK
    if not source.exists():
        raise FileNotFoundError(f"No source deck found: {SLIDE2_CANDIDATE_DECK} or {SOURCE_UPDATE_DECK}")
    prs = Presentation(str(source))
    if len(prs.slides) < 7:
        raise ValueError(f"Expected at least seven slides in {source}; found {len(prs.slides)}")
    builders = {
        2: build_slide3,
        3: build_slide4,
        4: build_slide5,
        5: build_slide6,
        6: build_slide7,
    }
    for zero_index, builder in builders.items():
        slide = prs.slides[zero_index]
        clear_slide(slide)
        builder(slide)
    prs.save(OUT_DECK)


def audit_deck() -> None:
    prs = Presentation(str(OUT_DECK))
    rows = []
    roles = {
        3: "log signal, lithology, and core/NMR calibration scaffold",
        4: "simplified ML architecture",
        5: "equation transformations",
        6: "four-well evidence review board",
        7: "stability context map",
    }
    for slide_number in range(3, 8):
        slide = prs.slides[slide_number - 1]
        shape_count = len(slide.shapes)
        picture_count = sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
        text_count = sum(1 for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip())
        line_count = sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.LINE)
        one_picture_full_slide_like = picture_count == 1 and shape_count <= 3
        rows.append(
            {
                "slide_number": slide_number,
                "slide_role": roles[slide_number],
                "shape_count": shape_count,
                "picture_shapes": picture_count,
                "text_shapes": text_count,
                "line_or_connector_shapes": line_count,
                "raster_only_flag": "yes" if one_picture_full_slide_like else "no",
                "manual_editability_status": "editable_candidate",
                "notes": "Source images may remain raster, but labels, callouts, captions, lanes, and equation cards are native PowerPoint objects.",
            }
        )
    with AUDIT_CSV.open("w", newline="", encoding="utf-8") as fh:
        writer = csv.DictWriter(fh, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        writer.writerows(rows)


def write_source_notes() -> None:
    SOURCE_NOTES.write_text(
        "\n".join(
            [
                "# Editable Slides 3-7 Source And Claim Notes",
                "",
                "Generated: 2026-06-19",
                "",
                "## Scope",
                "",
                "This package executes delegated Prompts 18-20 as an editable candidate deck section. It does not replace the active V5.5 mentor deck until reviewed.",
                "",
                "## Public-Safe Boundary",
                "",
                "- Slide 3 uses a schematic well-log scaffold because reviewed four-well DOE exports are not committed.",
                "- Slide 6 uses anonymous Well A-D placeholders until names, locations, core/NMR, and lithology evidence are public-safe for display.",
                "- No approved rows, private depths, row-level predictions, fitted models, trained metrics, or sensitive identifiers are included.",
                "",
                "## Slide Claim Controls",
                "",
                "- Slide 3: log signals are complementary evidence cues only; caliper/washout is QC only.",
                "- Slide 4: target-only fields supervise training/validation and bypass X predictors.",
                "- Slide 5: equations are transformations or screening context; Archie-style saturation logic remains a reference card until role and unit policy are verified.",
                "- Slide 6: the evidence board is a review checklist, not a result table.",
                "- Slide 7: the map is public stability-admissibility context, not hydrate proof, occurrence, saturation, producibility, or ranking.",
                "",
                "## Source Anchors",
                "",
                "- `docs/APPROVED_DATA_SCHEMA_COVERAGE_AND_MODEL_ARCHITECTURE_PLAN.md` for schema roles and leakage barrier.",
                "- `docs/FIRST_MODEL_EXPERIMENT_PLAN_2026-06-15.md` for occurrence/saturation split and validation-before-metrics logic.",
                "- `docs/DOE_THREE_DATASET_ML_PIPELINE_RUNBOOK_2026-06-16.md` for approved-runtime and output-boundary rules.",
                "- `docs/STABILITY_CALCULATION_PLAN.md` for pressure-temperature context and stability-not-proof guardrails.",
                "- `data/public_ml_products/public_parameter_evidence_registry_2026-06-16.csv` for parameter family directions and caveats.",
                "- `docs/project_blueprints/presentation_assets/website_well_maps_2026_06_18/unified_north_slope_slide_export_callout_space_2026_06_18.png` for Slide 7 map image.",
                "",
            ]
        ),
        encoding="utf-8",
    )


def scan_claim_text() -> None:
    prs = Presentation(str(OUT_DECK))
    forbidden = ["project website", "P-T gate", "final prediction", "validated prediction"]
    text_blob = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    hits = [term for term in forbidden if term.lower() in text_blob.lower()]
    if hits:
        raise ValueError(f"Forbidden slide wording found: {hits}")


def main() -> None:
    build_previews()
    build_deck()
    audit_deck()
    write_source_notes()
    scan_claim_text()
    print(f"Wrote {OUT_DECK}")
    print(f"Wrote {CONTACT_SHEET}")
    print(f"Wrote {AUDIT_CSV}")
    print(f"Wrote {SOURCE_NOTES}")
    for slide_num, path in PREVIEW_FILES.items():
        print(f"Wrote slide {slide_num}: {path}")


if __name__ == "__main__":
    main()
