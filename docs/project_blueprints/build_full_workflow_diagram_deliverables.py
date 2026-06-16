from __future__ import annotations

import csv
import hashlib
import math
from io import BytesIO
from pathlib import Path

from docx import Document
from docx.enum.section import WD_ORIENT
from docx.shared import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
BLUEPRINT_DIR = ROOT / "docs" / "project_blueprints"
ASSET_DIR = BLUEPRINT_DIR / "presentation_assets" / "v5_4_corrected_2026_06_16"
SOURCE_DECK = BLUEPRINT_DIR / "CURRENT_GMAIL_VISUAL_REVISION_9_SLIDE_North_Slope_Gas_Hydrate_Slides_2026-06-11.pptx"
OUT_DECK = BLUEPRINT_DIR / "V5_4_CORRECTED_North_Slope_Gas_Hydrate_ML_Workflow_Slides_2026-06-16.pptx"
OUT_DOCX = BLUEPRINT_DIR / "V5_4_CORRECTED_North_Slope_Gas_Hydrate_ML_Workflow_Companion_2026-06-16.docx"
OUT_CONTACT_SHEET = ASSET_DIR / "v5_4_corrected_contact_sheet.png"
PUBLIC_PRODUCTS = ROOT / "data" / "public_stability_products"
PUBLIC_ML_PRODUCTS = ROOT / "data" / "public_ml_products"
WEBSITE_CAPTURE_DIR = BLUEPRINT_DIR / "presentation_assets" / "v5_3_website_captures"
REFERENCE_IMAGE_DIR = ROOT / "references" / "presentation-revision-2026-06-11" / "images"
V52_ASSET_DIR = BLUEPRINT_DIR / "presentation_assets" / "full_workflow_diagram_2026_06_15"
PROCESSING_ASSET_DIR = BLUEPRINT_DIR / "presentation_assets" / "processing_revisions_2026_06_11"
PARAMETER_EVIDENCE_REGISTRY = PUBLIC_ML_PRODUCTS / "public_parameter_evidence_registry_2026-06-16.csv"
PHASE_CURVE_METHANE_5PPT = PUBLIC_PRODUCTS / "phase_curve_methane_5ppt_sir2008_csmhyd_digitized_v1.csv"

W, H = 1600, 900
EXPANDED_W, EXPANDED_H = 5200, 3000
NETWORK_W, NETWORK_H = 2600, 1500
NAVY = (11, 35, 48)
DEEP = (5, 19, 27)
TEAL = (18, 124, 139)
ICE = (82, 185, 214)
ICE_LIGHT = (223, 244, 248)
GREEN = (37, 154, 113)
GREEN_LIGHT = (224, 246, 236)
AMBER = (219, 159, 52)
AMBER_LIGHT = (250, 239, 211)
RED = (200, 62, 65)
RED_LIGHT = (252, 229, 230)
BLUE = (55, 133, 203)
BLUE_LIGHT = (228, 241, 253)
PURPLE = (111, 105, 190)
PURPLE_LIGHT = (236, 234, 250)
MUTED = (82, 103, 112)
LIGHT = (244, 249, 250)
LINE = (181, 210, 217)
WHITE = (255, 255, 255)
LOCKED_SLIDE_HASHES: dict[int, str] = {}


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        "arialbd.ttf" if bold else "arial.ttf",
        "segoeuib.ttf" if bold else "segoeui.ttf",
        "calibrib.ttf" if bold else "calibri.ttf",
    ]
    roots = [Path("C:/Windows/Fonts"), Path("/usr/share/fonts/truetype/dejavu")]
    for name in candidates:
        for root in roots:
            path = root / name
            if path.exists():
                return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def wrap_lines(draw: ImageDraw.ImageDraw, text: str, size: int, width: int, bold: bool = False) -> list[str]:
    f = font(size, bold)
    lines: list[str] = []
    for raw in text.splitlines():
        if not raw:
            lines.append("")
            continue
        current = ""
        for word in raw.split():
            test = f"{current} {word}".strip()
            if draw.textlength(test, font=f) <= width or not current:
                current = test
            else:
                lines.append(current)
                current = word
        if current:
            lines.append(current)
    return lines


def text(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    value: str,
    size: int,
    fill: tuple[int, int, int] = NAVY,
    bold: bool = False,
    width: int | None = None,
    align: str = "left",
    gap: int = 5,
) -> int:
    f = font(size, bold)
    if width is None:
        draw.text(xy, value, font=f, fill=fill)
        return xy[1] + size + gap
    y = xy[1]
    for line in wrap_lines(draw, value, size, width, bold):
        if align == "center":
            x = xy[0] + width // 2 - int(draw.textlength(line, font=f) // 2)
        elif align == "right":
            x = xy[0] + width - int(draw.textlength(line, font=f))
        else:
            x = xy[0]
        draw.text((x, y), line, font=f, fill=fill)
        y += size + gap
    return y


def card(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    fill: tuple[int, int, int] = WHITE,
    outline: tuple[int, int, int] = LINE,
    radius: int = 18,
    width: int = 2,
) -> None:
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def arrow(
    draw: ImageDraw.ImageDraw,
    start: tuple[int, int],
    end: tuple[int, int],
    fill: tuple[int, int, int] = TEAL,
    width: int = 4,
    label: str | None = None,
    label_offset: tuple[int, int] = (0, 0),
) -> None:
    draw.line((start, end), fill=fill, width=width)
    angle = math.atan2(end[1] - start[1], end[0] - start[0])
    length = 15
    points = [
        end,
        (int(end[0] - length * math.cos(angle - 0.45)), int(end[1] - length * math.sin(angle - 0.45))),
        (int(end[0] - length * math.cos(angle + 0.45)), int(end[1] - length * math.sin(angle + 0.45))),
    ]
    draw.polygon(points, fill=fill)
    if label:
        mx = (start[0] + end[0]) // 2 + label_offset[0]
        my = (start[1] + end[1]) // 2 + label_offset[1]
        f = font(12, True)
        pad = 5
        w = int(draw.textlength(label, font=f)) + pad * 2
        card(draw, (mx - w // 2, my - 14, mx + w // 2, my + 12), fill=WHITE, outline=fill, radius=9, width=1)
        draw.text((mx - w // 2 + pad, my - 10), label, font=f, fill=fill)


def pill(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    label: str,
    fill: tuple[int, int, int],
    text_fill: tuple[int, int, int] = NAVY,
) -> None:
    card(draw, box, fill=fill, outline=fill, radius=14, width=1)
    f = font(13, True)
    x = box[0] + (box[2] - box[0]) // 2 - int(draw.textlength(label, font=f) // 2)
    y = box[1] + (box[3] - box[1] - 13) // 2 - 2
    draw.text((x, y), label, font=f, fill=text_fill)


def node(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    heading: str,
    lines: list[str],
    fill: tuple[int, int, int],
    accent: tuple[int, int, int],
    tag: str | None = None,
) -> None:
    card(draw, box, fill=fill, outline=accent, radius=18, width=2)
    draw.rounded_rectangle((box[0], box[1], box[0] + 10, box[3]), radius=10, fill=accent)
    text(draw, (box[0] + 24, box[1] + 15), heading, 18, accent, True, width=box[2] - box[0] - 42)
    y = box[1] + 48
    for item in lines:
        y = text(draw, (box[0] + 26, y), item, 13, NAVY, width=box[2] - box[0] - 45, gap=4)
    if tag:
        pill(draw, (box[0] + 24, box[3] - 34, box[2] - 24, box[3] - 10), tag, WHITE, accent)


def canvas(dark: bool = False) -> Image.Image:
    img = Image.new("RGB", (W, H), DEEP if dark else WHITE)
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 18, H), fill=ICE if dark else TEAL)
    return img


def title(draw: ImageDraw.ImageDraw, heading: str, subheading: str, dark: bool = False) -> None:
    text(draw, (58, 38), heading, 38, WHITE if dark else NAVY, True)
    text(draw, (61, 88), subheading, 17, (196, 229, 235) if dark else MUTED, width=1340)


def footer(draw: ImageDraw.ImageDraw, value: str, dark: bool = False) -> None:
    y = H - 50
    draw.line((54, y - 8, W - 54, y - 8), fill=(73, 96, 105) if dark else LINE, width=2)
    text(draw, (58, y), value, 12, (168, 195, 204) if dark else MUTED, width=1420)


def read_summary(name: str) -> dict[str, str]:
    path = PUBLIC_PRODUCTS / name
    if not path.exists():
        return {}
    with path.open(newline="", encoding="utf-8") as fh:
        return {row["metric"]: row["value"] for row in csv.DictReader(fh)}


def fmt_count(summary: dict[str, str], key: str, fallback: str) -> str:
    value = summary.get(key, fallback)
    try:
        number = float(value)
    except ValueError:
        return value
    if number.is_integer():
        return f"{int(number):,}"
    return f"{number:,.1f}"


def summaries() -> dict[str, str]:
    well = read_summary("north_slope_well_stability_context_summary_2026-06-14.csv")
    temp = read_summary("g10015_temperature_profile_summary_2026-06-14.csv")
    model = read_summary("stability_temperature_model_summary_2026-06-14.csv")
    screen = read_summary("stability_screen_summary_2026-06-14_methane_5ppt_v1.csv")
    features = read_summary("public_ml_feature_scaffold_summary_2026-06-15.csv")
    return {
        "wells": fmt_count(well, "Arctic Slope public wells", "8,084"),
        "ggd_controls": "43",
        "profiles": fmt_count(temp, "G10015 profiles", "184"),
        "codes": fmt_count(temp, "Unique well codes", "24"),
        "hydrate_aus": "3",
        "temp_matches": fmt_count(features, "Rows with matched temperature profile", "483"),
        "temp_rows": fmt_count(model, "Temperature model rows", "16,168"),
        "temp_calculated": fmt_count(model, "Calculated key depths", "919"),
        "temp_extrapolated": fmt_count(model, "Extrapolated key depths", "387"),
        "temp_blocked": fmt_count(model, "Blocked key depths", "15,249"),
        "screen_rows": fmt_count(screen, "Screen rows", "8,084"),
        "screen_calculated": fmt_count(screen, "Calculated stability intervals", "22"),
        "screen_no_interval": fmt_count(screen, "No stable interval found", "8"),
        "screen_blocked": fmt_count(screen, "Blocked rows", "8,054"),
        "training_ready": fmt_count(features, "Rows training-ready for occurrence/saturation ML", "0"),
        "occurrence_labels": fmt_count(features, "Rows with validated hydrate occurrence labels", "0"),
        "approved_visible": "about 3 / 71",
    }


def save(img: Image.Image, name: str) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / name
    img.save(path, quality=94)
    return path


def dashed_line(
    draw: ImageDraw.ImageDraw,
    start: tuple[int, int],
    end: tuple[int, int],
    fill: tuple[int, int, int],
    width: int = 3,
    dash: int = 18,
    gap: int = 12,
) -> None:
    dx, dy = end[0] - start[0], end[1] - start[1]
    length = math.hypot(dx, dy)
    if length == 0:
        return
    ux, uy = dx / length, dy / length
    distance = 0.0
    while distance < length:
        line_end = min(distance + dash, length)
        draw.line(
            (
                int(start[0] + ux * distance),
                int(start[1] + uy * distance),
                int(start[0] + ux * line_end),
                int(start[1] + uy * line_end),
            ),
            fill=fill,
            width=width,
        )
        distance += dash + gap


def workflow_slide_summary_panel(values: dict[str, str]) -> Path:
    """Build a readable 16:9 summary for slide use; the poster stays separate."""
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    title(
        draw,
        "Full Project ML Workflow V5.2",
        "Mentor-scale map: public source context and approved OSL data feed a leakage-safe path to occurrence and saturation outputs.",
    )

    def arrow_between(start: tuple[int, int], end: tuple[int, int], color=TEAL, width: int = 4, dashed: bool = False) -> None:
        if dashed:
            dashed_line(draw, start, end, color, width=width, dash=20, gap=13)
        else:
            draw.line((start, end), fill=color, width=width)
        angle = math.atan2(end[1] - start[1], end[0] - start[0])
        size = 15
        head = [
            end,
            (int(end[0] - size * math.cos(angle - 0.45)), int(end[1] - size * math.sin(angle - 0.45))),
            (int(end[0] - size * math.cos(angle + 0.45)), int(end[1] - size * math.sin(angle + 0.45))),
        ]
        draw.polygon(head, fill=color)

    def icon(box: tuple[int, int, int, int], kind: str, accent) -> None:
        x1, y1, x2, y2 = box
        if kind == "sources":
            for idx, color in enumerate([ICE_LIGHT, BLUE_LIGHT, GREEN_LIGHT]):
                ox, oy = idx * 11, idx * 9
                card(draw, (x1 + ox, y1 + oy, x1 + 82 + ox, y1 + 58 + oy), fill=color, outline=accent, radius=8, width=1)
                draw.rectangle((x1 + 14 + ox, y1 + 16 + oy, x1 + 68 + ox, y1 + 21 + oy), fill=accent)
                draw.rectangle((x1 + 14 + ox, y1 + 32 + oy, x1 + 56 + ox, y1 + 37 + oy), fill=(151, 183, 191))
            for px, py, color in [(x2 - 78, y1 + 58, TEAL), (x2 - 36, y1 + 30, GREEN), (x2 - 18, y2 - 18, AMBER)]:
                draw.ellipse((px - 7, py - 7, px + 7, py + 7), fill=color, outline=WHITE, width=2)
            draw.line((x2 - 78, y1 + 58, x2 - 36, y1 + 30, x2 - 18, y2 - 18), fill=LINE, width=3)
        elif kind == "stability":
            ax1, ay1, ax2, ay2 = x1 + 10, y1 + 8, x2 - 10, y2 - 10
            draw.line((ax1, ay2, ax2, ay2), fill=NAVY, width=2)
            draw.line((ax1, ay2, ax1, ay1), fill=NAVY, width=2)
            phase = [(ax1 + 4, ay2 - 12), (ax1 + 48, ay2 - 36), (ax1 + 98, ay2 - 70), (ax2 - 8, ay1 + 13)]
            temp = [(ax1 + 6, ay1 + 10), (ax1 + 46, ay1 + 50), (ax1 + 88, ay2 - 32), (ax2 - 20, ay2 - 18)]
            draw.line(phase, fill=AMBER, width=4)
            draw.line(temp, fill=BLUE, width=4)
        elif kind == "features":
            colors = [GREEN, TEAL, BLUE, PURPLE]
            for idx, color in enumerate(colors):
                tx = x1 + idx * 35
                draw.rectangle((tx, y1 + 8, tx + 24, y2 - 8), outline=(219, 233, 237), width=2)
                pts = []
                for j in range(18):
                    yy = y1 + 15 + j * ((y2 - y1 - 30) / 17)
                    xx = tx + 5 + int(14 * (0.5 + 0.4 * math.sin(j * 0.7 + idx)))
                    pts.append((xx, int(yy)))
                draw.line(pts, fill=color, width=2)
            mx = x2 - 78
            for r in range(5):
                for c in range(4):
                    draw.rectangle((mx + c * 15, y1 + 20 + r * 15, mx + c * 15 + 12, y1 + 32 + r * 15), fill=ICE_LIGHT, outline=WHITE)
        elif kind == "model":
            layers = [(x1 + 20, 4, GREEN), (x1 + 70, 5, PURPLE), (x1 + 120, 4, PURPLE), (x2 - 26, 2, AMBER)]
            pts_by_layer = []
            for lx, count, color in layers:
                pts = [(lx, int(y1 + 15 + i * ((y2 - y1 - 30) / max(1, count - 1)))) for i in range(count)]
                pts_by_layer.append(pts)
            for a_layer, b_layer in zip(pts_by_layer, pts_by_layer[1:], strict=False):
                for a in a_layer[::2]:
                    for b in b_layer[::2]:
                        draw.line((a, b), fill=(207, 221, 226), width=1)
            for pts, (_, _, color) in zip(pts_by_layer, layers, strict=False):
                for px, py in pts:
                    draw.ellipse((px - 8, py - 8, px + 8, py + 8), fill=WHITE, outline=color, width=3)
        elif kind == "outputs":
            card(draw, (x1 + 2, y1 + 10, x1 + 70, y2 - 8), fill=BLUE_LIGHT, outline=BLUE, radius=8, width=1)
            for idx, h in enumerate([32, 46, 25, 56]):
                bx = x1 + 16 + idx * 12
                draw.rectangle((bx, y2 - 18 - h, bx + 7, y2 - 18), fill=BLUE)
            card(draw, (x1 + 88, y1 + 10, x2 - 2, y2 - 8), fill=GREEN_LIGHT, outline=GREEN, radius=8, width=1)
            pts = [(x1 + 104, y2 - 22), (x1 + 128, y2 - 50), (x1 + 155, y2 - 38), (x2 - 12, y1 + 28)]
            draw.line(pts, fill=GREEN, width=3)
            for px, py in pts:
                draw.ellipse((px - 5, py - 5, px + 5, py + 5), fill=GREEN, outline=WHITE, width=2)

    def step_card(
        box: tuple[int, int, int, int],
        heading: str,
        rows: list[str],
        accent,
        fill,
        kind: str,
        tag: str,
    ) -> None:
        card(draw, box, fill=fill, outline=accent, radius=18, width=2)
        draw.rectangle((box[0] + 18, box[1] + 18, box[0] + 28, box[1] + 64), fill=accent)
        text(draw, (box[0] + 42, box[1] + 18), heading, 21, accent, True, width=box[2] - box[0] - 58)
        icon((box[0] + 36, box[1] + 62, box[2] - 34, box[1] + 155), kind, accent)
        y = box[1] + 170
        for row in rows:
            y = text(draw, (box[0] + 32, y), row, 16, NAVY, width=box[2] - box[0] - 64, gap=5)
        pill(draw, (box[0] + 32, box[3] - 42, box[2] - 32, box[3] - 14), tag, WHITE, accent)

    top_y = 160
    cards = [
        ((58, top_y, 323, 555), "Source And Schema Controls", ["Public: DNR wells, GGD223, G10015, USGS AU.", "Approved later: LAS, core, NMR, workbook labels.", "Variable fingerprints preserve units and roles."], TEAL, ICE_LIGHT, "sources", "public vs OSL boundary"),
        ((368, top_y, 633, 555), "Stability Context", ["Hydrostatic pressure + G10015 temperature.", "Methane 5 ppt phase lookup.", "Outputs status, interval, confidence, and blocked reason."], BLUE, BLUE_LIGHT, "stability", "admissibility only"),
        ((678, top_y, 943, 555), "Feature Engineering", ["Measured logs and core context.", "Depth is alignment/context unless approved.", "Caliper coverage comes before washout filtering."], GREEN, GREEN_LIGHT, "features", "X allowed only"),
        ((988, top_y, 1253, 555), "Leakage-Safe ML", ["Target labels bypass predictors.", "Split before 0-1 scaling and selection.", "Baselines first; tree/boosting second; ANN third."], PURPLE, PURPLE_LIGHT, "model", "split before fit"),
        ((1298, top_y, 1563, 555), "Reviewed Outputs", ["Occurrence probability P(hydrate).", "Saturation estimate Sh pred.", "Uncertainty, QC, mimic, caveat, and public-safe export review."], AMBER, AMBER_LIGHT, "outputs", "validated later"),
    ]

    for box, heading, rows, accent, fill, kind, tag in cards:
        step_card(box, heading, rows, accent, fill, kind, tag)
    for left_box, right_box, color in [
        (cards[0][0], cards[1][0], TEAL),
        (cards[1][0], cards[2][0], BLUE),
        (cards[2][0], cards[3][0], GREEN),
        (cards[3][0], cards[4][0], PURPLE),
    ]:
        arrow_between((left_box[2] + 6, (left_box[1] + left_box[3]) // 2), (right_box[0] - 8, (right_box[1] + right_box[3]) // 2), color)

    # Target-only rail and public/runtime boundary.
    target_box = (988, 588, 1253, 665)
    card(draw, target_box, fill=RED_LIGHT, outline=RED, radius=16, width=2)
    text(draw, (target_box[0] + 24, target_box[1] + 13), "Target-only labels", 20, RED, True)
    text(draw, (target_box[0] + 24, target_box[1] + 41), "Sgh, Sh, NMR SAT, phase labels -> Y only", 15, NAVY, width=215)
    arrow_between((target_box[2], 626), (1340, 626), RED, width=3, dashed=True)
    text(draw, (1288, 585), "training / validation overlay", 15, RED, True, width=230)

    boundary = (58, 692, 1563, 795)
    card(draw, boundary, fill=LIGHT, outline=LINE, radius=18, width=2)
    text(draw, (92, 718), "Data boundary", 20, TEAL, True)
    text(
        draw,
        (262, 709),
        "GitHub/Streamlit shows public methods, source-backed scaffold counts, diagrams, and public-safe summaries. OSL or the approved runtime loads real logs, core/NMR, labels, feature rows, model fitting, and reviewed outputs.",
        21,
        NAVY,
        True,
        width=1220,
        gap=5,
    )

    stats = (
        f"Current public screen: {values['screen_rows']} rows, {values['screen_calculated']} calculated admissibility intervals, "
        f"{values['screen_no_interval']} no-stable-interval rows, {values['screen_blocked']} blocked rows."
    )
    footer(draw, stats + " Stability is context or mask only, not hydrate proof or saturation.")
    return save(img, "full_project_ml_workflow_flowchart.png")


def full_workflow_panel() -> Path:
    values = summaries()
    img = Image.new("RGB", (EXPANDED_W, EXPANDED_H), WHITE)
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 34, EXPANDED_H), fill=TEAL)
    draw.rectangle((0, 0, EXPANDED_W, 214), fill=(246, 251, 252))
    draw.line((84, 214, EXPANDED_W - 84, 214), fill=LINE, width=4)
    text(draw, (92, 48), "North Slope Gas Hydrate ML Workflow V5.2", 64, NAVY, True)
    text(
        draw,
        (96, 126),
        "Visual source-to-model architecture: data controls, stability physics, log-derived features, target-safe ML, validation, and reviewed outputs.",
        31,
        MUTED,
        width=4300,
    )

    boxes: dict[str, tuple[int, int, int, int]] = {}
    col_x = [100, 1120, 2140, 3160, 4180]
    col_w = 900
    section_top = 280
    section_bottom = 2740

    def pill_text(box: tuple[int, int, int, int], label: str, fill, accent, size: int = 24) -> None:
        card(draw, box, fill=fill, outline=accent, radius=18, width=2)
        f = font(size, True)
        x = box[0] + (box[2] - box[0]) // 2 - int(draw.textlength(label, font=f) // 2)
        draw.text((x, box[1] + 9), label, font=f, fill=accent)

    def section(idx: int, label: str, accent) -> None:
        x = col_x[idx]
        card(draw, (x, section_top, x + col_w, section_bottom), fill=(249, 252, 252), outline=(214, 230, 234), radius=28, width=2)
        draw.rectangle((x + 28, section_top + 28, x + 42, section_top + 82), fill=accent)
        text(draw, (x + 60, section_top + 28), label, 31, accent, True, width=col_w - 90, gap=5)

    def node(
        key: str,
        col: int,
        y: int,
        h: int,
        heading: str,
        rows: list[str],
        accent,
        fill=WHITE,
        tag: str | None = None,
        formula: bool = False,
    ) -> None:
        x = col_x[col] + 28
        box = (x, y, x + col_w - 56, y + h)
        boxes[key] = box
        card(draw, box, fill=fill, outline=accent, radius=20, width=3)
        draw.rounded_rectangle((box[0], box[1], box[0] + 13, box[3]), radius=10, fill=accent)
        y2 = text(draw, (box[0] + 32, box[1] + 16), heading, 28, accent, True, width=box[2] - box[0] - 55, gap=5)
        y2 += 4
        body_size = 22 if not formula else 21
        for row in rows:
            if formula:
                card(draw, (box[0] + 32, y2, box[2] - 22, y2 + 38), fill=(247, 251, 252), outline=(224, 235, 238), radius=8, width=1)
                y2 = text(draw, (box[0] + 48, y2 + 7), row, body_size, NAVY, True, width=box[2] - box[0] - 82, gap=2)
                y2 += 9
            else:
                y2 = text(draw, (box[0] + 34, y2), row, body_size, NAVY, width=box[2] - box[0] - 62, gap=5)
        if tag:
            pill_text((box[0] + 32, box[3] - 45, box[2] - 32, box[3] - 10), tag, WHITE, accent, 21)

    def left(key: str) -> tuple[int, int]:
        b = boxes[key]
        return (b[0], (b[1] + b[3]) // 2)

    def right(key: str) -> tuple[int, int]:
        b = boxes[key]
        return (b[2], (b[1] + b[3]) // 2)

    def top(key: str) -> tuple[int, int]:
        b = boxes[key]
        return ((b[0] + b[2]) // 2, b[1])

    def bottom(key: str) -> tuple[int, int]:
        b = boxes[key]
        return ((b[0] + b[2]) // 2, b[3])

    def arrow_head(end: tuple[int, int], prev: tuple[int, int], fill) -> None:
        angle = math.atan2(end[1] - prev[1], end[0] - prev[0])
        length = 26
        points = [
            end,
            (int(end[0] - length * math.cos(angle - 0.45)), int(end[1] - length * math.sin(angle - 0.45))),
            (int(end[0] - length * math.cos(angle + 0.45)), int(end[1] - length * math.sin(angle + 0.45))),
        ]
        draw.polygon(points, fill=fill)

    def dashed_segment(p1: tuple[int, int], p2: tuple[int, int], fill, width: int) -> None:
        dx, dy = p2[0] - p1[0], p2[1] - p1[1]
        length = math.hypot(dx, dy)
        if length == 0:
            return
        ux, uy = dx / length, dy / length
        dash, gap = 28, 16
        distance = 0.0
        while distance < length:
            start = min(distance, length)
            end = min(distance + dash, length)
            draw.line(
                (
                    int(p1[0] + ux * start),
                    int(p1[1] + uy * start),
                    int(p1[0] + ux * end),
                    int(p1[1] + uy * end),
                ),
                fill=fill,
                width=width,
            )
            distance += dash + gap

    def label(text_value: str, xy: tuple[int, int], accent) -> None:
        f = font(22, True)
        pad_x, pad_y = 14, 8
        w = int(draw.textlength(text_value, font=f)) + pad_x * 2
        h = 38
        card(draw, (xy[0] - w // 2, xy[1] - h // 2, xy[0] + w // 2, xy[1] + h // 2), fill=WHITE, outline=accent, radius=12, width=2)
        draw.text((xy[0] - w // 2 + pad_x, xy[1] - h // 2 + pad_y - 2), text_value, font=f, fill=accent)

    def arrow(points: list[tuple[int, int]], fill=TEAL, width: int = 6, dashed: bool = False, label_text: str | None = None, label_xy: tuple[int, int] | None = None) -> None:
        for start, end in zip(points, points[1:], strict=False):
            if dashed:
                dashed_segment(start, end, fill, width)
            else:
                draw.line((start, end), fill=fill, width=width, joint="curve")
        arrow_head(points[-1], points[-2], fill)
        if label_text and label_xy:
            label(label_text, label_xy, fill)

    def vchain(keys: list[str], color) -> None:
        for a, b in zip(keys, keys[1:], strict=False):
            arrow([bottom(a), top(b)], color, width=5)

    # Lanes.
    section(0, "1. Sources and schema", TEAL)
    section(1, "2. Stability context", BLUE)
    section(2, "3. Feature engineering", GREEN)
    section(3, "4. Leakage-safe ML", PURPLE)
    section(4, "5. Validation and exports", AMBER)

    # Source/schema lane.
    node("public_sources", 0, 390, 170, "Public context bundle", ["Alaska DNR wells + GGD223 controls", "G10015 profiles + USGS hydrate AUs", "methane 5 ppt phase-curve source"], TEAL, ICE_LIGHT)
    node("approved_inputs", 0, 630, 170, "Approved OSL inputs later", ["LAS and CSV logs", "core, NMR, workbook labels", "authorized runtime only"], GREEN, GREEN_LIGHT)
    node("role_registry", 0, 870, 180, "Variable fingerprint registry", ["original header + unit shown", "role + X permission + leakage risk"], TEAL)
    node("unit_depth", 0, 1130, 180, "Units and depth axis", ["units visible beside source headers", "depth = alignment/context unless approved"], GREEN)
    node("qc_gate", 0, 1390, 180, "Caliper-first QC gate", ["check caliper/CAL1 coverage first", "missing coverage -> QC flag, not filter"], AMBER, AMBER_LIGHT, "fail closed")
    node(
        "coverage",
        0,
        1650,
        230,
        "Current coverage",
        [
            f"{values['wells']} wells | {values['ggd_controls']} GGD223 controls",
            f"{values['profiles']} G10015 profiles | {values['hydrate_aus']} hydrate AUs",
            f"{values['temp_matches']} temp matches | {values['approved_visible']} visible",
            f"{values['screen_calculated']} intervals | {values['screen_no_interval']} no-stable | {values['screen_blocked']} blocked",
        ],
        BLUE,
        BLUE_LIGHT,
        "schema status",
    )
    node("guardrails", 0, 1930, 300, "Project guardrails", ["stability = admissibility/context only", "not proof, occurrence, saturation, or sweet spots", "top/base/thickness are conditional, not final proof"], RED, RED_LIGHT)

    # Stability lane.
    node("au_gate", 1, 390, 152, "Spatial context", ["inside USGS hydrate AU", "nearest GGD223 permafrost control"], BLUE, BLUE_LIGHT)
    node("depth_basis", 1, 620, 152, "Depth basis", ["TrueVertic preferred", "DrillerTot fallback flagged"], TEAL, ICE_LIGHT)
    node("pressure_eq", 1, 850, 188, "Pressure equation", ["P_abs = P_surface + rho_w*g*z_m/1e6", "units: MPa, meters"], BLUE, WHITE, formula=True)
    node("temperature_eq", 1, 1120, 188, "Temperature model", ["T_model(z) = G10015 interpolation/extrapolation", "flag missing coverage"], BLUE, WHITE, formula=True)
    node("phase_eq", 1, 1390, 188, "Phase-boundary lookup", ["T_eq = f(P_abs, CH4, salinity)", "baseline = methane, 5 ppt"], AMBER, AMBER_LIGHT, formula=True)
    node("stable_test", 1, 1660, 188, "Stability test", ["stable_candidate = T_model <= T_eq", "conditional top/base/thickness"], BLUE, BLUE_LIGHT, formula=True)
    node("stability_context", 1, 1930, 220, "Stability output to ML", ["status | interval | source confidence", "context/mask/confidence/caveat only", "blocked reason retained"], BLUE, WHITE, "not proof")

    # Feature engineering lane.
    node("measured_logs", 2, 390, 170, "Measured log families", ["GR, RHOB, Rt", "Vp, Vs, NMRPHI", "caliper where available"], TEAL, ICE_LIGHT)
    node("lithology_features", 2, 640, 170, "Lithology/reservoir", ["Vsh = (GR-GR_clean)/(GR_sh-GR_clean)", "reservoir-quality flags"], GREEN, WHITE, formula=True)
    node("porosity_features", 2, 890, 170, "Porosity/NMR", ["phi_D = (rho_ma-RHOB)/(rho_ma-rho_f)", "NMR separation when available"], GREEN, GREEN_LIGHT, formula=True)
    node("sonic_features", 2, 1140, 198, "Sonic and elastic", ["Vp = 304.8/DT ; Vs = 304.8/DTS", "AI = RHOB*Vp | Vp/Vs", "mu-rho, lambda-rho"], PURPLE, PURPLE_LIGHT, formula=True)
    node("resistivity_features", 2, 1420, 170, "Resistivity checks", ["Rt features and Archie-style baselines", "baseline/check, not automatic label"], AMBER, AMBER_LIGHT)
    node("context_features", 2, 1670, 170, "Context features", ["AU | permafrost | stability status", "caveats | confidence | blocked reasons"], BLUE, BLUE_LIGHT)
    node("feature_matrix", 2, 1920, 230, "Feature matrix", ["X allowed = measured + derived + QC + context", "stability enters as context/mask only", "targets excluded before modeling"], GREEN, GREEN_LIGHT, "predictors only")

    # ML runtime lane.
    node("target_registry", 3, 390, 210, "Target registry", ["Sgh, S_h, Sh, NMR_SAT", "Hydrate Saturation, Swr", "phase/core labels"], RED, RED_LIGHT, "Y only")
    node("leakage_barrier", 3, 660, 190, "Leakage barrier", ["labels bypass X allowed", "occurrence evidence is target/validation only", "core/NMR/log/seismic evidence"], RED, WHITE)
    node("split", 3, 920, 170, "Validation split", ["whole well / compartment / geography", "split before preprocessing"], BLUE, BLUE_LIGHT)
    node("preprocess", 3, 1170, 198, "Train-only preprocessing", ["0-1 scaling after split only", "fit imputation/selection on train"], PURPLE, PURPLE_LIGHT)
    node("baselines", 3, 1440, 170, "Baseline models", ["physics/simple baselines first", "check if ML beats transparent rules"], AMBER, AMBER_LIGHT)
    node("candidates", 3, 1690, 170, "Candidate models", ["tree/boosting second", "ANN/Keras third; adapters require approval"], PURPLE)
    node("model_heads", 3, 1940, 210, "Two model heads", ["occurrence classifier -> P(hydrate)", "saturation regressor -> Sh pred", "linked, separate tasks"], PURPLE, PURPLE_LIGHT)

    # Output lane.
    node("occurrence", 4, 390, 170, "Occurrence classifier", ["outputs P(hydrate)", "precision, recall, calibration"], BLUE, BLUE_LIGHT)
    node("saturation", 4, 640, 170, "Saturation regressor", ["outputs Sh pred", "RMSE, MAE, R2 where labels exist"], GREEN, GREEN_LIGHT)
    node("metrics", 4, 890, 198, "Validation metrics", ["held-out wells or compartments", "compare only to approved labels", "review calibration and residuals"], PURPLE, PURPLE_LIGHT)
    node("residuals", 4, 1170, 230, "Residual and mimic review", ["by well, depth, lithology, QC", "source confidence and mimic risk", "mimics: shale, ice, gas, cement, washout"], AMBER, AMBER_LIGHT)
    node("exports", 4, 1480, 198, "Approved runtime exports later", ["probability, Sh pred, uncertainty", "reason flags and blocked reasons"], GREEN, GREEN_LIGHT)
    node("public_review", 4, 1760, 210, "Public-safe outputs now", ["diagrams, counts, schema, caveats", "website shows method/readiness", "not final ML results"], TEAL, ICE_LIGHT)
    node("mentor_decisions", 4, 2050, 230, "Mentor decisions", ["methane 5 ppt vs scenarios", "target authority | validation split", "missing G10015 | adapters | public outputs"], RED, RED_LIGHT)

    # Vertical lane flows.
    arrow([bottom("approved_inputs"), top("role_registry")], GREEN, width=5)
    vchain(["role_registry", "unit_depth", "qc_gate"], GREEN)
    vchain(["au_gate", "depth_basis", "pressure_eq", "temperature_eq", "phase_eq", "stable_test", "stability_context"], BLUE)
    vchain(["measured_logs", "lithology_features", "porosity_features", "sonic_features", "resistivity_features", "context_features", "feature_matrix"], GREEN)
    vchain(["target_registry", "leakage_barrier"], RED)
    vchain(["split", "preprocess", "baselines", "candidates", "model_heads"], PURPLE)
    vchain(["metrics", "residuals", "exports", "public_review", "mentor_decisions"], AMBER)

    # Cross-lane flows.
    bus01 = col_x[1] - 62
    bus12 = col_x[2] - 62
    bus23 = col_x[3] - 62
    bus34a = col_x[4] - 92
    bus34b = col_x[4] - 48
    lower_bus = section_bottom - 500
    arrow([right("public_sources"), (bus01, right("public_sources")[1]), (bus01, left("au_gate")[1]), left("au_gate")], TEAL, label_text="public context", label_xy=(bus01, 350))
    arrow(
        [right("qc_gate"), (bus01, right("qc_gate")[1]), (bus01, lower_bus), (bus12, lower_bus), (bus12, left("measured_logs")[1]), left("measured_logs")],
        GREEN,
        label_text="clean curves",
        label_xy=((bus01 + bus12) // 2, lower_bus),
    )
    arrow([right("stability_context"), (bus12, right("stability_context")[1]), (bus12, left("context_features")[1]), left("context_features")], BLUE, label_text="mask/confidence", label_xy=(bus12, 1790))
    arrow([right("feature_matrix"), (bus23, right("feature_matrix")[1]), (bus23, left("split")[1]), left("split")], GREEN, label_text="X allowed", label_xy=(bus23, 1450))
    ml_bus_1 = bus34a
    ml_bus_2 = bus34b
    arrow(
        [right("leakage_barrier"), (ml_bus_2, right("leakage_barrier")[1]), (ml_bus_2, right("model_heads")[1]), right("model_heads")],
        RED,
        dashed=True,
        label_text="Y labels only",
        label_xy=(ml_bus_2, 1340),
    )
    arrow(
        [right("model_heads"), (ml_bus_1, right("model_heads")[1]), (ml_bus_1, left("occurrence")[1]), left("occurrence")],
        BLUE,
        label_text="classification",
        label_xy=(ml_bus_1, 800),
    )
    arrow(
        [right("model_heads"), (ml_bus_2, right("model_heads")[1]), (ml_bus_2, left("saturation")[1]), left("saturation")],
        GREEN,
        label_text="regression",
        label_xy=(ml_bus_2, 980),
    )
    arrow(
        [right("target_registry"), (ml_bus_2, right("target_registry")[1]), (ml_bus_2, left("residuals")[1]), left("residuals")],
        RED,
        dashed=True,
        label_text="validation overlay",
        label_xy=(ml_bus_2, 1360),
    )

    def visual_shell(col: int, heading: str, accent, draw_fn) -> None:
        x = col_x[col] + 54
        y = 2300
        box = (x, y, x + col_w - 108, y + 400)
        card(draw, box, fill=WHITE, outline=accent, radius=22, width=3)
        text(draw, (x + 24, y + 18), heading, 26, accent, True, width=col_w - 156)
        draw_fn((x + 24, y + 70, x + col_w - 132, y + 360))

    def draw_source_visual(box: tuple[int, int, int, int]) -> None:
        x1, y1, x2, y2 = box
        for i, color in enumerate([ICE_LIGHT, BLUE_LIGHT, GREEN_LIGHT, AMBER_LIGHT]):
            bx = x1 + i * 34
            by = y1 + i * 22
            card(draw, (bx, by, bx + 230, by + 120), fill=color, outline=TEAL, radius=12, width=2)
            draw.rectangle((bx + 18, by + 22, bx + 200, by + 30), fill=TEAL)
            draw.rectangle((bx + 18, by + 48, bx + 172, by + 56), fill=(157, 190, 199))
            draw.rectangle((bx + 18, by + 74, bx + 142, by + 82), fill=(157, 190, 199))
        map_x = x1 + 390
        card(draw, (map_x, y1 + 8, x2 - 10, y2 - 12), fill=(247, 251, 252), outline=LINE, radius=16, width=2)
        pts = [(map_x + 50, y1 + 190), (map_x + 100, y1 + 105), (map_x + 190, y1 + 88), (map_x + 270, y1 + 125), (map_x + 315, y1 + 208)]
        draw.line(pts, fill=LINE, width=5)
        for idx, (px, py) in enumerate(pts):
            color = [TEAL, BLUE, GREEN, AMBER, PURPLE][idx]
            draw.ellipse((px - 13, py - 13, px + 13, py + 13), fill=color, outline=WHITE, width=3)
        text(draw, (map_x + 34, y2 - 54), "source packages + spatial controls", 18, MUTED, True, width=x2 - map_x - 50)

    def draw_stability_visual(box: tuple[int, int, int, int]) -> None:
        x1, y1, x2, y2 = box
        axis_left, axis_top, axis_right, axis_bottom = x1 + 70, y1 + 34, x2 - 70, y2 - 44
        draw.line((axis_left, axis_bottom, axis_right, axis_bottom), fill=NAVY, width=4)
        draw.line((axis_left, axis_bottom, axis_left, axis_top), fill=NAVY, width=4)
        text(draw, (axis_left + 120, axis_bottom + 10), "temperature", 16, MUTED, True)
        text(draw, (x1 + 5, axis_top + 70), "depth / pressure", 16, MUTED, True, width=80)
        curve = []
        for i in range(80):
            t = i / 79
            x = int(axis_left + t * (axis_right - axis_left))
            y = int(axis_top + (0.16 + 0.78 * (t**1.75)) * (axis_bottom - axis_top))
            curve.append((x, y))
        draw.line(curve, fill=AMBER, width=7)
        temp_path = [(axis_left + 12, axis_top + 12), (axis_left + 90, axis_top + 92), (axis_left + 190, axis_top + 160), (axis_left + 290, axis_bottom - 18)]
        draw.line(temp_path, fill=BLUE, width=6)
        for px, py in temp_path:
            draw.ellipse((px - 10, py - 10, px + 10, py + 10), fill=BLUE, outline=WHITE, width=3)
        card(draw, (x1 + 345, y1 + 34, x2 - 26, y1 + 122), fill=BLUE_LIGHT, outline=BLUE, radius=14, width=2)
        text(draw, (x1 + 368, y1 + 56), "T_model <= T_eq", 24, BLUE, True, width=x2 - x1 - 400)
        text(draw, (x1 + 368, y1 + 91), "admissible, not proof", 17, MUTED, True, width=x2 - x1 - 400)

    def draw_feature_visual(box: tuple[int, int, int, int]) -> None:
        x1, y1, x2, y2 = box
        track_names = ["GR", "Rt", "RHOB", "Vp", "QC"]
        track_w = 76
        for i, name in enumerate(track_names):
            tx = x1 + i * (track_w + 18)
            card(draw, (tx, y1 + 10, tx + track_w, y2 - 14), fill=(247, 251, 252), outline=LINE, radius=10, width=2)
            text(draw, (tx + 8, y1 + 22), name, 16, [GREEN, TEAL, BLUE, PURPLE, AMBER][i], True, width=track_w - 14, align="center")
            pts = []
            for j in range(26):
                yy = y1 + 58 + j * ((y2 - y1 - 95) / 25)
                xx = tx + 12 + int((track_w - 24) * (0.5 + 0.38 * math.sin(j * 0.74 + i)))
                pts.append((xx, int(yy)))
            draw.line(pts, fill=[GREEN, TEAL, BLUE, PURPLE, AMBER][i], width=3)
        matrix_x = x1 + 510
        cell = 28
        for r in range(7):
            for c in range(6):
                value = (r * 31 + c * 17) % 100
                color = (230 - value // 3, 246 - value // 6, 236 + min(value // 8, 15))
                draw.rectangle(
                    (matrix_x + c * cell, y1 + 54 + r * cell, matrix_x + (c + 1) * cell - 3, y1 + 54 + (r + 1) * cell - 3),
                    fill=color,
                    outline=WHITE,
                )
        text(draw, (matrix_x - 12, y1 + 15), "X allowed matrix", 19, GREEN, True, width=220, align="center")
        text(draw, (matrix_x - 20, y2 - 40), "features + QC + context", 16, MUTED, True, width=240, align="center")

    def draw_ml_visual(box: tuple[int, int, int, int]) -> None:
        x1, y1, x2, y2 = box
        layers = [(x1 + 80, 5, GREEN), (x1 + 240, 7, PURPLE), (x1 + 400, 6, PURPLE), (x1 + 560, 2, AMBER)]
        nodes: list[list[tuple[int, int]]] = []
        for x, count, color in layers:
            pts = []
            for i in range(count):
                yy = int(y1 + 48 + i * ((y2 - y1 - 96) / max(1, count - 1)))
                pts.append((x, yy))
            nodes.append(pts)
        for a_layer, b_layer in zip(nodes, nodes[1:], strict=False):
            for a in a_layer:
                for b in b_layer:
                    draw.line((a, b), fill=(211, 222, 226), width=1)
        for pts, (_, _, color) in zip(nodes, layers, strict=False):
            for px, py in pts:
                r = 15
                draw.ellipse((px - r, py - r, px + r, py + r), fill=WHITE, outline=color, width=4)
        card(draw, (x2 - 145, y1 + 35, x2 - 10, y1 + 115), fill=BLUE_LIGHT, outline=BLUE, radius=14, width=2)
        text(draw, (x2 - 126, y1 + 54), "P(hydrate)", 18, BLUE, True, width=98, align="center")
        card(draw, (x2 - 145, y2 - 118, x2 - 10, y2 - 38), fill=GREEN_LIGHT, outline=GREEN, radius=14, width=2)
        text(draw, (x2 - 126, y2 - 99), "Sh_pred", 18, GREEN, True, width=98, align="center")
        dashed_segment((x1 + 20, y2 - 22), (x2 - 30, y2 - 22), RED, 4)
        text(draw, (x1 + 36, y2 - 55), "target-only supervision rail", 17, RED, True, width=360)

    def draw_output_visual(box: tuple[int, int, int, int]) -> None:
        x1, y1, x2, y2 = box
        card(draw, (x1 + 10, y1 + 10, x1 + 245, y2 - 18), fill=BLUE_LIGHT, outline=BLUE, radius=16, width=2)
        for i, h in enumerate([112, 65, 145, 90, 125]):
            bx = x1 + 38 + i * 37
            draw.rectangle((bx, y2 - 48 - h, bx + 24, y2 - 48), fill=BLUE)
        text(draw, (x1 + 34, y1 + 28), "validation metrics", 19, BLUE, True, width=180)
        card(draw, (x1 + 300, y1 + 10, x2 - 10, y2 - 18), fill=GREEN_LIGHT, outline=GREEN, radius=16, width=2)
        pts = [(x1 + 342, y2 - 54), (x1 + 390, y2 - 118), (x1 + 455, y2 - 92), (x1 + 525, y2 - 160), (x1 + 600, y2 - 128)]
        draw.line(pts, fill=GREEN, width=5)
        for px, py in pts:
            draw.ellipse((px - 10, py - 10, px + 10, py + 10), fill=GREEN, outline=WHITE, width=3)
        text(draw, (x1 + 334, y1 + 28), "per-well outputs", 19, GREEN, True, width=220)
        text(draw, (x1 + 334, y1 + 58), "probability, Sh, uncertainty, caveats", 16, MUTED, True, width=300)

    visual_shell(0, "Public + OSL source packages", TEAL, draw_source_visual)
    visual_shell(1, "Pressure-temperature stability check", BLUE, draw_stability_visual)
    visual_shell(2, "Logs become matrix columns", GREEN, draw_feature_visual)
    visual_shell(3, "Leakage-safe model runtime", PURPLE, draw_ml_visual)
    visual_shell(4, "Validated outputs, not claims", AMBER, draw_output_visual)

    # Status and legend.
    card(draw, (100, 2800, 5100, 2925), fill=(245, 249, 250), outline=LINE, radius=24, width=2)
    legend = [
        ("source/schema", TEAL),
        ("stability context", BLUE),
        ("allowed predictors", GREEN),
        ("target-only", RED),
        ("model runtime", PURPLE),
        ("review/export", AMBER),
    ]
    lx = 140
    for label_value, color in legend:
        draw.rounded_rectangle((lx, 2832, lx + 34, 2866), radius=8, fill=color)
        text(draw, (lx + 48, 2831), label_value, 22, NAVY, True)
        lx += 640 if label_value == "allowed predictors" else 520
    text(
        draw,
        (140, 2878),
        f"Current public status: {values['wells']} scaffold wells, {values['ggd_controls']} GGD223 controls, "
        f"{values['profiles']} G10015 profiles, {values['hydrate_aus']} hydrate AUs, {values['temp_matches']} temperature matches, "
        f"{values['screen_calculated']} baseline admissibility intervals, {values['screen_no_interval']} no-stable rows, "
        f"{values['screen_blocked']} blocked screen rows, and {values['approved_visible']} approved datasets visible for schema design only. "
        "These are workflow constraints, not hydrate proof or saturation results.",
        23,
        MUTED,
        width=4700,
    )

    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    expanded_path = ASSET_DIR / "full_project_ml_workflow_flowchart_expanded.png"
    img.save(expanded_path, quality=94)
    return workflow_slide_summary_panel(values)


def ml_network_detail_panel() -> Path:
    img = Image.new("RGB", (NETWORK_W, NETWORK_H), WHITE)
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 24, NETWORK_H), fill=PURPLE)
    draw.rectangle((0, 0, NETWORK_W, 145), fill=(246, 251, 252))
    draw.line((60, 145, NETWORK_W - 60, 145), fill=LINE, width=3)
    text(draw, (70, 32), "ML Runtime Detail V5.2", 48, NAVY, True)
    text(
        draw,
        (72, 91),
        "Readable architecture view: fingerprints and feature families build X allowed; target labels stay on a separate Y-only rail for training and validation.",
        23,
        MUTED,
        width=2240,
    )

    def arrow_line(start: tuple[int, int], end: tuple[int, int], color, width: int = 4, dashed: bool = False) -> None:
        if dashed:
            dashed_line(draw, start, end, color, width=width, dash=22, gap=14)
        else:
            draw.line((start, end), fill=color, width=width)
        angle = math.atan2(end[1] - start[1], end[0] - start[0])
        size = 18
        pts = [
            end,
            (int(end[0] - size * math.cos(angle - 0.45)), int(end[1] - size * math.sin(angle - 0.45))),
            (int(end[0] - size * math.cos(angle + 0.45)), int(end[1] - size * math.sin(angle + 0.45))),
        ]
        draw.polygon(pts, fill=color)

    def poly_arrow(points: list[tuple[int, int]], color, width: int = 4, dashed: bool = False) -> None:
        for start, end in zip(points, points[1:], strict=False):
            if dashed:
                dashed_line(draw, start, end, color, width=width, dash=22, gap=14)
            else:
                draw.line((start, end), fill=color, width=width)
        end = points[-1]
        prev = points[-2]
        angle = math.atan2(end[1] - prev[1], end[0] - prev[0])
        size = 18
        pts = [
            end,
            (int(end[0] - size * math.cos(angle - 0.45)), int(end[1] - size * math.sin(angle - 0.45))),
            (int(end[0] - size * math.cos(angle + 0.45)), int(end[1] - size * math.sin(angle + 0.45))),
        ]
        draw.polygon(pts, fill=color)

    def small_card(
        box: tuple[int, int, int, int],
        heading: str,
        lines: list[str],
        accent,
        fill=WHITE,
        title_size: int = 21,
        body_size: int = 18,
    ) -> None:
        card(draw, box, fill=fill, outline=accent, radius=16, width=2)
        draw.rounded_rectangle((box[0], box[1], box[0] + 10, box[3]), radius=9, fill=accent)
        y = text(draw, (box[0] + 25, box[1] + 14), heading, title_size, accent, True, width=box[2] - box[0] - 44, gap=4)
        for line in lines:
            y = text(draw, (box[0] + 25, y + 2), line, body_size, NAVY, width=box[2] - box[0] - 48, gap=3)

    def lane_header(box: tuple[int, int, int, int], label: str, accent) -> None:
        card(draw, box, fill=(247, 251, 252), outline=accent, radius=12, width=2)
        text(draw, (box[0] + 12, box[1] + 12), label, 18, accent, True, width=box[2] - box[0] - 24, align="center")

    def draw_matrix(box: tuple[int, int, int, int]) -> None:
        x1, y1, x2, y2 = box
        card(draw, box, fill=WHITE, outline=GREEN, radius=14, width=2)
        text(draw, (x1 + 15, y1 + 14), "X allowed matrix", 22, GREEN, True, width=x2 - x1 - 30, align="center")
        text(draw, (x1 + 28, y1 + 48), "rows = depth-aligned samples, columns = approved inputs", 16, MUTED, True, width=x2 - x1 - 56, align="center")
        start_x, start_y = x1 + 54, y1 + 92
        cell = 25
        for r in range(9):
            for c in range(8):
                value = (r * 29 + c * 17) % 100
                fill = (231 - value // 5, 247 - value // 9, 237 + min(value // 10, 12))
                draw.rectangle(
                    (start_x + c * cell, start_y + r * cell, start_x + (c + 1) * cell - 3, start_y + (r + 1) * cell - 3),
                    fill=fill,
                    outline=WHITE,
                )
        labels = ["logs", "physics", "QC", "context"]
        lx = x1 + 35
        for idx, label_value in enumerate(labels):
            pill(draw, (lx + idx * 76, y2 - 52, lx + 64 + idx * 76, y2 - 22), label_value, [ICE_LIGHT, GREEN_LIGHT, AMBER_LIGHT, BLUE_LIGHT][idx], [TEAL, GREEN, AMBER, BLUE][idx])

    def draw_log_strip(box: tuple[int, int, int, int]) -> None:
        x1, y1, x2, y2 = box
        card(draw, box, fill=(247, 251, 252), outline=LINE, radius=14, width=2)
        text(draw, (x1 + 12, y1 + 12), "depth-aligned log strip", 17, MUTED, True, width=x2 - x1 - 24, align="center")
        names = [("GR", GREEN), ("Rt", TEAL), ("RHOB", BLUE), ("Vp", PURPLE), ("QC", AMBER)]
        track_w = 52
        gap = 10
        for idx, (name, color) in enumerate(names):
            tx = x1 + 25 + idx * (track_w + gap)
            draw.rectangle((tx, y1 + 58, tx + track_w, y2 - 22), outline=(219, 233, 237), width=2)
            text(draw, (tx + 4, y1 + 35), name, 13, color, True, width=track_w - 8, align="center")
            pts = []
            for j in range(30):
                yy = y1 + 68 + j * ((y2 - y1 - 100) / 29)
                xx = tx + 6 + int((track_w - 12) * (0.5 + 0.38 * math.sin(j * 0.62 + idx)))
                pts.append((xx, int(yy)))
            draw.line(pts, fill=color, width=3)

    def draw_model(box: tuple[int, int, int, int]) -> tuple[tuple[int, int], tuple[int, int]]:
        x1, y1, x2, y2 = box
        card(draw, box, fill=WHITE, outline=PURPLE, radius=18, width=2)
        text(draw, (x1 + 22, y1 + 16), "Candidate model", 24, PURPLE, True, width=x2 - x1 - 44)
        text(draw, (x1 + 22, y1 + 48), "baseline first, then tree or ANN after leakage controls pass", 17, MUTED, True, width=x2 - x1 - 44)
        layer_specs = [(x1 + 70, 6, GREEN), (x1 + 170, 7, PURPLE), (x1 + 270, 6, PURPLE), (x1 + 370, 2, AMBER)]
        layers: list[list[tuple[int, int]]] = []
        for lx, count, color in layer_specs:
            pts = [(lx, int(y1 + 125 + i * ((y2 - y1 - 190) / max(1, count - 1)))) for i in range(count)]
            layers.append(pts)
        edge = (202, 216, 222)
        for a_layer, b_layer in zip(layers, layers[1:], strict=False):
            for idx, a in enumerate(a_layer):
                targets = [b_layer[min(len(b_layer) - 1, idx)], b_layer[min(len(b_layer) - 1, idx + 1)]]
                for b in targets:
                    draw.line((a, b), fill=edge, width=1)
        for pts, (_, _, color) in zip(layers, layer_specs, strict=False):
            for px, py in pts:
                r = 16 if len(pts) > 2 else 21
                draw.ellipse((px - r, py - r, px + r, py + r), fill=WHITE, outline=color, width=4)
                if len(pts) == 2:
                    draw.ellipse((px - 7, py - 7, px + 7, py + 7), fill=color)
        text(draw, (x1 + 48, y2 - 48), "linked occurrence and saturation heads", 18, AMBER, True, width=x2 - x1 - 96, align="center")
        return layers[-1][0], layers[-1][1]

    # Lane headers.
    lane_header((70, 168, 515, 218), "1. Feature families", TEAL)
    lane_header((565, 168, 965, 218), "2. X allowed", GREEN)
    lane_header((1015, 168, 1435, 218), "3. Split and controls", BLUE)
    lane_header((1485, 168, 1935, 218), "4. Model candidate", PURPLE)
    lane_header((1985, 168, 2525, 218), "5. Heads and review", AMBER)

    feature_cards = [
        ("Measured logs", ["GR, RHOB, Rt, Vp, Vs, NMRPHI"], TEAL, ICE_LIGHT),
        ("Derived physics", ["Vsh, porosity, AI, elastic terms"], GREEN, GREEN_LIGHT),
        ("QC and alignment", ["caliper coverage first; missing-QC flag"], AMBER, AMBER_LIGHT),
        ("Stability context", ["context/mask/confidence only", "not occurrence or saturation"], BLUE, BLUE_LIGHT),
        ("Core and NMR context", ["calibration later, not public rows"], PURPLE, PURPLE_LIGHT),
    ]
    feature_midpoints: list[tuple[int, int]] = []
    for idx, (heading, lines, accent, fill) in enumerate(feature_cards):
        y = 250 + idx * 125
        box = (85, y, 500, y + 86)
        small_card(box, heading, lines, accent, fill, title_size=20, body_size=17)
        feature_midpoints.append((box[2], (box[1] + box[3]) // 2))

    draw_log_strip((585, 250, 945, 475))
    draw_matrix((585, 535, 945, 860))
    for midpoint in feature_midpoints:
        arrow_line(midpoint, (585, min(max(midpoint[1], 290), 820)), GREEN, width=3)
    arrow_line((765, 475), (765, 535), GREEN, width=4)

    small_card((1035, 250, 1415, 390), "Validation split", ["whole well / compartment / geography", "split before preprocessing"], BLUE, BLUE_LIGHT)
    small_card((1035, 440, 1415, 610), "Train-only preprocessing", ["0-1 scaling after split only", "fit imputation and selection on train"], PURPLE, PURPLE_LIGHT)
    small_card((1035, 660, 1415, 805), "Baseline gate", ["transparent physics/simple rules first", "tree/boosting then ANN/Keras"], AMBER, AMBER_LIGHT)
    arrow_line((945, 695), (1035, 695), GREEN)
    arrow_line((1225, 390), (1225, 440), BLUE)
    arrow_line((1225, 610), (1225, 660), PURPLE)

    out_occurrence, out_saturation = draw_model((1495, 250, 1915, 805))
    arrow_line((1415, 733), (1495, 733), PURPLE)

    small_card((1995, 250, 2505, 385), "Occurrence head", ["P(hydrate), calibrated class", "not stability proof"], BLUE, BLUE_LIGHT)
    small_card((1995, 450, 2505, 585), "Saturation head", ["Sh pred with uncertainty", "compare only to approved labels"], GREEN, GREEN_LIGHT)
    small_card((1995, 650, 2505, 800), "Validation", ["held-out wells or compartments", "approved labels only", "metrics by QC/source confidence"], PURPLE, PURPLE_LIGHT)
    small_card((1995, 865, 2505, 1015), "Reviewed output package", ["probability, Sh pred, uncertainty", "reason flags, blocked reasons", "public-safe summaries"], TEAL, ICE_LIGHT)
    arrow_line(out_occurrence, (1995, 318), BLUE)
    arrow_line(out_saturation, (1995, 517), GREEN)
    arrow_line((2250, 585), (2250, 650), PURPLE)
    arrow_line((2250, 800), (2250, 865), TEAL)

    # Target-only rail.
    rail = (70, 1060, 2530, 1270)
    card(draw, rail, fill=(255, 247, 247), outline=RED, radius=20, width=2)
    text(draw, (105, 1088), "Target-only rail: labels supervise training and validation, but never enter X allowed", 29, RED, True, width=2200)
    target_boxes = [
        ((110, 1148, 540, 1238), "Target registry", ["Sgh, S_h, Sh, NMR_SAT", "Hydrate Saturation, Swr, phase/core labels"], RED, RED_LIGHT),
        ((605, 1148, 985, 1238), "Leakage barrier", ["Y only", "removed before feature matrix"], RED, WHITE),
        ((1050, 1148, 1430, 1238), "Allowed use", ["loss, calibration, validation overlay", "not predictor columns"], PURPLE, PURPLE_LIGHT),
        ((1495, 1148, 1875, 1238), "Mentor decisions", ["target authority", "split and phase-curve policy"], AMBER, AMBER_LIGHT),
        ((1940, 1148, 2490, 1238), "Current status", ["architecture only: no trained metrics, no saturation output", "stability can be context or mask only"], BLUE, BLUE_LIGHT),
    ]
    for box, heading, lines, accent, fill in target_boxes:
        small_card(box, heading, lines, accent, fill, title_size=18, body_size=15)
    poly_arrow([(790, 1060), (790, 930), (1705, 930), (1705, 805)], RED, width=4, dashed=True)
    poly_arrow([(1235, 1060), (1235, 1028), (1810, 1028), (1810, 828), (2250, 828), (2250, 800)], RED, width=4, dashed=True)
    text(draw, (815, 914), "Y labels feed loss or validation only", 19, RED, True, width=520)

    card(draw, (70, 1320, 2530, 1405), fill=(245, 249, 250), outline=LINE, radius=18, width=2)
    text(
        draw,
        (102, 1342),
        "Guardrail: this is architecture only. Missing-log adapters are optional and mentor-approved; real training waits for approved labels, units, split policy, and target authority.",
        22,
        MUTED,
        True,
        width=2380,
    )

    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / "ml_pipeline_network_detail_v5.png"
    img.save(path, quality=94)
    return path


def project_cover_slide(values: dict[str, str]) -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    text(draw, (62, 58), "North Slope Gas Hydrate ML Workflow V5.2", 52, NAVY, True, width=1220)
    text(
        draw,
        (66, 126),
        "Project status deck: public scaffold, stability-admissibility context, approved-data intake, and leakage-safe ML architecture.",
        23,
        MUTED,
        width=1260,
    )

    # Quiet North Slope layer sketch.
    frame = (930, 145, 1518, 660)
    card(draw, frame, fill=(246, 251, 252), outline=LINE, radius=22, width=2)
    for idx, color in enumerate([(187, 225, 238), (209, 236, 220), (236, 222, 190)]):
        y = 212 + idx * 78
        pts = [
            (frame[0] + 34, y + 34),
            (frame[0] + 136, y - 2),
            (frame[0] + 274, y + 24),
            (frame[0] + 410, y - 14),
            (frame[2] - 34, y + 26),
        ]
        draw.line(pts, fill=color, width=30, joint="curve")
    draw.polygon(
        [(frame[0] + 60, 510), (frame[0] + 205, 416), (frame[0] + 378, 458), (frame[2] - 68, 394), (frame[2] - 28, 632), (frame[0] + 28, 632)],
        fill=(157, 105, 163),
        outline=(118, 78, 133),
    )
    draw.polygon(
        [(frame[0] + 74, 475), (frame[0] + 218, 398), (frame[0] + 389, 435), (frame[2] - 82, 375), (frame[2] - 68, 414), (frame[0] + 380, 470), (frame[0] + 210, 438), (frame[0] + 78, 516)],
        fill=(45, 156, 126),
        outline=(45, 156, 126),
    )
    for px, py in [(1000, 318), (1085, 268), (1195, 310), (1325, 250), (1430, 298)]:
        draw.line((px, py - 38, px, py + 170), fill=(82, 103, 112), width=3)
        draw.ellipse((px - 9, py - 9, px + 9, py + 9), fill=TEAL, outline=WHITE, width=2)
    text(draw, (958, 166), "Public regional context -> approved runtime analysis", 21, TEAL, True, width=520)

    goal_box = (70, 230, 850, 410)
    card(draw, goal_box, fill=ICE_LIGHT, outline=TEAL, radius=22, width=2)
    text(draw, (104, 260), "Goal", 25, TEAL, True)
    text(
        draw,
        (104, 300),
        "Use public-safe methods and approved-runtime data later to predict hydrate occurrence and saturation without leaking target labels or private rows.",
        25,
        NAVY,
        True,
        width=690,
    )

    stats = [
        ("public wells", values["wells"], TEAL),
        ("GGD223 controls", values["ggd_controls"], BLUE),
        ("G10015 profiles", values["profiles"], GREEN),
        ("baseline intervals", values["screen_calculated"], AMBER),
    ]
    x = 70
    for label_value, count, color in stats:
        card(draw, (x, 455, x + 178, 575), fill=WHITE, outline=color, radius=18, width=2)
        text(draw, (x + 20, 480), count, 31, color, True, width=138, align="center")
        text(draw, (x + 18, 525), label_value, 15, MUTED, True, width=142, align="center")
        x += 205

    path_cards = [
        ((70, 635, 430, 760), "Public GitHub / Streamlit", ["methods, counts, schemas", "diagrams and caveats"], TEAL, ICE_LIGHT),
        ((470, 635, 830, 760), "OSL / Approved Runtime", ["real LAS, core, NMR later", "training and validation"], PURPLE, PURPLE_LIGHT),
        ((870, 635, 1230, 760), "Current Guardrail", ["stability = admissibility", "not proof or saturation"], RED, RED_LIGHT),
    ]
    for box, heading, rows, accent, fill in path_cards:
        node(draw, box, heading, rows, fill, accent)
    arrow(draw, (430, 698), (470, 698), TEAL)
    arrow(draw, (830, 698), (870, 698), PURPLE)
    footer(draw, "V5.2 refresh: the expanded architecture map and ML runtime detail are now included inside the slide deck.")
    return save(img, "slide_01_project_cover_v5_2.png")


def diagram_reference_slide(source_path: Path, filename: str, footer_note: str) -> Path:
    img = Image.new("RGB", (W, H), WHITE)
    if not source_path.exists():
        raise FileNotFoundError(source_path)
    with Image.open(source_path).convert("RGB") as source:
        scale = min((W - 40) / source.width, H / source.height)
        new_w = int(source.width * scale)
        new_h = int(source.height * scale)
        resized = source.resize((new_w, new_h), Image.Resampling.LANCZOS)
        img.paste(resized, ((W - new_w) // 2, (H - new_h) // 2))
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, H - 36, W, H), fill=(246, 251, 252))
    text(draw, (58, H - 29), footer_note, 13, MUTED, True, width=1480)
    return save(img, filename)


def variable_fingerprint_decisions_slide() -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    title(
        draw,
        "Variable Fingerprint Before X Allowed",
        "V5.2 decision boxes: headers keep their units and roles before any feature enters the approved-runtime matrix.",
    )

    cards = [
        ((70, 170, 505, 360), "Fingerprint fields", ["original header + unit", "normalized name + role", "allowed in X? + leakage risk"], TEAL, ICE_LIGHT),
        ((585, 170, 1020, 360), "Depth policy", ["alignment/context axis", "not automatically normalized predictor", "well/depth split comes first"], BLUE, BLUE_LIGHT),
        ((1100, 170, 1535, 360), "Train-only scaling", ["numeric predictors use 0-1 scaling", "fit only on training wells", "freeze for validation/test"], PURPLE, PURPLE_LIGHT),
        ((70, 430, 505, 620), "Caliper gate", ["check caliper/CAL1 coverage first", "washout filter only if supported", "otherwise missing-QC flag"], AMBER, AMBER_LIGHT),
        ((585, 430, 1020, 620), "X allowed", ["measured logs", "derived physics features", "QC/context/stability caveats"], GREEN, GREEN_LIGHT),
        ((1100, 430, 1535, 620), "Y-only labels", ["Sgh, S_h, Sh, NMR_SAT", "Hydrate Saturation, Swr", "phase/core labels"], RED, RED_LIGHT),
    ]
    for box, heading, rows, accent, fill in cards:
        node(draw, box, heading, rows, fill, accent)

    arrow(draw, (505, 265), (585, 265), TEAL)
    arrow(draw, (1020, 265), (1100, 265), BLUE)
    arrow(draw, (505, 525), (585, 525), AMBER)
    arrow(draw, (1020, 525), (1100, 525), RED)

    card(draw, (70, 700, 1535, 795), fill=LIGHT, outline=LINE, radius=18, width=2)
    text(draw, (104, 727), "Operational validator", 20, TEAL, True)
    text(
        draw,
        (328, 720),
        "dashboard/approved_data_intake.py and 01_pipeline/validate_approved_data_headers.py audit headers only, report unknowns/leakage/blocked reasons, and never expose approved row values.",
        23,
        NAVY,
        True,
        width=1140,
    )
    footer(draw, "Decision rule: target columns and occurrence labels supervise/validate the model; they do not become predictors.")
    return save(img, "slide_06_variable_fingerprint_decisions.png")


def model_decision_boxes_slide() -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    title(
        draw,
        "Modeling Decisions: Linked Outputs, Separate Claims",
        "The first approved-runtime experiment stays leakage-safe and baseline-controlled before ANN/Keras candidates are tested.",
    )

    cards = [
        ((70, 170, 505, 365), "Occurrence classifier", ["P(hydrate) or class", "core/NMR/log/seismic evidence", "not measured by stability"], BLUE, BLUE_LIGHT),
        ((585, 170, 1020, 365), "Saturation regressor", ["Sh_pred only where labels exist", "Sgh/S_h/Sh/NMR_SAT authority needed", "fraction vs percent required"], GREEN, GREEN_LIGHT),
        ((1100, 170, 1535, 365), "Validation split", ["whole well / compartment / geography", "split before preprocessing", "random depth-row final split rejected"], AMBER, AMBER_LIGHT),
        ((70, 445, 505, 640), "Model order", ["baselines first", "tree/boosting second", "ANN/Keras third"], PURPLE, PURPLE_LIGHT),
        ((585, 445, 1020, 640), "Missing-log strategy", ["alternate log combinations first", "Vp/RHOB adapters optional", "mentor approval + validation required"], RED, RED_LIGHT),
        ((1100, 445, 1535, 640), "Research anchors", ["Chong 2024: ANN pattern", "Singh: MTE/IGS case wells", "MDPI: adapter background"], TEAL, ICE_LIGHT),
    ]
    for box, heading, rows, accent, fill in cards:
        node(draw, box, heading, rows, fill, accent)

    arrow(draw, (505, 268), (585, 268), TEAL, label="linked")
    arrow(draw, (1020, 268), (1100, 268), AMBER)
    arrow(draw, (505, 543), (585, 543), PURPLE)
    arrow(draw, (1020, 543), (1100, 543), TEAL)

    card(draw, (70, 714, 1535, 805), fill=(255, 247, 247), outline=RED, radius=18, width=2)
    text(draw, (105, 738), "Guardrail", 21, RED, True)
    text(
        draw,
        (262, 732),
        "No occurrence probabilities, saturation predictions, trained metrics, or sweet-spot outputs exist yet. Public slides show architecture and readiness only.",
        24,
        NAVY,
        True,
        width=1190,
    )
    footer(draw, "Stability can be context, mask, confidence, caveat, or blocked reason only if approved by mentor.")
    return save(img, "slide_08_model_decision_boxes.png")


def overview_slide_panel() -> Path:
    # Reuse the readable workflow image as the main slide panel.
    return full_workflow_panel()


def flow_slide(
    filename: str,
    heading: str,
    subheading: str,
    steps: list[tuple[str, list[str], tuple[int, int, int], tuple[int, int, int]]],
    guardrail: str,
    footer_text: str,
) -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    title(draw, heading, subheading)
    start_x = 70
    y = 210
    box_w = 260
    gap = 35
    previous: tuple[int, int, int, int] | None = None
    for idx, (step_heading, lines, fill, accent) in enumerate(steps):
        x = start_x + idx * (box_w + gap)
        box = (x, y, x + box_w, y + 250)
        node(draw, box, step_heading, lines, fill, accent)
        if previous:
            arrow(draw, (previous[2] + 8, y + 125), (box[0] - 8, y + 125), accent)
        previous = box
    card(draw, (74, 595, 1528, 736), fill=LIGHT, outline=LINE, radius=18)
    text(draw, (108, 623), "Guardrail", 20, RED, True)
    text(draw, (250, 618), guardrail, 26, NAVY, True, width=1170)
    footer(draw, footer_text)
    return save(img, filename)


def slide_04_inputs_boundary() -> Path:
    return flow_slide(
        "slide_04_inputs_boundary_zoom.png",
        "The Inputs Split Into Public Context And Approved Runtime Data",
        "The diagram separates what can be shown publicly from what must stay inside OSL or the approved environment.",
        [
            ("Public sources", ["DNR wells", "GGD223 / G10015", "USGS AU + phase curve"], ICE_LIGHT, TEAL),
            ("Public counts", ["8,084 wells", "43 controls / 184 profiles", "3 AUs / 483 temp matches"], BLUE_LIGHT, BLUE),
            ("Approved runtime", ["LAS / CSV logs", "core and NMR", "workbook labels"], GREEN_LIGHT, GREEN),
            ("Schema controls", ["original headers", "about 3 / 71 visible", "roles before aliases"], WHITE, TEAL),
            ("Boundary review", ["no approved rows", "no restricted IDs", "no trained models"], RED_LIGHT, RED),
        ],
        "The public side explains the method. The approved side runs the real log/core analysis.",
        "Use this as the slide-deck explanation for why GitHub/Streamlit and OSL both exist.",
    )


def slide_05_stability_branch() -> Path:
    return flow_slide(
        "slide_05_stability_branch_zoom.png",
        "Stability Connects Into ML As Context, Not As Proof",
        "The pressure-temperature screen is a guarded branch that feeds context, masks, confidence, and caveats.",
        [
            ("Depth basis", ["TrueVertic preferred", "DrillerTot fallback", "public depth caveat"], WHITE, TEAL),
            ("Pressure", ["P_abs = P_surface + rho_w*g*z_m/1e6", "hydrostatic assumption", "not measured pressure"], ICE_LIGHT, BLUE),
            ("Temperature", ["T_model(z) from G10015", "interpolate / extrapolate", "block missing coverage"], BLUE_LIGHT, BLUE),
            ("Phase curve", ["T_eq = f(P_abs, CH4, salinity)", "baseline methane 5 ppt", "scenario capable"], AMBER_LIGHT, AMBER),
            ("Stability status", ["T_model <= T_eq", "top/base conditional", "not hydrate proof"], RED_LIGHT, RED),
        ],
        "A stable interval can only say hydrate is admissible under assumptions; logs/core decide occurrence and saturation later.",
        "Current guarded run: 8,084 screen rows, 22 calculated admissibility intervals, 8 no-stable-interval rows, 8,054 blocked rows.",
    )


def slide_06_features() -> Path:
    return flow_slide(
        "slide_06_features_zoom.png",
        "Well Logs Become Feature Blocks Only After QC And Units",
        "The ML table is built from measured curves, derived physics features, and context flags with source provenance.",
        [
            ("Measured logs", ["GR, RHOB, Rt", "Vp, Vs, NMRPHI", "caliper where available"], WHITE, TEAL),
            ("Normalize", ["feet vs meters", "density / sonic / porosity", "raw values retained"], ICE_LIGHT, BLUE),
            ("QC + alignment", ["washout / caliper", "missingness / outliers", "core-log depth offsets"], AMBER_LIGHT, AMBER),
            ("Derived features", ["Vsh, density porosity", "AI, Vp/Vs", "lambda-rho, mu-rho"], GREEN_LIGHT, GREEN),
            ("Context flags", ["AU / permafrost", "stability status", "confidence and caveats"], BLUE_LIGHT, BLUE),
        ],
        "High resistivity, low GR, high velocity, or stability context cannot become standalone hydrate labels.",
        "This follows the science-to-ML ladder: stability context -> reservoir quality -> hydrate response.",
    )


def slide_07_leakage_modeling() -> Path:
    return flow_slide(
        "slide_07_leakage_modeling_zoom.png",
        "The Leakage Barrier Controls What The Model Is Allowed To Learn",
        "Targets supervise or validate the model; they do not enter the feature table as predictors.",
        [
            ("Target registry", ["Sgh, S_h, Sh, NMR_SAT", "Hydrate Saturation, Swr", "phase/core labels"], RED_LIGHT, RED),
            ("Allowed inputs", ["measured logs", "derived features", "QC/context flags"], GREEN_LIGHT, GREEN),
            ("Split first", ["whole wells / compartments", "geographic holdout option", "before preprocessing"], ICE_LIGHT, BLUE),
            ("Train controls", ["train-only transforms", "physics/simple baselines first", "tree/ANN after gates"], PURPLE_LIGHT, PURPLE),
            ("Two heads", ["P(hydrate)", "Sh pred", "linked but separate"], WHITE, TEAL),
        ],
        "The answer columns cannot leak into the input side, even if they are useful for calibration or validation.",
        "ML source anchors: Chong et al. 2022, Singh et al. 2021, Chong et al. 2024; project results wait for approved validation.",
    )


def slide_08_outputs_validation() -> Path:
    return flow_slide(
        "slide_08_outputs_validation_zoom.png",
        "The Final Product Is A Reviewed Output Package, Not One Score",
        "The approved runtime should export interpretable outputs and only public-safe reductions after review.",
        [
            ("Runtime later", ["probability track", "Sh pred with uncertainty", "reason flags"], BLUE_LIGHT, BLUE),
            ("Website now", ["method readiness", "schema coverage", "guarded stability counts"], GREEN_LIGHT, GREEN),
            ("Structural Explorer", ["per-well curve concept", "map/layer views", "review-only summaries"], AMBER_LIGHT, AMBER),
            ("Validation", ["approved labels only", "well/compartment holdout", "residual and mimic review"], PURPLE_LIGHT, PURPLE),
            ("Exports", ["plots and tables", "reason / blocked flags", "public-safe summary"], WHITE, TEAL),
        ],
        "Do not present public stability status, screenshots, or comparative-source metrics as this project's validated ML results.",
        "Outputs remain approved-runtime products until the public-safe boundary review decides what can be shown.",
    )


def slide_09_status_decisions() -> Path:
    values = summaries()
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    title(draw, "V5.2 Status And Mentor Decisions", "The project has a defensible workflow map, validator, and handoff contract; it does not yet have approved-data hydrate predictions.")

    columns = [
        (
            "Complete",
            [
                "public / OSL boundary",
                "8,084-well public scaffold",
                "schema registry + V5.2 deck/doc",
            ],
            GREEN_LIGHT,
            GREEN,
        ),
        (
            "Calculated",
            [
                f"{values['temp_rows']} temperature rows",
                f"{values['screen_calculated']} baseline admissibility intervals",
                f"{values['screen_no_interval']} sufficient-input no-interval rows",
            ],
            BLUE_LIGHT,
            BLUE,
        ),
        (
            "Blocked / future",
            [
                f"{values['screen_blocked']} blocked screen rows",
                "approved logs/core/NMR rows",
                "validated ML outputs",
            ],
            RED_LIGHT,
            RED,
        ),
    ]
    x = 72
    for heading, rows, fill, accent in columns:
        card(draw, (x, 170, x + 455, 620), fill=fill, outline=accent, radius=22, width=2)
        text(draw, (x + 28, 198), heading, 28, accent, True)
        y = 258
        for row in rows:
            draw.ellipse((x + 28, y + 5, x + 40, y + 17), fill=accent)
            y = text(draw, (x + 54, y), row, 19, NAVY, width=360, gap=12)
        x += 505

    card(draw, (72, 675, 1530, 790), fill=LIGHT, outline=LINE, radius=18)
    text(draw, (105, 704), "Next decisions", 20, AMBER, True)
    text(
        draw,
        (280, 698),
        "methane 5 ppt only vs scenarios; target authority; validation split; missing G10015 policy; stability as context/mask only; missing-log adapters; acceptable public website outputs before validation.",
        22,
        NAVY,
        True,
        width=1160,
    )
    footer(draw, "Blocked means the workflow is refusing to overclaim without required inputs; it does not mean no hydrate.")
    return save(img, "slide_09_status_and_next_decisions.png")


def preserve_source_slide_panel(slide_number: int, name: str) -> Path:
    if not SOURCE_DECK.exists():
        raise FileNotFoundError(SOURCE_DECK)
    prs = Presentation(SOURCE_DECK)
    slide = prs.slides[slide_number - 1]
    pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
    if len(pictures) != 1:
        raise ValueError(f"Source slide {slide_number} should have exactly one raster panel, found {len(pictures)}")
    blob = pictures[0].image.blob
    with Image.open(BytesIO(blob)) as image:
        if image.size != (W, H):
            raise ValueError(f"Source slide {slide_number} image size is {image.size}, expected {(W, H)}")
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / name
    path.write_bytes(blob)
    LOCKED_SLIDE_HASHES[slide_number] = hashlib.sha256(blob).hexdigest()
    return path


def build_panels(network_path: Path) -> list[Path]:
    values = summaries()
    overview = overview_slide_panel()
    expanded = ASSET_DIR / "full_project_ml_workflow_flowchart_expanded.png"
    return [
        project_cover_slide(values),
        preserve_source_slide_panel(2, "slide_02_locked_from_current_gmail_deck.png"),
        overview,
        diagram_reference_slide(
            expanded,
            "slide_04_expanded_architecture_map.png",
            "Expanded architecture reference inside the deck; slides 5, 6, 8, and 9 provide readable zoom explanations.",
        ),
        slide_05_stability_branch(),
        variable_fingerprint_decisions_slide(),
        diagram_reference_slide(
            network_path,
            "slide_07_ml_runtime_detail.png",
            "ML runtime detail inside the deck: X allowed, Y-only rail, split controls, output heads, validation, and reviewed outputs.",
        ),
        model_decision_boxes_slide(),
        slide_09_status_decisions(),
    ]


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        element = shape._element  # noqa: SLF001 - python-pptx has no public remove API.
        element.getparent().remove(element)


def rebuild_deck(panel_paths: list[Path]) -> Path:
    prs = Presentation(SOURCE_DECK)
    if len(prs.slides) != 9:
        raise ValueError(f"Expected 9 slides in {SOURCE_DECK}, found {len(prs.slides)}")
    if len(panel_paths) != 9:
        raise ValueError(f"Expected 9 panels, found {len(panel_paths)}")
    for slide, panel in zip(prs.slides, panel_paths, strict=True):
        clear_slide(slide)
        slide.shapes.add_picture(str(panel), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(OUT_DECK)
    return OUT_DECK


def verify_deck(path: Path) -> None:
    prs = Presentation(path)
    if len(prs.slides) != 9:
        raise ValueError(f"Expected 9 slides, found {len(prs.slides)}")
    for idx, slide in enumerate(prs.slides, start=1):
        pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
        if len(pictures) != 1:
            raise ValueError(f"Slide {idx} should have exactly one raster panel, found {len(pictures)}")
        blob = pictures[0].image.blob
        with Image.open(BytesIO(blob)) as image:
            if image.size != (W, H):
                raise ValueError(f"Slide {idx} image size is {image.size}, expected {(W, H)}")
        if idx in LOCKED_SLIDE_HASHES and hashlib.sha256(blob).hexdigest() != LOCKED_SLIDE_HASHES[idx]:
            raise ValueError(f"Slide {idx} no longer matches the locked source panel")


def build_contact_sheet(panel_paths: list[Path]) -> Path:
    thumb_w, thumb_h = 480, 270
    margin = 34
    label_h = 34
    gap = 26
    sheet_w = margin * 2 + thumb_w * 3 + gap * 2
    sheet_h = margin * 2 + (thumb_h + label_h) * 3 + gap * 2
    img = Image.new("RGB", (sheet_w, sheet_h), WHITE)
    draw = ImageDraw.Draw(img)
    for idx, path in enumerate(panel_paths, start=1):
        row = (idx - 1) // 3
        col = (idx - 1) % 3
        x = margin + col * (thumb_w + gap)
        y = margin + row * (thumb_h + label_h + gap)
        panel = Image.open(path).convert("RGB").resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        card(draw, (x - 4, y - 4, x + thumb_w + 4, y + thumb_h + label_h), fill=(248, 251, 252), outline=LINE, radius=14, width=2)
        img.paste(panel, (x, y))
        draw.rectangle((x, y, x + thumb_w, y + thumb_h), outline=(160, 190, 198), width=2)
        text(draw, (x + 14, y + thumb_h + 9), f"Slide {idx}", 18, NAVY, True, width=thumb_w - 28)
    OUT_CONTACT_SHEET.parent.mkdir(parents=True, exist_ok=True)
    img.save(OUT_CONTACT_SHEET, quality=94)
    return OUT_CONTACT_SHEET


def apply_doc_style(document: Document) -> None:
    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(10.5)
    for style_name, size in [("Title", 22), ("Heading 1", 16), ("Heading 2", 13)]:
        style = document.styles[style_name]
        style.font.name = "Arial"
        style.font.size = Pt(size)
        style.font.bold = True


def build_word_companion(diagram_path: Path, network_path: Path) -> Path:
    document = Document()
    apply_doc_style(document)
    section = document.sections[0]
    section.orientation = WD_ORIENT.LANDSCAPE
    section.page_width, section.page_height = section.page_height, section.page_width
    section.top_margin = Inches(0.55)
    section.bottom_margin = Inches(0.55)
    section.left_margin = Inches(0.55)
    section.right_margin = Inches(0.55)

    props = document.core_properties
    props.title = "V5.2 North Slope Gas Hydrate Full ML Workflow Companion"
    props.subject = "V5.2 explanation of the public, OSL, stability, intake-validator, and ML workflow"
    props.author = "North Slope Gas Hydrates project"

    document.add_heading("V5.2 North Slope Gas Hydrate Full ML Workflow Companion", level=0)
    document.add_paragraph(
        "This companion page explains the V5.2 diagram package requested for the slide refresh. "
        "The package connects source and schema controls, OpenScienceLab and approved-runtime inputs, "
        "the guarded pressure-temperature stability branch, well-log physics equations, the leakage "
        "barrier, occurrence classification, saturation regression, validation, and reviewed exports. "
        "The refreshed slide deck includes both complex architecture diagrams inside the presentation: "
        "the expanded full workflow map and the ML runtime/network detail."
    )

    values = summaries()
    document.add_heading("Project Status In Plain Language", level=1)
    document.add_paragraph(
        "The project is building a defensible North Slope workflow for gas hydrate occurrence "
        "classification and saturation regression. The public GitHub/Streamlit side explains "
        "source-backed methods, public scaffold counts, diagrams, and public-safe outputs. "
        "OpenScienceLab and the later approved runtime are where raw source bundles, approved "
        "well logs, core/NMR information, target labels, feature engineering, model fitting, "
        "validation, and reviewed outputs belong."
    )
    document.add_paragraph(
        f"The public scaffold currently has {values['wells']} wells, {values['ggd_controls']} GGD223 "
        f"permafrost controls, {values['profiles']} G10015 profiles, {values['hydrate_aus']} hydrate "
        f"AUs, {values['temp_matches']} temperature-profile matches, and {values['approved_visible']} "
        "approved datasets visible for schema design only."
    )

    document.add_heading("Three Visual Levels", level=1)
    for item in [
        "Slide 3 is the readable mentor-scale overview for a live presentation.",
        "Slide 4 embeds the expanded poster as the detailed architecture map with sources, equations, gates, features, targets, validation, outputs, and caveats.",
        "Slide 7 embeds the ML runtime detail focused on X allowed, the Y-only label rail, split/preprocess/model controls, output heads, and validation.",
        "Slides 5, 6, 8, and 9 are readable zoom slides that explain stability context, variable fingerprints, modeling decisions, and mentor decisions without cramming poster text into one slide.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("Current ML Architecture Decisions", level=1)
    for item in [
        "Occurrence classification and saturation regression are linked but separate outputs. Occurrence uses approved evidence labels; saturation regression uses an approved saturation target field only after authority and units are confirmed.",
        "Numeric predictors use train-only 0-1 scaling after a whole-well, compartment, or geographic split. Depth remains the alignment and context axis unless a mentor explicitly approves depth as a predictor.",
        "The model ladder is baseline models first, tree or boosting models second, and ANN/Keras third. A candidate ANN is an architecture option, not a trained project result.",
        "Stability may enter X allowed only as context, mask, confidence, caveat, or blocked reason if approved. It cannot create occurrence labels, saturation targets, or validated predictions.",
        "Caliper coverage is checked before washout filtering. If caliper, CAL1, or differential caliper coverage is insufficient, the runtime carries a missing-QC flag instead of silently filtering rows.",
        "Missing-log adapters for Vp or RHOB remain optional and mentor-approved; alternate log combinations are the default lower-risk strategy until adapters are validated on North Slope data.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("Variable Fingerprint And Intake Validator", level=1)
    document.add_paragraph(
        "Before a source column can enter X allowed, it needs a variable fingerprint: original header, "
        "unit, normalized name, role, normalized-or-raw status, feature-matrix permission, leakage risk, "
        "public-safe display status, and unresolved mentor question. This keeps origin headers such as "
        "DEPTH, GR, RHOB, Rt, Sgh, S_h, Sh, and NMR_SAT visible instead of hiding them behind aliases."
    )
    document.add_paragraph(
        "The public-safe validator in dashboard/approved_data_intake.py and the CLI runner in "
        "01_pipeline/validate_approved_data_headers.py audit headers only. They can report recognized "
        "headers, unknown headers, missing required fields, target-only leakage, unresolved units, "
        "caliper-gate status, missing-log adapter status, occurrence/saturation target authority, "
        "blocked reasons, and training readiness without reading or printing approved row values."
    )

    document.add_heading("Research Source Anchors", level=1)
    document.add_paragraph(
        "Chong et al. 2024 / USGS supports the architecture pattern of ANN-style hydrate occurrence "
        "classification and saturation prediction from well-log inputs. It classifies hydrate occurrence "
        "types and predicts saturation in a marine hydrate setting, so this project uses it as a model-"
        "architecture anchor rather than as North Slope performance evidence: "
        "https://pubs.usgs.gov/publication/70250169"
    )
    document.add_paragraph(
        "Occurrence is treated here as a target or validation label, not something measured or proven "
        "by the pressure-temperature stability screen. In the approved runtime, occurrence evidence can "
        "come from core or pressure-core observations, NMR/core-derived saturation, validated multi-log "
        "interpretation, or documented seismic indicators. Each occurrence label needs source, interval, "
        "confidence, and caveat metadata before model training or validation."
    )
    document.add_paragraph(
        "The Singh/NETL source identifies Well-MTE as the Mt. Elbert Stratigraphic Test Well and "
        "Well-IGS as the Ignik Sikumi Test Well in the Alaska North Slope / Eileen Gas Hydrate Trend "
        "case-study context. The public workbook headers MTE_refined and IGS_refined should remain "
        "blue mentor questions until approved metadata confirms whether they are raw/refined processing "
        "stages: https://www.osti.gov/servlets/purl/1893637"
    )
    document.add_paragraph(
        "The MDPI missing-log study supports the general plausibility of using machine learning to "
        "estimate missing Vp and bulk density logs in marine hydrate settings. It does not automatically "
        "validate missing-log adapters for North Slope permafrost sediments, so any Vp/RHOB adapter must "
        "be optional, mentor-approved, and validated before use: https://www.mdpi.com/1996-1073/16/23/7709"
    )

    document.add_heading("What Is Complete Outside Stability", level=1)
    for item in [
        "The public/runtime data boundary is documented and reflected in the V5.2 workflow diagrams.",
        "The approved-data schema coverage matrix preserves original headers and separates measured inputs, derived features, QC/alignment fields, calibration/reference fields, target-only fields, and unresolved fields.",
        "The target registry and leakage guardrails keep Sgh, S_h, Sh, NMR_SAT, hydrate-saturation fields, Swr, S_wr, and phase labels out of the predictor matrix.",
        "The Analyze Hydrates page includes Schema Coverage & Architecture, Public ML Readiness, and Target Registry & Leakage views without exposing approved rows or model metrics.",
        "The future ML architecture is defined as a leakage-safe path to linked occurrence classification and saturation regression, but final training and performance reporting remain blocked until approved labels and complete rows are available.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("What Stability Currently Does", level=1)
    document.add_paragraph(
        f"The baseline public stability screen has {values['screen_rows']} rows, "
        f"{values['screen_calculated']} calculated pressure-temperature admissibility intervals, "
        f"{values['screen_no_interval']} sufficient-input no-stable-interval rows, and "
        f"{values['screen_blocked']} blocked rows. It uses the cited methane 5 ppt phase lookup, "
        "hydrostatic absolute pressure, public G10015/GGD223 temperature and permafrost context, "
        "and source-control caveats."
    )
    for item in [
        "It can serve the ML workflow as context, a mask, a confidence field, a caveat field, or a reason flag.",
        "It does not prove hydrate occurrence, estimate saturation, define final top/base/thickness for approved interpretation, rank sweet spots, or replace log/core evidence.",
        "Blocked rows mean the calculation refused to overclaim without required inputs; they do not mean no hydrate.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("How Approved OSL Data Feed The Future ML Pipeline", level=1)
    for item in [
        "Approved LAS/CSV/core/NMR and workbook inputs are loaded only inside OSL or the approved runtime.",
        "Original headers, units, sheet/source identity, and depth references are preserved before alias mapping.",
        "Measured log families and approved derived physics equations build X allowed after unit checks, QC, and depth alignment.",
        "Target-only occurrence and saturation fields bypass X allowed and are used only for target mapping, calibration, and validation overlays.",
        "Occurrence evidence can come from core or pressure-core observations, NMR/core-derived saturation, validated log interpretation, or regional seismic indicators, but it remains target/validation evidence.",
        "Validation uses complete wells, compartments, or geographic holdouts selected before preprocessing, tuning, or model fitting.",
        "Reviewed outputs become occurrence probability, saturation estimate, uncertainty/QC/reason flags, plots, tables, GIS links, and public-safe summaries only after boundary review.",
    ]:
        document.add_paragraph(item, style="List Number")

    document.add_heading("Slide 3 Mentor Overview", level=1)
    document.add_paragraph("This version stays presentation-readable and leaves detailed evidence to the poster, runtime-detail slide, and later zoom slides.")
    document.add_picture(str(diagram_path), width=Inches(9.9))

    expanded_path = ASSET_DIR / "full_project_ml_workflow_flowchart_expanded.png"
    if expanded_path.exists():
        document.add_heading("Expanded Poster Detail", level=1)
        document.add_paragraph(
            "This poster is the complete architecture reference. It includes the public source counts, "
            "approved OSL inputs later, stability equations, target-only labels, leakage controls, "
            "validation expectations, public-safe outputs, and mentor decision points."
        )
        document.add_picture(str(expanded_path), width=Inches(9.9))

    document.add_heading("ML Architecture Detail", level=1)
    document.add_paragraph(
        "The companion architecture visual is embedded in slide 7 and expands the modeling block into feature groups, "
        "QC, the X allowed matrix, whole-well split controls, train-only preprocessing, "
        "a simplified candidate model, two output heads, and the target-only rail used "
        "for training labels and validation overlays."
    )
    document.add_picture(str(network_path), width=Inches(9.9))

    document.add_heading("How To Read The Diagram", level=1)
    for item in [
        "The slide-sized summary separates source/schema controls, stability context, feature engineering, leakage-safe ML, and reviewed outputs.",
        "The expanded poster shows the full OSL and approved-runtime path: raw source rebuilds, approved logs/core/NMR later, pressure-temperature stability context, feature engineering, modeling, and validation.",
        "The equation blocks show how pressure, temperature, phase-boundary lookup, lithology, porosity, velocity, elastic, and saturation-baseline calculations become context fields, features, or validation checks.",
        "The stability branch feeds context, masks, confidence labels, and caveats into the ML workflow. It does not become hydrate proof or a saturation label.",
        "The leakage barrier keeps S_h, Sgh, NMR_SAT, phase labels, and final ranks out of predictor features unless a field is proven to be an independent measured input.",
        "The final ML design keeps occurrence classification and saturation regression linked but separate.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("Current Status Shown In The Figure", level=1)
    status_table = document.add_table(rows=1, cols=3)
    status_table.style = "Table Grid"
    headers = ["Complete", "Calculated", "Blocked / Future"]
    for i, header in enumerate(headers):
        run = status_table.rows[0].cells[i].paragraphs[0].add_run(header)
        run.bold = True
    row = status_table.add_row().cells
    row[0].text = (
        "Public/OSL boundary; phase curve; pressure model; temperature-model logic; "
        f"{values['wells']} public wells; {values['ggd_controls']} controls; "
        f"{values['profiles']} profiles; guarded writer."
    )
    row[1].text = (
        f"{values['temp_calculated']} temperature key depths, "
        f"{values['temp_extrapolated']} extrapolated key depths, and "
        f"{values['screen_calculated']} baseline stability-admissibility intervals."
    )
    row[2].text = (
        f"{values['temp_blocked']} temperature key-depth rows, "
        f"{values['screen_blocked']} screen rows, approved logs/core/NMR, "
        "trained ML metrics, saturation outputs, hydrate proof, and sweet-spot ranking."
    )

    document.add_heading("Required Guardrail Language", level=1)
    document.add_paragraph(
        "This is a workflow and admissibility-screen diagram. It does not claim final hydrate stability, "
        "hydrate proof, hydrate saturation, validated ML performance, producibility, or sweet spots."
    )

    document.add_heading("Mentor Decisions Needed", level=1)
    for item in [
        "Phase-curve policy: keep methane 5 ppt as the only official baseline, or add a clearly labeled scenario table?",
        "Target authority: which occurrence and saturation labels are official when Sgh, S_h, Sh, NMR_SAT, hydrate saturation, or phase labels differ?",
        "Validation split: should final testing use whole-well holdout, compartment holdout, geographic holdout, or a staged combination?",
        "Temperature handling: when G10015 is missing, should rows stay blocked, use nearest-control proxy tiers, or use explicit scenario-only gradients?",
        "ML use of stability: is the stability screen allowed as context, confidence, reason flag, or mask only, never as an occurrence label?",
        "Public website outputs: which diagrams, counts, schema, caveat views, and readiness views are acceptable before approved model validation?",
    ]:
        document.add_paragraph(item, style="List Number")

    OUT_DOCX.parent.mkdir(parents=True, exist_ok=True)
    document.save(OUT_DOCX)
    return OUT_DOCX


# ---------------------------------------------------------------------------
# V5.3 mentor-facing rebuild
# ---------------------------------------------------------------------------


def v53_image(path: Path, crop: tuple[int, int, int, int] | None = None) -> Image.Image | None:
    if not path.exists():
        return None
    source = Image.open(path).convert("RGB")
    if crop:
        source = source.crop(crop)
    return source


def v53_paste(
    base: Image.Image,
    path: Path,
    box: tuple[int, int, int, int],
    crop: tuple[int, int, int, int] | None = None,
    mode: str = "cover",
) -> bool:
    source = v53_image(path, crop)
    if source is None:
        return False
    x1, y1, x2, y2 = box
    bw, bh = x2 - x1, y2 - y1
    if mode == "cover":
        source_ratio = source.width / source.height
        box_ratio = bw / bh
        if source_ratio > box_ratio:
            new_w = int(source.height * box_ratio)
            left = (source.width - new_w) // 2
            source = source.crop((left, 0, left + new_w, source.height))
        else:
            new_h = int(source.width / box_ratio)
            top = (source.height - new_h) // 2
            source = source.crop((0, top, source.width, top + new_h))
        resized = source.resize((bw, bh), Image.Resampling.LANCZOS)
        base.paste(resized, (x1, y1))
        return True
    scale = min(bw / source.width, bh / source.height)
    nw, nh = int(source.width * scale), int(source.height * scale)
    resized = source.resize((nw, nh), Image.Resampling.LANCZOS)
    base.paste(resized, (x1 + (bw - nw) // 2, y1 + (bh - nh) // 2))
    return True


def v53_caption(draw: ImageDraw.ImageDraw, xy: tuple[int, int], value: str, width: int) -> None:
    text(draw, xy, value, 11, MUTED, width=width, gap=3)


def v53_section_label(draw: ImageDraw.ImageDraw, xy: tuple[int, int], label: str, color=TEAL) -> None:
    f = font(12, True)
    label_w = int(draw.textlength(label, font=f)) + 24
    card(draw, (xy[0], xy[1], xy[0] + label_w, xy[1] + 30), fill=(246, 251, 252), outline=color, radius=8, width=1)
    draw.text((xy[0] + 12, xy[1] + 8), label, font=f, fill=color)


def v53_panel_title(draw: ImageDraw.ImageDraw, heading: str, subheading: str) -> None:
    text(draw, (60, 38), heading, 38, NAVY, True)
    text(draw, (63, 91), subheading, 17, MUTED, width=1390, gap=4)


def v53_footer(draw: ImageDraw.ImageDraw, source: str) -> None:
    footer(draw, source)


def v53_placeholder_image(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], label: str, accent=TEAL) -> None:
    card(draw, box, fill=LIGHT, outline=accent, radius=10, width=2)
    text(draw, (box[0] + 20, box[1] + 24), label, 18, accent, True, width=box[2] - box[0] - 40, align="center")


def v53_hydrate_structure(draw: ImageDraw.ImageDraw, center: tuple[int, int], label: str, guest: str, color) -> None:
    cx, cy = center
    r = 48
    pts = []
    for i in range(10):
        angle = -math.pi / 2 + i * 2 * math.pi / 10
        pts.append((cx + int(math.cos(angle) * r), cy + int(math.sin(angle) * r)))
    draw.line(pts + [pts[0]], fill=ICE, width=3)
    for px, py in pts:
        draw.ellipse((px - 7, py - 7, px + 7, py + 7), fill=ICE_LIGHT, outline=ICE, width=2)
    draw.ellipse((cx - 31, cy - 31, cx + 31, cy + 31), fill=color, outline=WHITE, width=3)
    text(draw, (cx - 40, cy - 13), guest, 21, WHITE, True, width=80, align="center")
    text(draw, (cx - 90, cy + 64), label, 15, NAVY, True, width=180, align="center")


def v53_phase_plot(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    title_text: str = "P-T stability concept",
    show_interval: bool = False,
) -> None:
    card(draw, box, fill=WHITE, outline=LINE, radius=10, width=2)
    x1, y1, x2, y2 = box
    text(draw, (x1 + 22, y1 + 16), title_text, 18, NAVY, True)
    ax = (x1 + 58, y1 + 66, x2 - 32, y2 - 52)
    draw.line((ax[0], ax[3], ax[2], ax[3]), fill=MUTED, width=3)
    draw.line((ax[0], ax[3], ax[0], ax[1]), fill=MUTED, width=3)
    draw.text((ax[2] - 10, ax[3] + 12), "T", font=font(13, True), fill=MUTED)
    draw.text((ax[0] - 28, ax[1] - 10), "P", font=font(13, True), fill=MUTED)
    stable_poly = [
        (ax[0] + 2, ax[3] - 10),
        (ax[0] + 85, ax[3] - 84),
        (ax[0] + 192, ax[3] - 122),
        (ax[2] - 86, ax[3] - 145),
        (ax[2] - 25, ax[3] - 20),
        (ax[0] + 2, ax[3] - 10),
    ]
    draw.polygon(stable_poly, fill=(213, 242, 247))
    phase = [(ax[0] + 16, ax[3] - 30), (ax[0] + 85, ax[3] - 88), (ax[0] + 160, ax[3] - 118), (ax[2] - 16, ax[3] - 142)]
    temp = [(ax[0] + 52, ax[3] - 8), (ax[0] + 125, ax[3] - 47), (ax[0] + 208, ax[3] - 92), (ax[2] - 46, ax[3] - 136)]
    draw.line(phase, fill=TEAL, width=5)
    draw.line(temp, fill=RED, width=4)
    text(draw, (ax[0] + 130, ax[1] + 62), "admissible", 17, TEAL, True)
    text(draw, (ax[2] - 158, ax[3] - 58), "too warm", 13, RED, True)
    if show_interval:
        draw.rectangle((ax[0] + 93, ax[1] + 84, ax[0] + 177, ax[3] - 32), outline=GREEN, width=4)
        text(draw, (ax[0] + 82, ax[1] + 50), "interval where\nT_model <= T_eq", 13, GREEN, True, width=150, align="center")


def v53_slide_cover(values: dict[str, str]) -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, W, H), fill=(248, 252, 253))
    draw.rectangle((0, 0, 20, H), fill=TEAL)

    text(draw, (64, 72), "North Slope Gas Hydrate\nOccurrence And Saturation ML", 46, NAVY, True, width=760, gap=8)
    text(
        draw,
        (68, 206),
        "A public-safe, source-backed workflow for moving from regional context and approved well-log evidence to future occurrence probability and saturation estimates.",
        23,
        MUTED,
        width=705,
        gap=8,
    )

    card(draw, (66, 356, 720, 552), fill=WHITE, outline=LINE, radius=10, width=2)
    text(draw, (96, 382), "Public build guardrail", 20, RED, True)
    text(
        draw,
        (96, 420),
        "This deck explains workflow readiness. Validated occurrence and saturation predictions wait for approved data, whole-well validation, and mentor review.",
        19,
        NAVY,
        True,
        width=570,
        gap=6,
    )

    stats = [
        (values["wells"], "public wells"),
        (values["profiles"], "G10015 profiles"),
        (values["temp_matches"], "temp-profile matches"),
        (values["screen_calculated"], "admissible intervals"),
    ]
    x = 68
    for number, label in stats:
        card(draw, (x, 580, x + 154, 686), fill=WHITE, outline=LINE, radius=10, width=2)
        text(draw, (x + 15, 604), number, 25, TEAL, True, width=124, align="center")
        text(draw, (x + 15, 642), label, 12, MUTED, True, width=124, align="center")
        x += 168

    map_box = (818, 84, 1536, 462)
    card(draw, (map_box[0] - 6, map_box[1] - 6, map_box[2] + 6, map_box[3] + 34), fill=WHITE, outline=LINE, radius=10, width=2)
    if not v53_paste(img, WEBSITE_CAPTURE_DIR / "02_explore_regional_map.png", map_box, crop=(440, 424, 832, 630)):
        v53_placeholder_image(draw, map_box, "Project website: North Slope public well map")
    v53_caption(draw, (map_box[0], map_box[3] + 9), "Current Streamlit public regional layer overview: well distribution, assessment units, seismic footprints.", map_box[2] - map_box[0])

    struct_box = (818, 528, 1536, 780)
    card(draw, (struct_box[0] - 6, struct_box[1] - 6, struct_box[2] + 6, struct_box[3] + 34), fill=WHITE, outline=LINE, radius=10, width=2)
    if not v53_paste(img, WEBSITE_CAPTURE_DIR / "03_explore_3d_structure.png", struct_box, crop=(398, 330, 1140, 690)):
        v53_placeholder_image(draw, struct_box, "Project website: structural context")
    v53_caption(draw, (struct_box[0], struct_box[3] + 9), "Current Streamlit structural stack preview: public horizons and wells as spatial context only.", struct_box[2] - struct_box[0])

    v53_footer(draw, "Sources: project Streamlit public website captures; public scaffold summaries; no approved/private rows shown.")
    return save(img, "slide_01_cover_v5_3.png")


def v53_slide_context() -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Gas Hydrate + North Slope Context",
        "Definition, hydrate structure type, gas chemistry, and Alaska North Slope setting before model architecture.",
    )

    sem_box = (60, 150, 660, 575)
    card(draw, (sem_box[0] - 6, sem_box[1] - 6, sem_box[2] + 6, sem_box[3] + 70), fill=WHITE, outline=LINE, radius=10, width=2)
    sem_path = REFERENCE_IMAGE_DIR / "usgs_gas_hydrate_crystals_sem_public_domain.jpg"
    if not v53_paste(img, sem_path, sem_box):
        v53_placeholder_image(draw, sem_box, "USGS SEM image unavailable")
    draw.rectangle((sem_box[0], sem_box[1], sem_box[2], sem_box[1] + 70), fill=(11, 35, 48))
    text(draw, (sem_box[0] + 28, sem_box[1] + 17), "USGS SEM: gas hydrate crystals", 25, WHITE, True)
    draw.rectangle((sem_box[0], sem_box[3] - 84, sem_box[2], sem_box[3]), fill=(11, 35, 48))
    text(draw, (sem_box[0] + 26, sem_box[3] - 68), "Ice-like crystalline solid: water cages trap gas molecules.", 20, WHITE, True, width=540)
    v53_caption(draw, (sem_box[0], sem_box[3] + 14), "Source visual: USGS Gas Hydrate Crystals image, public domain.", sem_box[2] - sem_box[0])

    card(draw, (705, 150, 1005, 575), fill=WHITE, outline=LINE, radius=10, width=2)
    text(draw, (728, 173), "Structure type matters", 22, NAVY, True)
    v53_hydrate_structure(draw, (780, 278), "Structure I", "CH4", AMBER)
    v53_hydrate_structure(draw, (930, 278), "Structure II", "C3+", PURPLE)
    v53_hydrate_structure(draw, (855, 435), "Structure H", "mix", TEAL)
    text(draw, (728, 516), "Methane baseline now; C2-C4 and mixed-gas scenarios shift the phase boundary later.", 15, NAVY, True, width=250)

    v53_phase_plot(draw, (1040, 150, 1534, 430), "Methane 5 ppt baseline")
    v53_caption(draw, (1048, 438), "Redrawn concept after USGS SIR 2008-5175 Figure 1A. Stability is necessary, not proof.", 465)

    ns_box = (1040, 515, 1534, 735)
    card(draw, (ns_box[0] - 6, ns_box[1] - 6, ns_box[2] + 6, ns_box[3] + 62), fill=WHITE, outline=LINE, radius=10, width=2)
    if not v53_paste(img, WEBSITE_CAPTURE_DIR / "02_explore_regional_map.png", ns_box, crop=(440, 424, 832, 630)):
        v53_placeholder_image(draw, ns_box, "North Slope public website map")
    text(draw, (1062, 748), "Permafrost-associated North Slope hydrates: Eileen / Mt. Elbert / Ignik Sikumi context where source-supported.", 15, NAVY, True, width=430)

    v53_footer(
        draw,
        "Sources: USGS Gas Hydrate Crystals; USGS SIR 2008-5175; Sloan & Koh hydrate structure descriptions; USGS/DOE North Slope hydrate context; project Streamlit capture.",
    )
    return save(img, "slide_02_gas_hydrate_north_slope_context_v5_3.png")


def v53_range_row(
    draw: ImageDraw.ImageDraw,
    y: int,
    label: str,
    left_label: str,
    right_label: str,
    favorable: tuple[float, float],
    caution: tuple[float, float] | None,
    note: str,
    color=GREEN,
    target: bool = False,
) -> None:
    x0, x1 = 344, 1092
    row_h = 64
    card(draw, (62, y - 10, 1534, y + row_h), fill=(255, 247, 247) if target else WHITE, outline=RED if target else LINE, radius=8, width=2)
    text(draw, (88, y + 5), label, 18, RED if target else NAVY, True, width=210)
    if target:
        text(draw, (344, y + 8), "target-only rail", 18, RED, True, width=180)
        text(draw, (560, y + 8), note, 17, NAVY, True, width=830)
        return
    draw.line((x0, y + 28, x1, y + 28), fill=(220, 234, 238), width=13)
    draw.text((x0 - 6, y + 43), "0", font=font(11, True), fill=MUTED)
    draw.text((x1 - 6, y + 43), "1", font=font(11, True), fill=MUTED)
    draw.text((x0, y + 5), left_label, font=font(13, True), fill=MUTED)
    draw.text((x1 - int(draw.textlength(right_label, font=font(13, True))), y + 5), right_label, font=font(13, True), fill=MUTED)
    fav_x0 = x0 + int((x1 - x0) * favorable[0])
    fav_x1 = x0 + int((x1 - x0) * favorable[1])
    draw.line((fav_x0, y + 28, fav_x1, y + 28), fill=color, width=14)
    draw.polygon([(fav_x1 + 10, y + 28), (fav_x1 - 8, y + 18), (fav_x1 - 8, y + 38)], fill=color)
    if caution:
        cx0 = x0 + int((x1 - x0) * caution[0])
        cx1 = x0 + int((x1 - x0) * caution[1])
        draw.line((cx0, y + 28, cx1, y + 28), fill=AMBER, width=10)
    text(draw, (1140, y + 3), note, 15, NAVY, width=360, gap=4)


def v53_slide_parameter_ranges() -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Parameters And Expected Hydrate Ranges",
        "Relative 0-1 directions only: exact thresholds wait for source-specific units, approved sheets, and mentor-reviewed labels.",
    )
    text(draw, (78, 128), "feature side", 15, GREEN, True)
    text(draw, (1224, 128), "interpretation note", 15, MUTED, True)
    y = 166
    rows = [
        ("Gamma ray / GR", "clean", "shale", (0.02, 0.28), (0.70, 0.96), "lower GR = cleaner sand; higher GR = shale/clay mimic", GREEN),
        ("Deep resistivity / Rt", "wet", "resistive", (0.68, 0.96), (0.03, 0.25), "higher Rt can support hydrate; alone is not proof", RED),
        ("RHOB + density porosity", "tight", "porous", (0.42, 0.76), None, "porosity/storage context; density alone is not target", GREEN),
        ("NMRPHI / separation", "mobile water", "separation", (0.58, 0.88), (0.05, 0.35), "density-vs-NMR separation can matter; high mobile water trends wet", BLUE),
        ("Vp, Vs, Vp/Vs, AI", "soft", "stiff", (0.62, 0.92), (0.08, 0.35), "hydrate/cement can stiffen; gas and poor consolidation mimic", PURPLE),
        ("Caliper / QC", "in-gauge", "washout", (0.00, 0.22), (0.70, 0.96), "QC strip only: washout means caution, not hydrate", AMBER),
        ("P-T stability", "outside", "inside", (0.56, 0.90), (0.02, 0.34), "admissible under assumptions; not occurrence", TEAL),
    ]
    for label, left, right, fav, caution, note, color in rows:
        v53_range_row(draw, y, label, left, right, fav, caution, note, color)
        y += 76
    v53_range_row(
        draw,
        y,
        "Occurrence + saturation",
        "",
        "",
        (0, 0),
        None,
        "Sgh / S_h / Sh / NMR_SAT / hydrate saturation / phase labels stay Y-only and never enter X inputs.",
        RED,
        target=True,
    )
    v53_footer(draw, "Sources: WELL_LOG_REQUIREMENTS_MAP, science-to-ML ladder, approved-data intake spec; directional unless source/data support numeric thresholds.")
    return save(img, "slide_03_parameter_ranges_v5_3.png")


def v53_slide_workflow(values: dict[str, str]) -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Simplified Project Workflow For A Non-ML Audience",
        "Public context explains where hydrate could be physically plausible; approved labeled wells later teach the model occurrence and saturation patterns.",
    )
    steps = [
        ("1. Public context", ["wells", "maps", "permafrost", "phase curve"], TEAL, ICE_LIGHT),
        ("2. Approved", ["logs", "core", "NMR", "labels"], PURPLE, PURPLE_LIGHT),
        ("3. Clean + align", ["depth", "units", "QC"], BLUE, BLUE_LIGHT),
        ("4. Physics", ["stability", "admissible", "caveats"], AMBER, AMBER_LIGHT),
        ("5. ML learns", ["labeled wells", "whole-well split", "baselines"], GREEN, GREEN_LIGHT),
        ("6. Two outputs", ["P(hydrate)", "Sh estimate", "uncertainty"], TEAL, ICE_LIGHT),
        ("7. Review + public", ["validation", "summaries", "no rows"], RED, RED_LIGHT),
    ]
    box_w, gap, x = 198, 22, 68
    centers = []
    for idx, (head, rows, accent, fill) in enumerate(steps):
        box = (x, 184, x + box_w, 390)
        node(draw, box, head, rows, fill, accent)
        centers.append((box[2], (box[1] + box[3]) // 2, box[0]))
        if idx:
            arrow(draw, (x - 19, 287), (x - 5, 287), accent, width=3)
        x += box_w + gap

    map_box = (78, 455, 610, 735)
    card(draw, (map_box[0] - 6, map_box[1] - 6, map_box[2] + 6, map_box[3] + 42), fill=WHITE, outline=LINE, radius=10, width=2)
    if not v53_paste(img, WEBSITE_CAPTURE_DIR / "02_explore_regional_map.png", map_box, crop=(440, 424, 832, 630)):
        v53_placeholder_image(draw, map_box, "Website map capture")
    v53_caption(draw, (map_box[0], map_box[3] + 12), "Public context capture: well distribution and regional layer concept, without approved rows.", map_box[2] - map_box[0])

    card(draw, (650, 455, 1090, 735), fill=WHITE, outline=LINE, radius=10, width=2)
    text(draw, (684, 486), "What machine learning means here", 22, NAVY, True, width=370)
    text(draw, (684, 536), "The model is not guessing from a map. It learns relationships between approved input logs and mentor-approved labels, then gets tested on wells it did not train on.", 20, NAVY, True, width=360, gap=7)

    card(draw, (1130, 455, 1530, 735), fill=(255, 247, 247), outline=RED, radius=10, width=2)
    text(draw, (1162, 486), "Target labels bypass X", 22, RED, True)
    text(draw, (1162, 532), "Sgh, S_h, Sh, NMR_SAT, hydrate saturation, Swr/S_wr, and phase labels supervise training and validation only.", 20, NAVY, True, width=320, gap=7)
    dashed_line(draw, (1328, 455), (1328, 390), RED, width=3, dash=12, gap=8)
    arrow(draw, (1328, 390), (1328, 312), RED, width=3)

    v53_footer(draw, f"Website now: {values['wells']} public wells and public scaffold counts. Approved rows and trained ML outputs stay out of public artifacts.")
    return save(img, "slide_04_simplified_workflow_v5_3.png")


def v53_track(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], label: str, color, phase: float, good: tuple[int, int]) -> None:
    x1, y1, x2, y2 = box
    card(draw, box, fill=WHITE, outline=LINE, radius=8, width=2)
    text(draw, (x1 + 12, y1 + 10), label, 16, NAVY, True)
    ax = (x1 + 46, y1 + 48, x2 - 35, y2 - 28)
    draw.line((ax[0], ax[1], ax[0], ax[3]), fill=(205, 223, 229), width=2)
    draw.rectangle((ax[0] - 12, good[0], ax[2] + 12, good[1]), fill=(230, 248, 239), outline=GREEN, width=2)
    pts = []
    for i in range(92):
        t = i / 91
        yy = ax[1] + int(t * (ax[3] - ax[1]))
        xx = ax[0] + int((0.45 + 0.34 * math.sin(t * 8.0 + phase) + 0.12 * math.sin(t * 22.0 + phase)) * (ax[2] - ax[0]))
        pts.append((xx, yy))
    draw.line(pts, fill=color, width=4)


def v53_slide_evidence_panel() -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Parameter Evidence Visuals",
        "Conceptual normalized tracks show stacked evidence. No single parameter proves hydrate.",
    )
    track_top, track_bottom = 170, 680
    tracks = [
        ((70, track_top, 285, track_bottom), "GR", AMBER, 0.1, (270, 380)),
        ((310, track_top, 525, track_bottom), "Rt", RED, 1.1, (245, 350)),
        ((550, track_top, 765, track_bottom), "Vp / Vs / AI", PURPLE, 2.2, (235, 380)),
        ((790, track_top, 1005, track_bottom), "NMR vs density", BLUE, 3.1, (330, 465)),
        ((1030, track_top, 1245, track_bottom), "P-T", TEAL, 4.4, (260, 430)),
    ]
    for args in tracks:
        v53_track(draw, *args)
    labels = [
        (118, 715, "clean", GREEN),
        (350, 715, "resistive", RED),
        (590, 715, "stiff", PURPLE),
        (822, 715, "separation", BLUE),
        (1082, 715, "stable", TEAL),
    ]
    for lx, ly, label, color in labels:
        pill(draw, (lx, ly, lx + 128, ly + 34), label, color, WHITE)

    card(draw, (1285, 170, 1530, 680), fill=LIGHT, outline=LINE, radius=10, width=2)
    text(draw, (1310, 200), "Evidence stack", 23, NAVY, True, width=190)
    stack = [("clean", GREEN), ("resistive", RED), ("stiff", PURPLE), ("separation", BLUE), ("stable", TEAL)]
    y = 260
    for label, color in stack:
        pill(draw, (1312, y, 1488, y + 38), label, color, WHITE)
        y += 52
    text(draw, (1312, 545), "stacked evidence -> stronger candidate", 18, NAVY, True, width=175, align="center")
    text(draw, (1312, 620), "one alone -> weak or ambiguous", 17, RED, True, width=175, align="center")

    v53_footer(draw, "Conceptual normalized tracks only; no approved/private row data and no project hydrate predictions shown.")
    return save(img, "slide_05_parameter_evidence_visuals_v5_3.png")


def v53_slide_stability(values: dict[str, str]) -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Stability Physics",
        "Pressure, temperature, salinity, and gas chemistry define admissibility under assumptions. They do not prove occurrence.",
    )
    v53_phase_plot(draw, (72, 168, 770, 705), "Methane 5 ppt P-T gate", show_interval=True)

    card(draw, (830, 168, 1528, 365), fill=WHITE, outline=LINE, radius=10, width=2)
    text(draw, (866, 200), "One clean equation line", 23, NAVY, True)
    text(draw, (866, 256), "P_abs(z) = P_surface + rho_w * g * z_m / 1e6", 27, TEAL, True, width=605)
    text(draw, (866, 310), "Depth increases pressure; pressure and gas/salinity control T_eq.", 18, NAVY, width=570)

    cards = [
        ((830, 405, 1054, 590), "Inputs", ["depth basis", "temperature model", "phase curve"], TEAL, ICE_LIGHT),
        ((1074, 405, 1298, 590), "Decision", ["T_model <= T_eq", "interval candidate", "confidence flag"], AMBER, AMBER_LIGHT),
        ((1318, 405, 1528, 590), "Limits", ["not proof", "not saturation", "not final"], RED, RED_LIGHT),
    ]
    for box, head, rows, accent, fill in cards:
        node(draw, box, head, rows, fill, accent)
    arrow(draw, (1054, 498), (1074, 498), TEAL)
    arrow(draw, (1298, 498), (1318, 498), AMBER)

    card(draw, (830, 636, 1528, 745), fill=(255, 247, 247), outline=RED, radius=10, width=2)
    text(draw, (862, 662), "Guardrail", 20, RED, True)
    text(draw, (1000, 655), f"Current public run reports {values['screen_calculated']} baseline admissibility intervals, not hydrate occurrence or saturation.", 24, NAVY, True, width=470)

    v53_footer(draw, "Sources: USGS SIR 2008-5175 methane 5 ppt phase boundary; project public stability products; redrawn schematic.")
    return save(img, "slide_06_stability_physics_v5_3.png")


def v53_slide_ml_architecture() -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Machine Learning Architecture For Beginners",
        "Inputs enter the feature matrix; target labels travel on a separate rail and supervise two future outputs.",
    )

    input_x = 190
    hidden1_x = 530
    hidden2_x = 790
    output_x = 1140
    node_y = [220, 305, 390, 475, 560]
    input_labels = ["GR / Rt", "porosity", "Vp / Vs / AI", "QC flags", "P-T context"]
    for y, label in zip(node_y, input_labels, strict=True):
        draw.ellipse((input_x - 27, y - 27, input_x + 27, y + 27), fill=WHITE, outline=TEAL, width=4)
        text(draw, (58, y - 12), label, 16, NAVY, True, width=110, align="right")
    for x, count in [(hidden1_x, 6), (hidden2_x, 5)]:
        for i in range(count):
            y = 210 + i * (390 // max(1, count - 1))
            draw.ellipse((x - 22, y - 22, x + 22, y + 22), fill=WHITE, outline=PURPLE, width=4)
    for y in node_y:
        for i in range(6):
            hy = 210 + i * (390 // 5)
            draw.line((input_x + 27, y, hidden1_x - 22, hy), fill=(218, 230, 235), width=1)
    for i in range(6):
        y1 = 210 + i * (390 // 5)
        for j in range(5):
            y2 = 210 + j * (390 // 4)
            draw.line((hidden1_x + 22, y1, hidden2_x - 22, y2), fill=(218, 230, 235), width=1)
    for y in [318, 468]:
        for j in range(5):
            hy = 210 + j * (390 // 4)
            draw.line((hidden2_x + 22, hy, output_x - 30, y), fill=(218, 230, 235), width=1)
    for out_y, label, color in [(318, "Occurrence\nP(hydrate)", BLUE), (468, "Saturation\nSh estimate", GREEN)]:
        card(draw, (output_x - 30, out_y - 62, output_x + 245, out_y + 62), fill=WHITE, outline=color, radius=10, width=3)
        text(draw, (output_x, out_y - 34), label, 22, color, True, width=210, align="center", gap=5)

    card(draw, (90, 670, 760, 775), fill=LIGHT, outline=LINE, radius=10, width=2)
    text(draw, (120, 694), "Model ladder", 20, NAVY, True)
    text(draw, (300, 690), "physics/rules baseline -> tree/boosting baseline -> ANN/Keras candidate", 23, NAVY, True, width=400)

    card(draw, (850, 670, 1530, 775), fill=(255, 247, 247), outline=RED, radius=10, width=2)
    text(draw, (880, 694), "Target-only rail", 20, RED, True)
    text(draw, (1070, 688), "occurrence labels + Sgh / S_h / Sh / NMR_SAT supervise training and validation; never enter X inputs.", 22, NAVY, True, width=400)
    dashed_line(draw, (1218, 670), (1218, 530), RED, width=3, dash=12, gap=8)
    arrow(draw, (1218, 530), (1218, 468), RED, width=3)

    text(draw, (340, 126), "Split by whole well or geography before scaling", 19, AMBER, True, width=600, align="center")
    v53_footer(draw, "Architecture anchor: Chong et al. 2024 / USGS supports ANN-style occurrence and saturation modeling; this project has no trained metrics yet.")
    return save(img, "slide_07_ml_architecture_beginner_v5_3.png")


def v53_slide_validation() -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Validation, Uncertainty, And Outputs",
        "The reviewed output package must separate future prediction forms from validated project results.",
    )

    card(draw, (70, 160, 555, 405), fill=WHITE, outline=LINE, radius=10, width=2)
    text(draw, (102, 190), "Whole-well split", 22, NAVY, True)
    for i in range(36):
        col = i % 12
        row = i // 12
        x = 115 + col * 32
        y = 250 + row * 44
        color = GREEN if i % 5 else AMBER
        if i in {3, 9, 16, 28, 31, 34, 35}:
            color = BLUE
        draw.ellipse((x - 9, y - 9, x + 9, y + 9), fill=color, outline=WHITE, width=2)
    pill(draw, (108, 355, 238, 389), "train", GREEN, WHITE)
    pill(draw, (258, 355, 388, 389), "validation", AMBER, NAVY)
    pill(draw, (408, 355, 522, 389), "blind", BLUE, WHITE)

    card(draw, (595, 160, 1045, 405), fill=WHITE, outline=LINE, radius=10, width=2)
    text(draw, (625, 190), "Output form only", 22, NAVY, True)
    draw.line((655, 322, 990, 322), fill=LINE, width=3)
    pts = [(655, 312), (705, 290), (755, 274), (805, 305), (860, 250), (920, 236), (990, 260)]
    draw.line(pts, fill=BLUE, width=5)
    text(draw, (660, 340), "P(hydrate)", 17, BLUE, True)
    draw.rectangle((660, 240, 985, 292), outline=GREEN, width=3)
    text(draw, (760, 245), "review interval", 15, GREEN, True)

    card(draw, (1085, 160, 1530, 405), fill=WHITE, outline=LINE, radius=10, width=2)
    text(draw, (1115, 190), "Saturation estimate", 22, NAVY, True)
    ax = (1130, 245, 1490, 342)
    draw.rectangle(ax, outline=LINE, width=2)
    for x in range(ax[0] + 15, ax[2] - 15, 28):
        yy = ax[1] + 54 + int(20 * math.sin(x / 26))
        draw.line((x, yy - 24, x, yy + 24), fill=(178, 220, 205), width=4)
        draw.ellipse((x - 4, yy - 4, x + 4, yy + 4), fill=GREEN)
    text(draw, (1160, 356), "Sh_pred + uncertainty", 17, GREEN, True)

    bottom_cards = [
        ((70, 470, 408, 720), "Review against", ["core / NMR labels", "well-log labels", "residual plots", "mentor intervals"], TEAL, ICE_LIGHT),
        ((438, 470, 776, 720), "Mimic flags", ["shale / clay", "free gas", "cementation", "washout"], AMBER, AMBER_LIGHT),
        ((806, 470, 1144, 720), "Public-safe later", ["aggregate summaries", "maps", "method plots", "no approved rows"], GREEN, GREEN_LIGHT),
        ((1174, 470, 1530, 720), "Never claim early", ["no fake predictions", "no private rows", "no final proof", "no fake metrics"], RED, RED_LIGHT),
    ]
    for box, head, rows, accent, fill in bottom_cards:
        node(draw, box, head, rows, fill, accent)

    v53_footer(draw, "Validation policy: target labels and held-out wells decide performance after approved data arrive; current slides show workflow only.")
    return save(img, "slide_08_validation_uncertainty_outputs_v5_3.png")


def v53_slide_status(values: dict[str, str]) -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Current Status And Mentor Decisions",
        "Completed public workflow pieces are separated from blocked/future approved-data modeling steps.",
    )

    status_cards = [
        (
            (70, 160, 520, 470),
            "Complete",
            ["public / OSL boundary", "Streamlit public delivery surface", "public well map + structural explorer", "ML intake and leakage skeleton"],
            GREEN,
            GREEN_LIGHT,
        ),
        (
            (575, 160, 1025, 470),
            "Calculated public scaffold",
            [f"{values['wells']} public wells", "43 GGD223 controls", "184 G10015 profiles", f"{values['temp_matches']} temp-profile matches", f"{values['screen_calculated']} methane 5 ppt intervals"],
            BLUE,
            BLUE_LIGHT,
        ),
        (
            (1080, 160, 1530, 470),
            "Blocked / future",
            ["approved logs/core/NMR rows", "authoritative target labels", "blind validation split", "trained occurrence/saturation models", "public prediction release"],
            RED,
            RED_LIGHT,
        ),
    ]
    for box, head, rows, accent, fill in status_cards:
        node(draw, box, head, rows, fill, accent)

    card(draw, (70, 535, 1530, 785), fill=WHITE, outline=BLUE, radius=10, width=2)
    text(draw, (100, 562), "Mentor questions to keep in blue", 24, BLUE, True)
    qs = [
        "Methane 5 ppt baseline only, or add clearly labeled gas-chemistry scenarios?",
        "Which saturation label is authoritative when Sgh, S_h, Sh, NMR_SAT, or hydrate saturation differ?",
        "Should occurrence come from phase labels, saturation threshold, or mentor-reviewed intervals?",
        "MTE = Mt. Elbert and IGS = Ignik Sikumi; confirm *_refined workbook-stage metadata.",
        "Confirm whole-well blind validation policy, including the final holdout wells once all wells arrive.",
        "Are missing-log adapters allowed, or should models use only observed log combinations first?",
    ]
    y = 610
    for q in qs:
        draw.ellipse((105, y + 6, 117, y + 18), fill=BLUE)
        y = text(draw, (132, y), q, 18, NAVY, width=1320, gap=5)
    v53_footer(draw, "Status language: readiness and admissibility only; no final hydrate proof, saturation predictions, or ML metrics.")
    return save(img, "slide_09_status_mentor_decisions_v5_3.png")


def v53_draw_simplified_workflow(values: dict[str, str]) -> Path:
    return v53_slide_workflow(values)


def v53_draw_expanded_architecture(values: dict[str, str]) -> Path:
    img = Image.new("RGB", (EXPANDED_W, EXPANDED_H), WHITE)
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 34, EXPANDED_H), fill=TEAL)
    draw.rectangle((0, 0, EXPANDED_W, 230), fill=(246, 251, 252))
    draw.line((86, 230, EXPANDED_W - 86, 230), fill=LINE, width=4)
    text(draw, (95, 56), "North Slope Gas Hydrate ML Workflow V5.3", 66, NAVY, True)
    text(draw, (100, 137), "Audience-facing architecture with public context, approved runtime data, stability admissibility, leakage barrier, two-output ML, and review gates.", 31, MUTED, width=4450)

    cols = [
        ("Public Context", TEAL),
        ("Approved Runtime Later", PURPLE),
        ("QC + Feature Matrix", GREEN),
        ("Leakage-Safe Modeling", BLUE),
        ("Review + Public Summaries", AMBER),
    ]
    col_x = [100, 1120, 2140, 3160, 4180]
    col_w = 900
    for i, (heading, accent) in enumerate(cols):
        x = col_x[i]
        card(draw, (x, 300, x + col_w, 2715), fill=(249, 252, 252), outline=(215, 232, 236), radius=24, width=2)
        draw.rectangle((x + 30, 330, x + 44, 390), fill=accent)
        text(draw, (x + 68, 326), heading, 33, accent, True, width=col_w - 110)

    boxes: dict[str, tuple[int, int, int, int]] = {}

    def big_node(key: str, col: int, y: int, h: int, heading: str, rows: list[str], accent, fill=WHITE) -> None:
        x = col_x[col] + 42
        box = (x, y, x + col_w - 84, y + h)
        boxes[key] = box
        card(draw, box, fill=fill, outline=accent, radius=18, width=3)
        draw.rounded_rectangle((box[0], box[1], box[0] + 14, box[3]), radius=9, fill=accent)
        y2 = text(draw, (box[0] + 38, box[1] + 22), heading, 29, accent, True, width=box[2] - box[0] - 70, gap=6)
        y2 += 8
        for row in rows:
            y2 = text(draw, (box[0] + 42, y2), row, 22, NAVY, width=box[2] - box[0] - 82, gap=6)

    big_node("pub", 0, 455, 300, "Public GIS + source controls", [f"{values['wells']} public wells", "USGS hydrate AUs; DNR wells", "GGD223 permafrost + G10015 temperature", "source-backed captions and caveats"], TEAL, ICE_LIGHT)
    big_node("stab", 0, 840, 340, "Stability-admissibility branch", ["P_abs(z) = P_surface + rho_w*g*z_m/1e6", "T_model(z) from public temperature profiles", "methane 5 ppt phase lookup", "context/mask/caveat only"], BLUE, BLUE_LIGHT)
    big_node("website", 0, 1285, 260, "Website now", ["2D map and structural explorer", "schema readiness and target leakage views", "public-safe method diagrams"], GREEN, GREEN_LIGHT)
    big_node("approved", 1, 455, 350, "Approved logs/core/NMR later", ["preserve sheet/file names and original headers", "visible subset: about 3 / 71 datasets", "headers sufficient for schema design", "not enough for final training or metrics"], PURPLE, PURPLE_LIGHT)
    big_node("labels", 1, 900, 360, "Target-only labels", ["Sgh, S_h, Sh, NMR_SAT", "Hydrate Saturation, Swr, S_wr", "occurrence/phase/core labels", "supervise training and validation only"], RED, RED_LIGHT)
    big_node("unit", 2, 455, 335, "Unit + QC normalization", ["depth units retained; depth aligns rows", "density, sonic/velocity, porosity checks", "caliper coverage before washout filtering", "numeric predictors scaled 0-1 after split"], GREEN, GREEN_LIGHT)
    big_node("features", 2, 880, 395, "Feature families", ["GR / clean-sand proxy", "resistivity transforms", "density porosity and NMR separation", "Vp, Vs, Vp/Vs, AI, elastic terms where units support", "optional stability/context fields"], TEAL, ICE_LIGHT)
    big_node("barrier", 2, 1390, 260, "Leakage barrier", ["target columns bypass X_allowed", "unknown variables stay unresolved", "adapters require approval"], RED, RED_LIGHT)
    big_node("split", 3, 455, 300, "Whole-well split first", ["train/validation/test by well or geography", "freeze holdout before preprocessing", "avoid random depth-row final split"], BLUE, BLUE_LIGHT)
    big_node("models", 3, 850, 395, "Model ladder", ["simple physics/rules baseline", "tree or gradient boosting tabular baseline", "ANN/Keras candidate after baselines", "model must beat simple baselines"], PURPLE, PURPLE_LIGHT)
    big_node("outputs", 3, 1370, 300, "Two future outputs", ["occurrence probability P(hydrate)", "saturation regression Sh_pred", "uncertainty and reason flags"], GREEN, GREEN_LIGHT)
    big_node("validate", 4, 455, 395, "Validation against Y-only labels", ["core/NMR/log labels", "calibration and residual review", "mimic review: shale, free gas, cement, washout", "no fake metrics"], AMBER, AMBER_LIGHT)
    big_node("public", 4, 955, 380, "Public-safe after review", ["aggregate maps and summaries", "method plots and caveat counts", "no approved/private rows", "no final prediction release without mentor approval"], TEAL, ICE_LIGHT)
    big_node("decisions", 4, 1445, 350, "Mentor decisions", ["phase-curve scenarios", "target authority", "occurrence label policy", "MTE/IGS refined metadata", "blind validation wells", "missing-log adapters"], BLUE, BLUE_LIGHT)

    def connect(a: str, b: str, color=TEAL, dashed: bool = False) -> None:
        ba, bb = boxes[a], boxes[b]
        start = (ba[2] + 12, (ba[1] + ba[3]) // 2)
        end = (bb[0] - 12, (bb[1] + bb[3]) // 2)
        if dashed:
            dashed_line(draw, start, end, color, width=5, dash=22, gap=15)
            arrow(draw, (end[0] - 5, end[1]), end, color, width=1)
        else:
            arrow(draw, start, end, color, width=5)

    connect("pub", "approved", TEAL)
    connect("stab", "features", BLUE)
    connect("approved", "unit", PURPLE)
    connect("labels", "validate", RED, dashed=True)
    connect("unit", "split", GREEN)
    connect("features", "barrier", TEAL)
    connect("barrier", "models", RED)
    connect("split", "models", BLUE)
    connect("models", "outputs", PURPLE)
    connect("outputs", "validate", GREEN)
    connect("validate", "public", AMBER)
    connect("public", "decisions", TEAL)

    card(draw, (120, 2240, 5080, 2645), fill=(255, 247, 247), outline=RED, radius=22, width=3)
    text(draw, (165, 2290), "Guardrail contract", 36, RED, True)
    text(draw, (610, 2278), "Stability = admissible under assumptions, not occurrence or saturation. Target-only fields never enter predictors. Public outputs stay aggregate/method-level until approved validation and mentor review.", 39, NAVY, True, width=4200, gap=10)

    draw.line((90, 2850, EXPANDED_W - 90, 2850), fill=LINE, width=4)
    text(draw, (100, 2885), "Sources: project base, approved-data intake spec, first model experiment plan, USGS SIR 2008-5175, Chong et al. 2024 / USGS, OSTI/NETL MTE-IGS context, MDPI missing-log adapter literature.", 25, MUTED, width=4900)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / "full_project_ml_workflow_flowchart_expanded.png"
    img.save(path, quality=94)
    return path


def ml_network_detail_panel() -> Path:
    img = Image.new("RGB", (NETWORK_W, NETWORK_H), WHITE)
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 26, NETWORK_H), fill=TEAL)
    draw.rectangle((0, 0, NETWORK_W, 160), fill=(246, 251, 252))
    text(draw, (70, 40), "ML Runtime Detail V5.3", 48, NAVY, True)
    text(draw, (74, 104), "Feature matrix, target-only rail, split-first preprocessing, baselines, ANN candidate, and validation gates.", 27, MUTED, width=2300)

    groups = [
        ("Measured inputs", ["GR", "Rt / RES", "RHOB", "NMRPHI", "Vp / Vs", "caliper"], TEAL, 110, 260),
        ("Derived features", ["Vsh / clean sand", "density porosity", "log Rt", "Vp/Vs", "AI / elastic", "stability context"], GREEN, 110, 635),
        ("QC + context", ["unit flags", "depth alignment", "washout status", "missingness", "AU/permafrost", "caveats"], AMBER, 110, 1010),
    ]
    for heading, rows, accent, x, y in groups:
        card(draw, (x, y, x + 430, y + 285), fill=WHITE, outline=accent, radius=18, width=3)
        text(draw, (x + 30, y + 25), heading, 30, accent, True)
        yy = y + 78
        for row in rows:
            pill(draw, (x + 32, yy, x + 188, yy + 34), row, (246, 251, 252), accent)
            yy += 39

    matrix = (720, 410, 1100, 890)
    card(draw, matrix, fill=LIGHT, outline=GREEN, radius=18, width=4)
    text(draw, (785, 445), "X_allowed", 40, GREEN, True, width=250, align="center")
    for r in range(13):
        for c in range(8):
            shade = (215 + (r + c) % 3 * 8, 235, 237)
            draw.rectangle((760 + c * 38, 525 + r * 24, 792 + c * 38, 542 + r * 24), fill=shade, outline=WHITE)
    text(draw, (760, 850), "0-1 numeric scaling after split", 24, NAVY, True, width=310, align="center")
    for _, _, accent, x, y in groups:
        arrow(draw, (x + 430, y + 142), (matrix[0] - 16, (matrix[1] + matrix[3]) // 2), accent, width=5)

    split_box = (1235, 295, 1580, 520)
    card(draw, split_box, fill=BLUE_LIGHT, outline=BLUE, radius=18, width=3)
    text(draw, (1268, 328), "Split first", 34, BLUE, True)
    text(draw, (1268, 390), "whole well / geography\nthen fit preprocessing", 27, NAVY, True, width=280, gap=8)
    arrow(draw, (matrix[2] + 16, 620), (split_box[0] - 16, 405), GREEN, width=5)

    baseline_box = (1235, 620, 1580, 910)
    card(draw, baseline_box, fill=PURPLE_LIGHT, outline=PURPLE, radius=18, width=3)
    text(draw, (1268, 650), "Model ladder", 34, PURPLE, True)
    text(draw, (1268, 712), "rules baseline\n-> tree / boosting\n-> ANN / Keras", 28, NAVY, True, width=280, gap=8)
    arrow(draw, (split_box[0] + 170, split_box[3] + 18), (baseline_box[0] + 170, baseline_box[1] - 18), BLUE, width=5)

    ann = (1735, 430, 2100, 820)
    card(draw, ann, fill=WHITE, outline=PURPLE, radius=18, width=3)
    text(draw, (1790, 454), "ANN candidate", 32, PURPLE, True)
    layers = [(1775, 5, TEAL), (1860, 6, PURPLE), (1945, 5, PURPLE), (2030, 2, AMBER)]
    layer_pts = []
    for x, count, color in layers:
        pts = [(x, 525 + int(i * 230 / max(1, count - 1))) for i in range(count)]
        layer_pts.append((pts, color))
    for (pts_a, _), (pts_b, _) in zip(layer_pts, layer_pts[1:], strict=False):
        for a in pts_a:
            for b in pts_b:
                draw.line((a, b), fill=(220, 230, 234), width=1)
    for pts, color in layer_pts:
        for px, py in pts:
            draw.ellipse((px - 14, py - 14, px + 14, py + 14), fill=WHITE, outline=color, width=4)
    arrow(draw, (baseline_box[2] + 18, 765), (ann[0] - 18, 625), PURPLE, width=5)

    out1 = (2240, 370, 2520, 560)
    out2 = (2240, 690, 2520, 880)
    for box, label, color in [(out1, "Occurrence\nP(hydrate)", BLUE), (out2, "Saturation\nSh_pred", GREEN)]:
        card(draw, box, fill=WHITE, outline=color, radius=18, width=4)
        text(draw, (box[0] + 30, box[1] + 48), label, 34, color, True, width=220, align="center", gap=8)
        arrow(draw, (ann[2] + 18, (ann[1] + ann[3]) // 2), (box[0] - 18, (box[1] + box[3]) // 2), color, width=5)

    rail = (680, 1080, 2520, 1350)
    card(draw, rail, fill=RED_LIGHT, outline=RED, radius=22, width=4)
    text(draw, (720, 1120), "Y-only target rail", 38, RED, True)
    text(draw, (1110, 1110), "occurrence labels, Sgh, S_h, Sh, NMR_SAT, Hydrate Saturation, Swr/S_wr, phase labels", 34, NAVY, True, width=1220)
    text(draw, (1110, 1200), "Targets supervise training and validation but never join X_allowed.", 32, RED, True, width=1120)
    dashed_line(draw, (1880, rail[1]), (1880, out1[3]), RED, width=5, dash=24, gap=16)
    dashed_line(draw, (1880, rail[1]), (1880, out2[3]), RED, width=5, dash=24, gap=16)

    text(draw, (80, 1435), "No final training, metrics, occurrence prediction, or saturation result is claimed until approved labels and whole-well validation are available.", 25, MUTED, width=2380)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / "ml_pipeline_network_detail_v5.png"
    img.save(path, quality=94)
    return path


def v53_appendix_slide(source_path: Path, filename: str, heading: str, subheading: str) -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(draw, heading, subheading)
    box = (64, 142, 1536, 812)
    card(draw, (box[0] - 8, box[1] - 8, box[2] + 8, box[3] + 8), fill=WHITE, outline=LINE, radius=10, width=2)
    if source_path.exists():
        v53_paste(img, source_path, box, mode="contain")
    else:
        v53_placeholder_image(draw, box, "Reference image unavailable")
    v53_footer(draw, "Technical appendix slide included intact for detailed review; use main slides for audience explanation.")
    return save(img, filename)


def build_panels(network_path: Path) -> list[Path]:
    values = summaries()
    expanded = v53_draw_expanded_architecture(values)
    panels = [
        v53_slide_cover(values),
        v53_slide_context(),
        v53_slide_parameter_ranges(),
        v53_draw_simplified_workflow(values),
        v53_slide_evidence_panel(),
        v53_slide_stability(values),
        v53_slide_ml_architecture(),
        v53_slide_validation(),
        v53_slide_status(values),
        v53_appendix_slide(
            expanded,
            "slide_10_expanded_architecture_reference_v5_3.png",
            "Appendix A: Expanded Workflow Architecture",
            "The complex full-workflow reference remains intact for technical discussion after the audience slides.",
        ),
        v53_appendix_slide(
            network_path,
            "slide_11_ml_runtime_detail_reference_v5_3.png",
            "Appendix B: ML Runtime Detail",
            "The detailed runtime diagram keeps X inputs, Y-only targets, model order, and validation gates together.",
        ),
    ]
    return panels


def rebuild_deck(panel_paths: list[Path]) -> Path:
    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 5143500
    blank = prs.slide_layouts[6]
    for panel in panel_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(panel), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(OUT_DECK)
    return OUT_DECK


def verify_deck(path: Path) -> None:
    prs = Presentation(path)
    if len(prs.slides) < 9:
        raise ValueError(f"Expected at least 9 slides, found {len(prs.slides)}")
    for idx, slide in enumerate(prs.slides, start=1):
        pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
        if len(pictures) != 1:
            raise ValueError(f"Slide {idx} should have exactly one raster panel, found {len(pictures)}")
        blob = pictures[0].image.blob
        with Image.open(BytesIO(blob)) as image:
            if image.size != (W, H):
                raise ValueError(f"Slide {idx} image size is {image.size}, expected {(W, H)}")


def build_contact_sheet(panel_paths: list[Path]) -> Path:
    thumb_w, thumb_h = 480, 270
    margin = 34
    label_h = 34
    gap = 26
    columns = 3
    rows = math.ceil(len(panel_paths) / columns)
    sheet_w = margin * 2 + thumb_w * columns + gap * (columns - 1)
    sheet_h = margin * 2 + (thumb_h + label_h) * rows + gap * (rows - 1)
    img = Image.new("RGB", (sheet_w, sheet_h), WHITE)
    draw = ImageDraw.Draw(img)
    for idx, path in enumerate(panel_paths, start=1):
        row = (idx - 1) // columns
        col = (idx - 1) % columns
        x = margin + col * (thumb_w + gap)
        y = margin + row * (thumb_h + label_h + gap)
        panel = Image.open(path).convert("RGB").resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        card(draw, (x - 4, y - 4, x + thumb_w + 4, y + thumb_h + label_h), fill=(248, 251, 252), outline=LINE, radius=10, width=2)
        img.paste(panel, (x, y))
        draw.rectangle((x, y, x + thumb_w, y + thumb_h), outline=(160, 190, 198), width=2)
        text(draw, (x + 14, y + thumb_h + 9), f"Slide {idx}", 18, NAVY, True, width=thumb_w - 28)
    OUT_CONTACT_SHEET.parent.mkdir(parents=True, exist_ok=True)
    img.save(OUT_CONTACT_SHEET, quality=94)
    return OUT_CONTACT_SHEET


def build_word_companion(diagram_path: Path, network_path: Path) -> Path:
    values = summaries()
    document = Document()
    apply_doc_style(document)
    section = document.sections[0]
    section.top_margin = Inches(0.6)
    section.bottom_margin = Inches(0.6)
    section.left_margin = Inches(0.65)
    section.right_margin = Inches(0.65)

    props = document.core_properties
    props.title = "V5.3 North Slope Gas Hydrate ML Workflow Companion"
    props.subject = "Mentor-facing scientific and ML workflow explanation"
    props.author = "North Slope Gas Hydrates project"

    document.add_heading("V5.3 North Slope Gas Hydrate ML Workflow Companion", level=0)
    document.add_paragraph(
        "This companion explains the refreshed mentor-facing deck. It describes the public-safe workflow, "
        "the gas-hydrate science context, the parameter evidence logic, the guarded stability branch, and the "
        "future leakage-safe machine-learning architecture for occurrence classification and saturation regression. "
        "It does not report trained project metrics or final hydrate predictions."
    )

    document.add_heading("1. Project Purpose", level=1)
    document.add_paragraph(
        "The project is building a source-backed North Slope workflow that can later use approved well logs, core, "
        "NMR, and mentor-approved labels to predict gas hydrate occurrence and saturation. The current public work "
        "is readiness and method architecture: it documents what will enter the feature matrix, what stays target-only, "
        "how physics context is used, and how validation should be performed."
    )

    document.add_heading("2. Public vs OSL Boundary", level=1)
    document.add_paragraph(
        "The GitHub and Streamlit side can show public-source context, public scaffold counts, diagrams, source-ledger "
        "logic, schema coverage, and caveats. Approved runtime or OSL space is reserved for raw/approved well-log rows, "
        "core/NMR values, sensitive source files, model fitting, trained artifacts, and reviewed outputs."
    )
    document.add_paragraph(
        f"The public scaffold currently shows {values['wells']} public wells, 43 GGD223 permafrost controls, "
        f"{values['profiles']} G10015 temperature profiles, {values['temp_matches']} temperature-profile matches, "
        f"and {values['screen_calculated']} baseline methane 5 ppt admissibility intervals."
    )

    document.add_heading("3. Gas Hydrate Science And North Slope Context", level=1)
    document.add_paragraph(
        "Gas hydrate is an ice-like crystalline solid in which water cages trap guest gas molecules, most commonly "
        "methane. North Slope hydrate work is permafrost-associated and tied to public context such as hydrate "
        "assessment units, permafrost depth, public wells, and case-study areas including Mt. Elbert and Ignik Sikumi "
        "where supported by source documentation."
    )

    document.add_heading("4. Why Gas Chemistry And Structure Type Matter", level=1)
    document.add_paragraph(
        "Structure I is the methane-dominant baseline used for the current public stability screen. Structure II can be "
        "stabilized by larger molecules such as propane or isobutane, and Structure H involves larger hydrocarbons plus "
        "methane. Because gas composition shifts the pressure-temperature boundary, the current deck labels methane 5 ppt "
        "as the baseline and keeps variable chemistry as a later scenario capability."
    )

    document.add_heading("5. Parameter Evidence Table", level=1)
    table = document.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    for i, header in enumerate(["Parameter family", "Hydrate-compatible direction", "Opposite / mimic", "Caveat"]):
        run = table.rows[0].cells[i].paragraphs[0].add_run(header)
        run.bold = True
    rows = [
        ("GR", "Lower GR, cleaner sand/reservoir quality.", "Higher GR can mean shale/clay and can confuse response.", "Directional only unless thresholds are source-supported."),
        ("Rt / resistivity", "Higher resistivity can support resistive pore fill.", "Lower resistivity trends saline water or wet sand.", "Resistivity alone is not proof."),
        ("RHOB / density porosity", "Porosity/storage context and density-derived porosity.", "Tight or low-quality intervals are less favorable.", "Density is an input or derived feature, not a target."),
        ("NMRPHI / NMR separation", "Density-vs-NMR separation can be meaningful because hydrate may reduce mobile-water signal.", "High mobile water tends water-bearing.", "NMR saturation-like fields stay target-only."),
        ("Vp, Vs, Vp/Vs, AI", "Hydrate-bearing or cemented sand can be stiffer with higher velocity/impedance.", "Free gas, cement, or poor consolidation can mimic or reverse signals.", "Use only when units and logs support the calculation."),
        ("Caliper / QC", "In-gauge hole supports trust in logs.", "Washout or bad hole means caution.", "QC only; not a hydrate signal."),
        ("Pressure-temperature stability", "Inside methane 5 ppt phase boundary is admissible under assumptions.", "Outside is not stable under that assumption.", "Not occurrence and not saturation."),
    ]
    for row_values in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row_values):
            cells[i].text = value

    document.add_heading("6. Stability Method", level=1)
    document.add_paragraph(
        "The public stability branch combines a depth basis, hydrostatic absolute pressure, a public temperature model, "
        "and a methane 5 ppt phase-boundary lookup. The pressure equation shown in the deck is "
        "P_abs(z) = P_surface + rho_w * g * z_m / 1e6. The calculation identifies intervals where the modeled "
        "temperature is at or below the equilibrium temperature for the chosen phase curve."
    )
    document.add_paragraph(
        "The methane 5 ppt baseline is retained as the official public baseline. Scenario chemistry can be added later "
        "only with labeled sources and mentor approval. Stability means physically admissible under assumptions; it is "
        "not hydrate proof, occurrence, saturation, producibility, or final top/base/thickness."
    )

    document.add_heading("7. ML Workflow", level=1)
    for item in [
        "X inputs: measured logs, derived features, QC/context flags, and optional stability/context fields if approved.",
        "Y-only targets: occurrence labels, Sgh, S_h, Sh, NMR_SAT, hydrate saturation, Swr/S_wr, and interpreted phase labels.",
        "Two outputs: future occurrence classification as P(hydrate) and future saturation regression as Sh_pred.",
        "Model order: simple physics/rules baseline, tree or gradient boosting baseline, then ANN/Keras candidate.",
        "Validation: split by whole well, compartment, or geography before preprocessing; avoid final random depth-row splits.",
        "Leakage guardrail: target labels supervise training and validation but never enter the feature matrix.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("8. What The Website Currently Shows", level=1)
    for item in [
        "2D public regional map / North Slope layer overview.",
        "Structural explorer preview with public horizons and wells.",
        "Public scaffold counts and stability-admissibility readiness.",
        "Schema coverage and target-leakage architecture views.",
        "Public-safe diagrams without approved/private data rows.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("9. Complete vs Future Work", level=1)
    status_table = document.add_table(rows=1, cols=3)
    status_table.style = "Table Grid"
    for i, header in enumerate(["Complete", "Calculated", "Future / blocked"]):
        run = status_table.rows[0].cells[i].paragraphs[0].add_run(header)
        run.bold = True
    cells = status_table.add_row().cells
    cells[0].text = "Public/OSL boundary; website surface; public well map; structural explorer; first ML intake and leakage architecture."
    cells[1].text = f"{values['wells']} public wells; 43 controls; {values['profiles']} profiles; {values['temp_matches']} temp-profile matches; {values['screen_calculated']} methane 5 ppt admissibility intervals."
    cells[2].text = "Approved log/core/NMR rows; target authority; final blind validation wells; trained occurrence/saturation models; reviewed public prediction release."

    document.add_heading("10. Mentor Decision List", level=1)
    for item in [
        "Methane 5 ppt baseline only, or add clearly labeled variable gas-chemistry scenarios?",
        "Which saturation label is authoritative when Sgh, S_h, Sh, NMR_SAT, or hydrate saturation differ?",
        "Should occurrence labels come from phase labels, saturation threshold, or mentor-reviewed intervals?",
        "MTE is Mt. Elbert and IGS is Ignik Sikumi; confirm whether *_refined tabs are processing stages by workbook metadata.",
        "Which wells become blind validation once the full well inventory is available?",
        "Are missing-log adapters allowed, or should the first models use only observed log combinations?",
    ]:
        document.add_paragraph(item, style="List Number")

    document.add_heading("Research / Source Anchors", level=1)
    for item in [
        "USGS Gas Hydrate Crystals image: https://www.usgs.gov/media/images/gas-hydrate-crystals",
        "USGS SIR 2008-5175 North Slope stability/prospect method: https://pubs.usgs.gov/sir/2008/5175/",
        "Chong et al. 2024 / USGS ANN occurrence and saturation architecture anchor: https://pubs.usgs.gov/publication/70250169",
        "OSTI / NETL MTE and IGS context: https://www.osti.gov/servlets/purl/1893637",
        "MDPI missing-log adapter background: https://www.mdpi.com/1996-1073/16/23/7709",
    ]:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("Deck Figures", level=1)
    document.add_paragraph("Audience workflow slide:")
    if diagram_path.exists():
        document.add_picture(str(diagram_path), width=Inches(6.7))
    expanded_path = ASSET_DIR / "full_project_ml_workflow_flowchart_expanded.png"
    if expanded_path.exists():
        document.add_paragraph("Expanded technical workflow:")
        document.add_picture(str(expanded_path), width=Inches(6.7))
    if network_path.exists():
        document.add_paragraph("ML runtime detail:")
        document.add_picture(str(network_path), width=Inches(6.7))

    document.add_heading("Guardrail Language", level=1)
    document.add_paragraph(
        "Use 'admissible under assumptions,' 'candidate context,' 'confidence/caveat,' and 'validation required.' "
        "Do not call the stability scaffold hydrate proof, final stability, occurrence, saturation, or trained ML performance."
    )

    OUT_DOCX.parent.mkdir(parents=True, exist_ok=True)
    document.save(OUT_DOCX)
    return OUT_DOCX


def v54_registry_rows() -> dict[str, dict[str, str]]:
    if not PARAMETER_EVIDENCE_REGISTRY.exists():
        return {}
    with PARAMETER_EVIDENCE_REGISTRY.open(newline="", encoding="utf-8") as fh:
        return {row["parameter_family"]: row for row in csv.DictReader(fh)}


def v54_range(row: dict[str, str] | None, fallback: tuple[float, float]) -> tuple[float, float]:
    if not row:
        return fallback
    try:
        return (float(row["hydrate_window_norm_start"]), float(row["hydrate_window_norm_end"]))
    except (KeyError, TypeError, ValueError):
        return fallback


def v54_copy_authority_panel(
    source_path: Path,
    filename: str,
    heading: str | None = None,
    subheading: str | None = None,
) -> Path:
    if source_path.exists():
        img = Image.open(source_path).convert("RGB").resize((W, H), Image.Resampling.LANCZOS)
    else:
        img = canvas(False)
        draw = ImageDraw.Draw(img)
        v53_panel_title(draw, heading or "Source panel unavailable", subheading or "Expected authority visual was not found.")
        v53_placeholder_image(draw, (120, 180, 1480, 710), str(source_path))
    if heading:
        draw = ImageDraw.Draw(img)
        draw.rectangle((0, 0, W, 92), fill=(245, 250, 251))
        draw.rectangle((0, 0, 18, H), fill=TEAL)
        text(draw, (58, 24), heading, 34, NAVY, True)
        if subheading:
            text(draw, (61, 67), subheading, 15, MUTED, width=1410)
    return save(img, filename)


def v54_source_label(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], label: str) -> None:
    x1, y1, x2, y2 = box
    draw.rectangle((x1, y2 - 24, x2, y2), fill=(11, 35, 48))
    text(draw, (x1 + 10, y2 - 21), label, 10, WHITE, True, width=x2 - x1 - 20, gap=1)


def v54_phase_curve(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int]) -> None:
    card(draw, box, fill=WHITE, outline=LINE, radius=10, width=2)
    x1, y1, x2, y2 = box
    text(draw, (x1 + 22, y1 + 16), "Methane 5 ppt stability curve", 20, NAVY, True, width=x2 - x1 - 44)
    rows: list[tuple[float, float]] = []
    if PHASE_CURVE_METHANE_5PPT.exists():
        with PHASE_CURVE_METHANE_5PPT.open(newline="", encoding="utf-8") as fh:
            for row in csv.DictReader(fh):
                try:
                    rows.append((float(row["equilibrium_temperature_c"]), float(row["pressure_mpa_absolute"])))
                except (KeyError, TypeError, ValueError):
                    continue
    ax = (x1 + 60, y1 + 78, x2 - 38, y2 - 58)
    draw.line((ax[0], ax[3], ax[2], ax[3]), fill=MUTED, width=3)
    draw.line((ax[0], ax[3], ax[0], ax[1]), fill=MUTED, width=3)
    draw.text((ax[2] - 78, ax[3] + 14), "Temperature C", font=font(12, True), fill=MUTED)
    draw.text((ax[0] - 46, ax[1] - 10), "MPa", font=font(12, True), fill=MUTED)
    for i in range(1, 4):
        gx = ax[0] + i * (ax[2] - ax[0]) // 4
        gy = ax[3] - i * (ax[3] - ax[1]) // 4
        draw.line((gx, ax[1], gx, ax[3]), fill=(226, 238, 241), width=1)
        draw.line((ax[0], gy, ax[2], gy), fill=(226, 238, 241), width=1)
    if rows:
        temps = [r[0] for r in rows]
        pressures = [r[1] for r in rows]
        t_min, t_max = min(temps) - 1.0, max(temps) + 1.0
        p_min, p_max = min(pressures) - 0.4, max(pressures) + 0.4

        def project(t: float, p: float) -> tuple[int, int]:
            px = ax[0] + int((t - t_min) / (t_max - t_min) * (ax[2] - ax[0]))
            py = ax[3] - int((p - p_min) / (p_max - p_min) * (ax[3] - ax[1]))
            return px, py

        pts = [project(t, p) for t, p in rows]
        if len(pts) > 1:
            draw.line(pts, fill=TEAL, width=4)
        for pt in pts[:: max(1, len(pts) // 14)]:
            draw.ellipse((pt[0] - 3, pt[1] - 3, pt[0] + 3, pt[1] + 3), fill=TEAL)
        text(draw, (ax[0] + 22, ax[1] + 16), "T_model <= T_eq is admissible context", 14, TEAL, True, width=210)
        text(draw, (ax[0] + 20, ax[3] - 34), f"{t_min:.0f}", 10, MUTED)
        text(draw, (ax[2] - 20, ax[3] - 34), f"{t_max:.0f}", 10, MUTED)
    else:
        text(draw, (ax[0] + 28, ax[1] + 58), "Digitized phase curve CSV missing", 17, RED, True, width=280)
    v54_source_label(draw, box, "Lee et al. 2008 USGS SIR 2008-5175 Fig. 1A, digitized methane 5 ppt curve")


def v54_slide_personal_opener() -> Path:
    return v54_copy_authority_panel(
        V52_ASSET_DIR / "slide_01_locked_from_current_gmail_deck.png",
        "slide_01_personal_about_me_v5_4.png",
    )


def v54_slide_context() -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Gas Hydrate And North Slope Context",
        "Source-backed hydrate visuals only: USGS image, public map capture, and digitized methane 5 ppt stability context.",
    )

    sem_box = (66, 145, 642, 520)
    card(draw, (sem_box[0] - 6, sem_box[1] - 6, sem_box[2] + 6, sem_box[3] + 40), fill=WHITE, outline=LINE, radius=10, width=2)
    sem_path = REFERENCE_IMAGE_DIR / "usgs_gas_hydrate_crystals_sem_public_domain.jpg"
    if not v53_paste(img, sem_path, sem_box):
        v53_placeholder_image(draw, sem_box, "USGS hydrate SEM image unavailable")
    draw.rectangle((sem_box[0], sem_box[1], sem_box[2], sem_box[1] + 48), fill=(11, 35, 48))
    text(draw, (sem_box[0] + 22, sem_box[1] + 12), "Gas hydrate crystals", 21, WHITE, True)
    v54_source_label(draw, sem_box, "USGS public domain Gas Hydrate Crystals SEM")

    type_box = (66, 600, 642, 790)
    card(draw, type_box, fill=WHITE, outline=LINE, radius=10, width=2)
    text(draw, (92, 622), "Hydrate structure types", 22, NAVY, True)
    type_rows = [
        ("I", "methane-dominant hydrate", TEAL),
        ("II", "larger guests and mixed gases", PURPLE),
        ("H", "larger hydrocarbon guests", AMBER),
    ]
    tx = 92
    for label, desc, color in type_rows:
        card(draw, (tx, 670, tx + 148, 740), fill=(248, 252, 253), outline=color, radius=8, width=2)
        text(draw, (tx + 14, 683), f"Structure {label}", 17, color, True, width=120, align="center")
        text(draw, (tx + 12, 713), desc, 10, NAVY, True, width=124, align="center", gap=2)
        tx += 166
    text(draw, (92, 746), "Current public method: methane 5 ppt baseline; gas-composition scenarios are future mentor policy.", 11, MUTED, True, width=520)

    v54_phase_curve(draw, (690, 145, 1530, 458))

    map_box = (690, 540, 1530, 780)
    card(draw, (map_box[0] - 6, map_box[1] - 6, map_box[2] + 6, map_box[3] + 40), fill=WHITE, outline=LINE, radius=10, width=2)
    if not v53_paste(img, WEBSITE_CAPTURE_DIR / "02_explore_regional_map.png", map_box, crop=(360, 330, 1180, 670)):
        v53_placeholder_image(draw, map_box, "North Slope public map capture")
    draw.rectangle((map_box[0], map_box[1], map_box[2], map_box[1] + 42), fill=(11, 35, 48))
    text(draw, (map_box[0] + 20, map_box[1] + 10), "North Slope public context", 19, WHITE, True)
    v54_source_label(draw, map_box, "Project Streamlit public map capture, regional context only")

    footer(
        draw,
        "Sources: USGS Gas Hydrate Crystals; Lee et al. 2008 USGS SIR 2008-5175; Sloan/Koh hydrate structure context; project Streamlit capture. No generic AI cage diagram.",
    )
    return save(img, "slide_02_source_context_v5_4.png")


def v54_parameter_bar(
    draw: ImageDraw.ImageDraw,
    y: int,
    label: str,
    left: str,
    right: str,
    good: tuple[float, float],
    direction: str,
    opposite: str,
    mimic: str,
    role: str,
    color: tuple[int, int, int],
    target: bool = False,
) -> None:
    row_box = (58, y, 1542, y + 70)
    card(draw, row_box, fill=(255, 248, 248) if target else WHITE, outline=RED if target else LINE, radius=8, width=2)
    text(draw, (82, y + 13), label, 16, RED if target else NAVY, True, width=190)
    if target:
        text(draw, (298, y + 15), "Y-only target rail", 16, RED, True, width=190)
        text(draw, (520, y + 15), "Labels supervise training and validation, but never enter X_allowed.", 15, NAVY, True, width=610)
        pill(draw, (1260, y + 18, 1495, y + 48), role, RED_LIGHT, RED)
        return
    ax0, ax1 = 300, 900
    draw.line((ax0, y + 36, ax1, y + 36), fill=(220, 234, 238), width=13)
    gx0 = ax0 + int((ax1 - ax0) * good[0])
    gx1 = ax0 + int((ax1 - ax0) * good[1])
    draw.line((gx0, y + 36, gx1, y + 36), fill=color, width=15)
    draw.line((gx0, y + 24, gx0, y + 48), fill=color, width=3)
    draw.polygon([(gx1 + 12, y + 36), (gx1 - 8, y + 24), (gx1 - 8, y + 48)], fill=color)
    text(draw, (ax0, y + 9), left, 11, MUTED, True, width=150)
    text(draw, (ax1 - 150, y + 9), right, 11, MUTED, True, width=150, align="right")
    text(draw, (926, y + 8), direction, 13, color, True, width=130)
    text(draw, (1068, y + 8), opposite, 12, NAVY, True, width=160)
    text(draw, (1235, y + 8), mimic, 12, MUTED, True, width=170)
    pill(draw, (1414, y + 20, 1516, y + 50), role, ICE_LIGHT if role != "Y-only" else RED_LIGHT, TEAL if role != "Y-only" else RED)


def v54_slide_parameter_ranges() -> Path:
    registry = v54_registry_rows()
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Parameters And Expected Hydrate Ranges",
        "Working normalized screening envelopes only: direction, opposite meaning, mimic risk, and ML role.",
    )
    headers = [("family", 82), ("range direction", 300), ("supports", 926), ("opposite", 1068), ("mimic / mask", 1235), ("role", 1420)]
    for label, x in headers:
        text(draw, (x, 125), label, 12, MUTED, True)

    combined_stiff = (
        min(
            v54_range(registry.get("Compressional velocity stiffness"), (0.58, 0.88))[0],
            v54_range(registry.get("Vp/Vs ratio"), (0.25, 0.55))[0],
            v54_range(registry.get("Acoustic impedance and elastic contrast"), (0.55, 0.88))[0],
        ),
        max(
            v54_range(registry.get("Compressional velocity stiffness"), (0.58, 0.88))[1],
            v54_range(registry.get("Vp/Vs ratio"), (0.25, 0.55))[1],
            v54_range(registry.get("Acoustic impedance and elastic contrast"), (0.55, 0.88))[1],
        ),
    )
    rows = [
        ("GR clean sand", "clean", "shaly", v54_range(registry.get("Gamma ray clean-sand proxy"), (0.00, 0.35)), "clean sand", "shale/clay", "radioactive minerals", "input", GREEN),
        ("Density / porosity", "tight", "porous", v54_range(registry.get("Porosity and density support"), (0.45, 0.75)), "porous", "tight/shale", "gas, washout", "input", BLUE),
        ("Deep resistivity", "conductive", "resistive", v54_range(registry.get("Deep resistivity response"), (0.65, 0.98)), "resistive", "wet/saline", "ice/free gas/tight", "input", RED),
        ("NMR separation", "mobile water", "separation", v54_range(registry.get("NMR porosity and NMR-density separation"), (0.60, 0.90)), "separation", "water/clay", "processing/depth", "input/Y guard", PURPLE),
        ("Vp/Vs/AI elastic", "soft", "stiff", combined_stiff, "stiff contrast", "soft/gas", "compaction/cement", "derived", PURPLE),
        ("Caliper QC", "in gauge", "washout", v54_range(registry.get("Caliper and bad-hole QC"), (0.00, 0.25)), "trust gate", "bad hole", "tool standoff", "QC", AMBER),
        ("Stability context", "outside", "inside", v54_range(registry.get("Pressure-temperature stability"), (0.60, 1.00)), "admissible", "unstable", "gas/depth/temp", "context", TEAL),
    ]
    y = 154
    for row in rows:
        v54_parameter_bar(draw, y, *row)
        y += 76
    v54_parameter_bar(
        draw,
        y,
        "Occurrence / saturation",
        "predictor",
        "target",
        (0, 0),
        "Y-only",
        "leakage",
        "fake performance",
        "Y-only",
        RED,
        target=True,
    )
    footer(draw, "Source: data/public_ml_products/public_parameter_evidence_registry_2026-06-16.csv. Ranges are working envelopes unless source-locked; none prove hydrate.")
    return save(img, "slide_03_parameter_ranges_v5_4.png")


def v54_slide_full_workflow() -> Path:
    return v54_copy_authority_panel(
        V52_ASSET_DIR / "slide_04_expanded_architecture_map.png",
        "slide_04_full_complex_project_workflow_v5_4.png",
        "Full Complex Project Workflow V5.4",
        "Public sources, OSL later inputs, variable fingerprints, gates, stability context, features, leakage-safe ML, dual heads, validation, and exports.",
    )


def v54_behavior_card(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    label: str,
    tag: str,
    good: tuple[float, float],
    why: str,
    opposite: str,
    mimic: str,
    color: tuple[int, int, int],
) -> None:
    card(draw, box, fill=WHITE, outline=color, radius=10, width=2)
    x1, y1, x2, y2 = box
    pill(draw, (x1 + 16, y1 + 14, x1 + 132, y1 + 43), tag, ICE_LIGHT if color != RED else RED_LIGHT, color)
    text(draw, (x1 + 150, y1 + 14), label, 18, NAVY, True, width=x2 - x1 - 168)
    ax0, ax1 = x1 + 34, x2 - 34
    ay = y1 + 80
    draw.line((ax0, ay, ax1, ay), fill=(221, 235, 238), width=12)
    gx0 = ax0 + int((ax1 - ax0) * good[0])
    gx1 = ax0 + int((ax1 - ax0) * good[1])
    draw.line((gx0, ay, gx1, ay), fill=color, width=14)
    draw.line((gx0, ay - 13, gx0, ay + 13), fill=color, width=3)
    draw.polygon([(gx1 + 12, ay), (gx1 - 7, ay - 11), (gx1 - 7, ay + 11)], fill=color)
    text(draw, (x1 + 26, y1 + 106), f"Why: {why}", 13, NAVY, True, width=x2 - x1 - 52)
    text(draw, (x1 + 26, y1 + 149), f"Opposite: {opposite}", 12, MUTED, True, width=(x2 - x1) // 2 - 34)
    text(draw, (x1 + (x2 - x1) // 2 + 6, y1 + 149), f"Mimic: {mimic}", 12, RED if mimic else MUTED, True, width=(x2 - x1) // 2 - 34)


def v54_slide_parameter_behavior() -> Path:
    registry = v54_registry_rows()
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Why Parameter Ranges Matter",
        "The ranges are working screening envelopes: useful because they encode physics, risky because several non-hydrate cases can mimic them.",
    )
    cards = [
        ("GR", "clean sand", v54_range(registry.get("Gamma ray clean-sand proxy"), (0.00, 0.35)), "lower radiation supports cleaner reservoir host", "shaly, clay-prone interval", "low GR alone", GREEN),
        ("Porosity", "porous", v54_range(registry.get("Porosity and density support"), (0.45, 0.75)), "pore volume is needed before hydrate can fill pores", "tight/compacted or shale", "gas/washout", BLUE),
        ("Rt", "resistive", v54_range(registry.get("Deep resistivity response"), (0.65, 0.98)), "hydrate can replace conductive pore water", "wet/saline sand", "ice/free gas/tight", RED),
        ("Vp/Vs/AI", "stiff", (0.50, 0.88), "hydrate can stiffen the frame and raise elastic contrast", "soft gas or poor consolidation", "cement/lithology", PURPLE),
        ("NMR", "separation", v54_range(registry.get("NMR porosity and NMR-density separation"), (0.60, 0.90)), "mobile-water response can separate from density porosity", "mobile water or clay-bound water", "processing/depth", TEAL),
        ("Stability", "admissible", v54_range(registry.get("Pressure-temperature stability"), (0.60, 1.00)), "P-T window is required physical context", "too warm or too shallow", "gas chemistry/temp", AMBER),
        ("Caliper", "washout", v54_range(registry.get("Caliper and bad-hole QC"), (0.00, 0.25)), "in-gauge hole protects density, sonic, and resistivity trust", "bad hole, downweight or block", "tool standoff", AMBER),
        ("Labels", "target-only", (0.00, 0.00), "occurrence and saturation are Y labels for supervision", "using as predictors creates leakage", "fake metrics", RED),
    ]
    positions = [
        (60, 145, 420, 340),
        (440, 145, 800, 340),
        (820, 145, 1180, 340),
        (1200, 145, 1540, 340),
        (60, 382, 420, 577),
        (440, 382, 800, 577),
        (820, 382, 1180, 577),
        (1200, 382, 1540, 577),
    ]
    for box, card_data in zip(positions, cards):
        v54_behavior_card(draw, box, *card_data)

    card(draw, (60, 640, 1540, 790), fill=(248, 252, 253), outline=LINE, radius=10, width=2)
    text(draw, (88, 665), "Interpretation rule", 23, NAVY, True)
    text(
        draw,
        (88, 705),
        "A hydrate-compatible direction becomes useful only when the reservoir gate, pore volume, hydrate-response logs, QC, stability context, and leakage-safe label policy agree. A single curve is never proof.",
        23,
        NAVY,
        True,
        width=1370,
        gap=8,
    )
    footer(draw, "Source: public parameter evidence registry plus ML logic ladder. All thresholds are working envelopes unless source-locked and unit-checked.")
    return save(img, "slide_05_parameter_behavior_v5_4.png")


def v54_equation_chip(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    label: str,
    formula: str,
    note: str,
    color: tuple[int, int, int],
) -> None:
    card(draw, box, fill=WHITE, outline=color, radius=9, width=2)
    x1, y1, x2, _ = box
    text(draw, (x1 + 16, y1 + 12), label, 15, color, True, width=x2 - x1 - 32)
    text(draw, (x1 + 16, y1 + 42), formula, 13, NAVY, True, width=x2 - x1 - 32, gap=3)
    text(draw, (x1 + 16, y1 + 78), note, 10, MUTED, True, width=x2 - x1 - 32, gap=2)


def v54_slide_equations_unit_gate() -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Equations, Feature Engineering, And Unit Gate",
        "Equation features enter X_allowed only after source, unit, depth, QC, and leakage checks.",
    )

    gate_box = (58, 145, 350, 790)
    card(draw, gate_box, fill=(248, 252, 253), outline=TEAL, radius=10, width=2)
    text(draw, (84, 170), "Unit and leakage gate", 22, TEAL, True, width=240)
    gate_steps = [
        ("1", "source header preserved"),
        ("2", "units visible or mapped"),
        ("3", "depth axis aligned"),
        ("4", "caliper/QC first"),
        ("5", "derived formula provenance"),
        ("6", "Y-only fields removed"),
        ("7", "X_allowed matrix"),
    ]
    y = 228
    for num, label in gate_steps:
        draw.ellipse((86, y, 116, y + 30), fill=TEAL)
        text(draw, (96, y + 5), num, 13, WHITE, True)
        text(draw, (130, y + 4), label, 14, NAVY, True, width=180)
        if num != "7":
            arrow(draw, (101, y + 36), (101, y + 62), TEAL, width=3)
        y += 74

    chips = [
        ("GR clean/shale", "GR_clean = low-GR reservoir proxy", "input/gate; not hydrate by itself", GREEN),
        ("Density porosity", "phi_D = (rho_ma - RHOB)/(rho_ma - rho_f)", "requires matrix/fluid assumptions", BLUE),
        ("Resistivity", "log_Rt = log10(Rt); Archie only if a,m,n,Rw approved", "hydrate support, not proof", RED),
        ("Vp from sonic", "Vp = 304800 / DT_us_per_ft", "unit conversion must be explicit", PURPLE),
        ("Vs where available", "Vs from shear sonic or approved source", "missingness flag if absent", PURPLE),
        ("Vp/Vs", "VpVs = Vp / Vs", "derived crossplot feature", PURPLE),
        ("Acoustic impedance", "AI = RHOB * Vp", "inherits density and sonic QC", TEAL),
        ("Elastic moduli", "G = rho*Vs^2; K = rho*(Vp^2 - 4Vs^2/3)", "unit-consistent physics only", TEAL),
        ("lambda-rho / mu-rho", "lambda_rho = rho*(Vp^2 - 2Vs^2); mu_rho = rho*Vs^2", "elastic contrast, not label", TEAL),
        ("NMR-density separation", "sep = phi_D - NMRPHI", "NMR_SAT/Sgh/Sh remain Y-only", AMBER),
    ]
    x_positions = [390, 760]
    y0 = 145
    for idx, chip in enumerate(chips):
        col = idx % 2
        row = idx // 2
        box = (x_positions[col], y0 + row * 128, x_positions[col] + 335, y0 + row * 128 + 104)
        v54_equation_chip(draw, box, *chip)

    right_box = (1150, 145, 1540, 790)
    card(draw, right_box, fill=(248, 252, 253), outline=AMBER, radius=10, width=2)
    text(draw, (1178, 170), "Stability context chip", 22, AMBER, True, width=320)
    v54_equation_chip(
        draw,
        (1178, 230, 1512, 338),
        "Hydrostatic pressure",
        "P_abs(z) = P_surface + rho_w*g*z_m/1e6",
        "assumption, not measured reservoir pressure",
        AMBER,
    )
    v54_equation_chip(
        draw,
        (1178, 370, 1512, 478),
        "Phase check",
        "admissible = T_model <= T_eq(P, CH4, 5 ppt)",
        "context/mask only; no occurrence claim",
        AMBER,
    )
    card(draw, (1178, 535, 1512, 700), fill=WHITE, outline=RED, radius=9, width=2)
    text(draw, (1204, 560), "Leakage stop", 20, RED, True, width=280)
    text(draw, (1204, 600), "Sgh, S_h, Sh, NMR_SAT, hydrate saturation, Swr/S_wr, and phase labels stay out of X_allowed.", 16, NAVY, True, width=270, gap=5)
    footer(draw, "Sources: processing slide 6 direction, WELL_LOG_REQUIREMENTS_MAP, approved intake spec, stability calculation plan, and parameter evidence registry.")
    return save(img, "slide_06_equations_feature_unit_gate_v5_4.png")


def v54_slide_ml_runtime() -> Path:
    return v54_copy_authority_panel(
        V52_ASSET_DIR / "slide_07_ml_runtime_detail.png",
        "slide_07_complex_ml_runtime_architecture_v5_4.png",
        "Complex ML Runtime Architecture V5.4",
        "Measured logs, derived physics, QC, stability/core context, X_allowed, whole-well split, train-only preprocessing, baselines, ANN/Keras candidate, dual heads, validation, and Y-only rail.",
    )


def v54_validation_lane(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    heading: str,
    body: str,
    color: tuple[int, int, int],
) -> None:
    card(draw, box, fill=WHITE, outline=color, radius=10, width=2)
    x1, y1, x2, _ = box
    text(draw, (x1 + 20, y1 + 18), heading, 20, color, True, width=x2 - x1 - 40)
    text(draw, (x1 + 20, y1 + 58), body, 16, NAVY, True, width=x2 - x1 - 40, gap=5)


def v54_slide_validation_outputs() -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Validation, Uncertainty, And Outputs",
        "Future outputs are designed now, but no trained metrics, occurrence predictions, or saturation predictions are reported.",
    )
    lanes = [
        ((70, 155, 430, 295), "1. Whole-well split", "Hold out complete wells or compartments before any preprocessing. Random depth-row splits are not final validation.", BLUE),
        ((470, 155, 830, 295), "2. Train-only preprocessing", "Scaling, imputation, feature selection, and model choice are fit on training wells only.", PURPLE),
        ((870, 155, 1230, 295), "3. Planned checks", "Confusion, calibration, residuals, and source/QC confidence plots are placeholders until approved labels exist.", AMBER),
        ((70, 345, 430, 505), "Occurrence output", "Future P(hydrate) and class label with calibration and reason flags. No occurrence map yet.", TEAL),
        ((470, 345, 830, 505), "Saturation output", "Future Sh_pred regression with residual review against approved target labels only. No saturation estimate yet.", GREEN),
        ((870, 345, 1230, 505), "Uncertainty and mimic flags", "Flag sparse logs, missing NMR/core, washout, high-resistivity mimics, elastic ambiguity, and stability assumptions.", RED),
    ]
    for box, heading, body, color in lanes:
        v54_validation_lane(draw, box, heading, body, color)
    for start, end, color in [
        ((430, 225), (470, 225), BLUE),
        ((830, 225), (870, 225), PURPLE),
        ((1030, 295), (300, 345), AMBER),
        ((1030, 295), (700, 345), AMBER),
        ((1030, 295), (1030, 345), AMBER),
    ]:
        arrow(draw, start, end, color, width=4)

    package_box = (1265, 155, 1530, 705)
    card(draw, package_box, fill=(248, 252, 253), outline=TEAL, radius=10, width=2)
    text(draw, (1290, 180), "Reviewed output package", 21, TEAL, True, width=215)
    package_items = [
        "probability table",
        "Sh prediction table",
        "uncertainty flags",
        "mimic flags",
        "calibration/residual review",
        "locked reasons",
        "public-safe later maps/tables",
    ]
    y = 238
    for item in package_items:
        draw.rounded_rectangle((1290, y, 1510, y + 36), radius=8, fill=WHITE, outline=LINE, width=1)
        text(draw, (1306, y + 9), item, 13, NAVY, True, width=185)
        y += 54
    arrow(draw, (1230, 425), (1265, 425), TEAL, width=4)

    card(draw, (70, 630, 1230, 760), fill=(255, 248, 248), outline=RED, radius=10, width=2)
    text(draw, (98, 655), "No-results guardrail", 22, RED, True)
    text(draw, (98, 695), "The deck shows validation architecture and output shape only. It does not report trained metrics, confusion values, occurrence predictions, saturation predictions, or hydrate proof.", 20, NAVY, True, width=1070, gap=6)
    footer(draw, "Sources: FIRST_MODEL_EXPERIMENT_PLAN, first model output schema, target registry, leakage guardrails, and V5.2 ML runtime detail.")
    return save(img, "slide_08_validation_uncertainty_outputs_v5_4.png")


def v54_slide_status_decisions(values: dict[str, str]) -> Path:
    img = canvas(False)
    draw = ImageDraw.Draw(img)
    v53_panel_title(
        draw,
        "Status, Mentor Decisions, And Next Actions",
        "The project is ready for approved-data header audit and OSL runtime mapping; model training has not started.",
    )

    metrics = [
        (values["wells"], "public scaffold wells"),
        (values["profiles"], "G10015 profiles"),
        (values["temp_matches"], "temp-profile matches"),
        (values["screen_calculated"], "admissibility intervals"),
        (values["approved_visible"], "approved datasets visible"),
        ("0", "trained ML results"),
    ]
    metric_positions = [
        (70, 150),
        (338, 150),
        (70, 265),
        (338, 265),
        (70, 380),
        (338, 380),
    ]
    for idx, (number, label) in enumerate(metrics):
        x, y = metric_positions[idx]
        card(draw, (x, y, x + 235, y + 88), fill=WHITE, outline=LINE, radius=10, width=2)
        color = RED if label == "trained ML results" else TEAL
        text(draw, (x + 18, y + 18), number, 25, color, True, width=199, align="center")
        text(draw, (x + 18, y + 54), label, 12, MUTED, True, width=199, align="center")

    card(draw, (70, 505, 605, 740), fill=(248, 252, 253), outline=TEAL, radius=10, width=2)
    text(draw, (98, 528), "Current status", 22, TEAL, True)
    status_items = [
        "stability screen has 22 calculated intervals, not proof",
        "only about 3/71 approved datasets visible for schema design",
        "ML architecture/readiness scaffold exists",
        "no training, metrics, occurrence maps, or saturation outputs yet",
    ]
    sy = 572
    for item in status_items:
        text(draw, (108, sy), f"- {item}", 15, NAVY, True, width=455, gap=3)
        sy += 38

    card(draw, (645, 150, 1530, 555), fill=WHITE, outline=AMBER, radius=10, width=2)
    text(draw, (675, 178), "Mentor decisions to lock", 24, AMBER, True)
    decisions = [
        "authoritative saturation target field",
        "occurrence label policy",
        "blind validation wells or compartments",
        "missing-log adapter policy",
        "phase curve / gas composition scenario policy",
        "caliper QC policy",
    ]
    dx, dy = 675, 230
    for idx, item in enumerate(decisions, start=1):
        draw.ellipse((dx, dy + 2, dx + 28, dy + 30), fill=AMBER)
        text(draw, (dx + 10, dy + 6), str(idx), 12, WHITE, True)
        text(draw, (dx + 42, dy + 4), item, 17, NAVY, True, width=350)
        if idx == 3:
            dx, dy = 1110, 230
        else:
            dy += 82

    card(draw, (645, 600, 1530, 740), fill=(255, 248, 248), outline=RED, radius=10, width=2)
    text(draw, (675, 630), "Next action", 24, RED, True)
    text(draw, (675, 675), "Run approved-data header audit and runtime mapping in OSL. Keep GitHub/Streamlit public-safe: schema, counts, diagrams, guardrails, and no approved/private rows.", 18, NAVY, True, width=800, gap=5)
    footer(draw, "Counts from public scaffold summaries. Stability is admissibility only; outputs are planned until approved labels and validation are complete.")
    return save(img, "slide_09_status_mentor_decisions_v5_4.png")


def v54_build_panels() -> list[Path]:
    values = summaries()
    return [
        v54_slide_personal_opener(),
        v54_slide_context(),
        v54_slide_parameter_ranges(),
        v54_slide_full_workflow(),
        v54_slide_parameter_behavior(),
        v54_slide_equations_unit_gate(),
        v54_slide_ml_runtime(),
        v54_slide_validation_outputs(),
        v54_slide_status_decisions(values),
    ]


def v54_build_word_companion(panel_paths: list[Path], contact_sheet: Path) -> Path:
    values = summaries()
    document = Document()
    apply_doc_style(document)
    section = document.sections[0]
    section.top_margin = Inches(0.6)
    section.bottom_margin = Inches(0.6)
    section.left_margin = Inches(0.65)
    section.right_margin = Inches(0.65)

    props = document.core_properties
    props.title = "V5.4 Corrected North Slope Gas Hydrate ML Workflow Companion"
    props.subject = "Corrected mentor-facing slide companion"
    props.author = "North Slope Gas Hydrates project"

    document.add_heading("V5.4 Corrected North Slope Gas Hydrate ML Workflow Companion", level=0)
    document.add_paragraph(
        "This companion explains the corrected nine-slide mentor deck. V5.4 restores the personal opener and the "
        "complex workflow/runtime diagrams, replaces AI-looking hydrate visuals with source-backed figures, and "
        "keeps the project guardrails explicit. It does not report hydrate proof, trained metrics, occurrence "
        "predictions, saturation predictions, or approved/private data rows."
    )

    document.add_heading("Current Public State", level=1)
    document.add_paragraph(
        f"The public scaffold currently tracks {values['wells']} Arctic Slope public wells, {values['profiles']} "
        f"G10015 temperature profiles, {values['temp_matches']} temperature-profile matches, and "
        f"{values['screen_calculated']} methane 5 ppt calculated stability-admissibility intervals. Only about "
        "3/71 approved datasets are visible for schema design. The ML scaffold exists, but no model training has started."
    )

    document.add_heading("Hydrate Structures And North Slope Methane Context", level=1)
    document.add_paragraph(
        "Gas hydrate is an ice-like crystalline solid where water cages trap guest gas molecules. Structure I is the "
        "methane-dominant hydrate structure used as the current public baseline. Structure II can host larger guest "
        "molecules and mixed gases, and Structure H involves larger hydrocarbon guests. Because gas composition shifts "
        "the stability boundary, V5.4 treats methane 5 ppt as the current public method and leaves mixed-gas scenarios "
        "as a mentor policy decision."
    )

    document.add_heading("Stability Is Admissibility Only", level=1)
    document.add_paragraph(
        "The stability branch compares a public temperature model against a digitized methane 5 ppt phase curve under "
        "a hydrostatic pressure assumption. Intervals where modeled temperature is at or below the equilibrium "
        "temperature are physically admissible under those assumptions. This is not hydrate occurrence, saturation, "
        "producibility, final stability top/base/thickness, or proof."
    )

    document.add_heading("Why Log Responses Are Useful But Ambiguous", level=1)
    document.add_paragraph(
        "Low gamma ray helps identify cleaner reservoir-quality sand, but clean sand alone is not hydrate. High deep "
        "resistivity can support hydrate interpretation in clean porous sand, but free gas, ice, tight rock, salinity, "
        "cement, invasion, and bad borehole conditions can mimic or mask the response. Density, porosity, NMR, velocity, "
        "Vp/Vs, acoustic impedance, and elastic attributes are useful only after unit, provenance, depth, QC, and "
        "lithology checks."
    )

    document.add_heading("Leakage-Safe ML Design", level=1)
    document.add_paragraph(
        "The ML architecture keeps measured logs, derived physics, QC flags, alignment fields, and approved context in "
        "X_allowed after source and unit gates. Occurrence labels, Sgh, S_h, Sh, NMR_SAT, hydrate saturation, Swr/S_wr, "
        "and interpreted phase labels are Y-only targets. They supervise training and validation but never enter the "
        "feature matrix."
    )

    document.add_heading("Separate Outputs", level=1)
    document.add_paragraph(
        "The future model has two separate outputs: occurrence classification as P(hydrate) and saturation regression "
        "as Sh_pred. Both require approved labels, whole-well or grouped validation, train-only preprocessing, and "
        "mentor review before public release. V5.4 shows planned output shapes and review flags only, not results."
    )

    slide_notes = [
        (
            "Slide 1 - Personal/about-me opener",
            "Restores the original Gmail-style personal opener and keeps the project title, personal imagery, and old spine.",
        ),
        (
            "Slide 2 - Gas hydrate and North Slope context",
            "Uses source-backed visuals only: USGS hydrate crystals, the digitized USGS SIR methane 5 ppt curve, and the public website map capture. Hydrate structures are explained in text rather than with generic cage art.",
        ),
        (
            "Slide 3 - Parameters and expected hydrate ranges",
            "Shows only parameter families, normalized working envelopes, opposite meaning, mimic or mask risk, and ML role from the public parameter evidence registry.",
        ),
        (
            "Slide 4 - Full complex project workflow",
            "Restores the V5.2 complex workflow architecture as the main visual, including public sources, OSL later inputs, fingerprint registry, gates, feature engineering, leakage-safe ML, dual heads, validation, and exports.",
        ),
        (
            "Slide 5 - Why parameter ranges matter",
            "Explains why each family is scientifically useful and where it can fail. The point is evidence combination, not single-log proof.",
        ),
        (
            "Slide 6 - Equations, feature engineering, and unit gate",
            "Moves away from a pressure-temperature-only slide and centers the equations, unit checks, QC gates, and leakage stop that decide whether derived features enter X_allowed.",
        ),
        (
            "Slide 7 - Complex ML architecture",
            "Restores the complex runtime diagram: measured logs, derived physics, QC, stability/core context, X_allowed, grouped split, train-only preprocessing, baselines, ANN/Keras candidate, dual heads, validation, reviewed output package, and Y-only rail.",
        ),
        (
            "Slide 8 - Validation, uncertainty, and outputs",
            "Connects whole-well validation, planned calibration and residual review, uncertainty and mimic flags, occurrence classification, saturation regression, and the reviewed output package. Metrics are intentionally absent.",
        ),
        (
            "Slide 9 - Status, mentor decisions, and next steps",
            "Summarizes current counts, unresolved mentor decisions, and the immediate OSL action: approved-data header audit and runtime mapping.",
        ),
    ]
    document.add_heading("Slide-By-Slide Companion", level=1)
    for idx, (heading, body) in enumerate(slide_notes):
        document.add_heading(heading, level=2)
        document.add_paragraph(body)
        if idx < len(panel_paths) and panel_paths[idx].exists():
            document.add_picture(str(panel_paths[idx]), width=Inches(6.7))

    document.add_heading("Contact Sheet", level=1)
    if contact_sheet.exists():
        document.add_picture(str(contact_sheet), width=Inches(6.7))

    document.add_heading("Source And Guardrail Anchors", level=1)
    anchors = [
        "USGS Gas Hydrate Crystals public-domain image.",
        "Lee et al. 2008 USGS SIR 2008-5175 methane 5 ppt phase curve digitized product.",
        "data/public_ml_products/public_parameter_evidence_registry_2026-06-16.csv.",
        "docs/FULL_PROJECT_ML_WORKFLOW_DIAGRAM.md.",
        "docs/APPROVED_DATA_INTAKE_SPEC_2026-06-15.md and target leakage guardrails.",
        "docs/FIRST_MODEL_EXPERIMENT_PLAN_2026-06-15.md and first model output schema.",
    ]
    for item in anchors:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("Language To Preserve", level=1)
    for item in [
        "Use 'admissible under assumptions' rather than hydrate proof.",
        "Use 'planned validation/output shape' rather than trained metrics.",
        "Use 'future occurrence classification' and 'future saturation regression' rather than predictions.",
        "Keep OpenScienceLab as the approved/private runtime workbench and GitHub/Streamlit as the public delivery surface.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    OUT_DOCX.parent.mkdir(parents=True, exist_ok=True)
    document.save(OUT_DOCX)
    return OUT_DOCX


def main() -> None:
    panels = v54_build_panels()
    if len(panels) != 9:
        raise ValueError(f"V5.4 must have exactly 9 slides, found {len(panels)}")
    deck = rebuild_deck(panels)
    verify_deck(deck)
    contact_sheet = build_contact_sheet(panels)
    docx = v54_build_word_companion(panels, contact_sheet)
    print(f"Wrote {len(panels)} V5.4 panels")
    print(f"Wrote {contact_sheet}")
    print(f"Wrote {deck}")
    print(f"Wrote {docx}")
    print(f"Wrote panels to {ASSET_DIR}")


if __name__ == "__main__":
    main()
