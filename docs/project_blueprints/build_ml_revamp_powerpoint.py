from __future__ import annotations

import csv
import textwrap
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from build_processing_slide_assets import build_assets as build_processing_assets


ROOT = Path(__file__).resolve().parents[2]
BLUEPRINT_DIR = ROOT / "docs" / "project_blueprints"
ASSET_DIR = BLUEPRINT_DIR / "presentation_assets"
BASE_PPTX = BLUEPRINT_DIR / "Drive_base_CURRENT_June_10_North_Slope_Gas_Hydrate_Reservoir_Characterization_Research_Overview_Slides.pptx"
OUT_PPTX = BLUEPRINT_DIR / "North_Slope_Gas_Hydrate_Reservoir_Characterization_Research_Overview.pptx"
PARAMETER_CSV = BLUEPRINT_DIR / "ml_parameter_effect_tree.csv"

PROFILE_PHOTO = ASSET_DIR / "rohan_profile_photo.jpg"
REGIONAL_CONTEXT = ASSET_DIR / "regional_3d_context.png"
EVIDENCE_STACK = ASSET_DIR / "subsurface_evidence_stack.png"
SYNTHETIC_LOG = ASSET_DIR / "synthetic_well_log_panel.png"
SWEET_SPOT = ASSET_DIR / "sweet_spot_ranking.png"
REVISION_ASSET_DIR = ROOT / "references" / "presentation-revision-2026-06-11" / "images"
STREAMLIT_STRUCTURAL = REVISION_ASSET_DIR / "project_streamlit_structural_explorer_v2.png"
HYDRATE_CRYSTALS = REVISION_ASSET_DIR / "usgs_gas_hydrate_crystals_sem_public_domain.jpg"
GAMMA_LOG_REF = REVISION_ASSET_DIR / "usgs_gamma_logs_public_domain.png"
CALIPER_LOG_REF = REVISION_ASSET_DIR / "usgs_caliper_logs_public_domain.png"
RESISTIVITY_LOG_REF = REVISION_ASSET_DIR / "usgs_fluid_resistivity_logs_public_domain.jpg"
PROCESSING_ASSET_DIR = ASSET_DIR / "processing_revisions_2026_06_11"
PROCESSING_SLIDES = {
    "title": PROCESSING_ASSET_DIR / "slide_01_about_me.png",
    "intro": PROCESSING_ASSET_DIR / "slide_02_hydrate_intro.png",
    "parameters": PROCESSING_ASSET_DIR / "slide_03_parameter_scaffold.png",
    "architecture": PROCESSING_ASSET_DIR / "slide_04_ml_architecture.png",
    "behavior": PROCESSING_ASSET_DIR / "slide_05_parameter_behavior.png",
    "geomechanics": PROCESSING_ASSET_DIR / "slide_06_geomechanics.png",
    "map": PROCESSING_ASSET_DIR / "slide_07_map_context.png",
    "results": PROCESSING_ASSET_DIR / "slide_08_results_plan.png",
    "conclusion": PROCESSING_ASSET_DIR / "slide_09_conclusion.png",
}


NAVY = RGBColor(9, 34, 49)
DEEP = RGBColor(4, 21, 33)
TEAL = RGBColor(22, 125, 141)
ICE = RGBColor(103, 208, 223)
GREEN = RGBColor(37, 165, 138)
AMBER = RGBColor(219, 165, 72)
RED = RGBColor(204, 75, 74)
INK = RGBColor(18, 52, 71)
MUTED = RGBColor(86, 105, 115)
LIGHT = RGBColor(244, 248, 249)
PANEL = RGBColor(231, 240, 242)
WHITE = RGBColor(255, 255, 255)
PURPLE = RGBColor(128, 139, 214)
BLUE = RGBColor(63, 138, 201)
CHARCOAL = RGBColor(33, 45, 53)
SAND = RGBColor(224, 195, 137)
SHALE = RGBColor(110, 120, 128)
DARK_PANEL = RGBColor(11, 39, 54)


def clear_deck(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
    for slide_id in list(slide_ids):
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        slide_ids.remove(slide_id)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_background(slide, color: RGBColor = WHITE) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def set_shape_fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color: RGBColor, width: float = 1.0, transparency: int = 0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float = 14,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font: str = "Aptos",
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    lines = str(text).splitlines() or [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(1)
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
    return box


def add_title(slide, title: str, subtitle: str | None = None, dark: bool = False) -> None:
    color = WHITE if dark else NAVY
    textbox(slide, 0.42, 0.18, 12.3, 0.38, title, size=24, color=color, bold=True)
    if subtitle:
        textbox(slide, 0.47, 0.55, 12.0, 0.28, subtitle, size=8.8, color=ICE if dark else MUTED)


def add_footer(slide, text: str, dark: bool = False) -> None:
    color = RGBColor(160, 178, 187) if dark else MUTED
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.45),
        Inches(7.12),
        Inches(12.88),
        Inches(7.12),
    )
    set_line(line, RGBColor(182, 199, 204) if not dark else RGBColor(71, 93, 105), 0.6)
    textbox(slide, 0.48, 7.16, 12.2, 0.2, text, size=6.5, color=color)


def add_chip(slide, x: float, y: float, w: float, text: str, fill: RGBColor = PANEL, color: RGBColor = INK):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.24))
    set_shape_fill(shape, fill)
    set_line(shape, fill, 0.5)
    textbox(slide, x + 0.04, y + 0.035, w - 0.08, 0.16, text, size=6.2, color=color, bold=True, align=PP_ALIGN.CENTER)
    return shape


def add_card(slide, x: float, y: float, w: float, h: float, fill: RGBColor = WHITE, line: RGBColor = PANEL):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_fill(card, fill)
    set_line(card, line, 0.9)
    return card


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    body: str = "",
    fill: RGBColor = WHITE,
    line: RGBColor = PANEL,
    accent: RGBColor = TEAL,
    title_size: float = 11,
    body_size: float = 8,
):
    add_card(slide, x, y, w, h, fill=fill, line=line)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    set_shape_fill(bar, accent)
    set_line(bar, accent, 0)
    textbox(slide, x + 0.16, y + 0.11, w - 0.25, 0.22, label, size=title_size, color=INK, bold=True)
    if body:
        textbox(slide, x + 0.16, y + 0.42, w - 0.25, h - 0.5, body, size=body_size, color=MUTED)


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = TEAL, width: float = 1.4) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    set_line(line, color, width)
    line.line.end_arrowhead = True


def add_shape(slide, shape_name: str, x: float, y: float, w: float, h: float, fill: RGBColor, line: RGBColor | None = None):
    kind = getattr(MSO_AUTO_SHAPE_TYPE, shape_name, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_fill(shape, fill)
    set_line(shape, line or fill, 0.8)
    return shape


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    if not path.exists():
        add_box(slide, x, y, w, h, "Missing asset", path.name, fill=LIGHT, accent=RED)
        return None

    with Image.open(path) as img:
        img_w, img_h = img.size
    image_ratio = img_w / img_h
    box_ratio = w / h
    if image_ratio >= box_ratio:
        width = w
        height = w / image_ratio
    else:
        height = h
        width = h * image_ratio
    return slide.shapes.add_picture(
        str(path),
        Inches(x + (w - width) / 2),
        Inches(y + (h - height) / 2),
        width=Inches(width),
        height=Inches(height),
    )


def add_processing_slide(slide, key: str) -> bool:
    path = PROCESSING_SLIDES[key]
    if not path.exists():
        return False
    add_picture_contain(slide, path, 0, 0, 13.333333, 7.5)
    return True


def short(text: str, limit: int) -> str:
    return textwrap.shorten(str(text).replace(";", ",").replace("  ", " "), width=limit, placeholder="...")


def read_parameters() -> list[dict[str, str]]:
    with PARAMETER_CSV.open(newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def draw_borehole_icon(slide, x: float, y: float, size: float, color: RGBColor = TEAL) -> None:
    for dx in (0.16, size - 0.2):
        wall = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + dx), Inches(y + 0.05), Inches(0.03), Inches(size - 0.1))
        set_shape_fill(wall, color)
        set_line(wall, color, 0)
    add_arrow(slide, x + 0.24, y + size / 2, x + size - 0.28, y + size / 2, color, 1.0)


def draw_crossplot_icon(slide, x: float, y: float, size: float, color: RGBColor = TEAL) -> None:
    axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.16), Inches(y + size - 0.14), Inches(x + size - 0.1), Inches(y + size - 0.14))
    set_line(axis, color, 1.0)
    axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.16), Inches(y + 0.12), Inches(x + 0.16), Inches(y + size - 0.14))
    set_line(axis, color, 1.0)
    for dx, dy, fill in ((0.28, 0.44, GREEN), (0.46, 0.28, AMBER), (0.62, 0.58, BLUE), (0.78, 0.35, RED)):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + dx), Inches(y + dy), Inches(0.09), Inches(0.09))
        set_shape_fill(dot, fill)
        set_line(dot, WHITE, 0.4)


def draw_layer_icon(slide, x: float, y: float, size: float, color: RGBColor = TEAL) -> None:
    colors = [RGBColor(210, 224, 228), RGBColor(188, 210, 214), RGBColor(232, 219, 184), RGBColor(193, 221, 214)]
    for i, fill in enumerate(colors):
        layer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.08), Inches(y + 0.1 + i * 0.18), Inches(size - 0.16), Inches(0.14))
        set_shape_fill(layer, fill)
        set_line(layer, WHITE, 0.3)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + size - 0.18), Inches(y + 0.1), Inches(x + size - 0.18), Inches(y + size - 0.12))
    set_line(line, color, 1.1)


def draw_nmr_icon(slide, x: float, y: float, size: float, color: RGBColor = TEAL) -> None:
    for offset, alpha in ((0.07, 0), (0.18, 20), (0.29, 35)):
        ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + offset), Inches(y + offset), Inches(size - 2 * offset), Inches(size - 2 * offset))
        set_shape_fill(ring, RGBColor(221, 241, 243), alpha)
        set_line(ring, color, 0.8)
    pore = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.36), Inches(y + 0.36), Inches(0.2), Inches(0.2))
    set_shape_fill(pore, BLUE)
    set_line(pore, BLUE, 0.5)


def draw_velocity_icon(slide, x: float, y: float, size: float, kind: str, color: RGBColor = TEAL) -> None:
    block = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, Inches(x + size * 0.52), Inches(y + size * 0.22), Inches(size * 0.28), Inches(size * 0.44))
    set_shape_fill(block, RGBColor(222, 235, 237))
    set_line(block, color, 0.8)
    wave = add_shape(slide, "WAVE", x + size * 0.08, y + size * 0.34, size * 0.58, size * 0.22, ICE, ICE)
    if kind == "vs":
        wave.rotation = 90
        add_arrow(slide, x + size * 0.18, y + size * 0.82, x + size * 0.72, y + size * 0.82, AMBER, 1.0)


def draw_parameter_icon(slide, key: str, x: float, y: float, size: float = 0.52, color: RGBColor = TEAL) -> None:
    container = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    set_shape_fill(container, RGBColor(232, 244, 246))
    set_line(container, RGBColor(190, 219, 224), 0.6)
    inset_x = x + 0.05
    inset_y = y + 0.05
    draw_size = size - 0.1

    if key == "resistivity":
        bolt = add_shape(slide, "LIGHTNING_BOLT", inset_x + 0.1, inset_y + 0.03, draw_size - 0.2, draw_size - 0.05, AMBER, AMBER)
        bolt.rotation = 5
    elif key == "nmr":
        draw_nmr_icon(slide, inset_x, inset_y, draw_size, color)
    elif key == "vp":
        draw_velocity_icon(slide, inset_x, inset_y, draw_size, "vp", color)
    elif key == "vs":
        draw_velocity_icon(slide, inset_x, inset_y, draw_size, "vs", color)
    elif key == "density":
        cube = add_shape(slide, "CUBE", inset_x + 0.08, inset_y + 0.06, draw_size - 0.16, draw_size - 0.12, GREEN, GREEN)
        cube.rotation = 2
    elif key == "gamma":
        sun = add_shape(slide, "SUN", inset_x + 0.08, inset_y + 0.08, draw_size - 0.16, draw_size - 0.16, AMBER, AMBER)
        sun.rotation = 15
    elif key == "caliper":
        draw_borehole_icon(slide, inset_x, inset_y, draw_size, color)
    elif key == "elastic":
        draw_crossplot_icon(slide, inset_x, inset_y, draw_size, color)
    elif key == "core":
        can = add_shape(slide, "CAN", inset_x + 0.15, inset_y + 0.08, draw_size - 0.3, draw_size - 0.14, BLUE, BLUE)
        can.rotation = 90
    elif key == "depth":
        draw_layer_icon(slide, inset_x, inset_y, draw_size, color)
    else:
        gear = add_shape(slide, "GEAR_6", inset_x + 0.1, inset_y + 0.1, draw_size - 0.2, draw_size - 0.2, color, color)
        gear.rotation = 20


def slide_title(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    set_shape_fill(band, LIGHT)
    set_line(band, LIGHT, 0)
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(7.15), Inches(7.5))
    set_shape_fill(left, WHITE)
    set_line(left, WHITE, 0)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.13), Inches(7.5))
    set_shape_fill(accent, TEAL)
    set_line(accent, TEAL, 0)

    textbox(slide, 0.55, 0.48, 6.1, 0.72, "Gas Hydrate Occurrence and Saturation Prediction", size=23, color=NAVY, bold=True)
    textbox(slide, 0.57, 1.35, 5.9, 0.48, "Permafrost sediments on the Alaska North Slope using physics-constrained AI/ML", size=14.5, color=INK, bold=True)
    textbox(slide, 0.6, 2.12, 5.9, 0.78, "Use approved well-log, NMR, and core-analysis data to predict hydrate occurrence and saturation, then explain each result as reservoir-characterization evidence.", size=11.8, color=MUTED)

    add_chip(slide, 0.58, 3.12, 1.7, "9 slides", TEAL, WHITE)
    add_chip(slide, 2.43, 3.12, 2.05, "source-backed", RGBColor(217, 232, 236), NAVY)
    add_chip(slide, 4.65, 3.12, 1.6, "runtime-safe", RGBColor(232, 221, 189), NAVY)

    add_box(
        slide,
        0.58,
        4.1,
        5.9,
        1.25,
        "Presenter",
        "Rohan Nanda\nLead Geologist | M.S. Geosciences, The University of Texas at Dallas\nAbout me: drawing, swimming, and running",
        fill=WHITE,
        line=RGBColor(211, 224, 228),
        accent=GREEN,
        title_size=12,
        body_size=9,
    )
    textbox(slide, 0.62, 6.48, 6.25, 0.28, "Public deck only: approved logs, restricted identifiers, runtime results, and trained models stay out of this file.", size=6.9, color=MUTED)

    add_picture_contain(slide, PROFILE_PHOTO, 7.42, 0.52, 5.1, 6.35)
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.42), Inches(0.52), Inches(5.1), Inches(6.35))
    frame.fill.background()
    set_line(frame, NAVY, 1.0)
    add_footer(slide, "Profile photo sourced from the existing project presentation; deck rebuilt locally before Drive handoff.")


def slide_need_for_ml(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "Why Hydrate Prediction Needs ML", "Well logs give indirect evidence; ML is used to combine weak, caveated signals without treating any single curve as proof.")

    add_picture_contain(slide, REGIONAL_CONTEXT, 0.62, 1.05, 5.8, 4.9)
    textbox(slide, 0.85, 6.03, 5.35, 0.28, "Regional and 3D context is planning context, not a substitute for well-level validation.", size=7.2, color=MUTED, align=PP_ALIGN.CENTER)

    steps = [
        ("Occurrence", "Where hydrate is geologically plausible and log responses agree.", TEAL),
        ("Saturation", "How much pore space is likely occupied, with uncertainty.", GREEN),
        ("Reservoir quality", "Whether hydrate sits in usable sand-rich intervals.", AMBER),
    ]
    for i, (label, body, color) in enumerate(steps):
        x = 6.78 + i * 2.05
        add_card(slide, x, 1.05, 1.78, 1.35, fill=WHITE, line=RGBColor(209, 224, 228))
        draw_parameter_icon(slide, ["depth", "resistivity", "core"][i], x + 0.1, 1.2, 0.5, color)
        textbox(slide, x + 0.68, 1.17, 0.95, 0.22, label, size=10, color=NAVY, bold=True)
        textbox(slide, x + 0.13, 1.7, 1.52, 0.45, body, size=6.5, color=MUTED)

    add_box(
        slide,
        6.78,
        2.78,
        5.93,
        1.0,
        "Central constraint",
        "Gas hydrate interpretation is a converging-evidence problem: resistivity, porosity, velocity, lithology, depth, and QC must agree before a prediction is trusted.",
        fill=LIGHT,
        accent=TEAL,
        title_size=11,
        body_size=7.8,
    )
    add_box(
        slide,
        6.78,
        4.04,
        2.78,
        1.15,
        "What ML adds",
        "Pattern learning across wells, interactions between curves, and repeatable uncertainty checks.",
        fill=WHITE,
        accent=GREEN,
        title_size=10.5,
        body_size=7.5,
    )
    add_box(
        slide,
        9.78,
        4.04,
        2.93,
        1.15,
        "What ML cannot replace",
        "Physics review, target provenance, complete-well validation, and geologic masking checks.",
        fill=WHITE,
        accent=RED,
        title_size=10.5,
        body_size=7.5,
    )
    textbox(slide, 6.9, 5.73, 5.65, 0.48, "Design rule: use ML to rank and calibrate evidence, not to turn one high-resistivity curve into a hydrate label.", size=10, color=INK, bold=True)
    add_footer(slide, "Sources: Chong et al. 2022; Lee and Collett 2011; Haines et al. 2022; project source matrix.")


def slide_source_stack(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "Source-to-Workflow Evidence Stack", "The deck separates public source evidence from authorized runtime data and uses sources only for public-safe methods.")

    add_picture_contain(slide, EVIDENCE_STACK, 0.58, 1.02, 3.25, 5.08)
    columns = [
        (
            "Public source layer",
            "Peer-reviewed papers\npublic reports\nproject synthesis docs\nheader/equation maps",
            TEAL,
        ),
        (
            "Interpretation layer",
            "Parameter meanings\nfeature equations\nknown caveats\nvalidation design",
            GREEN,
        ),
        (
            "Authorized runtime layer",
            "Approved logs\nrestricted identifiers\ntrained models\npopulated outputs",
            RED,
        ),
        (
            "Public deliverables",
            "Concept deck\nmethod docs\nscaffolded workflow\nno restricted outputs",
            AMBER,
        ),
    ]
    for i, (label, body, color) in enumerate(columns):
        x = 4.1 + i * 2.18
        add_box(slide, x, 1.07, 1.82, 3.35, label, body, fill=WHITE, line=RGBColor(209, 224, 228), accent=color, title_size=9.5, body_size=7.5)
        if i < len(columns) - 1:
            add_arrow(slide, x + 1.86, 2.78, x + 2.12, 2.78, color=RGBColor(135, 161, 170), width=1.0)

    add_box(
        slide,
        4.1,
        4.85,
        8.55,
        1.12,
        "Deck rule",
        "Every visual can show the workflow, equations, and parameter logic. It must not expose approved well logs, named restricted identifiers, trained models, populated runtime settings, or derived sensitive outputs.",
        fill=LIGHT,
        line=RGBColor(202, 218, 223),
        accent=NAVY,
        title_size=10.5,
        body_size=8,
    )
    textbox(slide, 4.22, 6.25, 8.3, 0.28, "This is why the PowerPoint uses sketches, generated/scaffold visuals, and source-backed definitions instead of copied restricted data.", size=7.4, color=MUTED)
    add_footer(slide, "Boundary follows docs/opensciencelab_runtime_layout.md and the project architecture map.")


def slide_parameter_grid(prs: Presentation, params: list[dict[str, str]]) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "Parameter Signal and Caveat Grid", "Each parameter has a measurement meaning, hydrate interpretation role, and caveat review. Caveats are not the parameter itself.")

    grid_params = params[:10]
    row_h = 1.14
    columns = [(0.42, grid_params[:5]), (6.82, grid_params[5:])]
    for col_x, rows in columns:
        for i, row in enumerate(rows):
            y = 1.02 + i * row_h
            add_card(slide, col_x, y, 6.1, 1.03, fill=WHITE, line=RGBColor(206, 222, 227))
            draw_parameter_icon(slide, row.get("icon_key", ""), col_x + 0.12, y + 0.17, 0.62)
            textbox(slide, col_x + 0.86, y + 0.07, 2.28, 0.18, short(row["parameter"], 38), size=7.7, color=NAVY, bold=True)
            textbox(slide, col_x + 3.18, y + 0.07, 2.8, 0.18, short(row.get("family", ""), 31), size=6.0, color=TEAL, bold=True, align=PP_ALIGN.RIGHT)
            body = (
                f"Measures: {short(row.get('measures', ''), 72)}\n"
                f"Hydrate signal: {short(row.get('primary_hydrate_effect', ''), 76)}\n"
                f"Caveats: {short(row.get('slide_caveat_summary', ''), 82)}\n"
                f"ML role: {short(row.get('model_role', ''), 78)}"
            )
            textbox(slide, col_x + 0.86, y + 0.28, 5.05, 0.67, body, size=5.55, color=CHARCOAL)

    textbox(slide, 0.48, 6.86, 12.35, 0.2, "Example correction: deep resistivity measures formation electrical response; the masking review is gas, ice, carbonate/cement, tight rock, shale correction, salinity, invasion, and bad-hole context.", size=6.4, color=MUTED)
    add_footer(slide, "Sources: parameter CSV; Chong et al. 2022; Lee and Collett 2011; Haines et al. 2022; header and equation maps.")


def equation_card(slide, x: float, y: float, w: float, label: str, equation: str, note: str, accent: RGBColor) -> None:
    add_card(slide, x, y, w, 0.64, fill=WHITE, line=RGBColor(207, 224, 228))
    add_chip(slide, x + 0.12, y + 0.1, 1.05, label, fill=accent, color=WHITE)
    textbox(slide, x + 1.28, y + 0.09, w - 1.42, 0.2, equation, size=8.1, color=NAVY, bold=True)
    textbox(slide, x + 1.28, y + 0.36, w - 1.42, 0.18, note, size=5.8, color=MUTED)


def slide_equations(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "Equations Become ML Features", "Feature engineering turns log curves into named variables before modeling; these are inputs or screens, not final labels.")

    equations = [
        ("Vshale", "Vsh = clip((GR - GR_clean) / (GR_shale - GR_clean), 0, 1)", "GR = gamma ray API; GR_clean and GR_shale are local references.", TEAL),
        ("Porosity", "phi_den = clip((rho_ma - RHOB) / (rho_ma - rho_fl), 0, 0.70)", "RHOB = bulk density; rho_ma = matrix density; rho_fl = fluid density.", GREEN),
        ("Sonic", "Vp = 304.8 / DT   |   Vs = 304.8 / DTS", "DT and DTS are slowness in us/ft; Vp and Vs are km/s.", BLUE),
        ("Ratios", "Vp/Vs = Vp / Vs   |   mu-rho = rho_b * Vs^2", "rho_b = bulk density; mu-rho tracks shear rigidity contrast.", PURPLE),
        ("AI", "AI = rho_b * Vp", "AI = acoustic impedance; check density and sonic units before comparing wells.", TEAL),
        ("Moduli", "lambda-rho = rho_b * (Vp^2 - 2Vs^2)   |   K = rho_b * (Vp^2 - 4/3Vs^2)", "Elastic products help separate hydrate stiffness, gas effects, and lithology.", AMBER),
        ("NMR sep", "Delta_NMR = phi_den - phi_nmr   |   NMR_H = clip(Delta_NMR / phi_den, 0, 1)", "phi_nmr = measured NMR porosity; target-derived NMR_SAT stays out of inputs.", GREEN),
        ("Screen", "H_proxy = clip(1 - sqrt(Sw_ref / (phi_den^2 * Rt)), 0, 1)", "Rt = deep resistivity; Sw_ref is a reference only. This is a screening proxy.", RED),
    ]
    for i, item in enumerate(equations):
        equation_card(slide, 0.6, 0.96 + i * 0.72, 7.32, *item)

    add_box(
        slide,
        8.25,
        1.02,
        4.55,
        2.0,
        "Variable dictionary",
        "GR: gamma ray\nRHOB/rho_b: density\nphi_den/phi_nmr: porosities\nDT/DTS: sonic slowness\nVp/Vs: velocities\nAI: acoustic impedance\nRt: deep resistivity\nK, lambda-rho, mu-rho: elastic attributes",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=NAVY,
        title_size=11,
        body_size=8.2,
    )
    add_box(
        slide,
        8.25,
        3.35,
        4.55,
        1.25,
        "Leakage rule",
        "Do not feed target-derived NMR_SAT, core-calibrated labels, or post-outcome interpretations back into the input feature set.",
        fill=WHITE,
        line=RGBColor(218, 226, 229),
        accent=RED,
        title_size=11,
        body_size=8,
    )
    add_box(
        slide,
        8.25,
        4.92,
        4.55,
        1.15,
        "Training rule",
        "Fit scalers, imputers, encoders, and feature selectors only on training wells, then apply them to held-out wells.",
        fill=WHITE,
        line=RGBColor(218, 226, 229),
        accent=GREEN,
        title_size=11,
        body_size=8,
    )
    add_footer(slide, "Equations mirror dashboard/runtime/feature_engineering.py and docs/ML_PARAMETER_TREE_AND_DECK_REVAMP_PLAN.md.")


def slide_ml_pipeline(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "ML Workflow and Model Ladder", "The ML section must explain what is trained, why it is used, and how errors are controlled.")

    stages = [
        ("1. Intake", "public header map\napproved runtime logs", TEAL),
        ("2. QC", "unit checks\nwashout flags\ndepth alignment", BLUE),
        ("3. Features", "equations\ninteractions\nmissingness flags", GREEN),
        ("4. Split", "complete-well groups\nno random-row leakage", RED),
        ("5. Model ladder", "baselines -> trees -> neural nets -> ensemble", PURPLE),
        ("6. Review", "metrics\ncaveats\nuncertainty map", AMBER),
    ]
    for i, (label, body, color) in enumerate(stages):
        x = 0.62 + i * 2.08
        add_box(slide, x, 1.15, 1.72, 1.42, label, body, fill=WHITE, line=RGBColor(207, 224, 228), accent=color, title_size=8.7, body_size=6.5)
        if i < len(stages) - 1:
            add_arrow(slide, x + 1.74, 1.86, x + 2.03, 1.86, color=RGBColor(138, 161, 169), width=1.0)

    add_box(
        slide,
        0.62,
        3.02,
        3.0,
        2.0,
        "Why ML",
        "Gas hydrate response is multivariate. Tree models and neural networks can learn nonlinear interactions between Rt, porosity, GR, Vp, Vs, density, and geologic context when target provenance is controlled.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=TEAL,
        title_size=11,
        body_size=8.4,
    )
    add_box(
        slide,
        3.88,
        3.02,
        4.0,
        2.0,
        "Model ladder",
        "Start with rules and linear/logistic baselines so the workflow has an interpretable reference. Add Random Forest/Gradient Boosting/XGBoost for nonlinear logs. Add ANN/Keras saturation models only when held-out wells improve.",
        fill=WHITE,
        line=RGBColor(218, 226, 229),
        accent=GREEN,
        title_size=11,
        body_size=8.1,
    )
    add_box(
        slide,
        8.15,
        3.02,
        4.55,
        2.0,
        "Known error modes",
        "Leakage from target-derived saturation, random-row overfit, class imbalance, missing NMR/shear sonic, poor depth match, bad-hole intervals, out-of-distribution lithology, and gas/ice/cement resistivity look-alikes.",
        fill=WHITE,
        line=RGBColor(218, 226, 229),
        accent=RED,
        title_size=11,
        body_size=8.1,
    )
    textbox(slide, 0.74, 5.55, 11.72, 0.58, "ML source anchors: Chong et al. (2022) provides the published ANN saturation reference; the two local ML drafts provide the project-specific classification/regression ladder, grouped-well validation, and error controls.", size=9.3, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: Chong et al. 2022; ML Research Paper Draft; Classification Methods Draft; project ML/equation docs.")


def decision_node(slide, x: float, y: float, w: float, h: float, label: str, body: str, accent: RGBColor) -> None:
    add_card(slide, x, y, w, h, fill=WHITE, line=RGBColor(204, 221, 226))
    textbox(slide, x + 0.12, y + 0.12, w - 0.24, 0.22, label, size=9.2, color=NAVY, bold=True)
    textbox(slide, x + 0.12, y + 0.42, w - 0.24, h - 0.52, body, size=6.9, color=MUTED)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.05))
    set_shape_fill(band, accent)
    set_line(band, accent, 0)


def slide_decision_map(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, DEEP)
    add_title(slide, "Detailed ML Decision Map", "The pipeline separates feature creation, target governance, model heads, and geologic review.", dark=True)

    nodes = [
        (0.55, 1.0, 2.25, 1.05, "Log intake", "GR, RHOB, Rt, DT, DTS, NMRPHI, caliper, depth", TEAL),
        (3.05, 1.0, 2.25, 1.05, "QC gates", "units, nulls, washout, depth shift, curve provenance", BLUE),
        (5.55, 1.0, 2.25, 1.05, "Feature store", "Vsh, phi_den, Vp, Vs, Vp/Vs, moduli, flags", GREEN),
        (8.05, 1.0, 2.25, 1.05, "Split policy", "group by well; transformations fit on training wells", RED),
        (10.55, 1.0, 2.25, 1.05, "Target registry", "occurrence class, saturation target, uncertainty source", AMBER),
    ]
    for node in nodes:
        decision_node(slide, *node)
    for x1 in (2.82, 5.32, 7.82, 10.32):
        add_arrow(slide, x1, 1.52, x1 + 0.22, 1.52, color=RGBColor(126, 155, 164), width=1.0)

    decision_node(slide, 1.05, 2.9, 3.2, 1.2, "Classification head", "predict hydrate/no-hydrate or phase class\nmetrics: precision, recall, F1, confusion matrix", TEAL)
    decision_node(slide, 5.08, 2.9, 3.2, 1.2, "Saturation regression head", "predict continuous hydrate saturation\nmetrics: MAE, RMSE, R2, calibration by well", GREEN)
    decision_node(slide, 9.1, 2.9, 3.2, 1.2, "Uncertainty and abstention", "flag missing curves, low confidence, OOD lithology, bad hole", RED)
    add_arrow(slide, 6.67, 2.08, 2.65, 2.9, color=RGBColor(126, 155, 164), width=1.0)
    add_arrow(slide, 8.18, 2.08, 6.65, 2.9, color=RGBColor(126, 155, 164), width=1.0)
    add_arrow(slide, 10.9, 2.08, 10.75, 2.9, color=RGBColor(126, 155, 164), width=1.0)

    add_box(
        slide,
        0.8,
        4.82,
        3.45,
        1.08,
        "Leakage barrier",
        "Target-derived variables, post-outcome interpretations, and fitted preprocessors never cross from validation/test wells into training inputs.",
        fill=RGBColor(239, 245, 247),
        line=RGBColor(73, 101, 112),
        accent=RED,
        title_size=10,
        body_size=7.4,
    )
    add_box(
        slide,
        4.92,
        4.82,
        3.45,
        1.08,
        "Geologic review",
        "Check gas, ice, shale, carbonate/cement, compaction, overburden, salinity, invasion, and bad-hole caveats.",
        fill=RGBColor(239, 245, 247),
        line=RGBColor(73, 101, 112),
        accent=AMBER,
        title_size=10,
        body_size=7.4,
    )
    add_box(
        slide,
        9.03,
        4.82,
        3.45,
        1.08,
        "Output",
        "Well-by-well prediction, confidence class, reason codes, and sweet-spot candidates for review.",
        fill=RGBColor(239, 245, 247),
        line=RGBColor(73, 101, 112),
        accent=GREEN,
        title_size=10,
        body_size=7.4,
    )
    add_footer(slide, "Architecture uses complete-well validation and public-safe feature logic; populated runtime outputs remain ignored.", dark=True)


def slide_errors(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "Errors, Validation, and Masking Review", "A crisp ML deck should show where the workflow can be wrong, not only where it can work.")

    quadrants = [
        ("Classification errors", "False positive: hydrate predicted where caveats explain the signal.\nFalse negative: hydrate missed when curves are absent/noisy.\nReview: confusion matrix, precision, recall, F1.", TEAL),
        ("Regression errors", "Saturation biased high by Rt or biased low by missing NMR.\nReview: MAE, RMSE, R2, residuals by well and interval.", GREEN),
        ("Data pipeline errors", "Random-row leakage, target leakage, unit conversion, depth mismatch, imputation artifacts, bad-hole intervals.\nReview: grouped splits and per-well diagnostics.", RED),
        ("Geologic caveats", "Free gas, permafrost/ice, shale, carbonate/cement, tight rock, salinity, invasion, compaction, stress, and overburden context.\nReview: reason codes.", AMBER),
    ]
    for i, (label, body, color) in enumerate(quadrants):
        x = 0.62 + (i % 2) * 6.2
        y = 1.1 + (i // 2) * 2.25
        add_box(slide, x, y, 5.76, 1.72, label, body, fill=WHITE, line=RGBColor(209, 224, 228), accent=color, title_size=11, body_size=8.1)

    add_picture_contain(slide, SYNTHETIC_LOG, 0.86, 5.76, 2.95, 0.75)
    textbox(slide, 4.2, 5.82, 8.25, 0.55, "Validation standard: final claims should be stated as calibrated predictions with caveats and confidence, not as direct observations from one log curve.", size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Public-safe diagnostics only; approved measured data and populated model outputs belong in authorized runtime folders.")


def slide_final_output(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "Final Output and Next Work", "The nine-slide deck ends with the decision product and the next gated inputs required before true runtime modeling.")

    add_picture_contain(slide, SWEET_SPOT, 0.62, 1.05, 4.6, 4.7)
    add_box(
        slide,
        5.52,
        1.08,
        3.28,
        1.38,
        "Decision product",
        "Candidate hydrate intervals ranked by evidence agreement, model confidence, and geologic caveat review.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=TEAL,
        title_size=11,
        body_size=8.4,
    )
    add_box(
        slide,
        9.12,
        1.08,
        3.08,
        1.38,
        "Do not overclaim",
        "The deck explains a workflow. It does not publish approved logs, restricted identifiers, trained models, or runtime predictions.",
        fill=WHITE,
        line=RGBColor(218, 226, 229),
        accent=RED,
        title_size=11,
        body_size=8.2,
    )

    next_items = [
        ("1", "Finalize workbook formulas and header mappings."),
        ("2", "Load approved log curves only inside ignored runtime folders."),
        ("3", "Confirm target provenance for saturation and phase labels."),
        ("4", "Run complete-well validation and error review."),
        ("5", "Promote only public-safe summaries into docs or slides."),
    ]
    for i, (num, body) in enumerate(next_items):
        y = 2.95 + i * 0.54
        add_shape(slide, "OVAL", 5.6, y, 0.32, 0.32, TEAL, TEAL)
        textbox(slide, 5.65, y + 0.06, 0.21, 0.1, num, size=7.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, 6.05, y + 0.04, 5.95, 0.22, body, size=9.2, color=INK)

    textbox(slide, 5.58, 6.18, 6.55, 0.48, "Final acceptance: exactly 9 slides, first slide profile photo present, slide 4 caveat language corrected, ML equations named, and ML errors/pipelines shown.", size=8.7, color=NAVY, bold=True)
    add_footer(slide, "Next action after this deck: upload/import the verified PPTX to Drive as the 9-slide revision.")


def slide_introduction(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "Introduction: What and Why Gas Hydrates", "A potential natural-gas resource that requires reservoir-scale characterization")

    add_picture_contain(slide, REGIONAL_CONTEXT, 0.62, 1.04, 4.65, 3.45)
    add_box(
        slide,
        0.62,
        4.72,
        4.65,
        1.05,
        "Core message",
        "Stability is necessary, but it is not proof of hydrate or producibility. The project tests whether logs, NMR, core evidence, and geologic context agree.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=TEAL,
        title_size=10.5,
        body_size=7.8,
    )

    intro = [
        ("Gas hydrate", "Methane is held inside crystalline water cages in sediment pore space.", TEAL),
        ("North Slope value", "Permafrost stability, documented hydrate-bearing sands, and large assessed resource potential.", GREEN),
        ("Characterization goal", "Translate resource potential into interval-scale occurrence, saturation, uncertainty, and sweet-spot evidence.", AMBER),
    ]
    for i, (label, body, color) in enumerate(intro):
        add_box(slide, 5.65 + i * 2.35, 1.05, 2.05, 1.42, label, body, fill=WHITE, line=RGBColor(209, 224, 228), accent=color, title_size=10, body_size=7.3)

    chain = [
        ("stability", "PT window"),
        ("reservoir", "sand + pore space"),
        ("gas charge", "supply + migration"),
        ("logs/core", "measured evidence"),
        ("saturation", "quantified target"),
    ]
    for i, (label, body) in enumerate(chain):
        x = 5.72 + i * 1.38
        add_card(slide, x, 3.18, 1.12, 0.78, fill=WHITE, line=RGBColor(209, 224, 228))
        textbox(slide, x + 0.06, 3.28, 1.0, 0.18, label, size=7.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, x + 0.07, 3.55, 0.98, 0.16, body, size=5.7, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(chain) - 1:
            add_arrow(slide, x + 1.13, 3.56, x + 1.34, 3.56, color=RGBColor(130, 154, 164), width=1.0)

    add_box(
        slide,
        5.65,
        4.58,
        6.9,
        1.18,
        "Why this matters for ML",
        "ML is useful only after the science separates occurrence, saturation, reservoir quality, and producibility. The model should rank converging evidence, not convert one high-resistivity interval into a label.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=NAVY,
        title_size=11,
        body_size=8.3,
    )
    add_footer(slide, "Sources: USGS 2019; DOE/NETL; Chong et al. 2022; Lee and Collett 2011; Haines et al. 2022.")


def slide_parameters_scaffold(prs: Presentation, params: list[dict[str, str]]) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "Parameters: Well-Log Scaffold", "Each curve keeps its measured property, hydrate use, caveat, and ML role")

    grid_params = params[:10]
    row_h = 1.04
    columns = [(0.42, grid_params[:5]), (6.82, grid_params[5:])]
    for col_x, rows in columns:
        for i, row in enumerate(rows):
            y = 0.93 + i * row_h
            add_card(slide, col_x, y, 6.1, 0.96, fill=WHITE, line=RGBColor(206, 222, 227))
            draw_parameter_icon(slide, row.get("icon_key", ""), col_x + 0.12, y + 0.15, 0.55)
            textbox(slide, col_x + 0.78, y + 0.06, 2.6, 0.16, short(row["parameter"], 42), size=7.3, color=NAVY, bold=True)
            textbox(slide, col_x + 3.45, y + 0.06, 2.4, 0.16, short(row.get("family", ""), 28), size=5.8, color=TEAL, bold=True, align=PP_ALIGN.RIGHT)
            body = (
                f"Measures: {short(row.get('measures', ''), 60)}\n"
                f"Hydrate use: {short(row.get('hydrate_positive_side', ''), 64)}\n"
                f"Caveats: {short(row.get('slide_caveat_summary', ''), 66)}\n"
                f"ML role: {short(row.get('model_role', ''), 66)}"
            )
            textbox(slide, col_x + 0.78, y + 0.24, 5.12, 0.66, body, size=4.75, color=CHARCOAL)

    add_box(
        slide,
        0.54,
        6.22,
        5.95,
        0.52,
        "Leakage rule",
        "Sgh, S_h, NMR_SAT, phase labels, and final rankings calibrate or score models; they cannot enter predictors.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=RED,
        title_size=8.6,
        body_size=5.9,
    )
    add_box(
        slide,
        6.83,
        6.22,
        5.95,
        0.52,
        "Scaffold rule",
        "Measured curves, derived equations, QC flags, alignment fields, labels, and outputs stay separated through validation.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=TEAL,
        title_size=8.6,
        body_size=5.9,
    )
    add_footer(slide, "Sources: recovered header map; parameter matrix; Chong et al. 2022; Lee and Collett 2011; Haines et al. 2022.")


def slide_ml_architecture_connected(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, DEEP)
    add_title(slide, "ML Methodology: Architecture", "Workflow adapted from gas-hydrate ML sources and approved-data constraints", dark=True)

    # Arrows are drawn first so the nodes read cleanly on top.
    top_y = 1.53
    for x1, x2, label in [
        (1.98, 2.32, "clean + align"),
        (3.72, 4.06, "derive physics"),
        (5.48, 5.82, "combine"),
        (7.34, 7.68, "split by well"),
        (9.22, 9.58, "train + test"),
    ]:
        add_arrow(slide, x1, top_y, x2, top_y, color=RGBColor(136, 167, 176), width=1.1)
        textbox(slide, x1 - 0.08, top_y + 0.1, 0.78, 0.14, label, size=4.9, color=ICE, align=PP_ALIGN.CENTER)

    nodes = [
        (0.55, 1.05, 1.45, 0.95, "Approved logs", "rho/phi, Rt, GR,\nVp, Vs, NMR,\ncaliper, core", TEAL),
        (2.35, 1.05, 1.45, 0.95, "QC gates", "units, depth,\nmissing rows,\ncaliper, outliers", BLUE),
        (4.15, 1.05, 1.45, 0.95, "Equations", "Vsh, phi_den,\nDelta_NMR, AI,\nVp/Vs, elastic", GREEN),
        (5.95, 1.05, 1.55, 0.95, "Feature table", "measured + derived;\ntrain-fit scaling;\nblocked targets", AMBER),
        (7.85, 1.05, 1.45, 0.95, "Split policy", "known wells split;\nprediction wells\nlocked", RED),
        (9.65, 1.05, 1.45, 0.95, "Model ladder", "rules/logit/Ridge\n-> trees/GBM\n-> Keras ANN", PURPLE),
    ]
    for node in nodes:
        decision_node(slide, *node)

    decision_node(slide, 1.05, 3.05, 3.0, 1.05, "Classification branch", "hydrate, gas, water,\ngood sand/no hydrate,\nnon-reservoir, expert review", TEAL)
    decision_node(slide, 5.05, 3.05, 3.0, 1.05, "Saturation regression", "NMR-density Sgh label\nMAE, RMSE, R2\ncalibration by band", GREEN)
    decision_node(slide, 9.05, 3.05, 3.0, 1.05, "Uncertainty review", "abstain if missing curves,\nbad hole, OOD lithology,\nor unresolved caveats", RED)
    add_arrow(slide, 10.38, 2.0, 2.55, 3.05, color=RGBColor(136, 167, 176), width=1.0)
    add_arrow(slide, 10.38, 2.0, 6.55, 3.05, color=RGBColor(136, 167, 176), width=1.0)
    add_arrow(slide, 11.1, 2.0, 10.55, 3.05, color=RGBColor(136, 167, 176), width=1.0)

    decision_node(slide, 0.95, 5.02, 2.9, 0.95, "Chong source anchor", "5 wells; >10k depth points;\nKeras ANN; min-max inputs;\n2-3 log combos tested", AMBER)
    decision_node(slide, 4.35, 5.02, 3.1, 0.95, "Leakage barrier", "NMR/core targets supervise;\npreprocessing fits only\non training wells", RED)
    decision_node(slide, 8.1, 5.02, 3.55, 0.95, "Prediction output", "class probability, Sh,\nreason codes, interval table,\nuncertainty, sweet-spot lanes", GREEN)
    add_arrow(slide, 3.85, 5.5, 4.35, 5.5, color=RED, width=1.2)
    add_arrow(slide, 7.45, 5.5, 8.1, 5.5, color=GREEN, width=1.2)

    textbox(slide, 0.65, 6.32, 12.05, 0.32, "Project upgrade: adapt the published ANN saturation workflow to approved North Slope data with complete-well validation, blocked targets, train-only transforms, and explicit geologic review gates.", size=8.0, color=ICE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: Classification Methods Draft; Chong et al. 2022; ML Project Reference Notes; runtime schemas and feature engineering.", dark=True)


def slide_why_parameters_and_models(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "ML Methodology: Why These Parameters", "Each parameter family earns a model role because it measures a different part of the hydrate system")

    headers = [
        (2.25, "Equation / feature"),
        (4.8, "Physical reason"),
        (7.25, "Caveat or error to check"),
        (9.86, "ML use"),
    ]
    for x, label in headers:
        textbox(slide, x, 0.88, 2.05, 0.16, label, size=6.4, color=TEAL, bold=True)

    lanes = [
        ("Lithology gate", "GR -> Vsh", "clean sand vs shale", "prevents shale false positives", "rules + trees", TEAL),
        ("Porosity/fluid", "RHOB + NMR -> phi_den, Delta_NMR", "pore volume and mobile-fluid loss", "Sgh label must not leak", "regression + ANN", GREEN),
        ("Electrical", "Rt + phi_den -> H_proxy", "hydrate raises resistivity by replacing conductive water", "checks gas, ice, salinity, cement", "tree/GBM interactions", AMBER),
        ("Elastic", "Vp, Vs, rho -> AI, Vp/Vs, lambda-rho, mu-rho", "stiffness and phase response", "checks compaction, cement, gas", "nonlinear models", PURPLE),
        ("QC/context", "caliper + depth + missingness + compartment", "data trust, stability, and structural context", "stops bad-hole, drift, and row leakage", "abstention + grouped split", RED),
    ]
    for i, (family, eqn, reason, caveat, model, color) in enumerate(lanes):
        y = 1.12 + i * 0.88
        add_card(slide, 0.55, y, 12.15, 0.72, fill=WHITE, line=RGBColor(207, 224, 228))
        add_chip(slide, 0.72, y + 0.18, 1.35, family, fill=color, color=WHITE)
        textbox(slide, 2.25, y + 0.14, 2.35, 0.18, eqn, size=7.8, color=NAVY, bold=True)
        textbox(slide, 4.8, y + 0.13, 2.25, 0.18, reason, size=6.7, color=INK)
        textbox(slide, 7.25, y + 0.13, 2.35, 0.18, caveat, size=6.7, color=MUTED)
        textbox(slide, 9.86, y + 0.13, 1.95, 0.18, model, size=6.7, color=INK, bold=True)
        add_arrow(slide, 2.03, y + 0.36, 2.2, y + 0.36, color=RGBColor(96, 126, 137), width=1.0)
        add_arrow(slide, 4.58, y + 0.36, 4.74, y + 0.36, color=RGBColor(96, 126, 137), width=1.0)
        add_arrow(slide, 7.05, y + 0.36, 7.18, y + 0.36, color=RGBColor(96, 126, 137), width=1.0)
        add_arrow(slide, 9.62, y + 0.36, 9.78, y + 0.36, color=RGBColor(96, 126, 137), width=1.0)

    add_box(
        slide,
        0.72,
        5.97,
        3.7,
        0.72,
        "Model ladder",
        "Rule baseline -> logit/Ridge -> RF/GBM -> Keras ANN only if held-out wells improve.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=TEAL,
        title_size=8.8,
        body_size=6.5,
    )
    add_box(
        slide,
        4.82,
        5.97,
        3.7,
        0.72,
        "Validation unit",
        "Fit imputer/scaler/min-max on training wells only; validate by held-out well or compartment.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=GREEN,
        title_size=8.8,
        body_size=6.5,
    )
    add_box(
        slide,
        8.92,
        5.97,
        3.7,
        0.72,
        "Output contract",
        "Report precision/recall/F1, MAE/RMSE/R2, calibration, residuals, reason codes, and review flags.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=PURPLE,
        title_size=8.8,
        body_size=6.5,
    )
    add_footer(slide, "Sources: Classification Methods Draft; Chong et al. 2022; ML Project Reference Notes; Lee and Collett 2011; Haines et al. 2022.")


def slide_geomechanical_feature_sketch(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "Geomechanical Feature Sketch", "Rock-physics parameters help evaluate hydrate-consistent stiffness without treating one curve as proof")

    add_box(slide, 0.62, 1.04, 2.35, 1.35, "Inputs", "rho_b, Vp, Vs\nGR, phi_den, phi_nmr\nRt, caliper, depth", fill=LIGHT, accent=TEAL, title_size=11, body_size=8.2)
    add_box(slide, 3.32, 1.04, 3.45, 1.35, "Feature equations", "AI = rho_b * Vp\nVp/Vs = Vp / Vs\nmu-rho = rho_b * Vs^2\nlambda-rho = rho_b * (Vp^2 - 2Vs^2)", fill=WHITE, accent=GREEN, title_size=11, body_size=7.6)
    add_box(slide, 7.12, 1.04, 2.45, 1.35, "Stiffness question", "Does elastic behavior support pore-filling hydrate, or a non-hydrate hard rock effect?", fill=LIGHT, accent=AMBER, title_size=11, body_size=8.1)
    add_box(slide, 9.92, 1.04, 2.55, 1.35, "Model use", "Elastic features become inputs and reason codes, not direct labels.", fill=WHITE, accent=PURPLE, title_size=11, body_size=8.1)
    add_arrow(slide, 2.98, 1.7, 3.3, 1.7, color=RGBColor(140, 162, 170), width=1.0)
    add_arrow(slide, 6.78, 1.7, 7.1, 1.7, color=RGBColor(140, 162, 170), width=1.0)
    add_arrow(slide, 9.58, 1.7, 9.9, 1.7, color=RGBColor(140, 162, 170), width=1.0)

    checks = [
        ("Hydrate", "higher stiffness + resistivity\nwith clean reservoir/NMR support", GREEN),
        ("Free gas", "lowers Vp more than Vs;\ncan raise Rt but weakens elastic fit", RED),
        ("Shale", "GR and bound water distort\nporosity and NMR response", AMBER),
        ("Ice/cement", "stiff and resistive without\nhydrate pore occupancy", PURPLE),
        ("Overburden", "compaction raises density\nand velocities across baselines", BLUE),
        ("Bad hole", "washout corrupts density,\nsonic, neutron, and resistivity", RED),
    ]
    for i, (label, body, color) in enumerate(checks):
        x = 0.8 + (i % 3) * 4.15
        y = 3.0 + (i // 3) * 1.15
        add_card(slide, x, y, 3.55, 0.82, fill=WHITE, line=RGBColor(207, 224, 228))
        add_chip(slide, x + 0.12, y + 0.16, 0.98, label, fill=color, color=WHITE)
        textbox(slide, x + 1.22, y + 0.16, 2.05, 0.28, body, size=6.6, color=INK)

    textbox(slide, 0.85, 5.85, 11.72, 0.5, "Purpose: integrate mechanical behavior with log response so high resistivity or high velocity becomes evidence to review, not a hydrate label by itself.", size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: local equation documents; Lee and Collett 2011; Haines et al. 2022; runtime feature engineering.")


def slide_3d_map_context(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "3D Map and Well Context", "Regional context constrains interpretation but does not replace log evidence")

    add_picture_contain(slide, REGIONAL_CONTEXT, 0.62, 1.05, 6.15, 4.85)
    add_box(slide, 7.15, 1.08, 2.35, 1.18, "Map use", "Show well positions, structure, data coverage, and candidate interval context.", fill=WHITE, accent=TEAL, title_size=11, body_size=8)
    add_box(slide, 9.85, 1.08, 2.55, 1.18, "Boundary", "Public deck can show public/synthetic context; approved identifiers and results stay runtime-only.", fill=WHITE, accent=RED, title_size=11, body_size=7.8)
    add_box(slide, 7.15, 2.72, 5.25, 1.12, "How it connects to ML", "Map layers explain why complete-well validation and out-of-distribution checks matter: wells can sit in different structural, burial, and reservoir contexts.", fill=LIGHT, accent=GREEN, title_size=11, body_size=8.2)
    add_box(slide, 7.15, 4.22, 5.25, 1.12, "Future replacement", "After approved-data execution, this slide can swap in real runtime-safe summary figures, not raw restricted logs or named sensitive identifiers.", fill=LIGHT, accent=AMBER, title_size=11, body_size=8.2)
    add_footer(slide, "Source anchor: project GIS atlas and future approved-runtime outputs.")


def slide_results_discussion_plan(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "Results and Discussion Plan", "Slides will be populated after approved-data execution")

    add_picture_contain(slide, SYNTHETIC_LOG, 0.58, 1.12, 3.3, 2.35)
    add_picture_contain(slide, SWEET_SPOT, 0.8, 3.82, 3.0, 2.15)

    panels = [
        ("Expected figures", "well-log panel\nstaged-gate table\nconfusion matrix\ncalibration + residual plots", TEAL),
        ("Data quality", "null/missingness table\ncaliper washout counts\nrow loss and outlier audit\nfeature drift by well", BLUE),
        ("Outputs", "probabilities\nSh regression\nreason codes\nsweet-spot lanes", AMBER),
        ("Review flags", "good sand/no hydrate\nexpert review\ncompartment mismatch\ncore-log depth offset", RED),
    ]
    for i, (label, body, color) in enumerate(panels):
        x = 4.38 + (i % 2) * 4.15
        y = 1.12 + (i // 2) * 2.0
        add_box(slide, x, y, 3.6, 1.38, label, body, fill=WHITE, line=RGBColor(209, 224, 228), accent=color, title_size=11, body_size=8.2)

    textbox(slide, 4.52, 5.64, 7.68, 0.48, "Results should explain the decision path: gate outcome, equations driving the class, data-quality state, evidence against competing explanations, and uncertainty after calibration.", size=9.4, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: Classification Methods Draft; USGS/DOE/NETL; Chong et al. 2022; ML Project Reference Notes; Lee and Collett 2011; Haines et al. 2022.")


def slide_conclusion(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_background(slide, WHITE)
    add_title(slide, "Conclusion", "Final project message")

    items = [
        ("Scientific value", "Hydrate occurrence, saturation, reservoir quality, uncertainty, and producibility are separated.", TEAL),
        ("ML value", "Models use physics-backed features, connected equation logic, and complete-well validation.", GREEN),
        ("Energy value", "Characterization supports future North Slope natural-gas resource evaluation and energy security.", AMBER),
    ]
    for i, (label, body, color) in enumerate(items):
        x = 0.82 + i * 4.1
        add_box(slide, x, 1.4, 3.55, 1.55, label, body, fill=WHITE, line=RGBColor(209, 224, 228), accent=color, title_size=12, body_size=9)

    add_box(
        slide,
        1.08,
        3.6,
        11.15,
        1.28,
        "Final message",
        "The goal is a defensible, explainable workflow for predicting gas hydrate occurrence and saturation from approved North Slope well-log, NMR, and core data.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=NAVY,
        title_size=13,
        body_size=11,
    )
    add_box(
        slide,
        1.08,
        5.32,
        11.15,
        0.82,
        "Next work",
        "Gather the full workbook/formulas, confirm saturation and phase-label provenance, run approved-data validation, then replace conceptual graphics with verified runtime figures.",
        fill=WHITE,
        line=RGBColor(218, 226, 229),
        accent=RED,
        title_size=10.5,
        body_size=8.2,
    )
    add_footer(slide, "Sources: USGS/DOE/NETL; Chong et al. 2022; Lee and Collett 2011; Haines et al. 2022; project header/equation docs.")


def add_visual_title(slide, title: str, subtitle: str | None = None, dark: bool = False) -> None:
    title_color = WHITE if dark else NAVY
    sub_color = RGBColor(178, 217, 224) if dark else MUTED
    textbox(slide, 0.52, 0.22, 12.2, 0.42, title, size=25, color=title_color, bold=True)
    if subtitle:
        textbox(slide, 0.55, 0.68, 12.0, 0.26, subtitle, size=9.2, color=sub_color)


def add_source_footer(slide, text: str, dark: bool = False) -> None:
    add_footer(slide, text, dark=dark)


def add_caption(slide, x: float, y: float, w: float, text: str, dark: bool = False) -> None:
    textbox(slide, x, y, w, 0.24, text, size=7.0, color=RGBColor(174, 190, 198) if dark else MUTED, align=PP_ALIGN.CENTER)


def add_symbol_chip(
    slide,
    x: float,
    y: float,
    symbol: str,
    name: str,
    color: RGBColor = TEAL,
    w: float = 1.16,
    h: float = 0.46,
    dark: bool = False,
) -> None:
    fill = RGBColor(25, 70, 84) if dark else RGBColor(232, 244, 246)
    line = RGBColor(91, 164, 177) if dark else RGBColor(190, 219, 224)
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_fill(chip, fill)
    set_line(chip, line, 0.7)
    textbox(slide, x + 0.05, y + 0.04, w - 0.1, 0.16, symbol, size=9.8, color=WHITE if dark else color, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, x + 0.05, y + 0.23, w - 0.1, 0.13, name, size=4.6, color=RGBColor(190, 214, 220) if dark else MUTED, align=PP_ALIGN.CENTER)


def add_icon_tile(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    symbol: str,
    name: str,
    key: str,
    measure: str,
    caveat: str,
    color: RGBColor,
) -> None:
    add_card(slide, x, y, w, h, fill=WHITE, line=RGBColor(206, 222, 227))
    draw_parameter_icon(slide, key, x + 0.18, y + 0.18, 0.68, color)
    textbox(slide, x + 0.98, y + 0.14, 0.72, 0.22, symbol, size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, x + 1.68, y + 0.15, w - 1.86, 0.20, name, size=9.4, color=NAVY, bold=True)
    textbox(slide, x + 0.98, y + 0.46, w - 1.16, 0.20, f"Measures: {measure}", size=6.0, color=INK)
    textbox(slide, x + 0.98, y + 0.69, w - 1.16, 0.20, f"Caveat: {caveat}", size=5.8, color=MUTED)


def draw_curve(slide, x: float, y: float, w: float, h: float, color: RGBColor, points: list[float], width: float = 1.5) -> None:
    if len(points) < 2:
        return
    coords = []
    for i, value in enumerate(points):
        px = x + max(0.03, min(0.97, value)) * w
        py = y + (i / (len(points) - 1)) * h
        coords.append((px, py))
    for (x1, y1), (x2, y2) in zip(coords[:-1], coords[1:]):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1),
            Inches(y1),
            Inches(x2),
            Inches(y2),
        )
        set_line(line, color, width)


def draw_mini_log(slide, x: float, y: float, w: float, h: float, label: str, color: RGBColor, points: list[float]) -> None:
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_fill(frame, RGBColor(246, 249, 250))
    set_line(frame, RGBColor(206, 222, 227), 0.55)
    for i in range(1, 4):
        grid = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + i * w / 4), Inches(y), Inches(x + i * w / 4), Inches(y + h))
        set_line(grid, RGBColor(224, 232, 235), 0.35)
    draw_curve(slide, x + 0.04, y + 0.10, w - 0.08, h - 0.18, color, points, 1.4)
    textbox(slide, x + 0.05, y + h - 0.20, w - 0.1, 0.14, label, size=5.3, color=MUTED, align=PP_ALIGN.CENTER)


def draw_hydrate_cage(slide, x: float, y: float, size: float) -> None:
    center_x = x + size / 2
    center_y = y + size / 2
    radius = size * 0.36
    nodes = []
    for i in range(12):
        angle = (i / 12) * 6.28318
        px = center_x + radius * (0.86 if i % 2 else 1.05) * __import__("math").cos(angle)
        py = center_y + radius * (0.72 if i % 2 else 1.0) * __import__("math").sin(angle)
        nodes.append((px, py))
    for i, (x1, y1) in enumerate(nodes):
        x2, y2 = nodes[(i + 1) % len(nodes)]
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        set_line(line, ICE, 0.95)
    for px, py in nodes:
        o = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(px - 0.045), Inches(py - 0.045), Inches(0.09), Inches(0.09))
        set_shape_fill(o, RGBColor(204, 237, 244))
        set_line(o, WHITE, 0.3)
    methane = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(center_x - 0.14), Inches(center_y - 0.14), Inches(0.28), Inches(0.28))
    set_shape_fill(methane, AMBER)
    set_line(methane, WHITE, 0.6)
    textbox(slide, center_x - 0.14, center_y - 0.035, 0.28, 0.08, "CH4", size=5.6, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def draw_about_icon(slide, kind: str, x: float, y: float, color: RGBColor) -> None:
    add_card(slide, x, y, 0.78, 0.78, fill=RGBColor(239, 247, 248), line=RGBColor(203, 225, 229))
    if kind == "draw":
        pencil = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, Inches(x + 0.23), Inches(y + 0.19), Inches(0.34), Inches(0.12))
        pencil.rotation = -30
        set_shape_fill(pencil, color)
        set_line(pencil, color, 0.5)
        draw_curve(slide, x + 0.18, y + 0.50, 0.42, 0.12, color, [0.1, 0.5, 0.25, 0.85], 1.0)
    elif kind == "gym":
        bar = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.18), Inches(y + 0.39), Inches(x + 0.60), Inches(y + 0.39))
        set_line(bar, color, 1.4)
        for dx in (0.12, 0.60):
            for j in range(2):
                plate = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + dx + j * 0.03), Inches(y + 0.30), Inches(0.03), Inches(0.18))
                set_shape_fill(plate, color)
                set_line(plate, color, 0)
    elif kind == "run":
        head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.34), Inches(y + 0.12), Inches(0.12), Inches(0.12))
        set_shape_fill(head, color)
        set_line(head, color, 0)
        for x1, y1, x2, y2 in [(0.40, 0.25, 0.34, 0.43), (0.34, 0.43, 0.23, 0.56), (0.34, 0.43, 0.55, 0.53), (0.36, 0.31, 0.20, 0.34), (0.37, 0.31, 0.55, 0.26)]:
            limb = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + x1), Inches(y + y1), Inches(x + x2), Inches(y + y2))
            set_line(limb, color, 1.25)
    else:
        for i in range(3):
            draw_curve(slide, x + 0.14, y + 0.26 + i * 0.13, 0.50, 0.08, color, [0.0, 0.30, 0.58, 1.0], 1.0)
        head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.50), Inches(y + 0.17), Inches(0.10), Inches(0.10))
        set_shape_fill(head, color)
        set_line(head, color, 0)


def v2_slide_title(prs: Presentation) -> None:
    slide = blank_slide(prs)
    if add_processing_slide(slide, "title"):
        return
    fill_background(slide, WHITE)
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(7.25), Inches(7.5))
    set_shape_fill(left, RGBColor(246, 250, 251))
    set_line(left, RGBColor(246, 250, 251), 0)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    set_shape_fill(accent, TEAL)
    set_line(accent, TEAL, 0)

    textbox(slide, 0.60, 0.62, 6.25, 1.05, "Gas Hydrate Occurrence and Saturation Prediction", size=26, color=NAVY, bold=True)
    textbox(slide, 0.63, 1.88, 5.95, 0.44, "Alaska North Slope permafrost reservoirs using physics-constrained AI/ML", size=14, color=INK, bold=True)
    textbox(slide, 0.64, 2.64, 5.75, 0.42, "Goal: combine approved well logs, NMR, core context, and public GIS without exposing runtime-only data.", size=10.4, color=MUTED)

    add_chip(slide, 0.64, 3.35, 1.62, "source-backed", TEAL, WHITE)
    add_chip(slide, 2.43, 3.35, 1.50, "9 slides", RGBColor(217, 232, 236), NAVY)
    add_chip(slide, 4.10, 3.35, 1.62, "runtime-safe", RGBColor(232, 221, 189), NAVY)

    textbox(slide, 0.66, 4.34, 1.2, 0.22, "About me", size=10.8, color=NAVY, bold=True)
    for i, (kind, label, color) in enumerate([("draw", "drawing", TEAL), ("gym", "gym", GREEN), ("run", "running", AMBER), ("swim", "swimming", BLUE)]):
        x = 0.68 + i * 1.28
        draw_about_icon(slide, kind, x, 4.76, color)
        textbox(slide, x - 0.05, 5.62, 0.88, 0.15, label, size=6.2, color=MUTED, align=PP_ALIGN.CENTER)

    add_picture_contain(slide, PROFILE_PHOTO, 7.58, 0.56, 4.95, 6.02)
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.58), Inches(0.56), Inches(4.95), Inches(6.02))
    frame.fill.background()
    set_line(frame, NAVY, 1.0)
    add_source_footer(slide, "Public deck only: approved logs, restricted identifiers, runtime outputs, and trained models remain outside this file.")


def v2_slide_introduction(prs: Presentation) -> None:
    slide = blank_slide(prs)
    if add_processing_slide(slide, "intro"):
        return
    fill_background(slide, WHITE)
    add_visual_title(slide, "Introduction: What and Why Gas Hydrates", "Hydrate is a pressure-temperature system, a reservoir question, and a log-interpretation problem.")

    add_picture_contain(slide, STREAMLIT_STRUCTURAL, 6.35, 1.03, 6.40, 4.35)
    add_caption(slide, 6.66, 5.45, 5.75, "Current Streamlit structural explorer: public wells, boundaries, and horizons")

    photo = add_picture_contain(slide, HYDRATE_CRYSTALS, 0.62, 1.06, 2.92, 1.72)
    if photo:
        photo.line.color.rgb = RGBColor(205, 219, 224)
        photo.line.width = Pt(0.8)
    draw_hydrate_cage(slide, 3.72, 1.05, 1.65)
    textbox(slide, 0.72, 2.93, 4.62, 0.36, "Methane sits inside an ice-like water cage in sediment pore space.", size=11.5, color=NAVY, bold=True)

    tiles = [
        ("What", "methane + water cages", TEAL, "nmr"),
        ("Where", "permafrost + marine sediments", BLUE, "depth"),
        ("Why", "large North Slope resource", GREEN, "layer"),
        ("Goal", "predict occurrence and saturation separately", AMBER, "elastic"),
    ]
    for i, (label, body, color, key) in enumerate(tiles):
        x = 0.64 + (i % 2) * 2.68
        y = 3.70 + (i // 2) * 1.08
        add_card(slide, x, y, 2.30, 0.82, fill=WHITE, line=RGBColor(208, 225, 229))
        draw_parameter_icon(slide, key, x + 0.12, y + 0.14, 0.48, color)
        textbox(slide, x + 0.70, y + 0.13, 1.38, 0.18, label, size=9.6, color=color, bold=True)
        textbox(slide, x + 0.70, y + 0.38, 1.42, 0.24, body, size=6.4, color=INK)

    textbox(slide, 6.52, 6.12, 6.05, 0.40, "Map context narrows the question. Hydrate evidence still comes from well logs, NMR/core calibration, and uncertainty review.", size=10.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_source_footer(slide, "Sources: USGS FAQ/Primer; USGS FS 2019-3037; DOE/NETL North Slope hydrate pages; project Streamlit asset.")


def v2_slide_parameters(prs: Presentation) -> None:
    slide = blank_slide(prs)
    if add_processing_slide(slide, "parameters"):
        return
    fill_background(slide, WHITE)
    add_visual_title(slide, "Parameters: Well-Log Scaffold", "Symbols and icons stay consistent through the deck; targets remain locked outside the input set.")

    tiles = [
        ("GR", "Gamma ray", "gamma", "natural radioactivity", "low GR is not hydrate", AMBER),
        ("Rt", "Deep resistivity", "resistivity", "electrical resistance", "gas/ice/cement can mimic", RED),
        ("RHOB | phi_D", "Density porosity", "density", "bulk density + pore volume", "mineralogy and washout shift it", GREEN),
        ("phi_NMR", "NMR porosity", "nmr", "pore-fluid signal", "target leakage if NMR_SAT is used", BLUE),
        ("Vp | Vs", "Sonic velocities", "vp", "P-wave and S-wave speed", "compaction/cement alter baseline", PURPLE),
        ("AI | Vp/Vs", "Elastic features", "elastic", "derived stiffness ratios", "inherits input errors", TEAL),
        ("CAL/DCAL", "Borehole QC", "caliper", "hole diameter/washout", "QC gate, not hydrate evidence", RED),
        ("z | P-T | Core", "Context + calibration", "depth", "depth, stability, core ties", "stability is necessary only", NAVY),
    ]
    for i, tile in enumerate(tiles):
        x = 0.56 + (i % 4) * 3.18
        y = 1.15 + (i // 4) * 1.56
        add_icon_tile(slide, x, y, 2.82, 1.16, *tile)

    add_picture_contain(slide, GAMMA_LOG_REF, 0.78, 4.70, 2.35, 1.20)
    add_picture_contain(slide, RESISTIVITY_LOG_REF, 3.28, 4.70, 2.35, 1.20)
    add_picture_contain(slide, CALIPER_LOG_REF, 5.78, 4.70, 2.35, 1.20)
    add_box(
        slide,
        8.55,
        4.72,
        3.92,
        1.17,
        "Locked target fields",
        "S_h, Sgh, NMR_SAT, phase labels, and final rankings supervise or score models; they do not become predictors.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=RED,
        title_size=10.5,
        body_size=7.1,
    )
    add_source_footer(slide, "Sources: WELL_LOG_REQUIREMENTS_MAP; parameter matrix; USGS gamma/caliper/resistivity images; Haines et al. 2022; NETL ML source.")


def v2_gate(slide, x: float, y: float, label: str, color: RGBColor, dark: bool = True) -> None:
    fill = RGBColor(18, 56, 70) if dark else WHITE
    line = RGBColor(76, 121, 136) if dark else RGBColor(206, 222, 227)
    add_card(slide, x, y, 1.48, 0.50, fill=fill, line=line)
    textbox(slide, x + 0.07, y + 0.13, 1.34, 0.14, label, size=5.9, color=WHITE if dark else INK, bold=True, align=PP_ALIGN.CENTER)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y + 0.46), Inches(1.48), Inches(0.04))
    set_shape_fill(bar, color)
    set_line(bar, color, 0)


def v2_slide_ml_architecture(prs: Presentation) -> None:
    slide = blank_slide(prs)
    if add_processing_slide(slide, "architecture"):
        return
    fill_background(slide, DEEP)
    add_visual_title(slide, "ML Methodology: Architecture", "A detailed evidence pipeline: parameters -> QC gates -> equations -> model branches -> reviewable outputs.", dark=True)

    input_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.52), Inches(1.08), Inches(1.55), Inches(4.55))
    set_shape_fill(input_panel, DARK_PANEL)
    set_line(input_panel, RGBColor(65, 107, 121), 0.8)
    textbox(slide, 0.64, 1.22, 1.30, 0.18, "Approved logs", size=8.2, color=ICE, bold=True, align=PP_ALIGN.CENTER)
    chips = [("GR", "lith"), ("Rt", "elec"), ("RHOB", "dens"), ("phi_NMR", "nmr"), ("Vp", "P"), ("Vs", "S"), ("CAL", "qc"), ("Core", "cal")]
    for i, (symbol, name) in enumerate(chips):
        add_symbol_chip(slide, 0.67, 1.62 + i * 0.43, symbol, name, TEAL, 1.24, 0.34, dark=True)

    gate_labels = [
        ("units + source", TEAL),
        ("depth align", BLUE),
        ("caliper QC", RED),
        ("missing/outlier", AMBER),
        ("sand gate", GREEN),
        ("leakage lock", RED),
    ]
    gate_ys = [1.18 + i * 0.58 for i in range(len(gate_labels))]
    for y in gate_ys:
        add_arrow(slide, 2.07, y + 0.25, 2.35, y + 0.25, RGBColor(105, 140, 153), 0.75)
        add_arrow(slide, 3.83, y + 0.25, 4.23, y + 0.25, RGBColor(105, 140, 153), 0.75)
    for i, (label, color) in enumerate(gate_labels):
        v2_gate(slide, 2.35, gate_ys[i], label, color)

    equations = [
        ("Vsh", "from GR"),
        ("phi_D", "from RHOB"),
        ("Delta_NMR-D", "phi_D - phi_NMR"),
        ("Vp/Vs", "Vp / Vs"),
        ("AI", "rho_b * Vp"),
        ("mu-rho", "rho_b * Vs^2"),
        ("lambda-rho", "rho_b * (Vp^2 - 2Vs^2)"),
    ]
    add_card(slide, 4.23, 1.08, 1.66, 3.56, fill=RGBColor(17, 58, 64), line=RGBColor(68, 119, 126))
    textbox(slide, 4.38, 1.21, 1.36, 0.18, "Equations", size=8.3, color=ICE, bold=True, align=PP_ALIGN.CENTER)
    for i, (eq, desc) in enumerate(equations):
        textbox(slide, 4.36, 1.56 + i * 0.40, 0.74, 0.15, eq, size=6.4, color=WHITE, bold=True)
        textbox(slide, 5.04, 1.57 + i * 0.40, 0.74, 0.15, desc, size=4.7, color=RGBColor(179, 205, 211))

    add_arrow(slide, 5.89, 2.08, 6.38, 2.08, RGBColor(121, 153, 165), 1.0)
    add_card(slide, 6.38, 1.20, 1.55, 1.48, fill=RGBColor(28, 63, 76), line=RGBColor(77, 123, 138))
    textbox(slide, 6.54, 1.38, 1.22, 0.18, "Feature table", size=8.2, color=ICE, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 6.54, 1.76, 1.22, 0.42, "train-only scaling\nmeasured + derived\nblocked targets", size=5.7, color=WHITE, align=PP_ALIGN.CENTER)

    add_arrow(slide, 7.93, 2.08, 8.42, 2.08, RGBColor(121, 153, 165), 1.0)
    add_card(slide, 8.42, 1.20, 1.42, 1.48, fill=RGBColor(45, 60, 88), line=RGBColor(102, 112, 154))
    textbox(slide, 8.55, 1.38, 1.16, 0.18, "Split policy", size=8.0, color=ICE, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 8.57, 1.76, 1.10, 0.38, "held-out wells\nno random rows", size=5.9, color=WHITE, align=PP_ALIGN.CENTER)

    model_ladder = [("Rules", TEAL), ("Logit/Ridge", GREEN), ("RF/GBM", AMBER), ("Keras ANN", PURPLE)]
    add_arrow(slide, 9.84, 2.08, 10.42, 2.08, RGBColor(121, 153, 165), 1.0)
    for i, (label, color) in enumerate(model_ladder):
        add_card(slide, 10.42, 1.12 + i * 0.48, 1.36, 0.35, fill=RGBColor(22, 56, 69), line=RGBColor(68, 112, 126))
        textbox(slide, 10.52, 1.21 + i * 0.48, 1.16, 0.12, label, size=5.8, color=color, bold=True, align=PP_ALIGN.CENTER)

    branches = [
        (0.98, "Classification", "hydrate | gas | water\nsand/no hydrate | shale\nexpert review", TEAL),
        (4.88, "Saturation regression", "S_h / Sgh prediction\nMAE, RMSE, R2\ncalibration by interval", GREEN),
        (8.78, "Uncertainty + outputs", "probability, reason codes\nreview flags\nsweet-spot lanes", AMBER),
    ]
    add_arrow(slide, 3.10, 4.58, 2.46, 4.72, RGBColor(121, 153, 165), 0.9)
    add_arrow(slide, 5.12, 4.64, 6.26, 4.72, RGBColor(121, 153, 165), 0.9)
    add_arrow(slide, 10.98, 3.02, 10.28, 4.72, RGBColor(121, 153, 165), 0.9)
    for x, label, body, color in branches:
        add_card(slide, x, 4.72, 3.0, 0.95, fill=RGBColor(16, 50, 63), line=RGBColor(73, 118, 132))
        textbox(slide, x + 0.16, 4.87, 2.66, 0.17, label, size=9.0, color=color, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, x + 0.20, 5.14, 2.58, 0.32, body, size=5.8, color=WHITE, align=PP_ALIGN.CENTER)

    rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.42), Inches(6.05), Inches(8.58), Inches(0.09))
    set_shape_fill(rail, RED)
    set_line(rail, RED, 0)
    textbox(slide, 2.55, 6.25, 8.3, 0.24, "Locked targets below this rail: S_h, Sgh, NMR_SAT, phase labels, final rankings. They supervise or score models; they do not become predictors.", size=7.5, color=ICE, bold=True, align=PP_ALIGN.CENTER)
    add_source_footer(slide, "Sources: Classification Methods Draft; Chong et al. 2022; NETL ML source; WELL_LOG_REQUIREMENTS_MAP; runtime feature engineering.", dark=True)


def v2_behavior_panel(slide, x: float, y: float, title: str, subtitle: str, color: RGBColor, curves: list[tuple[str, list[float], RGBColor]]) -> None:
    add_card(slide, x, y, 3.78, 1.72, fill=WHITE, line=RGBColor(205, 222, 227))
    textbox(slide, x + 0.18, y + 0.15, 3.35, 0.20, title, size=10.2, color=color, bold=True)
    textbox(slide, x + 0.18, y + 0.43, 3.30, 0.18, subtitle, size=6.7, color=MUTED)
    for i, (label, points, curve_color) in enumerate(curves):
        draw_mini_log(slide, x + 0.20 + i * 1.14, y + 0.78, 0.88, 0.72, label, curve_color, points)


def v2_slide_parameter_behavior(prs: Presentation) -> None:
    slide = blank_slide(prs)
    if add_processing_slide(slide, "behavior"):
        return
    fill_background(slide, WHITE)
    add_visual_title(slide, "ML Methodology: Why These Parameters", "The model learns behavior patterns, not isolated curve slogans.")

    v2_behavior_panel(
        slide,
        0.58,
        1.10,
        "Clean sand",
        "reservoir gate, not hydrate proof",
        GREEN,
        [("GR", [0.24, 0.28, 0.18, 0.26, 0.23], AMBER), ("phi", [0.58, 0.62, 0.65, 0.60, 0.63], GREEN), ("Rt", [0.30, 0.34, 0.32, 0.29, 0.33], RED)],
    )
    v2_behavior_panel(
        slide,
        4.78,
        1.10,
        "Hydrate in sand",
        "high Rt + NMR-density separation + stiffness",
        TEAL,
        [("Rt", [0.40, 0.82, 0.88, 0.76, 0.42], RED), ("NMR", [0.62, 0.28, 0.24, 0.31, 0.58], BLUE), ("Vp", [0.46, 0.72, 0.76, 0.70, 0.50], PURPLE)],
    )
    v2_behavior_panel(
        slide,
        8.98,
        1.10,
        "Shale",
        "high GR and bound water can mislead porosity",
        AMBER,
        [("GR", [0.68, 0.74, 0.82, 0.70, 0.76], AMBER), ("phi", [0.55, 0.58, 0.52, 0.57, 0.54], GREEN), ("Rt", [0.28, 0.34, 0.31, 0.36, 0.30], RED)],
    )
    v2_behavior_panel(
        slide,
        0.58,
        3.40,
        "Free gas",
        "Rt may rise; Vp weakens relative to Vs",
        RED,
        [("Rt", [0.42, 0.70, 0.78, 0.68, 0.38], RED), ("Vp", [0.62, 0.28, 0.24, 0.31, 0.58], BLUE), ("Vs", [0.50, 0.52, 0.54, 0.51, 0.53], PURPLE)],
    )
    v2_behavior_panel(
        slide,
        4.78,
        3.40,
        "Ice / cement / carbonate",
        "stiff and resistive mimic without hydrate proof",
        PURPLE,
        [("Rt", [0.46, 0.78, 0.82, 0.75, 0.48], RED), ("Vp", [0.52, 0.76, 0.80, 0.74, 0.55], BLUE), ("GR", [0.30, 0.28, 0.26, 0.29, 0.27], AMBER)],
    )
    v2_behavior_panel(
        slide,
        8.98,
        3.40,
        "Bad hole",
        "washout corrupts density, sonic, NMR, Rt",
        RED,
        [("CAL", [0.25, 0.30, 0.90, 0.88, 0.32], RED), ("RHOB", [0.52, 0.50, 0.20, 0.22, 0.48], GREEN), ("Vp", [0.56, 0.54, 0.30, 0.34, 0.55], BLUE)],
    )

    textbox(slide, 1.02, 6.02, 11.25, 0.40, "Interpretation rule: each parameter is evidence only after lithology, borehole quality, depth alignment, and competing explanations are checked.", size=10.6, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_source_footer(slide, "Sources: parameter matrix; Haines et al. 2022; NETL ML source; USGS gamma/caliper/resistivity references.")


def v2_slide_geomechanics(prs: Presentation) -> None:
    slide = blank_slide(prs)
    if add_processing_slide(slide, "geomechanics"):
        return
    fill_background(slide, WHITE)
    add_visual_title(slide, "Geomechanical Feature Sketch", "Rock-physics features test hydrate-consistent stiffness without turning one curve into a label.")

    for i, (symbol, name, color) in enumerate([("RHOB", "density", GREEN), ("Vp", "P wave", BLUE), ("Vs", "S wave", PURPLE), ("Rt", "resistivity", RED), ("GR", "lithology", AMBER), ("CAL", "QC", TEAL)]):
        add_symbol_chip(slide, 0.68, 1.18 + i * 0.62, symbol, name, color, 1.25, 0.42)

    rock = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, Inches(3.05), Inches(1.40), Inches(3.45), Inches(2.45))
    set_shape_fill(rock, RGBColor(218, 226, 219))
    set_line(rock, RGBColor(135, 158, 150), 1.2)
    for i in range(26):
        px = 3.28 + (i % 7) * 0.42
        py = 1.63 + (i // 7) * 0.46
        pore = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(px), Inches(py), Inches(0.16), Inches(0.11))
        set_shape_fill(pore, RGBColor(187, 211, 215) if i % 4 else RGBColor(235, 241, 242))
        set_line(pore, RGBColor(187, 211, 215), 0.2)
    draw_curve(slide, 2.25, 1.82, 1.15, 0.42, BLUE, [0.02, 0.25, 0.45, 0.70, 0.95], 2.0)
    add_arrow(slide, 2.78, 2.05, 3.12, 2.05, BLUE, 1.3)
    draw_curve(slide, 2.25, 2.83, 1.15, 0.42, PURPLE, [0.10, 0.88, 0.18, 0.92, 0.20], 2.0)
    add_arrow(slide, 2.78, 3.06, 3.12, 3.06, PURPLE, 1.3)
    textbox(slide, 3.58, 3.99, 2.36, 0.18, "P-wave + S-wave response through hydrate-bearing sediment", size=7.0, color=MUTED, align=PP_ALIGN.CENTER)

    equations = [
        ("AI", "rho_b * Vp"),
        ("Vp/Vs", "Vp / Vs"),
        ("mu-rho", "rho_b * Vs^2"),
        ("lambda-rho", "rho_b * (Vp^2 - 2Vs^2)"),
    ]
    for i, (label, eq) in enumerate(equations):
        y = 1.20 + i * 0.70
        add_card(slide, 7.05, y, 2.55, 0.50, fill=LIGHT, line=RGBColor(205, 222, 227))
        textbox(slide, 7.18, y + 0.12, 0.78, 0.13, label, size=8.3, color=TEAL, bold=True)
        textbox(slide, 7.98, y + 0.12, 1.38, 0.13, eq, size=7.3, color=INK, bold=True)

    add_box(
        slide,
        10.05,
        1.20,
        2.25,
        2.60,
        "Model use",
        "Elastic features become inputs and reason codes.\n\nThey are reviewed with Rt, GR, NMR, caliper, and depth context.",
        fill=WHITE,
        line=RGBColor(205, 222, 227),
        accent=PURPLE,
        title_size=12,
        body_size=8,
    )

    checks = [
        ("hydrate-supportive", GREEN),
        ("free gas", RED),
        ("shale", AMBER),
        ("ice/cement/carbonate", PURPLE),
        ("overburden", BLUE),
        ("bad hole", RED),
    ]
    for i, (label, color) in enumerate(checks):
        x = 0.78 + i * 2.05
        add_chip(slide, x, 5.32, 1.55, label, fill=color, color=WHITE)
    textbox(slide, 1.15, 6.12, 11.05, 0.36, "Takeaway: high resistivity or high velocity is evidence to review, not a hydrate label by itself.", size=12.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_source_footer(slide, "Sources: runtime feature engineering; ODP sonic-tools reference; Haines et al. 2022; parameter matrix.")


def v2_slide_map_context(prs: Presentation) -> None:
    slide = blank_slide(prs)
    if add_processing_slide(slide, "map"):
        return
    fill_background(slide, WHITE)
    add_visual_title(slide, "3D Map and Well Context", "Use the current Streamlit structural explorer as regional context, then route evidence back to logs.")

    add_picture_contain(slide, STREAMLIT_STRUCTURAL, 0.62, 1.05, 7.10, 4.72)
    add_caption(slide, 1.10, 5.84, 6.18, "Public Streamlit Structural Explorer: no approved runtime well-log values")

    steps = [
        ("Public map", "structure, public wells, boundaries", TEAL, "depth"),
        ("Runtime logs", "approved curves and core ties", GREEN, "core"),
        ("ML review", "probabilities, reason codes, flags", AMBER, "elastic"),
    ]
    for i, (label, body, color, key) in enumerate(steps):
        y = 1.28 + i * 1.40
        add_card(slide, 8.10, y, 3.70, 0.94, fill=WHITE, line=RGBColor(205, 222, 227))
        draw_parameter_icon(slide, key, 8.25, y + 0.18, 0.52, color)
        textbox(slide, 8.96, y + 0.16, 2.30, 0.18, label, size=10.5, color=color, bold=True)
        textbox(slide, 8.96, y + 0.47, 2.38, 0.18, body, size=6.8, color=INK)
        if i < 2:
            add_arrow(slide, 9.95, y + 1.00, 9.95, y + 1.30, RGBColor(126, 153, 163), 1.0)

    add_box(
        slide,
        8.10,
        5.45,
        3.70,
        0.72,
        "Boundary",
        "GIS constrains interpretation; it does not replace multi-log evidence.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=RED,
        title_size=9.5,
        body_size=6.8,
    )
    add_source_footer(slide, "Source anchor: current public Streamlit atlas and future approved-runtime summary outputs.")


def v2_slide_results_plan(prs: Presentation) -> None:
    slide = blank_slide(prs)
    if add_processing_slide(slide, "results"):
        return
    fill_background(slide, WHITE)
    add_visual_title(slide, "Results and Discussion Plan", "Final results should explain the decision path, not only show a prediction.")

    add_picture_contain(slide, SYNTHETIC_LOG, 0.62, 1.08, 4.95, 3.00)
    add_picture_contain(slide, SWEET_SPOT, 0.90, 4.50, 4.35, 1.55)

    outputs = [
        ("Gate table", "QC, stability, reservoir, phase evidence", TEAL, "caliper"),
        ("Model figures", "confusion, calibration, residuals", BLUE, "cross"),
        ("Interval output", "class probability + S_h", GREEN, "resistivity"),
        ("Review flags", "bad hole, shale, gas, OOD context", RED, "depth"),
    ]
    for i, (label, body, color, key) in enumerate(outputs):
        x = 6.15 + (i % 2) * 3.05
        y = 1.25 + (i // 2) * 1.55
        add_card(slide, x, y, 2.70, 1.04, fill=WHITE, line=RGBColor(205, 222, 227))
        draw_parameter_icon(slide, key, x + 0.16, y + 0.18, 0.52, color)
        textbox(slide, x + 0.82, y + 0.19, 1.52, 0.18, label, size=9.3, color=color, bold=True)
        textbox(slide, x + 0.82, y + 0.49, 1.62, 0.24, body, size=6.2, color=INK)

    add_box(
        slide,
        6.20,
        4.62,
        5.62,
        0.98,
        "Discussion lens",
        "Explain why an interval was classified, what evidence disagreed, and what uncertainty or review flag remains.",
        fill=LIGHT,
        line=RGBColor(204, 221, 226),
        accent=NAVY,
        title_size=11,
        body_size=8.1,
    )
    add_source_footer(slide, "Sources: Classification Methods Draft; USGS/DOE/NETL; Chong et al. 2022; Haines et al. 2022; project scaffold visuals.")


def v2_slide_conclusion(prs: Presentation) -> None:
    slide = blank_slide(prs)
    if add_processing_slide(slide, "conclusion"):
        return
    fill_background(slide, WHITE)
    add_visual_title(slide, "Conclusion", "A defensible workflow separates hydrate occurrence, saturation, reservoir quality, uncertainty, and producibility.")

    center = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(4.55), Inches(1.40), Inches(4.10), Inches(2.35))
    set_shape_fill(center, RGBColor(231, 243, 245))
    set_line(center, RGBColor(158, 205, 212), 1.1)
    textbox(slide, 5.10, 2.00, 3.00, 0.32, "Explainable hydrate prediction", size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 5.15, 2.55, 2.90, 0.30, "physics-backed features + complete-well validation + public/runtime boundary", size=8.2, color=MUTED, align=PP_ALIGN.CENTER)

    values = [
        (0.78, 1.30, "Science", "separate occurrence and saturation", TEAL, "nmr"),
        (9.65, 1.30, "ML", "transparent features and reason codes", GREEN, "elastic"),
        (0.78, 4.35, "Energy", "rank sweet spots for future evaluation", AMBER, "depth"),
        (9.65, 4.35, "Next", "confirm workbook labels and runtime figures", RED, "core"),
    ]
    for x, y, label, body, color, key in values:
        add_card(slide, x, y, 2.75, 1.10, fill=WHITE, line=RGBColor(205, 222, 227))
        draw_parameter_icon(slide, key, x + 0.16, y + 0.22, 0.52, color)
        textbox(slide, x + 0.85, y + 0.23, 1.50, 0.18, label, size=10.4, color=color, bold=True)
        textbox(slide, x + 0.85, y + 0.54, 1.58, 0.24, body, size=6.6, color=INK)
        add_arrow(slide, x + (2.75 if x < 4 else 0), y + 0.56, 4.70 if x < 4 else 8.62, 2.58 if y < 3 else 3.08, RGBColor(148, 171, 180), 0.8)

    textbox(slide, 1.50, 6.10, 10.34, 0.36, "Final message: predict hydrate occurrence and saturation only when public context, approved logs, target provenance, and validation all stay traceable.", size=11.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_source_footer(slide, "Sources: USGS/DOE/NETL; Chong et al. 2022; Haines et al. 2022; project header/equation docs.")


def build_deck() -> Presentation:
    build_processing_assets(ROOT)
    prs = Presentation(str(BASE_PPTX)) if BASE_PPTX.exists() else Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    clear_deck(prs)

    v2_slide_title(prs)
    v2_slide_introduction(prs)
    v2_slide_parameters(prs)
    v2_slide_ml_architecture(prs)
    v2_slide_parameter_behavior(prs)
    v2_slide_geomechanics(prs)
    v2_slide_map_context(prs)
    v2_slide_results_plan(prs)
    v2_slide_conclusion(prs)
    return prs


def main() -> None:
    prs = build_deck()
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"Wrote {OUT_PPTX} with {len(prs.slides)} slides.")


if __name__ == "__main__":
    main()
