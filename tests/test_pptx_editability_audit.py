from pathlib import Path
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

import importlib.util


MODULE_PATH = Path(__file__).resolve().parents[1] / "docs" / "project_blueprints" / "audit_pptx_editability.py"
SPEC = importlib.util.spec_from_file_location("audit_pptx_editability", MODULE_PATH)
audit_module = importlib.util.module_from_spec(SPEC)
assert SPEC.loader is not None
sys.modules[SPEC.name] = audit_module
SPEC.loader.exec_module(audit_module)


def test_editable_text_slide_passes(tmp_path: Path) -> None:
    deck_path = tmp_path / "editable.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    box.text = "Editable title"
    slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(2), Inches(1))
    presentation.save(deck_path)

    results = audit_module.audit_presentation(deck_path)

    assert len(results) == 1
    assert results[0].text_shape_count == 1
    assert results[0].passes_editability is True


def test_single_full_slide_picture_fails(tmp_path: Path) -> None:
    image_path = tmp_path / "panel.png"
    Image.new("RGB", (1280, 720), "white").save(image_path)

    deck_path = tmp_path / "flat.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    presentation.save(deck_path)

    results = audit_module.audit_presentation(deck_path)

    assert results[0].one_picture_full_slide_raster is True
    assert results[0].passes_editability is False
